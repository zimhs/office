import io
import os
import re
import sys
import html
import json
import glob
import time
import subprocess
import shutil
import tempfile
import urllib.parse
import requests
from bs4 import BeautifulSoup
import numpy as np
import pandas as pd
import streamlit as st
import streamlit.components.v1 as components
import plotly.graph_objects as go
import plotly.express as px
import datetime
from matplotlib.colors import LinearSegmentedColormap
import sys
# 유저님 맥북에 다트 부품이 확실하게 설치되어 있는 '진짜 주소'를 강제로 주입합니다!
sys.path.append("/Library/Frameworks/Python.framework/Versions/3.13/lib/python3.13/site-packages")

try:
    import OpenDartReader
except ImportError:
    OpenDartReader = None
# [UI 패치] 검색창 팝업 입력칸 바로 아래 딱 붙이기 (좌우 위치 이탈 방지)
st.markdown(
    """
    <style>
    /* 1. 팝업창의 높이만 깔끔하게 제한 (위치 계산은 스트림릿에게 맡겨서 입력칸에 딱 붙게 함) */
    div[data-baseweb="popover"] > div,
    ul[role="listbox"] {
        max-height: 35vh !important;
        overflow-y: auto !important;
    }
    
    /* 2. 고정바 안에서도 스트림릿이 "밑에 공간 넓다!"고 착각하게 만들어 무조건 아래로 열게 유도 */
    .main .block-container {
        padding-bottom: 60vh !important;
    }
    </style>
    """,
    unsafe_allow_html=True
)
# Tab3 히트맵 — 파스텔만 사용 (진한 네이비/브라운 제외)
_TAB3_CMAP_BLUE = LinearSegmentedColormap.from_list(
    "tab3_blue", ["#F8FAFC", "#DBEAFE", "#93C5FD"]
)
_TAB3_CMAP_GREEN = LinearSegmentedColormap.from_list(
    "tab3_green", ["#F8FAFC", "#DCFCE7", "#86EFAC"]
)
_TAB3_CMAP_ORANGE = LinearSegmentedColormap.from_list(
    "tab3_orange", ["#FFF7ED", "#FFEDD5", "#FDBA74"]
)
# OpenDartReader 강제 주입 및 인식 패치
import sys
import subprocess
import importlib
import streamlit as st

try:
    import OpenDartReader
except ImportError:
    try:
        # 1. 강제 설치
        subprocess.check_call([sys.executable, "-m", "pip", "install", "opendartreader"])
        # 2. 파이썬 안경 씌우기 (새로 설치된 부품 즉시 인식!)
        importlib.invalidate_caches()
        # 3. 다시 임포트
        import OpenDartReader
    except Exception:
        OpenDartReader = None
# 페이지 및 Styler 가동 한도 설정 (과부하 방지용)
pd.set_option("styler.render.max_elements", 2000000)
st.set_page_config(page_title="통합 영업 분석 대시보드", layout="wide", initial_sidebar_state="expanded")
# ==========================================
# 0. 로컬 파일 자동 저장을 위한 디렉토리 설정
# ==========================================
CACHE_DIR = "./uploaded_cache"
os.makedirs(CACHE_DIR, exist_ok=True)

# 맥·아이패드 공용 Streamlit Cloud (업무일지·캐시 연동)
# secrets/환경변수로 덮어쓸 수 있음
DASHBOARD_CLOUD_URL = (
    os.environ.get("DASHBOARD_CLOUD_URL")
    or "https://office-g8ryabkapprkpjmfwa5aypw.streamlit.app"
).strip().rstrip("/")

try:
    from drive_autoload import sync_drive_copy_into_cache, sync_cache_to_drive_copy
except Exception:  # pragma: no cover
    sync_drive_copy_into_cache = None  # type: ignore
    sync_cache_to_drive_copy = None  # type: ignore


def _is_local_macos() -> bool:
    """로컬 맥에서 streamlit run 중일 때만 True (Notes 자동화 가능)."""
    return sys.platform == "darwin" and not _is_streamlit_cloud()


def _is_streamlit_cloud() -> bool:
    """Streamlit Community Cloud 여부. 로컬 맥 Desktop은 False."""
    try:
        env = (os.environ.get("STREAMLIT_RUNTIME_ENVIRONMENT") or "").strip().lower()
        if env == "cloud":
            return True
    except Exception:
        pass
    try:
        cwd = os.path.abspath(os.getcwd())
        if cwd.startswith("/mount/src"):
            return True
    except Exception:
        pass
    return False


def _cloud_app_url() -> str:
    """공용 Cloud URL (secrets > env > 기본값)."""
    try:
        secrets = getattr(st, "secrets", None)
        if secrets is not None:
            u = str(secrets.get("dashboard_cloud_url", "") or "").strip()
            if u:
                return u.rstrip("/")
    except Exception:
        pass
    return DASHBOARD_CLOUD_URL


def _render_cloud_sync_banner() -> None:
    """로컬(localhost)에서만: 맥·아이패드는 같은 Cloud URL 쓰도록 안내."""
    url = _cloud_app_url()
    if _is_streamlit_cloud():
        st.sidebar.success("Cloud 공용 · 맥·아이패드 동일 저장소 (여기에 업로드)")
        st.sidebar.caption("새 파일은 이 사이드바에만 올리면 양쪽에서 같이 보입니다.")
        return
    st.sidebar.error(
        "로컬 실행 중입니다. 아이패드와 맞추려면 **Cloud**를 쓰세요."
    )
    try:
        st.sidebar.link_button(
            "▶ Cloud에서 열기 (권장)",
            url,
            use_container_width=True,
            type="primary",
        )
    except TypeError:
        st.sidebar.markdown(f"**[▶ Cloud에서 열기]({url})**")
    st.sidebar.caption(url)
    st.sidebar.caption(
        "로컬 업로드는 이 맥에만 남습니다. Drive 복사는 백업용입니다."
    )

# ==========================================
# 1. 상단 공백 최소화 및 사이드바 무손실 복구 CSS
# ==========================================
# OpenDartReader 임포트 (설치되지 않았을 경우를 대비한 예외 처리)
try:
    import OpenDartReader
except ImportError:
    OpenDartReader = None

# 페이지 및 Styler 가동 한도 설정 (과부하 방지용)
pd.set_option("styler.render.max_elements", 2000000)
st.set_page_config(page_title="통합 영업 분석 대시보드", layout="wide", initial_sidebar_state="expanded")

# ==========================================
# 0. 로컬 파일 자동 저장을 위한 디렉토리 설정
# ==========================================
CACHE_DIR = "./uploaded_cache"
os.makedirs(CACHE_DIR, exist_ok=True)


# ==========================================
# 1. 상단 공백 최소화 및 사이드바 무손실 복구 CSS
# ==========================================
# ==========================================
# 1. 상단 공백 최소화 및 사이드바 무손실 복구 CSS
# ==========================================
# ==========================================
# 1. 상단 공백 최소화 및 사이드바 무손실 복구 CSS
# ==========================================
def inject_custom_css():
    st.markdown(
        """
        <script>
            document.documentElement.lang = 'ko';
            document.documentElement.classList.add('notranslate');
        </script>
        <meta name="google" content="notranslate" />
        <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no" />
        <style>
            /* 기본 불필요 UI 숨김 */
            div[data-baseweb="select"] + div:has(span) { display: none !important; }
            div[data-testid="stMultiSelect"] [data-testid="stWidgetInstructions"] { display: none !important; }
            small[data-testid="stCaptionContainer"] { display: none !important; }
            
            html, body, .stApp {
                background-color: #F8FAFC !important;
                color: #1E293B !important;
                font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
                -webkit-tap-highlight-color: transparent;
            }
            
            /* ★★★ 상단 여백 최소화 (Streamlit 헤더는 그대로) ★★★ */
            section.main .block-container { 
                padding-top: 0.25rem !important; 
                padding-bottom: 0.75rem !important;
                padding-left: 1rem !important;
                padding-right: 1rem !important;
                max-width: 99% !important; 
            }

            section.main [data-testid="stVerticalBlock"]:first-child,
            section.main [data-testid="stVerticalBlockBorderWrapper"]:first-child {
                margin-top: 0 !important;
                padding-top: 0 !important;
            }

            /* 사이드바 — Streamlit 기본 동작 유지, 색상만 보조 */
            [data-testid="stSidebar"] {
                background-color: #F1F5F9 !important;
                border-right: 1px solid #E2E8F0;
                height: 100vh !important;
                position: -webkit-sticky !important;
                position: sticky !important;
                top: 0 !important;
            }
            [data-testid="stSidebar"] > div:first-child {
                overflow-y: auto !important;
                height: 100% !important;
            }
            [data-testid="stSidebar"] .block-container { 
                padding-top: 1.5rem !important; 
                padding-bottom: 15rem !important; /* 하단 짤림 방지 여유 쿠션 */
            }
            [data-testid="stSidebar"] p, [data-testid="stSidebar"] span, 
            [data-testid="stSidebar"] label, [data-testid="stSidebar"] .stMarkdown {
                color: #334155 !important;
            }

            div[data-testid="column"] { align-self: flex-start; }

            .metric-box {
                background: #FFFFFF;
                padding: 16px 20px;
                border-radius: 10px;
                border: 1px solid #E2E8F0;
                box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.05);
                margin-bottom: 8px;
            }
            .metric-label { color: #64748B; font-size: 13px; font-weight: 500; margin-bottom: 6px; }
            .metric-value { color: #0F172A; font-size: 20px; font-weight: 700; }
            
            .sub-header {
                color: #1E3A8A;
                font-size: 17px;
                font-weight: 700;
                margin-top: 0 !important;
                margin-bottom: 12px;
                border-left: 4px solid #1E3A8A;
                padding-left: 10px;
            }

            .dashboard-tabs-host-compact [role="tabpanel"] {
                padding-top: 8px !important;
            }

            .dashboard-tab-title-row {
                min-height: 44px !important;
                align-items: flex-start !important;
            }
            
            div[role="radiogroup"] {
                padding: 10px;
                background: white;
                border-radius: 8px;
                border: 1px solid #E2E8F0;
            }

            /* ===== 상단 통합 고정바 (필터 + 탭) ===== */
            #dashboard-sticky-spacer {
                width: 100% !important;
                height: var(--dashboard-fixed-bar-height, 108px) !important;
                margin: 0 !important;
                padding: 0 !important;
                pointer-events: none !important;
                flex-shrink: 0 !important;
            }

            #dashboard-top-shield {
                pointer-events: none !important;
            }

            section.main {
                overflow: visible !important;
            }

            .dashboard-tabs-host-compact {
                margin-top: 0 !important;
                padding-top: 0 !important;
                height: auto !important;
                max-height: none !important;
                overflow: visible !important;
            }

            /* tablist만 숨김 — tabpanel 직접 자식은 반드시 유지 */
            .dashboard-tabs-host-compact > div:has(> [role="tablist"]),
            .dashboard-tabs-host-compact > div.dashboard-tabs-list-shell:empty {
                display: none !important;
                height: 0 !important;
                min-height: 0 !important;
                margin: 0 !important;
                padding: 0 !important;
                overflow: hidden !important;
            }

            .dashboard-tabs-host-compact > div[role="tabpanel"],
            .dashboard-tabs-host-compact [role="tabpanel"] {
                display: block !important;
                height: auto !important;
                max-height: none !important;
                min-height: 0 !important;
                overflow: visible !important;
            }

            .dashboard-filter-sticky {
                position: fixed !important;
                top: var(--dashboard-bar-top, 2.75rem) !important;
                left: var(--dashboard-bar-left, 0px) !important;
                width: var(--dashboard-bar-width, 100%) !important;
                max-width: var(--dashboard-bar-width, 100%) !important;
                z-index: 990 !important;
                background-color: #FFFFFF !important;
                border: 1px solid #BFDBFE !important;
                border-radius: 0 0 6px 6px !important;
                padding: 1px 6px 0 6px !important;
                box-shadow: 0 2px 8px -2px rgba(15, 23, 42, 0.06) !important;
                box-sizing: border-box !important;
                overflow-x: hidden !important;
                overflow-y: visible !important;
                margin: 0 !important;
            }

            .dashboard-filter-sticky-touch {
                z-index: 999999 !important;
                top: 3.65rem !important;
                border: 2px solid #2563EB !important;
                border-top: 3px solid #2563EB !important;
                border-radius: 0 0 8px 8px !important;
                padding: 8px 10px 0 10px !important;
                box-shadow: 0 8px 14px -4px rgba(37, 99, 235, 0.18) !important;
                margin-top: 0 !important;
            }

            html.dashboard-touch-mode [data-testid="stHeader"],
            html.dashboard-touch-mode [data-testid="stToolbar"],
            html.dashboard-touch-mode [data-testid="stDecoration"] {
                z-index: 1000001 !important;
                position: relative !important;
            }

            html.dashboard-touch-mode [data-testid="stSidebar"],
            html.dashboard-touch-mode [data-testid="stSidebarBackdrop"],
            html.dashboard-touch-mode section[data-testid="stSidebar"] {
                z-index: 1000005 !important;
            }

            html.dashboard-touch-mode [data-testid="stAppViewContainer"],
            html.dashboard-touch-mode section.main,
            html.dashboard-touch-mode section.main .block-container {
                overflow: visible !important;
            }

            /* 고정바 내부 공백 최소화 */
            .dashboard-filter-sticky [data-testid="stVerticalBlock"],
            .dashboard-filter-sticky [data-testid="stHorizontalBlock"] {
                gap: 0.1rem !important;
            }
            .dashboard-filter-sticky [data-testid="column"] {
                padding-top: 0 !important;
                padding-bottom: 0 !important;
            }
            .dashboard-filter-sticky [data-testid="stWidgetLabel"],
            .dashboard-filter-sticky label {
                margin-bottom: 0 !important;
                padding-bottom: 0 !important;
                font-size: 11px !important;
                line-height: 1.15 !important;
                min-height: auto !important;
            }
            .dashboard-filter-sticky [data-testid="stTextInput"],
            .dashboard-filter-sticky [data-testid="stSelectbox"],
            .dashboard-filter-sticky [data-testid="stMultiSelect"] {
                margin-bottom: 0 !important;
                padding-bottom: 0 !important;
            }
            .dashboard-filter-sticky [data-testid="stTextInput"] > div > div,
            .dashboard-filter-sticky [data-baseweb="select"] > div {
                min-height: 1.85rem !important;
            }
            .dashboard-filter-sticky [data-testid="stMarkdownContainer"] {
                margin: 0 !important;
                padding: 0 !important;
            }

            div[data-testid="stVerticalBlockBorderWrapper"].dashboard-filter-sticky {
                padding-top: 2px !important;
                padding-bottom: 0 !important;
            }

            /* 고정바 하단 구분선 — 가늘고 연하게 */
            .dashboard-filter-sticky::after {
                content: '' !important;
                display: block !important;
                height: 1px !important;
                margin: 2px -6px 0 -6px !important;
                background: #CBD5E1 !important;
                box-shadow: none !important;
                border-radius: 0 !important;
            }

            /* 탭 패널 높이 — 대시보드 메인 탭만 */
            .dashboard-tabs-host-compact,
            .dashboard-tabs-host-compact > div,
            .dashboard-tabs-host-compact [role="tabpanel"] {
                overflow: visible !important;
                max-height: none !important;
                height: auto !important;
            }

            div[data-testid="stDataFrame"] > div {
                -webkit-overflow-scrolling: touch;
            }

            /* 탭: 필터 고정바 내부 하단, 세로모드 가로 스크롤 */
            .dashboard-filter-sticky [role="tablist"],
            .dashboard-tabs-in-filter {
                display: flex !important;
                flex-wrap: nowrap !important;
                overflow-x: auto !important;
                overflow-y: visible !important;
                -webkit-overflow-scrolling: touch !important;
                width: 100% !important;
                max-width: 100% !important;
                margin-top: 1px !important;
                margin-bottom: 0 !important;
                padding: 0 0 2px 0 !important;
                border-top: 1px solid #E2E8F0 !important;
                background-color: #FFFFFF !important;
                box-sizing: border-box !important;
                scrollbar-width: none;
                touch-action: pan-x !important;
            }

            .dashboard-filter-sticky [role="tab"] {
                flex: 0 0 auto !important;
                white-space: nowrap !important;
                min-width: -webkit-max-content !important;
                min-width: max-content !important;
                padding: 2px 8px 3px 8px !important;
                font-size: 12px !important;
            }

            .dashboard-filter-sticky [role="tablist"]::-webkit-scrollbar,
            .dashboard-tabs-in-filter::-webkit-scrollbar {
                display: none;
            }

            .dashboard-tabs-host-compact [role="tabpanel"] {
                padding-top: 4px !important;
            }

            .dashboard-tab-panel-head {
                padding-top: 2px !important;
                margin-bottom: 6px !important;
            }

            .dashboard-tabs-host-compact [role="tabpanel"]:not([hidden]) {
                padding-bottom: 48px !important;
            }

            /* Streamlit 기본 스크롤 복구 */
            [data-testid="stAppViewContainer"] {
                overflow-y: auto !important;
                overflow-x: hidden !important;
                -webkit-overflow-scrolling: touch !important;
            }

            .dashboard-filter-sticky [role="tab"] p,
            .dashboard-filter-sticky [role="tab"] span,
            .dashboard-filter-sticky [role="tab"] label {
                white-space: nowrap !important;
            }

            @media (max-width: 1024px) {
                .dashboard-filter-sticky {
                    padding: 2px 6px 0 6px !important;
                }

                .dashboard-filter-sticky [role="tab"] {
                    font-size: 10px !important;
                    padding: 2px 6px 3px 6px !important;
                }
            }

            @supports (top: env(safe-area-inset-top)) {
                .dashboard-filter-sticky {
                    padding-left: max(10px, env(safe-area-inset-left)) !important;
                    padding-right: max(10px, env(safe-area-inset-right)) !important;
                }
            }

            /* ==========================================
               [최종 팝업 방어 패치] 거대 하얀 박스 완벽 차단
               ========================================== */
            /* 1. 팝업창 가로 크기가 화면 전체를 덮는 현상 원천 차단 */
            div[data-baseweb="popover"] {
                max-width: 400px !important; 
            }
            
            /* 2. 스크롤 높이 제한 및 하단 공간(쿠션) 확보 */
            div[data-baseweb="popover"] > div,
            ul[role="listbox"] {
                max-height: 35vh !important;
                overflow-y: auto !important;
            }
            div[data-testid="stAppViewContainer"] {
                padding-bottom: 50vh !important;
            }
            
            /* 3. 쓸데없는 경고 툴팁 아예 숨김 처리 */
            div[data-baseweb="tooltip"] {
                display: none !important;
                opacity: 0 !important;
                pointer-events: none !important;
            }
        </style>
        """,
        unsafe_allow_html=True
    )
# ==========================================
# 2. 날짜 파싱 및 데이터 정규화 유틸리티
# ==========================================
def parse_date_series_robust(series, default_year="2026"):
    if series.empty:
        return pd.Series(pd.NaT, index=series.index)
    s_str = series.astype(str).str.strip()
    parsed = pd.Series(pd.NaT, index=series.index)
    digits = s_str.str.replace(r"\D", "", regex=True)
    cond_8 = (digits.str.len() == 8) & parsed.isna()
    if cond_8.any():
        parsed[cond_8] = pd.to_datetime(
            digits[cond_8], format="%Y%m%d", errors="coerce"
        )
    cond_6 = (digits.str.len() == 6) & parsed.isna()
    if cond_6.any():
        parsed[cond_6] = pd.to_datetime(
            "20" + digits[cond_6], format="%Y%m%d", errors="coerce"
        )
    remaining_mask = (
        parsed.isna() & (s_str != "") & (s_str != "nan") & (s_str != "None")
    )
    if remaining_mask.any():
        rem_series = s_str[remaining_mask]
        parts_df = rem_series.str.split(r"[-/.\s]+", expand=True)
        parts_df = parts_df.apply(lambda col: col.str.strip()).replace("", None)
        num_parts = parts_df.notna().sum(axis=1)
        cond_3 = num_parts >= 3
        if cond_3.any():
            sub = parts_df[cond_3]
            y = sub[0].astype(str).str.strip()
            y = np.where(y.str.len() == 2, "20" + y, y)
            m = sub[1].astype(str).str.strip().str.zfill(2)
            d = sub[2].astype(str).str.strip().str.zfill(2)
            dt_str = y + "-" + m + "-" + d
            parsed.loc[sub.index] = pd.to_datetime(
                dt_str, format="%Y-%m-%d", errors="coerce"
            )
        cond_2 = num_parts == 2
        if cond_2.any():
            sub = parts_df[cond_2]
            m = sub[0].astype(str).str.strip().str.zfill(2)
            d = sub[1].astype(str).str.strip().str.zfill(2)
            dt_str = str(default_year) + "-" + m + "-" + d
            parsed.loc[sub.index] = pd.to_datetime(
                dt_str, format="%Y-%m-%d", errors="coerce"
            )
    valid_range = (parsed >= pd.Timestamp("2000-01-01")) & (
        parsed <= pd.Timestamp("2099-12-31")
    )
    parsed[~valid_range] = pd.NaT
    return parsed
def normalize_items_vectorized(df):
    if "품목명" not in df.columns or df.empty:
        return df
    p_str = df["품목명"].astype(str)
    p_upper = p_str.str.upper().str.replace(" ", "")
    is_bulk = p_upper.str.contains("BULK", na=False) | p_str.str.contains("벌크", na=False)
    
    is_ar = is_bulk & (p_upper.str.contains("AR", na=False) | p_str.str.contains("아르곤|아르", na=False))
    is_co2 = is_bulk & (p_upper.str.contains("CO2", na=False) | p_str.str.contains("탄산", na=False))
    is_o2 = is_bulk & ~is_co2 & (p_upper.str.contains("O2", na=False) | p_str.str.contains("산소", na=False))
    is_n2 = is_bulk & (p_upper.str.contains("N2", na=False) | p_str.str.contains("질소", na=False))
    is_n2_liter = is_n2 & (p_upper.str.contains(r"\bL\b|LITER", regex=True, na=False) | p_str.str.contains("리터", na=False))
    df.loc[is_ar, "품목명"] = "AR (kg, Bulk)"
    df.loc[is_co2, "품목명"] = "CO2 (kg, Bulk)"
    df.loc[is_o2, "품목명"] = "O2 (kg, Bulk)"
    df.loc[is_n2, "품목명"] = "N2 (kg, Bulk)"
    return df
def get_exact_original_price(series):
    s = series[series > 0]
    if s.empty:
        return 0
    return s.mode().iloc[0]
def apply_forward_unit_price(unit_price_df, qty_df, years, all_months):
    """당월 출고 없음(0)이어도 직전 최종 변경 단가를 이월 표시."""
    if unit_price_df.empty:
        return unit_price_df
    filled = unit_price_df.copy()
    chron_cols = [
        f"{yr[2:]}년 {m}"
        for yr in sorted(years)
        for m in all_months
        if f"{yr[2:]}년 {m}" in filled.columns
    ]
    for idx in filled.index:
        last_price = 0.0
        for col in chron_cols:
            price = float(filled.at[idx, col]) if pd.notna(filled.at[idx, col]) else 0.0
            qty = 0.0
            if idx in qty_df.index and col in qty_df.columns:
                qty = float(qty_df.at[idx, col]) if pd.notna(qty_df.at[idx, col]) else 0.0
            if qty > 0 and price > 0:
                last_price = price
            elif last_price > 0 and (qty == 0 or price == 0):
                filled.at[idx, col] = last_price
    return filled


def build_unit_price_change_pivot(unit_price_df, years, all_months, item_names=None):
    """단가표에서 최초 적용·변동 월만 남긴 추이표 (tab3 품목 선택 시 사용).

    item_names가 있으면 해당 품목만. 열=변동 연월, 값=그 시점 단가.
    """
    if unit_price_df is None or unit_price_df.empty:
        return pd.DataFrame()
    src = unit_price_df
    if item_names:
        keep = [i for i in item_names if i in src.index]
        if not keep:
            return pd.DataFrame()
        src = src.loc[keep]
    yrs = sorted([str(y) for y in years if y])
    chron_cols = [
        f"{yr[2:]}년 {m}"
        for yr in yrs
        for m in all_months
        if f"{yr[2:]}년 {m}" in src.columns
    ]
    if not chron_cols:
        return pd.DataFrame()

    change_months = []
    seen_m = set()
    per_item = {}
    for idx in src.index:
        prev = None
        hits = {}
        for col in chron_cols:
            raw = src.at[idx, col]
            if pd.isna(raw):
                continue
            val = float(raw)
            if val <= 0:
                continue
            if prev is None or abs(val - prev) > 0.5:
                hits[col] = val
                if col not in seen_m:
                    seen_m.add(col)
                    change_months.append(col)
                prev = val
        if hits:
            per_item[idx] = hits
    if not per_item:
        return pd.DataFrame()

    chron_rank = {c: i for i, c in enumerate(chron_cols)}
    change_months = sorted(change_months, key=lambda c: chron_rank.get(c, 10**9))
    out = pd.DataFrame(index=list(per_item.keys()), columns=change_months, dtype=float)
    for idx, hits in per_item.items():
        for col, val in hits.items():
            out.at[idx, col] = val
    return out


@st.cache_data
def convert_dfs_to_excel(dfs_dict):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        sheets_written = False
        for sheet_name, (df, use_index) in dfs_dict.items():
            if not df.empty:
                if isinstance(df.columns, pd.MultiIndex):
                    df_to_save = df.copy()
                    df_to_save.columns = [
                        "_".join([str(c) for c in col if c])
                        for col in df_to_save.columns
                    ]
                    df_to_save.to_excel(
                        writer, sheet_name=sheet_name, index=use_index
                    )
                else:
                    df.to_excel(writer, sheet_name=sheet_name, index=use_index)
                sheets_written = True
        
        # [핵심] 저장할 데이터가 전혀 없어 시트가 0개일 때 발생하는 크래시 방지
        if not sheets_written:
            pd.DataFrame({"알림": ["저장할 데이터가 없습니다."]}).to_excel(writer, sheet_name="No Data", index=False)
            
    return output.getvalue()
def _ppt_format_value(val):
    if val is None or (isinstance(val, float) and pd.isna(val)):
        return ""
    if isinstance(val, (pd.Timestamp, datetime.datetime, datetime.date)):
        return val.strftime("%Y-%m-%d")
    if isinstance(val, (int, np.integer)):
        return f"{val:,}"
    if isinstance(val, (float, np.floating)):
        if abs(val) >= 100:
            return f"{val:,.0f}"
        return f"{val:,.2f}"
    text = str(val).replace("\n", " ")
    return text[:40] + ("…" if len(text) > 40 else "")
def _ppt_add_title_bar(slide, title, subtitle=""):
    from pptx.util import Inches, Pt
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.7))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.4), Inches(0.85), Inches(9.2), Inches(0.45))
        stf = sub.text_frame
        stf.text = subtitle
        stf.paragraphs[0].font.size = Pt(11)
def _ppt_add_bullets(slide, lines, top=1.35):
    from pptx.util import Inches, Pt
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(9.0), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.level = 0
def _ppt_add_dataframe(slide, df, max_rows=14, max_cols=8, top=1.25):
    from pptx.util import Inches
    if df is None or df.empty:
        _ppt_add_bullets(slide, ["표시할 데이터가 없습니다."], top=top)
        return
    view = df.copy()
    if isinstance(view.columns, pd.MultiIndex):
        view.columns = ["_".join(str(c) for c in col if c) for col in view.columns]
    view = view.iloc[:max_rows, :max_cols]
    rows = len(view) + 1
    cols = len(view.columns) + (1 if view.index.name or not isinstance(view.index, pd.RangeIndex) else 0)
    has_index = not isinstance(view.index, pd.RangeIndex) or view.index.name
    if not has_index:
        cols = len(view.columns)
    if cols == 0:
        _ppt_add_bullets(slide, ["표시할 데이터가 없습니다."], top=top)
        return
    left, width, height = Inches(0.35), Inches(9.3), Inches(0.28)
    table_shape = slide.shapes.add_table(rows, cols, left, Inches(top), width, height * rows)
    table = table_shape.table
    col_offset = 0
    if has_index:
        table.cell(0, 0).text = str(view.index.name or "")
        col_offset = 1
    for c, col_name in enumerate(view.columns):
        table.cell(0, c + col_offset).text = _ppt_format_value(col_name)
    for r in range(len(view)):
        row = view.iloc[r]
        if has_index:
            table.cell(r + 1, 0).text = _ppt_format_value(view.index[r])
        for c, col_name in enumerate(view.columns):
            table.cell(r + 1, c + col_offset).text = _ppt_format_value(row[col_name])
def _ppt_try_add_chart(slide, fig, top=1.2):
    from pptx.util import Inches
    if fig is None:
        return False
    try:
        img_bytes = fig.to_image(format="png", width=960, height=540, scale=1)
        slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.45), Inches(top), width=Inches(9.1))
        return True
    except Exception:
        return False
def _ppt_new_content_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])
def _ppt_tab_slide(prs, tab_title, section_title, df, chart_fig=None, bullets=None):
    slide = _ppt_new_content_slide(prs)
    _ppt_add_title_bar(slide, tab_title, section_title)
    row_top = 1.25
    if bullets:
        _ppt_add_bullets(slide, bullets, top=row_top)
        row_top = 2.15
    if chart_fig is not None and _ppt_try_add_chart(slide, chart_fig, top=row_top):
        return
    _ppt_add_dataframe(slide, df, top=row_top)
def _ppt_cover_subtitle(latest_update_str, selected_client, selected_staff_tuple):
    today_str = datetime.date.today().strftime("%Y-%m-%d")
    staff_label = ", ".join(selected_staff_tuple) if selected_staff_tuple else "전체"
    return (
        f"생성일: {today_str}\n"
        f"데이터 기준: {latest_update_str}\n"
        f"조회 거래처: {selected_client}\n"
        f"담당자: {staff_label}"
    )
def _ppt_save_presentation(prs):
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
@st.cache_data(show_spinner="PPT 파일 생성 중...")
def convert_dashboard_to_ppt(
    latest_update_str,
    selected_client,
    selected_staff_tuple,
    latest_month_str_total,
    tot_sales_val,
    cur_sales_val,
    mom_rate_total,
    avg_rate_total,
    pivot_m_total,
    pivot_m_client,
    client_item_qty_pivot,
    sales_p,
    qty_p,
    staff_pivot,
    df_detail,
    filtered_debt_df,
    df_base,
    df_tank,
    df_vaporizer,
    df_integrated,
    all_months_tuple,
    years_tuple,
    target_items_tuple,
):
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError("python-pptx 패키지가 필요합니다. pip install python-pptx") from exc
    all_months = list(all_months_tuple)
    years = list(years_tuple)
    target_items = list(target_items_tuple)
    prs = Presentation()
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "통합 영업 분석 대시보드"
    cover.placeholders[1].text = _ppt_cover_subtitle(latest_update_str, selected_client, selected_staff_tuple)
    # Tab1
    _ppt_tab_slide(
        prs, "📌 Tab1 · 영업 종합 요약", "주요 실적 지표", pivot_m_total,
        bullets=[
            f"총 누적 매출(VAT포함): {tot_sales_val:,.0f} 만원",
            f"최근 월 매출 ({latest_month_str_total}): {cur_sales_val:,.0f} 만원",
            f"전월 대비(MoM): {mom_rate_total:+.0f}%",
            f"월평균 대비 증감: {avg_rate_total:+.0f}%",
        ],
    )
    _ppt_tab_slide(
        prs, "📌 Tab1 · 영업 종합 요약", "연도별 월 매출 추이 (차트)", pivot_m_total,
        chart_fig=create_stacked_bar_chart(pivot_m_total, title_text=""),
    )
    _ppt_tab_slide(prs, "📌 Tab1 · 영업 종합 요약", "연도별 월 매출 (표)", pivot_m_total)
    if not df_base.empty and target_items:
        item_pivot = cached_get_item_pivot(
            df_base, target_items[0], "매출액 (만원)", all_months, years
        )
        _ppt_tab_slide(prs, "📌 Tab1 · 영업 종합 요약", f"주요 품목 분석 — {target_items[0]}", item_pivot)
    # Tab2
    tab2_title = f"🏢 Tab2 · 거래처 분석 [{selected_client}]"
    _ppt_tab_slide(
        prs, tab2_title, "연도별 월 매출 추이 (차트)", pivot_m_client,
        chart_fig=create_stacked_bar_chart(pivot_m_client, title_text="") if not pivot_m_client.empty else None,
    )
    _ppt_tab_slide(prs, tab2_title, "연도별 월 매출 (표)", pivot_m_client)
    _ppt_tab_slide(prs, tab2_title, "거래처별 품목 사용량", client_item_qty_pivot)
    # Tab3
    sales_view = sales_p * 1.1 / 10000 if not sales_p.empty else sales_p
    _ppt_tab_slide(prs, "📦 Tab3 · 품목 및 단가 분석", "품목별 매출액 (만원, VAT포함)", sales_view)
    _ppt_tab_slide(prs, "📦 Tab3 · 품목 및 단가 분석", "품목별 출고량", qty_p)
    # Tab4
    _ppt_tab_slide(prs, "👤 Tab4 · 담당자 & 상세내역", "담당자별 월 매출 (만원)", staff_pivot)
    detail_view = df_detail.sort_values(by="매출일_dt", ascending=False).head(20) if not df_detail.empty else df_detail
    _ppt_tab_slide(prs, "👤 Tab4 · 담당자 & 상세내역", "거래 상세 내역 (최신 20건)", detail_view)
    # Tab5
    _ppt_tab_slide(prs, "📌 Tab5 · 채권 관리", "채권 현황", filtered_debt_df)
    # Tab6
    if not df_base.empty:
        map_summary = (
            df_base.groupby("담당자")["거래처"].nunique().reset_index(name="거래처 수").sort_values("거래처 수", ascending=False)
        )
        addr_count = df_base["거래처"].nunique()
        _ppt_tab_slide(prs, "📍 Tab6 · 지도 분포", f"담당자별 거래처 수 (전체 {addr_count}곳)", map_summary)
    else:
        _ppt_tab_slide(prs, "📍 Tab6 · 지도 분포", "거래처 분포", pd.DataFrame())
    # Tab7
    _ppt_tab_slide(prs, "🏭 Tab7 · 설비 재고", "탱크 재고 현황", df_tank)
    _ppt_tab_slide(prs, "🏭 Tab7 · 설비 재고", "기화기 재고 현황", df_vaporizer)
    # Tab8
    _ppt_tab_slide(prs, "🛢️ Tab8 · 통합 탱크 재고", "통합 탱크 재고 현황", df_integrated)
    return _ppt_save_presentation(prs)
# ==========================================
# ★ 네이버 크롤링 + DART API 하이브리드 유틸리티 ★
# ==========================================
def _company_name_candidates(company_name):
    """DART/네이버 검색용 상호 후보 최적화 (무한로딩 폭탄 제거)"""
    raw = str(company_name or "").strip()
    cands = []

    def _add(x):
        x = re.sub(r"\s+", " ", str(x or "").strip())
        if x and len(x) > 1 and x not in cands:
            cands.append(x)

    # 1. 괄호 안의 내용(구상호) 및 법인 표기 완전 제거 -> 핵심 이름만 쏙 뽑기 (예: 라쿨)
    core_name = re.sub(r"\(.*?\)|\[.*?\]|\(주\)|주식회사|㈜|\(유\)|유한회사", "", raw).strip()
    _add(core_name)
    
    # 2. 원래 이름에서 괄호만 제거한 버전 (예: 라쿨 주식회사)
    no_bracket = re.sub(r"\(.*?\)|\[.*?\]", "", raw).strip()
    _add(no_bracket)

    # 핵심 키워드 딱 2개까지만 반환 (후보가 많아지면 수십번 API를 긁어 앱이 뻗는 현상 원천 차단)
    return cands[:2]


def _dart_first_corp_row(by_name):
    """company_by_name 결과(list/DataFrame)에서 첫 법인 dict 추출."""
    rows = _dart_corp_rows(by_name)
    return rows[0] if rows else None


def _dart_corp_rows(by_name):
    """company_by_name 결과 → dict 리스트."""
    if by_name is None:
        return []
    if isinstance(by_name, pd.DataFrame):
        if by_name.empty:
            return []
        return [r.to_dict() for _, r in by_name.iterrows()]
    if isinstance(by_name, list):
        return [r for r in by_name if isinstance(r, dict)]
    if isinstance(by_name, dict) and by_name.get("corp_code"):
        return [by_name]
    return []


def _loc_tokens_from_address(address):
    """주소에서 동명회사 구분용 지역 토큰 (시/군/구·읍/면/동 등)."""
    if not address or str(address).strip() in ("", "등록된 주소 정보가 없습니다."):
        return []
    a = re.sub(r"\(.*?\)|\[.*?\]", " ", str(address))
    tokens = []
    for m in re.finditer(
        r"((?:서울|부산|대구|인천|광주|대전|울산|세종)특별시|세종특별자치시|"
        r"경기(?:도)?|충남|충북|전남|전북|경남|경북|강원(?:도)?|제주(?:도)?)",
        a,
    ):
        tokens.append(m.group(1))
    for m in re.finditer(r"([가-힣]+(?:시|군|구))", a):
        tokens.append(m.group(1))
    for m in re.finditer(r"([가-힣]+(?:읍|면|동|리))", a):
        tokens.append(m.group(1))
    out = []
    for t in tokens:
        t = str(t).strip()
        if t and t not in out:
            out.append(t)
    return out[:8]


def _score_text_vs_loc_tokens(text, tokens):
    if not tokens:
        return 0
    blob = str(text or "")
    return sum(1 for t in tokens if t and t in blob)


# --- 공장등록(팩토리온) 공공데이터 — tab2 기업정보 보강용 ---
FACTORY_KEY_FILE = os.path.join(CACHE_DIR, "factory_api_key.txt")
FACTORY_API_PROD_URL = (
    "https://apis.data.go.kr/B550624/fctryRegistInfo/getFctryPrdctnService_v2"
)
FACTORY_API_LAND_URL = (
    "https://apis.data.go.kr/B550624/fctryRegistLndpclInfo/getFctryLndpclService"
)


def _load_factory_api_key():
    """Secrets → 로컬 파일 → session 순. 사이드바(DART)와 분리."""
    try:
        secrets = getattr(st, "secrets", None)
        if secrets is not None:
            for k in (
                "FACTORY_API_KEY",
                "DATA_GO_KR_SERVICE_KEY",
                "factory_api_key",
            ):
                try:
                    v = secrets.get(k) if hasattr(secrets, "get") else secrets[k]
                except Exception:
                    v = None
                if v and str(v).strip():
                    return str(v).strip()
    except Exception:
        pass
    try:
        if os.path.exists(FACTORY_KEY_FILE):
            with open(FACTORY_KEY_FILE, "r", encoding="utf-8") as f:
                v = f.read().strip()
                if v:
                    return v
    except Exception:
        pass
    return str(st.session_state.get("_factory_api_key") or "").strip()


def _persist_factory_api_key(key):
    key = str(key or "").strip()
    if not key:
        return
    st.session_state["_factory_api_key"] = key
    try:
        with open(FACTORY_KEY_FILE, "w", encoding="utf-8") as f:
            f.write(key)
    except Exception:
        pass


def _fmt_factory_area(val):
    """면적 표시 — 정수만."""
    if val is None or str(val).strip() in ("", "None", "nan"):
        return ""
    try:
        s = str(val).replace(",", "").strip()
        return f"{int(round(float(s))):,}㎡"
    except Exception:
        return str(val).strip()


def _fmt_factory_date(val):
    s = re.sub(r"\D", "", str(val or ""))
    if len(s) >= 8:
        return f"{s[:4]}-{s[4:6]}-{s[6:8]}"
    return str(val or "").strip()


def _factory_extract_items(payload):
    """공공데이터 JSON/XML-ish dict → item list."""
    if payload is None:
        return [], "응답 없음"
    if not isinstance(payload, dict):
        return [], "응답 형식 오류"

    # gateway error envelope
    gw = payload.get("OpenAPI_ServiceResponse") or {}
    if isinstance(gw, dict) and gw.get("cmmMsgHeader"):
        hdr = gw.get("cmmMsgHeader") or {}
        msg = hdr.get("returnAuthMsg") or hdr.get("errMsg") or "API 게이트웨이 오류"
        return [], str(msg)

    resp = payload.get("response") if isinstance(payload.get("response"), dict) else payload
    header = resp.get("header") if isinstance(resp, dict) else None
    body = resp.get("body") if isinstance(resp, dict) else None
    if isinstance(header, dict):
        code = str(header.get("resultCode") or "")
        if code and code not in ("00", "0", "000"):
            return [], str(header.get("resultMsg") or f"오류코드 {code}")
    if not isinstance(body, dict):
        return [], "데이터 없음"
    items = body.get("items")
    if items is None or items == "" or items == {}:
        return [], ""
    if isinstance(items, dict):
        item = items.get("item")
    else:
        item = items
    if item is None:
        return [], ""
    if isinstance(item, list):
        return [x for x in item if isinstance(x, dict)], ""
    if isinstance(item, dict):
        return [item], ""
    return [], "항목 파싱 실패"


def _factory_parse_response_text(text):
    """JSON 우선, 실패 시 XML → dict. 빈 본문은 안내 메시지로 반환."""
    raw = (text or "").strip()
    if not raw:
        return {
            "OpenAPI_ServiceResponse": {
                "cmmMsgHeader": {
                    "errMsg": "EMPTY_BODY",
                    "returnAuthMsg": (
                        "API 응답이 비어 있습니다. "
                        "일반 인증키·활용신청(승인)을 확인하세요."
                    ),
                }
            }
        }
    if raw[0] in "{[":
        try:
            return json.loads(raw)
        except Exception as e:
            return {
                "OpenAPI_ServiceResponse": {
                    "cmmMsgHeader": {
                        "errMsg": "JSON_PARSE_FAIL",
                        "returnAuthMsg": f"JSON 파싱 실패: {e}",
                    }
                }
            }
    if raw[0] == "<" or raw.lstrip().startswith("<?xml"):
        low = raw[:80].lower()
        if "<html" in low or "<!doctype" in low:
            snippet = re.sub(r"\s+", " ", raw)[:160]
            return {
                "OpenAPI_ServiceResponse": {
                    "cmmMsgHeader": {
                        "errMsg": "HTML_BODY",
                        "returnAuthMsg": f"HTML 응답(키/URL 확인): {snippet}",
                    }
                }
            }
        try:
            import xml.etree.ElementTree as ET

            root = ET.fromstring(raw)

            def _xml_to_obj(el):
                children = list(el)
                if not children:
                    return (el.text or "").strip()
                grouped = {}
                for ch in children:
                    grouped.setdefault(ch.tag, []).append(_xml_to_obj(ch))
                out = {}
                for tag, vals in grouped.items():
                    out[tag] = vals[0] if len(vals) == 1 else vals
                return out

            return {root.tag: _xml_to_obj(root)}
        except Exception as e:
            return {
                "OpenAPI_ServiceResponse": {
                    "cmmMsgHeader": {
                        "errMsg": "XML_PARSE_FAIL",
                        "returnAuthMsg": f"XML 파싱 실패: {e}",
                    }
                }
            }
    snippet = re.sub(r"\s+", " ", raw)[:160]
    return {
        "OpenAPI_ServiceResponse": {
            "cmmMsgHeader": {
                "errMsg": "NON_JSON_BODY",
                "returnAuthMsg": f"JSON/XML이 아닌 응답: {snippet}",
            }
        }
    }


def _factory_is_auth_error(err):
    s = str(err or "")
    return any(
        x in s
        for x in (
            "SERVICE_KEY",
            "등록되지 않은 서비스키",
            "UNAUTHORIZED",
            "인증키",
            "INVALID_SERVICE_KEY",
        )
    )


def _factory_http_get(url, api_key, extra_params, timeout=1.5):
    """data.go.kr 호출 (해외 IP 차단 무한로딩 방지 1.5초 컷 패치 완료)"""
    if st.session_state.get("factory_is_blocked"):
        return {"OpenAPI_ServiceResponse": {"cmmMsgHeader": {"errMsg": "BLOCKED", "returnAuthMsg": "EMPTY_BODY (해외IP 차단됨)"}}}

    raw_key = str(api_key or "").strip()
    decoded = urllib.parse.unquote(raw_key)
    encoded = urllib.parse.quote(decoded, safe="")

    base_extra = {}
    for k, v in (extra_params or {}).items():
        if v is not None and str(v).strip() != "":
            base_extra[k] = str(v).strip()

    attempts = [("params", decoded)]
    if encoded != decoded:
        attempts.append(("raw", encoded))
    if "%" in raw_key and raw_key not in (encoded, decoded):
        attempts.append(("raw", raw_key))

    last_payload = None
    last_auth_err = ""
    for mode, key in attempts:
        for resp_type in ("json", "xml"):
            q = {"pageNo": "1", "numOfRows": "20", "type": resp_type}
            q.update(base_extra)
            try:
                if mode == "params":
                    params = dict(q)
                    params["serviceKey"] = key
                    # 🚀 신분증(User-Agent) 장착 완료!
                    r = requests.get(url, params=params, headers={'User-Agent': 'Mozilla/5.0'}, timeout=timeout)
                else:
                    qs = urllib.parse.urlencode(q, doseq=True)
                    # 🚀 신분증(User-Agent) 장착 완료!
                    r = requests.get(f"{url}?serviceKey={key}&{qs}", headers={'User-Agent': 'Mozilla/5.0'}, timeout=timeout)
                
                payload = _factory_parse_response_text(r.text)
                last_payload = payload
                items, err = _factory_extract_items(payload)
                
                if items or not err:
                    return payload
                
                if _factory_is_auth_error(err):
                    last_auth_err = err
                    break
            except Exception as e:
                err_str = str(e).lower()
                if "timeout" in err_str or "connection" in err_str or "max retries" in err_str:
                    st.session_state.factory_is_blocked = True
                    return {"OpenAPI_ServiceResponse": {"cmmMsgHeader": {"errMsg": "TIMEOUT", "returnAuthMsg": "EMPTY_BODY (해외IP 차단/타임아웃)"}}}
                
                last_payload = {
                    "OpenAPI_ServiceResponse": {
                        "cmmMsgHeader": {
                            "errMsg": "REQUEST_FAIL",
                            "returnAuthMsg": str(e),
                        }
                    }
                }

    if last_auth_err:
        return {
            "OpenAPI_ServiceResponse": {
                "cmmMsgHeader": {
                    "errMsg": "SERVICE_KEY_ERROR",
                    "returnAuthMsg": f"등록되지 않은 서비스키 (원문: {last_auth_err})",
                }
            }
        }
    return last_payload or {
        "OpenAPI_ServiceResponse": {
            "cmmMsgHeader": {
                "errMsg": "NO_RESPONSE",
                "returnAuthMsg": "공장등록 API 응답을 받지 못했습니다.",
            }
        }
    }


def _factory_pick_best(items, loc_tokens, company_names):
    """주소 토큰 우선, 없으면 상호 유사도로 1건 선택. 주소 있는데 0점이면 거부."""
    if not items:
        return None, 0
    name_cores = []
    for n in company_names or []:
        core = re.sub(
            r"주식회사|㈜|\(주\)|유한회사|\(유\)|\s+",
            "",
            str(n or ""),
        )
        if core:
            name_cores.append(core)

    best = None
    best_score = -1
    for it in items:
        if not isinstance(it, dict):
            continue
        adres = str(it.get("rnAdres") or it.get("adres") or "")
        loc_score = _score_text_vs_loc_tokens(adres, loc_tokens)
        nm = str(it.get("cmpnyNm") or "")
        nm_core = re.sub(
            r"주식회사|㈜|\(주\)|유한회사|\(유\)|\s+",
            "",
            nm,
        )
        name_hit = 1 if any(c and (c in nm_core or nm_core in c) for c in name_cores) else 0
        score = loc_score * 10 + name_hit
        if score > best_score:
            best_score = score
            best = it
    if loc_tokens and best is not None:
        adres = str(best.get("rnAdres") or best.get("adres") or "")
        if _score_text_vs_loc_tokens(adres, loc_tokens) <= 0:
            return None, 0
    return best, best_score


def _factory_row_to_info(prod_row, land_row=None):
    row = dict(prod_row or {})
    if isinstance(land_row, dict):
        for k, v in land_row.items():
            if v is not None and str(v).strip() != "":
                row[k] = v
    area_land = _fmt_factory_area(row.get("fctryLndpclAr"))
    area_bldg = _fmt_factory_area(row.get("fctryDongBuldAr"))
    return {
        "ok": True,
        "company": str(row.get("cmpnyNm") or "").strip(),
        "ceo": str(row.get("rprsntvNm") or "").strip(),
        "address": str(row.get("rnAdres") or row.get("adres") or "").strip(),
        "industry": str(row.get("indutyNm") or "").strip(),
        "product": str(row.get("mainProductCn") or "").strip(),
        "land_area": area_land,
        "bldg_area": area_bldg,
        "zone": str(row.get("spfcSeCodeNm") or "").strip(),
        "admin": str(row.get("cvplChrgOrgnztNm") or "").strip(),
        "reg_date": _fmt_factory_date(row.get("frstFctryRegistDe")),
        "tel": str(row.get("cmpnyTelno") or "").strip(),
        "fax": str(row.get("cmpnyFxnum") or "").strip(),
        "homepage": str(row.get("hmpadr") or "").strip(),
        "employees": str(row.get("allEmplyCo") or "").strip(),
        "complex": str(row.get("irsttNm") or "").strip(),
        "manage_no": str(row.get("fctryManageNo") or "").strip(),
        "source": "한국산업단지공단 공장등록(팩토리온)",
        "source_url": "https://www.data.go.kr/data/15087611/openapi.do",
        "error": "",
    }


@st.cache_data(ttl=86400, show_spinner=False)
def fetch_factory_registry(company_name, address=None, api_key=None, _cache_v=4):
    """공장등록 생산정보+필지정보 조회. 상호+주소로 동명 오매칭 방지."""
    empty = {
        "ok": False,
        "company": "",
        "ceo": "",
        "address": "",
        "industry": "",
        "product": "",
        "land_area": "",
        "bldg_area": "",
        "zone": "",
        "admin": "",
        "reg_date": "",
        "tel": "",
        "fax": "",
        "homepage": "",
        "employees": "",
        "complex": "",
        "manage_no": "",
        "source": "한국산업단지공단 공장등록(팩토리온)",
        "source_url": "https://www.data.go.kr/data/15087611/openapi.do",
        "error": "",
    }
    key = str(api_key or "").strip()
    if not key:
        empty["error"] = "공장등록 API 키가 없습니다. (공공데이터포털 인증키)"
        return empty

    candidates = _company_name_candidates(company_name)
    loc_tokens = _loc_tokens_from_address(address)
    adres_hint = " ".join(loc_tokens[:2]).strip() if loc_tokens else ""

    def _search(url, names, with_adres=True):
        last_err = ""
        found = []
        for name in names:
            params = {"cmpnyNm": name}
            if with_adres and adres_hint:
                params["adres"] = adres_hint
            payload = _factory_http_get(url, key, params)
            items, err = _factory_extract_items(payload)
            if err and not items:
                last_err = err
                if any(
                    x in str(err)
                    for x in (
                        "SERVICE_KEY",
                        "등록되지 않은 서비스키",
                        "UNAUTHORIZED",
                        "만료",
                        "폐기",
                        "응답이 비어",
                        "EMPTY_BODY",
                        "SERVICE_KEY_ERROR",
                    )
                ):
                    return [], err, True  # fatal
                continue
            if items:
                found = items
                break
        return found, last_err, False

    # 1) 회사명(+지역) → 2) 회사명만 (지역 파라미터 때문에 0건 나는 경우 완화)
    prod_items, last_err, fatal = _search(
        FACTORY_API_PROD_URL, candidates[:5], with_adres=True
    )
    if fatal:
        empty["error"] = last_err
        return empty
    if not prod_items:
        prod_items, last_err2, fatal = _search(
            FACTORY_API_PROD_URL, candidates[:5], with_adres=False
        )
        last_err = last_err2 or last_err
        if fatal:
            empty["error"] = last_err
            return empty

    if not prod_items:
        land_items, last_err3, fatal = _search(
            FACTORY_API_LAND_URL, candidates[:4], with_adres=True
        )
        if fatal:
            empty["error"] = last_err3
            return empty
        if not land_items:
            land_items, last_err4, fatal = _search(
                FACTORY_API_LAND_URL, candidates[:4], with_adres=False
            )
            last_err3 = last_err4 or last_err3
            if fatal:
                empty["error"] = last_err3
                return empty
        best, _ = _factory_pick_best(land_items, loc_tokens, candidates)
        if best:
            return _factory_row_to_info(best, best)
        empty["error"] = (
            last_err
            or last_err3
            or (
                f"『{candidates[0] if candidates else company_name}』 "
                "공장등록 정보가 없습니다. "
                "(미등록·500㎡ 미만·상호 불일치 가능)"
            )
        )
        return empty

    best_prod, score = _factory_pick_best(prod_items, loc_tokens, candidates)
    # 주소 토큰 전부 거절됐을 때: 상호만 정확히 맞는 1건이면 허용
    if not best_prod and prod_items and loc_tokens:
        name_only = []
        for it in prod_items:
            nm = re.sub(
                r"주식회사|㈜|\(주\)|유한회사|\(유\)|\s+",
                "",
                str(it.get("cmpnyNm") or ""),
            )
            if any(
                c
                and (
                    c
                    == nm
                    or c in nm
                    or nm in c
                )
                for c in [
                    re.sub(
                        r"주식회사|㈜|\(주\)|유한회사|\(유\)|\s+",
                        "",
                        str(x or ""),
                    )
                    for x in candidates
                ]
            ):
                name_only.append(it)
        if len(name_only) == 1:
            best_prod = name_only[0]
        elif name_only:
            best_prod, _ = _factory_pick_best(name_only, [], candidates)

    if not best_prod:
        empty["error"] = "주소가 맞는 공장등록 정보가 없습니다. (동명 오매칭 방지)"
        return empty

    land_row = None
    manage_no = str(best_prod.get("fctryManageNo") or "").strip()
    land_params = {"cmpnyNm": best_prod.get("cmpnyNm") or candidates[0]}
    if manage_no:
        land_params["fctryManageNo"] = manage_no
    land_payload = _factory_http_get(FACTORY_API_LAND_URL, key, land_params)
    land_items, _ = _factory_extract_items(land_payload)
    if not land_items and not manage_no:
        land_payload = _factory_http_get(
            FACTORY_API_LAND_URL,
            key,
            {"cmpnyNm": best_prod.get("cmpnyNm") or candidates[0]},
        )
        land_items, _ = _factory_extract_items(land_payload)
    if land_items:
        if manage_no:
            for it in land_items:
                if str(it.get("fctryManageNo") or "").strip() == manage_no:
                    land_row = it
                    break
        if land_row is None:
            land_row, _ = _factory_pick_best(land_items, loc_tokens, candidates)

    return _factory_row_to_info(best_prod, land_row)


def _dart_pick_corp_by_address(dart, name, loc_tokens, max_check=2):
    """동명 법인 중 주소(adres)가 거래처 주소와 맞는 항목 우선 선택. (타임아웃 감지 즉시 영구 차단 패치)"""
    if st.session_state.get("dart_is_blocked"):
        return None, 0
        
    def _check_timeout(e):
        err_str = str(e).lower()
        if "timeout" in err_str or "connection" in err_str or "max retries" in err_str:
            st.session_state.dart_is_blocked = True
            return True
        return False

    rows = []
    try:
        if hasattr(dart, "company_by_name"):
            rows = _dart_corp_rows(dart.company_by_name(name))
    except Exception as e:
        if _check_timeout(e): return None, 0
        rows = []

    # corp_codes 보조 검색
    if len(rows) < 2 and hasattr(dart, "corp_codes") and dart.corp_codes is not None:
        try:
            cc = dart.corp_codes
            hit = cc[cc["corp_name"].astype(str).str.contains(re.escape(str(name)), na=False)]
            for _, r in hit.head(max_check).iterrows():
                d = r.to_dict()
                if d.get("corp_code") and not any(
                    str(x.get("corp_code")) == str(d.get("corp_code")) for x in rows
                ):
                    rows.append(d)
        except Exception:
            pass

    if st.session_state.get("dart_is_blocked"): return None, 0

    if not rows:
        code = None
        try:
            if hasattr(dart, "find_corp_code"):
                code = dart.find_corp_code(name)
        except Exception as e:
            if _check_timeout(e): return None, 0
            code = None
        if code:
            rows = [{"corp_code": str(code), "corp_name": name}]

    best = None
    best_score = -1
    for row in rows[:max_check]:
        if st.session_state.get("dart_is_blocked"): break
        code = str(row.get("corp_code") or "").strip()
        ci = row
        if code:
            try:
                fetched = dart.company(code)
                if isinstance(fetched, dict) and str(fetched.get("status", "000")) in ("000", "0", ""):
                    ci = fetched
            except Exception as e:
                if _check_timeout(e): break
        if not isinstance(ci, dict):
            continue
        if str(ci.get("status", "000")) not in ("000", "0", "", "None"):
            if ci.get("status") is not None and str(ci.get("status")) not in ("000", "0", ""):
                continue
        adres = str(ci.get("adres") or ci.get("address") or "")
        score = _score_text_vs_loc_tokens(adres, loc_tokens)
        if not loc_tokens and best is None:
            score = 0
        if score > best_score:
            best_score = score
            best = ci
            
    if loc_tokens and best_score <= 0:
        return None, 0
    return best, best_score


def _fmt_dart_amount(val):
    try:
        n = int(str(val).replace(",", "").replace(" ", ""))
        return f"{n:,} 원"
    except Exception:
        s = str(val).strip()
        return s if s else "정보 없음"


def _pick_fin_account(fin_state, exact_names, contains_names=None):
    if fin_state is None or fin_state.empty or "account_nm" not in fin_state.columns:
        return None
    acc = fin_state["account_nm"].astype(str)
    for nm in exact_names:
        rows = fin_state[acc == nm]
        if not rows.empty:
            return rows.iloc[0]
    for nm in contains_names or []:
        rows = fin_state[acc.str.contains(nm, na=False)]
        if not rows.empty:
            return rows.iloc[0]
    return None


def _prepare_opendart_corp_codes_cache(max_age_days=14):
    """OpenDartReader 일별 corpCode pickle 재다운로드 완화.
    최근 캐시가 있으면 오늘 파일명으로 재사용해 초기화만 빠르게(내용은 동일)."""
    docs = "docs_cache"
    try:
        os.makedirs(docs, exist_ok=True)
    except Exception:
        return
    today = datetime.date.today().strftime("%Y%m%d")
    today_path = os.path.join(docs, f"opendartreader_corp_codes_{today}.pkl")
    try:
        if os.path.exists(today_path) and os.path.getsize(today_path) > 1000:
            return
    except Exception:
        pass
    newest = None
    newest_mtime = 0.0
    for fp in glob.glob(os.path.join(docs, "opendartreader_corp_codes_*.pkl")):
        try:
            mt = os.path.getmtime(fp)
            if mt > newest_mtime and os.path.getsize(fp) > 1000:
                newest_mtime = mt
                newest = fp
        except Exception:
            continue
    if not newest:
        return
    age_days = (time.time() - newest_mtime) / 86400.0
    if age_days <= float(max_age_days):
        try:
            shutil.copy2(newest, today_path)
        except Exception:
            pass


@st.cache_resource(show_spinner="DART 법인목록 준비 중…")
def _make_opendart_reader(dart_api_key: str):
    """성공한 Reader만 캐시. 실패는 예외로 올려 None을 캐시하지 않음."""
    _prepare_opendart_corp_codes_cache(max_age_days=14)
    return OpenDartReader(str(dart_api_key).strip())


def get_opendart_reader(dart_api_key):
    """OpenDartReader 재사용. 연결 실패 시 None 반환 및 재시도 영구 차단 패치"""
    key = str(dart_api_key or "").strip()
    
    # [핵심] 이미 한 번이라도 차단당한 적이 있으면 아예 시도조차 하지 않고 즉시 0.1초 컷!
    if not key or OpenDartReader is None or st.session_state.get("dart_is_blocked"):
        return None
        
    import socket
    old_timeout = socket.getdefaulttimeout()
    socket.setdefaulttimeout(10.0) # 연결 대기 시간 1.5초로 강제 고정
    try:
        return _make_opendart_reader(key)
    except Exception as e:
        try:
            _make_opendart_reader.clear()
        except Exception:
            pass
        st.session_state["_opendart_last_error"] = str(e)
        # [핵심] 연결에 실패하면 '차단됨' 도장을 찍어서 다음 거래처부터는 절대 헛고생하지 않도록 방어
        st.session_state.dart_is_blocked = True 
        return None
    finally:
        socket.setdefaulttimeout(old_timeout)


@st.cache_data(show_spinner=False, max_entries=100)
def get_company_info_hybrid(company_name, dart_api_key=None, address=None, _cache_v=10):
    """기업정보 조회 (비상장/차단 무한로딩 원천 차단 패치)"""
    import socket
    
    import re
    clean_name = re.sub(r'\(.*?\)|\[.*?\]|주식회사|수식회사|㈜|\(주\)|주\)', '', str(company_name)).strip()
    if not clean_name:
        clean_name = str(company_name)
    candidates = [clean_name]
    loc_tokens = _loc_tokens_from_address(address)
    loc_hint = " ".join(loc_tokens[:3]).strip()

    info = {
        "ceo": "정보 없음", "industry": "정보 없음", "revenue": "정보 없음", "profit": "정보 없음",
        "clean_name": clean_name, "matched_name": "", "corp_code": "", "source": "정보 없음",
        "dart_error": "", "dart_ready": bool(dart_api_key) and OpenDartReader is not None,
        "lookup_hint": f"{clean_name} {loc_hint}".strip(),
    }
    dart_success = False

    if st.session_state.get("dart_is_blocked"):
        info["dart_error"] = "해외 IP 차단으로 DART 조회를 생략합니다."
    elif not dart_api_key:
        info["dart_error"] = "사이드바에 DART API 키가 없습니다."
    elif OpenDartReader is None:
        info["dart_error"] = ""
    else:
        old_timeout = socket.getdefaulttimeout()
        socket.setdefaulttimeout(10.0)
        try:
            dart = get_opendart_reader(dart_api_key)
            if dart is None:
                info["dart_error"] = f"DART 연결 실패: {st.session_state.get('_opendart_last_error')}"
            else:
                corp_code = None
                matched = None
                corp_info = None

                def _is_timeout(e):
                    err_str = str(e).lower()
                    # [핵심 1] Timeout뿐만 아니라 WAF 차단 페이지(HTML)로 인한 JSON 에러도 완벽하게 차단 감지!
                    if "timeout" in err_str or "connection" in err_str or "max retries" in err_str or "jsondecode" in err_str or "expecting value" in err_str:
                        st.session_state.dart_is_blocked = True
                        return True
                    return False

                # 1. 빠르고 가벼운 로컬 캐시에서 먼저 상호를 뒤져서 확실한 법인코드를 확보!
                for name in candidates:
                    if st.session_state.get("dart_is_blocked"): break
                    try:
                        picked, score = _dart_pick_corp_by_address(dart, name, loc_tokens)
                        if picked and picked.get("corp_code"):
                            corp_info = picked
                            corp_code = str(picked.get("corp_code")).strip()
                            matched = str(picked.get("corp_name") or name)
                            break
                    except Exception as e:
                        if _is_timeout(e): break

                # [핵심 2] DART에 없는 회사면 10번씩 삽질(API 호출)하지 않고 즉시 네이버로 0.1초 만에 탈출!
                if not corp_code:
                    info["dart_error"] = "DART 미등록 기업 (빠른 네이버 조회로 전환)"
                else:
                    if isinstance(corp_info, dict):
                        if corp_info.get("ceo_nm"): info["ceo"] = corp_info["ceo_nm"]
                        if corp_info.get("induty_nm"): info["industry"] = corp_info.get("induty_nm") or "정보 없음"
                        info["matched_name"] = corp_info.get("corp_name") or corp_info.get("stock_name") or matched
                        info["corp_code"] = corp_code

                    if not st.session_state.get("dart_is_blocked"):
                        this_year = datetime.date.today().year
                        fin_state = None
                        fin_year_used = None
                        # 재무제표는 최근 3년만, 헛고생 없이 딱 찾은 법인코드로만 찌르기
                        for yr in [this_year - 1, this_year - 2, this_year - 3]:
                            try:
                                fs = dart.finstate(corp_code, yr)
                                if fs is not None and not getattr(fs, "empty", True):
                                    fin_state = fs
                                    fin_year_used = yr
                                    break
                            except Exception as e:
                                if _is_timeout(e): break

                        if fin_state is not None:
                            sales_row = _pick_fin_account(fin_state, exact_names=["매출액", "수익(매출액)", "영업수익"], contains_names=["매출액"])
                            profit_row = _pick_fin_account(fin_state, exact_names=["영업이익", "영업이익(손실)"], contains_names=["영업이익"])
                            if sales_row is not None and "thstrm_amount" in sales_row.index:
                                info["revenue"] = _fmt_dart_amount(sales_row["thstrm_amount"])
                            if profit_row is not None and "thstrm_amount" in profit_row.index:
                                info["profit"] = _fmt_dart_amount(profit_row["thstrm_amount"])
                            
                            if info["revenue"] != "정보 없음" or info["profit"] != "정보 없음":
                                info["source"] = f"금융감독원 DART ({fin_year_used})"
                                dart_success = True
                            else:
                                info["source"] = "금융감독원 DART 기업개요"
                                info["dart_error"] = "재무제표 내 매출액/영업이익을 찾을 수 없습니다."
                        else:
                            info["source"] = "금융감독원 DART 기업개요"
                            info["dart_error"] = "최근 3년 내 구조화된 OpenAPI 재무제표가 없습니다."
                            
        except Exception as e:
            info["dart_error"] = str(e)
        finally:
            socket.setdefaulttimeout(old_timeout)

    # --- 네이버 검색 로직 그대로 유지 ---
    need_naver = (
        not dart_success
        or info["revenue"] == "정보 없음"
        or info["ceo"] == "정보 없음"
        or info["profit"] == "정보 없음"
    )
    
    if need_naver and not st.session_state.get("scraping_is_blocked"):
        naver_hit = False
        dart_matched = bool(info.get("matched_name") or info.get("corp_code"))
        name_limit = 2 if dart_matched else 4
        naver_queries = []
        for name in candidates[:name_limit]:
            if loc_hint:
                naver_queries.append(f"{name} {loc_hint} 기업정보")
            naver_queries.append(f"{name} 기업정보")
        
        ordered_q = []
        _seen_q = set()
        for q in naver_queries:
            if q not in _seen_q:
                _seen_q.add(q)
                ordered_q.append(q)
        
        headers = {
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        }
        
        for q in ordered_q:
            if st.session_state.get("scraping_is_blocked"): break
            try:
                url = "https://search.naver.com/search.naver?query=" + urllib.parse.quote(q)
                response = requests.get(url, headers=headers, timeout=1.5)
                if response.status_code != 200: continue
                soup = BeautifulSoup(response.text, "html.parser")
                texts = list(soup.stripped_strings)
                head_blob = " ".join(texts[:100])
                if loc_tokens and loc_hint and loc_hint.split()[0] in q:
                    if not any(t in head_blob for t in loc_tokens[:3]): continue
                for i, text in enumerate(texts):
                    if text in ["대표자", "대표자명"] and info["ceo"] == "정보 없음" and i + 1 < len(texts):
                        info["ceo"] = texts[i + 1]
                        naver_hit = True
                    elif text in ["업종", "산업(업종)"] and info["industry"] == "정보 없음" and i + 1 < len(texts):
                        info["industry"] = texts[i + 1]
                        naver_hit = True
                    elif text in ["매출액"] and info["revenue"] == "정보 없음" and i + 1 < len(texts):
                        info["revenue"] = texts[i + 1]
                        naver_hit = True
                    elif text in ["영업이익"] and info["profit"] == "정보 없음" and i + 1 < len(texts):
                        info["profit"] = texts[i + 1]
                        naver_hit = True
                
                if info["ceo"] != "정보 없음" and info["revenue"] != "정보 없음" and info["profit"] != "정보 없음":
                    naver_hit = True
                    
                if naver_hit:
                    if not dart_success:
                        info["source"] = "네이버 기업정보 요약" + (f" ({loc_hint})" if loc_hint else "")
                    elif info["revenue"] != "정보 없음" and "DART" not in str(info["source"]):
                        info["source"] = f"{info['source']} + 네이버"
                    break
            except Exception as e:
                if "timeout" in str(e).lower() or "connection" in str(e).lower():
                    st.session_state.scraping_is_blocked = True
                    break

    if not dart_success or info["industry"] == "정보 없음":
        try:
            job = enrich_company_from_job_portals(company_name, address=address, _cache_v=4)
            if job.get("industry") and info["industry"] == "정보 없음":
                info["industry"] = job["industry"]
            if job.get("source") and not dart_success:
                src = str(info.get("source") or "")
                if src in ("", "정보 없음"):
                    info["source"] = job["source"]
                elif job["source"] not in src:
                    info["source"] = f"{src} + {job['source']}"
            info["job_links"] = job.get("links") or {}
        except Exception:
            pass

    return info



@st.cache_data(show_spinner=False, max_entries=50)
def list_dart_audit_reports(corp_key, dart_api_key, years_back=4, _cache_v=1):
    """DART 외부감사(감사보고서) 목록 → [{name, date, rcept_no, url}, ...]"""
    if not corp_key or not dart_api_key or OpenDartReader is None:
        return []
    try:
        dart = get_opendart_reader(dart_api_key)
        if dart is None:
            return []
        end = datetime.date.today()
        start = datetime.date(end.year - int(years_back), 1, 1)
        # F = 외부감사관련
        used_kind_f = True
        df = dart.list(
            str(corp_key),
            start=start.strftime("%Y-%m-%d"),
            end=end.strftime("%Y-%m-%d"),
            kind="F",
            final=True,
        )
        if df is None or getattr(df, "empty", True):
            used_kind_f = False
            df = dart.list(
                str(corp_key),
                start=start.strftime("%Y-%m-%d"),
                end=end.strftime("%Y-%m-%d"),
                final=True,
            )
        if df is None or getattr(df, "empty", True):
            return []
        out = []
        for _, row in df.iterrows():
            nm = str(row.get("report_nm") or "")
            if not used_kind_f and "감사" not in nm:
                continue
            rcp = str(row.get("rcept_no") or row.get("rcp_no") or "").strip()
            if not rcp:
                continue
            out.append(
                {
                    "name": nm or "감사보고서",
                    "date": str(row.get("rcept_dt") or ""),
                    "rcept_no": rcp,
                    "url": f"https://dart.fss.or.kr/dsaf001/main.do?rcpNo={rcp}",
                }
            )
        out.sort(key=lambda x: x.get("date") or "", reverse=True)
        return out[:15]
    except Exception:
        return []


def _dart_fetch_url_html(url, timeout=10):
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
            "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        )
    }
    try:
        for u in (url, str(url).replace("http://", "https://")):
            r = requests.get(u, headers=headers, timeout=timeout)
            if r.status_code == 200 and len(r.text) > 200:
                return r.text
    except Exception:
        return ""
    return ""


def _parse_audit_note_company_overview(html_doc):
    """주석 '1. 회사의 개요' ~ 주요 주주 지분율까지만 구조화 추출."""
    out = {
        "overview_intro": "",
        "hq": "",
        "business": "",
        "ceo_note": "",
        "shareholders": [],
        "overview_ok": False,
    }
    if not html_doc:
        return out
    soup = BeautifulSoup(str(html_doc), "html.parser")
    text = soup.get_text("\n", strip=True)
    m = re.search(
        r"(?:^|\n)\s*(?:1\.\s*)?회사의\s*개요\s*\n([\s\S]+?)"
        r"(?=\n\s*2\.\s*중요한\s*회계|\n\s*2\.\s*회계|\n\s*2\.\s)",
        text,
    )
    if not m:
        m = re.search(r"회사의\s*개요\s*\n([\s\S]{80,1200})", text)
    if not m:
        return out
    block = m.group(1).strip()
    # 지분율 표 이후 회계방침이 섞이면 자름
    block = re.split(r"\n\s*2\.\s*", block, maxsplit=1)[0].strip()
    intro = re.split(r"\n\s*\(1\)", block, maxsplit=1)[0].strip()
    out["overview_intro"] = re.sub(r"\s+", " ", intro)[:400]
    for pat, key in (
        (r"\(1\)\s*본사\s*및\s*공장\s*소재지\s*[:：]\s*(.+)", "hq"),
        (r"\(2\)\s*주요\s*사업\s*내용\s*[:：]\s*(.+)", "business"),
        (r"\(3\)\s*대표이사\s*[:：]\s*(.+)", "ceo_note"),
    ):
        mm = re.search(pat, block)
        if mm:
            out[key] = re.sub(r"\s+", " ", mm.group(1)).strip()[:200]

    # HTML 표 우선
    for table in soup.find_all("table"):
        tplain = table.get_text(" ", strip=True)
        if "지분율" not in tplain and "주주" not in tplain:
            continue
        rows = []
        for tr in table.find_all("tr"):
            cells = [c.get_text(" ", strip=True) for c in tr.find_all(["td", "th"])]
            cells = [c for c in cells if c]
            if cells:
                rows.append(cells)
        data_rows = [r for r in rows if len(r) >= 5 and not re.search(r"주주명|소유주식", r[0])]
        for r in data_rows:
            name = r[0]
            if name in ("합계", "계"):
                continue
            out["shareholders"].append(
                {
                    "주주명": name,
                    "당기주식수": r[1] if len(r) > 1 else "",
                    "당기지분율": (r[2] + "%") if len(r) > 2 and r[2] and "%" not in r[2] else (r[2] if len(r) > 2 else ""),
                    "전기주식수": r[3] if len(r) > 3 else "",
                    "전기지분율": (r[4] + "%") if len(r) > 4 and r[4] and "%" not in r[4] else (r[4] if len(r) > 4 else ""),
                }
            )
        if out["shareholders"]:
            break

    # 텍스트 fallback (표 파싱 실패 시)
    if not out["shareholders"]:
        for sm in re.finditer(
            r"([가-힣A-Za-z0-9·\s]{2,40})\n([\d,]+)\n([\d.]+)\n([\d,]+)\n([\d.]+)",
            block,
        ):
            name = sm.group(1).strip()
            if name in ("합계", "계", "주주명") or "소유" in name:
                continue
            out["shareholders"].append(
                {
                    "주주명": name,
                    "당기주식수": sm.group(2),
                    "당기지분율": sm.group(3) + "%",
                    "전기주식수": sm.group(4),
                    "전기지분율": sm.group(5) + "%",
                }
            )

    out["overview_ok"] = bool(
        out["hq"] or out["business"] or out["ceo_note"] or out["shareholders"]
    )
    return out


def _extract_going_concern_issue(text):
    """독립된 감사인 보고서에서 '계속기업 관련 중요한 불확실성' 중점 이슈만 축약."""
    if not text:
        return "", False
    text_n = re.sub(r"[ \t]+", " ", text)
    m = re.search(
        r"계속기업\s*관련\s*중요한\s*불확실성\s*(.+?)"
        r"(?=재무제표에\s*대한\s*경영진|지배기구는\s*회사|"
        r"재무제표감사에\s*대한\s*감사인|이사의\s*책임|"
        r"감사인의\s*책임|핵심\s*감사|강조\s*사항|\Z)",
        text_n,
        flags=re.S,
    )
    if not m:
        m = re.search(
            r"(계속기업으로서의\s*존속능력에\s*유의적\s*의문을\s*제기.{20,500})",
            text_n,
        )
        if not m:
            return "", False
        block = m.group(1)
    else:
        block = m.group(1)
    block = re.sub(r"\s+", " ", block).strip()
    # 의견 비변형 문구는 이슈 뒤에 짧게 남김
    flag = True
    # 핵심 문장만: 유동성·존속·주석 언급 위주
    sents = re.split(r"(?<=다\.)\s+", block)
    keep = []
    for s in sents:
        s = s.strip()
        if not s:
            continue
        if re.search(
            r"유동부채|유동자산|존속|불확실|주석\s*\d+|자금|차입|손실|자본잠식|의문",
            s,
        ):
            keep.append(s)
        elif re.search(r"의견은\s*이\s*사항으로부터\s*영향을\s*받지", s):
            keep.append(s)
    issue = " ".join(keep) if keep else block
    if len(issue) > 480:
        issue = issue[:479] + "…"
    return issue, flag


@st.cache_data(show_spinner=False, max_entries=30)
def parse_dart_audit_report_summary(rcept_no, dart_api_key, _cache_v=5):
    """감사보고서에서 (1) 독립된 감사인 의견·계속기업 이슈
    (2) 주석 '회사의 개요'~지분율을 함께 추출한다."""
    empty = {
        "opinion": "확인 불가",
        "opinion_note": "본문에서 감사의견을 찾지 못했습니다.",
        "opinion_summary": "",
        "kam": "확인 불가",
        "going_concern": "확인 불가",
        "going_concern_issue": "",
        "going_concern_flag": False,
        "emphasis": "확인 불가",
        "revenue": "",
        "profit": "",
        "chapter": "",
        "overview_intro": "",
        "hq": "",
        "business": "",
        "ceo_note": "",
        "shareholders": [],
        "overview_ok": False,
    }
    if not rcept_no or not dart_api_key or OpenDartReader is None:
        return empty

    def _html_to_text(html_doc):
        return BeautifulSoup(str(html_doc), "html.parser").get_text("\n", strip=True)

    def _fetch_url(url):
        return _dart_fetch_url_html(url, timeout=10)

    try:
        dart = get_opendart_reader(dart_api_key)
        if dart is None:
            return empty
        chapter_title = ""
        text = ""

        # 1) 하위문서 중 '독립된 감사인…' / '감사의견' 챕터 우선
        try:
            subs = dart.sub_docs(str(rcept_no), match="독립된 감사인의 감사보고서")
            if subs is not None and not getattr(subs, "empty", True):
                # 제목 우선순위
                prefer = [
                    "독립된 감사인의 감사보고서",
                    "감사인의 보고서",
                    "감사보고서",
                ]
                chosen = None
                for pref in prefer:
                    hit = subs[subs["title"].astype(str).str.contains(pref, na=False)]
                    if not hit.empty:
                        chosen = hit.iloc[0]
                        break
                if chosen is None:
                    chosen = subs.iloc[0]
                chapter_title = str(chosen.get("title") or "")
                html_body = _fetch_url(str(chosen.get("url") or ""))
                if html_body:
                    text = _html_to_text(html_body)
        except Exception:
            pass

        # 2) 실패 시 전체 document
        if len(text) < 80:
            doc = dart.document(str(rcept_no))
            if doc:
                text = _html_to_text(doc)
                chapter_title = chapter_title or "감사보고서(전체)"

        if len(text) < 40:
            return empty

        text_one = re.sub(r"[ \t]+", " ", text)
        result = dict(empty)
        result["chapter"] = chapter_title or "독립된 감사인의 감사보고서"

        # —— 감사의견 문단 축약 (섹션 제목 이후)
        opinion_block = ""
        m_op_sec = re.search(
            r"(?:^|\n)\s*감사\s*의견\s*\n(.+?)(?=\n\s*감사의견\s*근거|\n\s*감사의견근거|"
            r"\n\s*핵심\s*감사|\n\s*계속기업|\n\s*강조\s*사항|\n\s*재무제표에\s*대한|"
            r"\n\s*이사의\s*책임|\n\s*감사인의\s*책임|\Z)",
            text,
            flags=re.S,
        )
        if m_op_sec:
            opinion_block = re.sub(r"\s+", " ", m_op_sec.group(1)).strip()
        else:
            # 챕터 전체에서 의견 문장만 추출
            m2 = re.search(
                r"(우리의\s*의견으로는.{20,500}?입니다\.)",
                text_one,
            )
            if m2:
                opinion_block = m2.group(1).strip()

        if opinion_block:
            # 너무 길면 앞 280자
            result["opinion_summary"] = (
                opinion_block[:279] + "…" if len(opinion_block) > 280 else opinion_block
            )

        # —— 의견 유형 판별 (문구·관용표현)
        # 주의: '적정의견' 단어가 없어도 "공정하게 표시하고 있습니다" = 대개 적정
        if re.search(r"의견\s*거절|거절\s*의견|의견을\s*거절", text_one):
            result["opinion"] = "의견거절"
            result["opinion_note"] = "자료 부족 등으로 판단을 내릴 수 없음 (위험 높음)"
        elif re.search(r"부적정\s*의견|부적정의견|부적정\s*의견을\s*표명", text_one):
            result["opinion"] = "부적정의견"
            result["opinion_note"] = "재무제표가 왜곡되었을 가능성 (매우 위험)"
        elif re.search(r"한정\s*의견|한정의견|한정\s*의견을\s*표명", text_one):
            result["opinion"] = "한정의견"
            result["opinion_note"] = "일부에만 이슈 — 주의 필요"
        elif re.search(
            r"적정\s*의견|적정의견|의견[:：]\s*적정|"
            r"우리의\s*의견으로는.{0,120}공정하게\s*표시|"
            r"중요성의\s*관점에서\s*공정하게\s*표시하고\s*있습니다",
            text_one,
        ):
            result["opinion"] = "적정의견"
            result["opinion_note"] = "기준에 맞게 작성됨 (가장 양호)"
        elif opinion_block:
            result["opinion"] = "의견 문단 확인"
            result["opinion_note"] = "유형 키워드는 약하나 감사의견 문단은 존재합니다."

        def _snip_after(patterns, max_chars=320):
            for p in patterns:
                m = re.search(p, text_one, flags=re.IGNORECASE)
                if not m:
                    continue
                snip = text_one[m.end() : m.end() + max_chars].strip()
                snip = re.sub(r"\s+", " ", snip)
                snip = re.split(
                    r"(?=핵심감사|강조사항|계속기업|재무제표|주석|독립된|감사인|이사의\s*책임|감사의견\s*근거)",
                    snip,
                    maxsplit=1,
                )[0].strip()
                if len(snip) > 20:
                    return (snip[: max_chars - 1] + "…") if len(snip) > max_chars else snip
            return ""

        kam = _snip_after(
            [
                r"핵심\s*감사\s*사항\s*[:：]?",
                r"Key\s*Audit\s*Matters?\s*[:：]?",
                r"핵심감사사항",
            ]
        )
        result["kam"] = kam or (
            "본문에 ‘핵심감사사항’ 문구가 없거나 추출되지 않았습니다."
            if not re.search(r"핵심\s*감사", text_one)
            else "핵심감사사항 제목은 있으나 상세 문단 추출에 실패했습니다."
        )

        gc_issue, gc_flag = _extract_going_concern_issue(text)
        result["going_concern_flag"] = bool(gc_flag)
        result["going_concern_issue"] = gc_issue
        if gc_flag and gc_issue:
            result["going_concern"] = "⚠️ 계속기업 관련 중요한 불확실성 언급"
        else:
            result["going_concern"] = "계속기업 불확실성 관련 문구를 찾지 못함"

        emp = _snip_after([r"강조\s*사항\s*[:：]?", r"강조사항"])
        if emp:
            result["emphasis"] = emp
        elif re.search(r"강조\s*사항", text_one):
            result["emphasis"] = "강조사항 제목은 있으나 상세 문단 추출에 실패했습니다."
        else:
            result["emphasis"] = "강조사항 문구를 찾지 못함"

        # 손익계산서 하위문서에서 매출/영업이익 보강
        try:
            subs_pl = dart.sub_docs(str(rcept_no), match="손익계산서")
            if subs_pl is not None and not getattr(subs_pl, "empty", True):
                hit = subs_pl[subs_pl["title"].astype(str).str.contains("손익", na=False)]
                row = hit.iloc[0] if not hit.empty else subs_pl.iloc[0]
                pl_html = _fetch_url(str(row.get("url") or ""))
                pl_text = _html_to_text(pl_html) if pl_html else ""
                pl_one = re.sub(r"\s+", " ", pl_text)
                for label, key in (("매출액", "revenue"), ("영업이익", "profit")):
                    if result.get(key):
                        continue
                    m = re.search(
                        rf"{label}\s*[:：]?\s*\(?\s*([-]?[\d,]+)\s*\)?",
                        pl_one,
                    )
                    if m:
                        raw = m.group(1).replace(",", "")
                        try:
                            result[key] = f"{int(raw):,} 원"
                        except Exception:
                            result[key] = m.group(1)
        except Exception:
            pass

        if not result.get("revenue") or not result.get("profit"):
            for label, key in (("매출액", "revenue"), ("영업이익", "profit")):
                if result.get(key):
                    continue
                m = re.search(
                    rf"{label}\s*[:：]?\s*\(?\s*([-]?[\d,]+)\s*\)?",
                    text_one,
                )
                if m:
                    raw = m.group(1).replace(",", "")
                    try:
                        result[key] = f"{int(raw):,} 원"
                    except Exception:
                        result[key] = m.group(1)

        # 주석: 1. 회사의 개요 ~ 지분율
        try:
            note_html = ""
            subs_note = dart.sub_docs(str(rcept_no), match="주석")
            if subs_note is not None and not getattr(subs_note, "empty", True):
                hit_n = subs_note[
                    subs_note["title"].astype(str).str.contains("주석", na=False)
                ]
                nrow = hit_n.iloc[0] if not hit_n.empty else subs_note.iloc[0]
                note_html = _fetch_url(str(nrow.get("url") or ""))
            if note_html:
                ov = _parse_audit_note_company_overview(note_html)
                result.update(ov)
        except Exception:
            pass

        return result
    except Exception as e:
        empty["opinion_note"] = f"본문 파싱 오류: {e}"
        return empty


# 하위 호환 별칭
def try_parse_dart_audit_finance(rcept_no, dart_api_key, _cache_v=1):
    s = parse_dart_audit_report_summary(rcept_no, dart_api_key)
    return {k: s[k] for k in ("revenue", "profit") if s.get(k)}


def _job_portal_query(company_name, address=None):
    """상호 + 주소(시/군/구·읍면동)로 채용포털 검색어 구성."""
    name = str(company_name or "").strip()
    name = re.sub(r"\(.*?\)|\[.*?\]", "", name).strip()
    name = re.sub(r"주식회사|㈜|\(주\)", "", name).strip() or str(company_name or "").strip()
    tokens = _loc_tokens_from_address(address)
    # 시/군/구 + 읍/면 정도만 (너무 길면 검색 노이즈)
    loc_parts = []
    for t in tokens:
        if t.endswith(("시", "군", "구")) and t not in loc_parts:
            loc_parts.append(t)
        elif t.endswith(("읍", "면")) and len(loc_parts) < 2 and t not in loc_parts:
            loc_parts.append(t)
    loc = " ".join(loc_parts[:2])
    q = f"{name} {loc}".strip() if loc else name
    return name, loc, q


def build_job_portal_links(company_name, address=None):
    """사람인·워크넷(고용24)·잡코리아 검색 링크. (상호+주소 병합)"""
    name, loc, q = _job_portal_query(company_name, address)
    qe = urllib.parse.quote(q)
    ne = urllib.parse.quote(name)
    return {
        "query": q,
        "name": name,
        "loc": loc,
        "saramin": f"https://www.saramin.co.kr/zf_user/search?searchType=search&searchword={qe}",
        "saramin_company": (
            f"https://www.saramin.co.kr/zf_user/search/company?searchword={qe}"
        ),
        "worknet": (
            f"https://www.work.go.kr/consltJobCarpa/srch/jobInfoSrch/srchWantedInfo.do"
            f"?searchKeyword={qe}"
        ),
        "work24": f"https://www.work24.go.kr/cm/c/c0100/c0101/companySearch.do?keyword={qe}",
        "jobkorea": f"https://www.jobkorea.co.kr/Search/?stext={qe}",
        "naver_map": f"https://map.naver.com/p/search/{urllib.parse.quote(q)}",
    }


@st.cache_data(show_spinner=False, max_entries=40)
def enrich_company_from_job_portals(company_name, address=None, _cache_v=4):
    """사람인/네이버 기업개요 스크래핑 (해외망 무한로딩 방지 1.5초 컷 패치)"""
    links = build_job_portal_links(company_name, address)
    out = {
        "links": links, "snippets": [], "industry": "", "size": "", 
        "summary": "", "source": "", "note": ""
    }
    
    # [핵심 패치 2] 한 번 차단된 걸 감지하면 사람인도 미련 없이 패스
    if st.session_state.get("scraping_is_blocked"):
        out["note"] = "해외망 차단으로 사람인/네이버 스크래핑 생략"
        return out

    headers = {
        "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
        "Accept-Language": "ko-KR,ko;q=0.9",
    }
    loc_tokens = _loc_tokens_from_address(address)
    saramin_urls = [links["saramin_company"]]
    if links.get("loc"):
        ne = urllib.parse.quote(links.get("name") or company_name)
        saramin_urls.append(f"https://www.saramin.co.kr/zf_user/search/company?searchword={ne}")
    
    try:
        for s_url in saramin_urls:
            if st.session_state.get("scraping_is_blocked"): break
            try:
                r = requests.get(s_url, headers=headers, timeout=1.5) # 1.5초 컷!
                if r.status_code != 200 or len(r.text) <= 500: continue
                soup = BeautifulSoup(r.text, "html.parser")
                texts = [t.strip() for t in soup.stripped_strings if t and len(t.strip()) > 1]
                head_blob = " ".join(texts[:120])
                if loc_tokens and links.get("loc") and links["loc"].split()[0] in urllib.parse.unquote(s_url.split("searchword=")[-1]):
                    if not any(t in head_blob for t in loc_tokens[:3]): continue
                for i, t in enumerate(texts):
                    if t in ("업종", "산업") and i + 1 < len(texts) and not out["industry"]:
                        cand = texts[i + 1]
                        if cand not in ("업종", "산업", "기업형태") and len(cand) < 40: out["industry"] = cand
                    if t in ("사원수", "직원수", "기업규모") and i + 1 < len(texts) and not out["size"]:
                        cand = texts[i + 1]
                        if len(cand) < 30: out["size"] = cand
                name_key = links.get("name") or ""
                for t in texts:
                    if name_key and name_key in t and 20 < len(t) < 180:
                        out["snippets"].append(t)
                        break
                if out["industry"] or out["size"] or out["snippets"]:
                    out["source"] = "사람인 기업검색(자동추출)"
                    out["summary"] = " · ".join(x for x in [out.get("industry"), out.get("size")] if x)
                    break
            except Exception as e:
                if "timeout" in str(e).lower() or "connection" in str(e).lower():
                    st.session_state.scraping_is_blocked = True
                    break
    except Exception as e:
        out["note"] = f"사람인 추출 실패: {e}"

    if not out["summary"] and not st.session_state.get("scraping_is_blocked"):
        try:
            q = links["query"] + " 기업"
            url = f"https://search.naver.com/search.naver?query={urllib.parse.quote(q)}"
            r = requests.get(url, headers=headers, timeout=1.5)
            if r.status_code == 200:
                soup = BeautifulSoup(r.text, "html.parser")
                texts = list(soup.stripped_strings)
                for i, t in enumerate(texts):
                    if t in ("업종", "산업(업종)") and i + 1 < len(texts) and not out["industry"]:
                        out["industry"] = texts[i + 1]
                    if t in ("사원수", "직원수") and i + 1 < len(texts) and not out["size"]:
                        out["size"] = texts[i + 1]
                if out["industry"] or out["size"]:
                    out["source"] = (out["source"] + " + " if out["source"] else "") + "네이버요약"
                    out["summary"] = " · ".join(x for x in [out.get("industry"), out.get("size")] if x)
        except Exception as e:
            if "timeout" in str(e).lower() or "connection" in str(e).lower():
                st.session_state.scraping_is_blocked = True

    if not out["source"]: out["note"] = out.get("note") or "검색 링크로만 제공합니다."
    return out

# ==========================================
# ★ 카카오 지오코딩 API ★
# ==========================================
@st.cache_data(show_spinner=False, max_entries=5000)
def get_lat_lon_kakao(company_name, address, rest_api_key):
    headers = {"Authorization": f"KakaoAK {rest_api_key}"}
    
    clean_addr = ""
    if address and address != "등록된 주소 정보가 없습니다.":
        clean_addr = re.sub(r'\(.*?\)|\[.*?\]', '', str(address))
        clean_addr = clean_addr.split(',')[0].strip()
        
    temp_name = re.sub(r'\(주\)|\(유\)|\(합\)|주식회사|㈜', '', str(company_name))
    match = re.search(r'\((.*?)\)', temp_name)
    
    clean_name = ""
    if match:
        clean_name = match.group(1).strip()
    else:
        clean_name = re.sub(r'^[zZ]', '', temp_name).strip()
    if clean_addr:
        try:
            res = requests.get("https://dapi.kakao.com/v2/local/search/address.json", headers=headers, params={"query": clean_addr}, timeout=3)
            if res.status_code == 200 and res.json().get('documents'):
                return float(res.json()['documents'][0]['y']), float(res.json()['documents'][0]['x'])
        except: pass
        
        try:
            res = requests.get("https://dapi.kakao.com/v2/local/search/keyword.json", headers=headers, params={"query": clean_addr}, timeout=3)
            if res.status_code == 200 and res.json().get('documents'):
                return float(res.json()['documents'][0]['y']), float(res.json()['documents'][0]['x'])
        except: pass
        
    if clean_name:
        try:
            res = requests.get("https://dapi.kakao.com/v2/local/search/keyword.json", headers=headers, params={"query": clean_name}, timeout=3)
            if res.status_code == 200 and res.json().get('documents'):
                return float(res.json()['documents'][0]['y']), float(res.json()['documents'][0]['x'])
        except: pass
        
    if company_name:
        try:
            res = requests.get("https://dapi.kakao.com/v2/local/search/keyword.json", headers=headers, params={"query": str(company_name)}, timeout=3)
            if res.status_code == 200 and res.json().get('documents'):
                return float(res.json()['documents'][0]['y']), float(res.json()['documents'][0]['x'])
        except: pass
    return None, None


# Tab6 전용: 프로세스 재시작 후에도 좌표 API를 다시 치지 않도록 디스크에 저장
KAKAO_GEOCODE_CACHE_PATH = os.path.join(CACHE_DIR, "kakao_geocode_cache.json")


def _kakao_geocode_disk_key(company_name, address):
    return f"{str(company_name).strip()}||{str(address or '').strip()}"


def _load_kakao_geocode_disk():
    path = KAKAO_GEOCODE_CACHE_PATH
    if not os.path.exists(path):
        return {}
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, dict) else {}
    except Exception:
        return {}


def _save_kakao_geocode_disk(cache_dict):
    try:
        os.makedirs(CACHE_DIR, exist_ok=True)
        with open(KAKAO_GEOCODE_CACHE_PATH, "w", encoding="utf-8") as f:
            json.dump(cache_dict, f, ensure_ascii=False)
    except Exception:
        pass


def get_lat_lon_kakao_disk(company_name, address, rest_api_key, disk_cache, dirty_flag):
    """디스크 캐시 우선 → 없으면 Kakao API(@st.cache_data) 호출."""
    key = _kakao_geocode_disk_key(company_name, address)
    if key in disk_cache:
        hit = disk_cache[key]
        if not hit or hit.get("lat") is None or hit.get("lon") is None:
            return None, None
        try:
            return float(hit["lat"]), float(hit["lon"])
        except Exception:
            return None, None
    lat, lon = get_lat_lon_kakao(company_name, address, rest_api_key)
    disk_cache[key] = {"lat": lat, "lon": lon}
    dirty_flag[0] = True
    return lat, lon


KAKAO_REST_API_KEY = "21a8c4d7312051598c2e05dba0b9c0c7"
@st.cache_data(show_spinner=False, max_entries=2000)
def kakao_place_search(query, rest_api_key=None, size=15):
    """주소·상호·명칭 통합검색 → 후보 목록.
    반환: [{lat, lon, label, place_name, address}, ...]
    """
    key = rest_api_key or KAKAO_REST_API_KEY
    q = str(query or "").strip()
    if not q:
        return []
    headers = {"Authorization": f"KakaoAK {key}"}
    cands = []
    seen = set()

    def _add(lat, lon, label, place_name="", address=""):
        try:
            lat_f, lon_f = float(lat), float(lon)
        except Exception:
            return
        sig = (round(lat_f, 5), round(lon_f, 5), str(label).strip())
        if sig in seen:
            return
        seen.add(sig)
        cands.append(
            {
                "lat": lat_f,
                "lon": lon_f,
                "label": str(label).strip() or q,
                "place_name": str(place_name or "").strip(),
                "address": str(address or "").strip(),
            }
        )

    # 1) 키워드(장소명·상호·지점)
    try:
        res = requests.get(
            "https://dapi.kakao.com/v2/local/search/keyword.json",
            headers=headers,
            params={"query": q, "size": min(int(size or 15), 15)},
            timeout=5,
        )
        if res.status_code == 200:
            for doc in res.json().get("documents") or []:
                place = str(doc.get("place_name") or "").strip()
                road = str(doc.get("road_address_name") or "").strip()
                jibun = str(doc.get("address_name") or "").strip()
                addr = road or jibun
                label = f"{place} · {addr}" if place and addr else (place or addr or q)
                _add(doc.get("y"), doc.get("x"), label, place, addr)
    except Exception:
        pass
    # 2) 주소 검색 (키워드에 없을 때 보강)
    try:
        res = requests.get(
            "https://dapi.kakao.com/v2/local/search/address.json",
            headers=headers,
            params={"query": q},
            timeout=5,
        )
        if res.status_code == 200:
            for doc in res.json().get("documents") or []:
                addr = doc.get("address_name") or q
                if doc.get("road_address"):
                    addr = doc["road_address"].get("address_name") or addr
                elif doc.get("address"):
                    addr = doc["address"].get("address_name") or addr
                _add(doc.get("y"), doc.get("x"), addr, "", addr)
    except Exception:
        pass
    return cands
def geocode_address_kakao(address, rest_api_key=None):
    """주소·상호·명칭 통합검색 → (lat, lon, 표기주소). 후보 1순위 사용."""
    cands = kakao_place_search(address, rest_api_key=rest_api_key, size=10)
    if not cands:
        return None, None, ""
    c0 = cands[0]
    return c0["lat"], c0["lon"], c0["label"]
def haversine_km(lat1, lon1, lat2, lon2):
    """두 좌표 직선거리(km) — Haversine."""
    r = 6371.0
    p1, p2 = np.radians(lat1), np.radians(lat2)
    dphi = np.radians(lat2 - lat1)
    dlmb = np.radians(lon2 - lon1)
    a = np.sin(dphi / 2) ** 2 + np.cos(p1) * np.cos(p2) * np.sin(dlmb / 2) ** 2
    return float(2 * r * np.arcsin(np.sqrt(a)))
# 한국도로공사 차종: 20톤 벌크로리 ≈ 5종(특수). 길찾기 API 통행료는 보통 1종(승용).
# 실무 근사 배수(1종 대비): 3종 1.5 / 4종 2.0 / 5종 2.5
TOLL_CLASS_5_MULT = 2.5
def _toll_for_bulk20t(toll_class1_won):
    """1종 통행료 → 20톤 벌크(5종) 편도 통행료(원, 10원 단위)."""
    raw = float(toll_class1_won or 0) * TOLL_CLASS_5_MULT
    return float(int(round(raw / 10.0) * 10))
@st.cache_data(show_spinner=False, max_entries=2000)
def kakao_route_from_coords(o_lat, o_lon, d_lat, d_lon, o_label, d_label, rest_api_key=None, _route_v=4):
    """좌표 기준 자동차 거리(km) + 통행료(원, 편도)."""
    _ = _route_v
    key = rest_api_key or KAKAO_REST_API_KEY
    o_label = str(o_label or "출발")
    d_label = str(d_label or "도착")
    try:
        o_lat, o_lon = float(o_lat), float(o_lon)
        d_lat, d_lon = float(d_lat), float(d_lon)
    except Exception:
        return {
            "ok": False, "km": 0.0, "toll": 0.0, "toll_class1": 0.0, "method": "",
            "origin_label": o_label, "dest_label": d_label,
            "message": "좌표가 올바르지 않습니다.",
        }
    try:
        headers = {"Authorization": f"KakaoAK {key}", "Content-Type": "application/json"}
        res = requests.get(
            "https://apis-navi.kakaomobility.com/v1/directions",
            headers=headers,
            params={
                "origin": f"{o_lon},{o_lat}",
                "destination": f"{d_lon},{d_lat}",
                "priority": "RECOMMEND",
                "car_hipass": "true",
            },
            timeout=8,
        )
        if res.status_code == 200:
            routes = res.json().get("routes") or []
            if routes:
                summary = routes[0].get("summary") or {}
                dist_m = float(summary.get("distance") or 0)
                fare = summary.get("fare") or {}
                toll_c1 = float(fare.get("toll") or 0)
                toll_5 = _toll_for_bulk20t(toll_c1)
                if dist_m > 0:
                    msg = (
                        f"카카오 길찾기 {dist_m/1000.0:.1f}km  ·  "
                        f"통행료 1종 {toll_c1:,.0f}원 → 20톤벌크(5종×{TOLL_CLASS_5_MULT:g}) {toll_5:,.0f}원/편도"
                    )
                    return {
                        "ok": True,
                        "km": dist_m / 1000.0,
                        "toll": toll_5,
                        "toll_class1": toll_c1,
                        "method": "kakao_mobility",
                        "origin_label": o_label,
                        "dest_label": d_label,
                        "message": msg,
                    }
    except Exception:
        pass
    straight = haversine_km(o_lat, o_lon, d_lat, d_lon)
    road_km = straight * 1.3
    return {
        "ok": True,
        "km": road_km,
        "toll": 0.0,
        "toll_class1": 0.0,
        "method": "haversine_x1.3",
        "origin_label": o_label,
        "dest_label": d_label,
        "message": (
            f"길찾기 API 미사용/실패 → 직선 {straight:.1f}km × 1.3 = 도로근사 {road_km:.1f}km  ·  "
            f"통행료는 수동 입력"
        ),
    }
@st.cache_data(show_spinner=False, max_entries=2000)
def kakao_route_distance_km(origin_addr, dest_addr, rest_api_key=None, _route_v=4):
    """출발·도착(주소/상호/명칭) → 자동차 거리(km) + 통행료(원, 편도)."""
    _ = _route_v
    key = rest_api_key or KAKAO_REST_API_KEY
    o_lat, o_lon, o_label = geocode_address_kakao(origin_addr, key)
    d_lat, d_lon, d_label = geocode_address_kakao(dest_addr, key)
    if o_lat is None or d_lat is None:
        return {
            "ok": False,
            "km": 0.0,
            "toll": 0.0,
            "toll_class1": 0.0,
            "method": "",
            "origin_label": o_label or str(origin_addr),
            "dest_label": d_label or str(dest_addr),
            "message": "출발/도착을 찾지 못했습니다. 주소·상호·명칭을 확인해 주세요.",
        }
    return kakao_route_from_coords(
        o_lat, o_lon, d_lat, d_lon, o_label, d_label, rest_api_key=key, _route_v=_route_v
    )
def _parse_ton_number(token):
    """톤수 숫자 토큰 → float. '4.9'/'4,9' 지원."""
    t = str(token or "").strip().replace(",", ".")
    t = re.sub(r"[^0-9.]", "", t)
    if not t or t == ".":
        return 0.0
    try:
        return float(t)
    except Exception:
        return 0.0
# 액화가스 내용적(L) → kg 환산 밀도 (대략값, 액화 기준 kg/L)
GAS_DENSITY_KG_PER_L = {
    "질소": 0.808,    # LN2
    "알곤": 1.395,    # LAr
    "산소": 1.141,    # LOX
    "탄산": 1.101,    # LCO2
    "수소": 0.0708,   # LH2
    "헬륨": 0.125,    # LHe
}
# 액화 kg → 기체 Nm³ 환산 (0℃, 1atm 대략값)
GAS_NM3_PER_KG = {
    "질소": 0.7996,
    "알곤": 0.5605,
    "산소": 0.6998,
    "탄산": 0.5090,
    "수소": 11.126,
    "헬륨": 5.596,
}
GAS_OPTIONS = list(GAS_DENSITY_KG_PER_L.keys())
def liters_to_tank_kg(liters, gas_name):
    """내용적(L) × 가스밀도 → 탱크 용량(kg)."""
    dens = float(GAS_DENSITY_KG_PER_L.get(str(gas_name), 0.808) or 0.808)
    return float(liters or 0) * dens
def tank_kg_to_nm3(tank_kg, gas_name):
    """탱크 용량(kg) → 기체 환산(Nm³)."""
    factor = float(GAS_NM3_PER_KG.get(str(gas_name), 0.0) or 0.0)
    return float(tank_kg or 0) * factor
def tank_liters_to_nm3(liters, gas_name):
    """내용적(L) → 기체 환산(Nm³) = L × 밀도 × Nm³/kg."""
    return tank_kg_to_nm3(liters_to_tank_kg(liters, gas_name), gas_name)
def nm3_per_h_to_kg_per_h(nm3_h, gas_name):
    """시간당 루베(Nm³/h) → kg/h."""
    factor = float(GAS_NM3_PER_KG.get(str(gas_name), 0.0) or 0.0)
    if factor <= 0:
        return 0.0
    return float(nm3_h or 0) / factor
def kg_per_h_to_nm3_per_h(kg_h, gas_name):
    """시간당 kg/h → 루베(Nm³/h)."""
    factor = float(GAS_NM3_PER_KG.get(str(gas_name), 0.0) or 0.0)
    return float(kg_h or 0) * factor
def gas_conversion_rows(liters, tank_kg, selected_gas):
    """선택 가스·전 가스 기체환산 표용 rows."""
    rows = []
    for g in GAS_OPTIONS:
        dens = GAS_DENSITY_KG_PER_L[g]
        nm3_per_kg = GAS_NM3_PER_KG[g]
        # 동일 내용적(L) 기준 비교 + 현재 탱크 kg 기준(선택 가스)도 표기
        kg_from_l = float(liters or 0) * dens
        nm3_from_l = kg_from_l * nm3_per_kg
        nm3_from_kg = float(tank_kg or 0) * nm3_per_kg if g == selected_gas else None
        rows.append(
            {
                "가스": g + (" ★" if g == selected_gas else ""),
                "밀도(kg/L)": dens,
                "Nm³/kg": nm3_per_kg,
                "내용적 기준 kg": round(kg_from_l, 1),
                "내용적 기준 Nm³": round(nm3_from_l, 1),
                "현재탱크 Nm³": round(nm3_from_kg, 1) if nm3_from_kg is not None else "—",
            }
        )
    return rows
def compute_tank_usage_cycle(
    tank_kg, hourly_usage_kg, operating_hours, fill_ratio=0.8, days_per_month=30.0
):
    """시간당 사용량·가동시간 → 탱크 사용주기.
    충전기준 = 탱크kg × 80% 소모 시 충전.
    사용주기(일) = 충전기준 ÷ (시간당사용량 × 일가동시간)
    월 사용량 = 일사용량 × 월가동일수
    """
    tank = float(tank_kg or 0)
    hourly = float(hourly_usage_kg or 0)
    hours = float(operating_hours or 0)
    days_m = max(0.0, float(days_per_month or 0))
    ratio = float(fill_ratio or 0) or 0.8
    charge_kg = tank * ratio
    daily_kg = hourly * hours
    monthly_kg = daily_kg * days_m
    if hourly <= 0 or charge_kg <= 0:
        return {
            "ok": False,
            "charge_kg": charge_kg,
            "daily_kg": daily_kg,
            "monthly_kg": monthly_kg,
            "days_per_month": days_m,
            "cycle_days": 0.0,
            "cycle_hours": 0.0,
            "fills_per_month": 0.0,
            "message": "시간당 사용량과 탱크 용량을 입력하세요.",
        }
    cycle_hours = charge_kg / hourly  # 가동시간 누적 기준
    cycle_days = (charge_kg / daily_kg) if daily_kg > 0 else 0.0
    fills_per_month = (days_m / cycle_days) if cycle_days > 0 else 0.0
    return {
        "ok": True,
        "charge_kg": charge_kg,
        "daily_kg": daily_kg,
        "monthly_kg": monthly_kg,
        "days_per_month": days_m,
        "cycle_days": cycle_days,
        "cycle_hours": cycle_hours,
        "fills_per_month": fills_per_month,
        "message": (
            f"충전기준 {charge_kg:,.0f}kg(탱크×{ratio*100:.0f}%) ÷ "
            f"일사용 {daily_kg:,.0f}kg = 사용주기 {cycle_days:,.0f}일"
        ),
    }
def parse_tank_capacity_kg(tank_value):
    """TANK 용량 → kg.
    - 숫자만 입력: kg로 그대로 사용 (예: 4900)
    - 구형식(4.9t, 35T*2 등): 톤→kg 환산 (하위호환)
    """
    if isinstance(tank_value, (int, float)) and not isinstance(tank_value, bool):
        return float(tank_value) if float(tank_value) > 0 else 0.0
    s = str(tank_value or "").strip()
    if not s:
        return 0.0
    for ch in ("，", "·", "．", "･", "。"):
        s = s.replace(ch, ".")
    s_norm = s.replace(",", ".")
    # 숫자만 → kg
    if re.fullmatch(r"\d+(?:\.\d+)?", s_norm):
        return float(s_norm)
    # 구형식: 톤 표기 → kg
    unit = r"(?:TON|톤|T|t)"
    pairs = re.findall(
        rf"(?<![0-9.])(\d+(?:[.,]\d+)?)\s*{unit}\s*[*xX×]\s*(\d+(?:[.,]\d+)?)",
        s,
    )
    if pairs:
        tons = sum(_parse_ton_number(a) * _parse_ton_number(b) for a, b in pairs)
        return float(tons * 1000.0)
    singles = re.findall(rf"(?<![0-9.])(\d+(?:[.,]\d+)?)\s*{unit}\b", s)
    if singles:
        tons = sum(_parse_ton_number(a) for a in singles)
        return float(tons * 1000.0)
    return 0.0
def compute_roundtrips_from_tank(monthly_usage_kg, tank_value, fill_ratio=0.8):
    """왕복횟수 = ceil(월평균공급량 ÷ (탱크용량kg × 0.8)).
    탱크(kg)의 80%가 소모되었을 때 1회 충전(왕복).
    """
    import math
    tank_kg = parse_tank_capacity_kg(tank_value)
    ratio = float(fill_ratio or 0) or 0.8
    charge_kg = tank_kg * ratio
    monthly = float(monthly_usage_kg or 0)
    if charge_kg <= 0:
        return {
            "ok": False,
            "roundtrips": 0.0,
            "roundtrips_exact": 0.0,
            "tank_tons": 0.0,
            "tank_kg": tank_kg,
            "usable_kg": 0.0,
            "fill_ratio": ratio,
            "message": "TANK 용량(kg)을 숫자로 입력하세요. 예: 4900",
        }
    exact = (monthly / charge_kg) if monthly > 0 else 0.0
    trips = int(math.ceil(exact - 1e-12)) if exact > 0 else 0
    return {
        "ok": True,
        "roundtrips": float(trips),
        "roundtrips_exact": exact,
        "tank_tons": tank_kg / 1000.0,
        "tank_kg": tank_kg,
        "usable_kg": charge_kg,
        "fill_ratio": ratio,
        "message": (
            f"탱크 {tank_kg:,.0f}kg의 {ratio*100:.0f}% 소모 시 충전 "
            f"= {charge_kg:,.0f}kg/회 → 월공급 {monthly:,.0f}kg ÷ 충전량 = {trips}회 "
            f"(정확 {exact:.2f}, 올림)"
        ),
    }
def profit_int_comma_input(label, *, key, value=0, help=None):
    """정수 입력 + 천단위 콤마 표시 (소수점 없음). Tab9 투자비 전용."""
    txt_key = f"{key}__comma"
    try:
        init_n = int(round(float(value or 0)))
    except Exception:
        init_n = 0
    if key not in st.session_state:
        st.session_state[key] = init_n
    if txt_key not in st.session_state:
        st.session_state[txt_key] = f"{int(st.session_state[key]):,}"

    def _sync():
        raw = str(st.session_state.get(txt_key, "") or "")
        digits = re.sub(r"[^\d]", "", raw)
        n = int(digits) if digits else 0
        st.session_state[key] = n
        st.session_state[txt_key] = f"{n:,}"

    st.text_input(label, key=txt_key, on_change=_sync, help=help)
    # 화면에 남아 있는 미커밋 값도 숫자로 해석
    raw_now = str(st.session_state.get(txt_key, "") or "")
    digits_now = re.sub(r"[^\d]", "", raw_now)
    if digits_now:
        st.session_state[key] = int(digits_now)
    return float(st.session_state.get(key, 0) or 0)
def compute_logistics_unit_cost(
    distance_km,
    fuel_price_per_l,
    fuel_efficiency_km_per_l,
    toll_fee=0.0,
    round_trips=1.0,
    load_kg=1.0,
):
    """표준 탱크로리 물류비(원/kg) 산식.
    연료비(편도) = 거리(km) × 유류비(원/L) ÷ 연비(km/L)
    편도비용 = 연료비 + 통행료
    왕복총비용 = 편도비용 × 왕복횟수
    물류비(원/kg) = 왕복총비용 ÷ 공급량(kg)
    (엑셀 P18과 동일: (KM×유류비/연비+통행료)×왕복÷KG)
    """
    eff = float(fuel_efficiency_km_per_l or 0) or 1.0
    kg = float(load_kg or 0) or 1.0
    fuel_one_way = float(distance_km or 0) * float(fuel_price_per_l or 0) / eff
    one_way = fuel_one_way + float(toll_fee or 0)
    total = one_way * float(round_trips or 0)
    per_kg = total / kg
    return {
        "fuel_one_way": fuel_one_way,
        "one_way_cost": one_way,
        "round_trip_cost": total,
        "per_kg": per_kg,
    }
DIESEL_PRICE_FILE = os.path.join(CACHE_DIR, "diesel_market_price.json")
def fetch_diesel_market_price():
    """전국 주유소 경유 평균가(원/L).
    소스: 네이버 금융 OIL_LO(경유) = 한국석유공사 Opinet 전국평균.
    20톤 벌크로리(경유) 물류비 산정용 시장 정가 기준.
    """
    url = "https://finance.naver.com/marketindex/oilDetail.naver?marketindexCd=OIL_LO"
    headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36"}
    try:
        res = requests.get(url, headers=headers, timeout=10)
        res.encoding = res.apparent_encoding or "euc-kr"
        soup = BeautifulSoup(res.text, "html.parser")
        nt = soup.select_one("p.no_today")
        price = None
        if nt:
            for b in nt.select(".blind"):
                b.decompose()
            spaced = re.sub(r"\s+", "", nt.get_text(" ", strip=True))
            num = re.search(r"([0-9]+(?:,[0-9]+)*(?:\.[0-9]+)?)", spaced)
            if num:
                price = float(num.group(1).replace(",", ""))
        if price is None or price < 500 or price > 5000:
            return {"ok": False, "price": 0.0, "asof": "", "source": "", "message": "경유 시세를 파싱하지 못했습니다."}
        asof = datetime.date.today().strftime("%Y-%m-%d")
        dm = re.search(r"(20\d{2}\.\d{2}\.\d{2})", res.text)
        if dm:
            asof = dm.group(1).replace(".", "-")
        return {
            "ok": True,
            "price": round(price, 2),
            "asof": asof,
            "source": "네이버금융·오피넷 전국 경유 평균",
            "message": f"경유 {price:,.2f} 원/L ({asof})",
        }
    except Exception as e:
        return {"ok": False, "price": 0.0, "asof": "", "source": "", "message": f"경유 시세 조회 실패: {e}"}
def load_diesel_price_cache():
    if os.path.exists(DIESEL_PRICE_FILE):
        try:
            with open(DIESEL_PRICE_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
            if isinstance(data, dict):
                return data
        except Exception:
            pass
    return {}
def save_diesel_price_cache(data):
    try:
        with open(DIESEL_PRICE_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception:
        pass
def get_diesel_price_monthly(force_refresh=False):
    """매월 시장가 자동 갱신(캐시). force_refresh면 즉시 재조회."""
    month = datetime.date.today().strftime("%Y-%m")
    cached = load_diesel_price_cache()
    if (
        not force_refresh
        and cached.get("ok")
        and cached.get("month") == month
        and float(cached.get("price") or 0) > 0
    ):
        return cached
    fetched = fetch_diesel_market_price()
    if fetched.get("ok"):
        fetched["month"] = month
        fetched["vehicle"] = "20톤 벌크로리 · 경유"
        save_diesel_price_cache(fetched)
        return fetched
    # 조회 실패 시 이전 캐시라도 사용
    if cached.get("ok") and float(cached.get("price") or 0) > 0:
        cached = dict(cached)
        cached["message"] = f"최신 조회 실패 → 캐시 사용 ({cached.get('asof','')})"
        return cached
    fetched["month"] = month
    return fetched
# ==========================================
# ★ 메모 생성 AppleScript ★
# ==========================================
def _client_note_inventory_html(client_name, df_integrated) -> str:
    inventory_html = ""
    if (
        df_integrated is not None
        and not df_integrated.empty
        and "거래처(사용처)/보관장소" in df_integrated.columns
    ):
        client_inv = df_integrated[
            df_integrated["거래처(사용처)/보관장소"]
            .astype(str)
            .str.contains(client_name, regex=False, na=False)
        ]
        if not client_inv.empty:
            inventory_html = "<h3>🛢️ 설치/보관 장비 현황 (통합 탱크 재고)</h3><ul>"
            for _, row in client_inv.iterrows():
                item = row.get("품목", "미상")
                status = row.get("사용구분", "")
                serial = row.get("일련(제조)번호", "S/N 없음")
                vol = row.get("저장부피(L)", "")
                weight = row.get("저장무게(kg)", "")

                vol_str = f"{vol}L" if pd.notna(vol) and str(vol).strip() != "" else ""
                weight_str = (
                    f"{weight}kg" if pd.notna(weight) and str(weight).strip() != "" else ""
                )
                cap_str = (
                    f" / 용량: {vol_str} {weight_str}".strip()
                    if vol_str or weight_str
                    else ""
                )

                inventory_html += (
                    f"<li><b>[{html.escape(str(item))}]</b> "
                    f"{html.escape(str(status))} "
                    f"(S/N: {html.escape(str(serial))}{html.escape(cap_str)})</li>"
                )
            inventory_html += "</ul><br>"
    return inventory_html


def _format_client_note_html(client_name, info, inventory_html="") -> str:
    search_q = info.get("lookup_hint") or info["clean_name"]
    encoded_name = urllib.parse.quote(search_q)
    title = html.escape(str(client_name))
    clean_name = html.escape(str(info["clean_name"]))
    return f"""<h1>{title}</h1>
<br>
<h3>📌 요약 기업 정보 (데이터 출처: {html.escape(str(info['source']))})</h3>
<ul>
    <li><b>대표자:</b> {html.escape(str(info['ceo']))}</li>
    <li><b>업종:</b> {html.escape(str(info['industry']))}</li>
    <li><b>매출액:</b> {html.escape(str(info['revenue']))}</li>
    <li><b>영업이익:</b> {html.escape(str(info['profit']))}</li>
</ul>
<br>
{inventory_html}
<h3>🔗 상세 정보 원클릭 검색</h3>
<ul>
    <li><a href="https://search.naver.com/search.naver?query={encoded_name} 기업정보">네이버에서 '{clean_name}' 재무정보 보기</a></li>
    <li><a href="https://www.saramin.co.kr/zf_user/search/company?searchword={encoded_name}">사람인에서 '{clean_name}' 기업/채용 검색</a></li>
    <li><a href="https://www.jobkorea.co.kr/Search/?stext={encoded_name}&tabType=corp">잡코리아에서 '{clean_name}' 기업 검색</a></li>
</ul>
<br>
<h3>📝 영업 및 특이사항</h3>
<p></p>
"""


def _format_client_note_plain(client_name, info, df_integrated=None) -> str:
    search_q = info.get("lookup_hint") or info["clean_name"]
    encoded_name = urllib.parse.quote(search_q)
    lines = [
        str(client_name),
        "",
        f"📌 요약 기업 정보 (데이터 출처: {info['source']})",
        f"- 대표자: {info['ceo']}",
        f"- 업종: {info['industry']}",
        f"- 매출액: {info['revenue']}",
        f"- 영업이익: {info['profit']}",
        "",
    ]
    if (
        df_integrated is not None
        and not df_integrated.empty
        and "거래처(사용처)/보관장소" in df_integrated.columns
    ):
        client_inv = df_integrated[
            df_integrated["거래처(사용처)/보관장소"]
            .astype(str)
            .str.contains(client_name, regex=False, na=False)
        ]
        if not client_inv.empty:
            lines.append("🛢️ 설치/보관 장비 현황 (통합 탱크 재고)")
            for _, row in client_inv.iterrows():
                item = row.get("품목", "미상")
                status = row.get("사용구분", "")
                serial = row.get("일련(제조)번호", "S/N 없음")
                vol = row.get("저장부피(L)", "")
                weight = row.get("저장무게(kg)", "")
                vol_str = f"{vol}L" if pd.notna(vol) and str(vol).strip() != "" else ""
                weight_str = (
                    f"{weight}kg" if pd.notna(weight) and str(weight).strip() != "" else ""
                )
                cap_str = (
                    f" / 용량: {vol_str} {weight_str}".strip()
                    if vol_str or weight_str
                    else ""
                )
                lines.append(f"- [{item}] {status} (S/N: {serial}{cap_str})")
            lines.append("")
    lines.extend(
        [
            "🔗 상세 정보 검색",
            f"- 네이버: https://search.naver.com/search.naver?query={encoded_name} 기업정보",
            f"- 사람인: https://www.saramin.co.kr/zf_user/search/company?searchword={encoded_name}",
            f"- 잡코리아: https://www.jobkorea.co.kr/Search/?stext={encoded_name}&tabType=corp",
            "",
            "📝 영업 및 특이사항",
            "",
        ]
    )
    return "\n".join(lines)


def prepare_client_note_export(
    client_name, dart_api_key, df_integrated=None, address=None
) -> tuple[str, str, str, str]:
    info = get_company_info_hybrid(client_name, dart_api_key, address=address)
    inventory_html = _client_note_inventory_html(client_name, df_integrated)
    body_html = _format_client_note_html(client_name, info, inventory_html)
    plain_text = _format_client_note_plain(client_name, info, df_integrated)
    safe_title = html.escape(str(client_name))
    full_html = (
        "<!DOCTYPE html>\n"
        '<html lang="ko"><head><meta charset="utf-8">'
        f"<title>{safe_title}</title></head>\n"
        f"<body>{body_html}</body></html>"
    )
    fname = re.sub(r'[\\/:*?"<>|]+', "_", str(client_name)).strip()[:60] or "note"
    return body_html, plain_text, full_html, f"{fname}_메모.html"


def _render_tab2_note_share_html(plain_text: str, title: str) -> None:
    js_title = json.dumps(str(title), ensure_ascii=False)
    js_text = json.dumps(str(plain_text), ensure_ascii=False)
    components.html(
        f"""
<div style="font-family:system-ui,-apple-system,sans-serif;font-size:14px;">
  <button id="tab2-copy-btn" style="width:100%;padding:10px 12px;margin-bottom:8px;cursor:pointer;border:1px solid #CBD5E1;border-radius:8px;background:#fff;">📋 클립보드에 복사</button>
  <button id="tab2-share-btn" style="width:100%;padding:10px 12px;cursor:pointer;border:1px solid #CBD5E1;border-radius:8px;background:#fff;">📤 공유 (iPad·iPhone 메모 등)</button>
  <p id="tab2-share-status" style="margin:8px 0 0;color:#64748B;font-size:12px;line-height:1.4;"></p>
</div>
<script>
const title = {js_title};
const text = {js_text};
const status = document.getElementById("tab2-share-status");
document.getElementById("tab2-copy-btn").onclick = async () => {{
  try {{
    await navigator.clipboard.writeText(text);
    status.textContent = "클립보드에 복사했습니다. 메모 앱에 붙여넣기 하세요.";
  }} catch (e) {{
    status.textContent = "복사 실패. Safari/Chrome에서 이 페이지 클립보드 권한을 확인해 주세요.";
  }}
}};
document.getElementById("tab2-share-btn").onclick = async () => {{
  if (navigator.share) {{
    try {{
      await navigator.share({{ title: title, text: text }});
      status.textContent = "공유 시트를 열었습니다. 메모를 선택하세요.";
    }} catch (e) {{
      if (!e || e.name !== "AbortError") {{
        status.textContent = "공유를 완료하지 못했습니다. 클립보드 복사를 사용해 주세요.";
      }}
    }}
  }} else {{
    status.textContent = "이 브라우저는 공유를 지원하지 않습니다. 클립보드 복사를 사용하세요.";
  }}
}};
</script>
""",
        height=130,
    )


def _build_client_note_html(client_name, dart_api_key, df_integrated=None, address=None) -> str:
    info = get_company_info_hybrid(client_name, dart_api_key, address=address)
    inventory_html = _client_note_inventory_html(client_name, df_integrated)
    return _format_client_note_html(client_name, info, inventory_html)


def _note_name_key(name: str) -> str:
    s = re.sub(r"\s+", "", str(name or "").strip().lower())
    s = re.sub(r"[.\u3002]+$", "", s)
    s = re.sub(r"주식회사|㈜|\(주\)|유한회사|\(유\)", "", s)
    return s


def _macos_notes_auth_hint(err: str) -> str:
    low = (err or "").lower()
    if (
        "not authorized" in low
        or "assistive" in low
        or "-1743" in (err or "")
        or "1002" in (err or "")
    ):
        return (
            " 시스템 설정 → 개인정보 보호 및 보안 → **자동화**에서 "
            "터미널/Cursor/Python에 **메모(Notes)** 제어를 허용해 주세요."
        )
    return ""


def _run_osascript_file(script_text: str, args: list[str], timeout: int = 60):
    script_path = ""
    try:
        with tempfile.NamedTemporaryFile(
            mode="w",
            suffix=".applescript",
            prefix="dash_notes_",
            delete=False,
            encoding="utf-8",
        ) as script_file:
            script_file.write(script_text.strip())
            script_path = script_file.name
        return subprocess.run(
            ["osascript", script_path, *args],
            capture_output=True,
            text=True,
            timeout=timeout,
        )
    finally:
        if script_path:
            try:
                os.unlink(script_path)
            except Exception:
                pass


def list_macos_client_notes(force: bool = False) -> tuple[list[dict], str]:
    """메모「거래처」폴더의 노트·하위폴더 목록. (로컬 맥 전용)"""
    if not _is_local_macos():
        return [], "로컬 맥에서만 메모 목록을 읽을 수 있습니다."
    cache = st.session_state.get("_macos_notes_index")
    now = time.time()
    if (
        not force
        and isinstance(cache, dict)
        and now - float(cache.get("ts") or 0) < 300
        and cache.get("items") is not None
    ):
        return list(cache["items"]), ""

    script = """
on run
    set outLines to {}
    tell application "Notes"
        repeat with acc in accounts
            try
                set accName to name of acc
                if exists folder "거래처" of acc then
                    set parentFolder to folder "거래처" of acc
                    repeat with n in (notes of parentFolder)
                        set end of outLines to "NOTE\t" & accName & "\t" & (name of n)
                    end repeat
                    try
                        repeat with f in (folders of parentFolder)
                            set fName to name of f
                            set end of outLines to "FOLDER\t" & accName & "\t" & fName
                            repeat with n in (notes of f)
                                set end of outLines to "SUBNOTE\t" & accName & "\t" & fName & "\t" & (name of n)
                            end repeat
                        end repeat
                    end try
                end if
            end try
        end repeat
    end tell
    set AppleScript's text item delimiters to linefeed
    return outLines as text
end run
"""
    try:
        result = _run_osascript_file(script, [], timeout=90)
    except subprocess.TimeoutExpired:
        return [], "메모 목록 조회 시간 초과."
    except Exception as e:
        return [], f"메모 목록 조회 오류: {e}"
    if result.returncode != 0:
        err = (result.stderr or result.stdout or "").strip()
        return [], f"메모 목록 조회 실패: {err or '알 수 없는 오류'}.{_macos_notes_auth_hint(err)}"

    items = []
    for line in (result.stdout or "").splitlines():
        parts = line.split("\t")
        if len(parts) < 3:
            continue
        kind, acc = parts[0].strip(), parts[1].strip()
        if kind == "NOTE" and len(parts) >= 3:
            items.append({"kind": "NOTE", "account": acc, "name": parts[2].strip(), "parent": ""})
        elif kind == "FOLDER" and len(parts) >= 3:
            items.append({"kind": "FOLDER", "account": acc, "name": parts[2].strip(), "parent": ""})
        elif kind == "SUBNOTE" and len(parts) >= 4:
            items.append(
                {
                    "kind": "SUBNOTE",
                    "account": acc,
                    "parent": parts[2].strip(),
                    "name": parts[3].strip(),
                }
            )
    st.session_state["_macos_notes_index"] = {"ts": now, "items": items}
    return items, ""


def match_macos_client_note(client_name: str, items: list[dict]) -> dict | None:
    """거래처명과 같은(또는 정규화 일치) 메모 노트/폴더를 고른다."""
    raw = str(client_name or "").strip()
    if not raw or raw == "전체 거래처":
        return None
    cands = [raw]
    for c in _company_name_candidates(raw):
        if c not in cands:
            cands.append(c)
    cand_set = {c.strip() for c in cands if c and c.strip()}
    key_set = {_note_name_key(c) for c in cand_set if _note_name_key(c)}

    def _score(item: dict, mode: str) -> tuple:
        kind_rank = {"NOTE": 0, "SUBNOTE": 1, "FOLDER": 2}.get(item.get("kind"), 9)
        return (kind_rank, len(str(item.get("name") or "")), str(item.get("name") or ""))

    for mode, pred in (
        ("exact", lambda it: str(it.get("name") or "").strip() == raw),
        ("candidate", lambda it: str(it.get("name") or "").strip() in cand_set),
        ("normalized", lambda it: _note_name_key(it.get("name") or "") in key_set),
    ):
        hits = [it for it in items if pred(it)]
        if hits:
            hits.sort(key=lambda it: _score(it, mode))
            return {**hits[0], "match": mode}

    return None


def _show_macos_matched_note(match: dict) -> tuple[bool, str, str]:
    """일치한 노트를 메모 앱에서 열고 plaintext를 반환."""
    kind_raw = str(match.get("kind") or "NOTE").upper()
    mode = {"FOLDER": "folder", "SUBNOTE": "sub", "NOTE": "top"}.get(kind_raw, "top")
    name = str(match.get("name") or "")
    parent = str(match.get("parent") or "")
    script = """
on run argv
    set targetMode to item 1 of argv
    set noteName to item 2 of argv
    set parentName to item 3 of argv
    tell application "Notes"
        activate
        repeat with acc in accounts
            try
                if exists folder "거래처" of acc then
                    set parentFolder to folder "거래처" of acc
                    if targetMode is "folder" then
                        if exists folder noteName of parentFolder then
                            set sf to folder noteName of parentFolder
                            if (count of notes of sf) > 0 then
                                set oneNote to note 1 of sf
                                show oneNote
                                return plaintext of oneNote
                            end if
                        end if
                    else if targetMode is "sub" then
                        if exists folder parentName of parentFolder then
                            set sf to folder parentName of parentFolder
                            repeat with oneNote in (notes of sf)
                                if name of oneNote is noteName then
                                    show oneNote
                                    return plaintext of oneNote
                                end if
                            end repeat
                        end if
                    else
                        repeat with oneNote in (notes of parentFolder)
                            if name of oneNote is noteName then
                                show oneNote
                                return plaintext of oneNote
                            end if
                        end repeat
                    end if
                end if
            end try
        end repeat
    end tell
    return ""
end run
"""
    try:
        result = _run_osascript_file(script, [mode, name, parent], timeout=45)
    except subprocess.TimeoutExpired:
        return False, "메모 앱 응답 시간 초과.", ""
    except Exception as e:
        return False, f"메모 열기 오류: {e}", ""
    if result.returncode != 0:
        err = (result.stderr or result.stdout or "").strip()
        return False, f"메모 열기 실패: {err or '알 수 없는 오류'}.{_macos_notes_auth_hint(err)}", ""
    body = (result.stdout or "").strip()
    label = parent + "/" + name if parent else name
    return True, f"메모「거래처」에서 '{label}' 노트를 열었습니다.", body


def open_macos_notes_folder(
    client_name, dart_api_key, df_integrated=None, address=None
) -> dict:
    empty = {
        "ok": False,
        "msg": "",
        "created": False,
        "matched_name": "",
        "body": "",
        "kind": "",
    }
    if not client_name or client_name == "전체 거래처":
        empty["msg"] = "사이드바에서 특정 거래처를 선택한 뒤 다시 시도해 주세요."
        return empty

    if not _is_local_macos():
        empty["msg"] = (
            "macOS 메모 연동은 **맥에서 로컬 실행**할 때만 동작합니다. "
            "바탕화면의 **dashboard_Local** 바로가기로 실행하세요."
        )
        return empty

    items, list_err = list_macos_client_notes()
    if list_err and not items:
        empty["msg"] = list_err
        return empty

    match = match_macos_client_note(client_name, items)
    if match:
        ok, msg, body = _show_macos_matched_note(match)
        return {
            "ok": ok,
            "msg": msg,
            "created": False,
            "matched_name": str(match.get("name") or ""),
            "body": body,
            "kind": str(match.get("kind") or ""),
        }

    note_content = _build_client_note_html(
        client_name, dart_api_key, df_integrated=df_integrated, address=address
    )
    body_path = ""
    try:
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".html", prefix="dash_note_", delete=False, encoding="utf-8"
        ) as body_file:
            body_file.write(note_content.strip())
            body_path = body_file.name

        applescript = """
on run argv
    set bodyPath to item 1 of argv
    set noteName to item 2 of argv
    set noteBody to do shell script "cat " & quoted form of bodyPath

    tell application "Notes"
        activate
        set targetFolderName to "거래처"
        set targetAcc to missing value

        repeat with acc in accounts
            try
                if exists folder targetFolderName of acc then
                    set targetAcc to acc
                    exit repeat
                end if
            end try
        end repeat
        if targetAcc is missing value then
            repeat with acc in accounts
                if name of acc is "iCloud" then
                    set targetAcc to acc
                    exit repeat
                end if
            end repeat
        end if
        if targetAcc is missing value then set targetAcc to first account

        if not (exists folder targetFolderName of targetAcc) then
            make new folder at targetAcc with properties {name:targetFolderName}
        end if

        set parentFolder to folder targetFolderName of targetAcc
        set newNote to make new note at parentFolder with properties {body:noteBody}
        show newNote
        return plaintext of newNote
    end tell
end run
"""
        result = _run_osascript_file(applescript, [body_path, str(client_name)], timeout=45)
        if result.returncode == 0:
            st.session_state.pop("_macos_notes_index", None)
            return {
                "ok": True,
                "msg": f"같은 이름의 노트가 없어 '{client_name}' 노트를 새로 만들었습니다.",
                "created": True,
                "matched_name": str(client_name),
                "body": (result.stdout or "").strip(),
                "kind": "NOTE",
            }
        err = (result.stderr or result.stdout or "").strip()
        empty["msg"] = f"메모 연동 실패: {err or '알 수 없는 오류'}.{_macos_notes_auth_hint(err)}"
        return empty
    except subprocess.TimeoutExpired:
        empty["msg"] = "메모 앱 응답 시간 초과. Notes가 실행 중인지 확인해 주세요."
        return empty
    except Exception as e:
        empty["msg"] = f"메모 연동 오류: {e}"
        return empty
    finally:
        if body_path:
            try:
                os.unlink(body_path)
            except Exception:
                pass
def create_stacked_bar_chart(pivot_df, title_text="", y_suffix="", y_format=",.0f"):
    fig = go.Figure()
    sorted_years = sorted(pivot_df.columns, key=lambda x: str(x))
    
    color_map = {
        '2020': '#0052CC',
        '2021': '#4C9AFF',
        '2022': '#FF2B2B',
        '2023': '#FF9999',
        '2024': '#00B894',
        '2025': '#55E6A5',
        '2026': '#FF9F1A',
        # 월×주 비교용
        '1주': '#93C5FD',
        '2주': '#34D399',
        '3주': '#FBBF24',
        '4주': '#F472B6',
    }
    for yr in sorted_years:
        col_name = str(yr)
        color = color_map.get(col_name, None)
        fig.add_trace(
            go.Bar(
                x=pivot_df.index,
                y=pivot_df[yr],
                name=col_name,
                marker_color=color,
                hovertemplate=f"%{{x}} ({col_name}): %{{y:{y_format}}}{y_suffix}<extra></extra>"
            )
        )
    layout_args = dict(
        barmode='stack',
        xaxis=dict(title=None, tickangle=0),
        yaxis=dict(title=None, gridcolor='#E2E8F0'),
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=-0.25,
            xanchor="center",
            x=0.5
        ),
        margin=dict(l=10, r=10, t=10, b=40), 
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        height=420
    )
    
    if title_text:
        layout_args['title'] = dict(text=title_text, font=dict(size=14, color="#334155"))
        layout_args['margin']['t'] = 40
        
    fig.update_layout(**layout_args)
    return fig
def create_grouped_bar_chart(pivot_df, title_text="", y_suffix="", y_format=",.0f"):
    """월별 × 연도 비교용 그룹 막대그래프 (스택 아님). 블루 톤으로 통일."""
    fig = go.Figure()
    sorted_years = sorted(pivot_df.columns, key=lambda x: str(x))
    # 연도별 블루 스케일 (연한 하늘→진한 네이비) — 표 Blues 그라데이션과 톤 맞춤
    _blue_scale = [
        "#DBEAFE", "#BFDBFE", "#93C5FD", "#60A5FA",
        "#3B82F6", "#2563EB", "#1D4ED8", "#1E3A8A",
    ]
    color_map = {
        str(yr): _blue_scale[i % len(_blue_scale)]
        for i, yr in enumerate(sorted_years)
    }
    for yr in sorted_years:
        col_name = str(yr)
        fig.add_trace(
            go.Bar(
                x=pivot_df.index,
                y=pivot_df[yr],
                name=col_name,
                marker_color=color_map.get(col_name),
                hovertemplate=f"%{{x}} ({col_name}): %{{y:{y_format}}}{y_suffix}<extra></extra>",
            )
        )
    layout_args = dict(
        barmode="group",
        xaxis=dict(title=None, tickangle=0),
        yaxis=dict(title=None, gridcolor="#E2E8F0"),
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=-0.25,
            xanchor="center",
            x=0.5,
        ),
        margin=dict(l=10, r=10, t=10, b=40),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        height=420,
    )
    if title_text:
        layout_args["title"] = dict(text=title_text, font=dict(size=14, color="#334155"))
        layout_args["margin"]["t"] = 40
    fig.update_layout(**layout_args)
    return fig
@st.cache_data
def cached_client_item_qty_pivot_two_years(df_client_filtered, col_keys_tuple):
    """거래처 품목×년월 출고량 — col_keys 순서 그대로 (당월→과거 역순)."""
    if df_client_filtered is None or df_client_filtered.empty:
        return pd.DataFrame()
    cols = [str(c) for c in (col_keys_tuple or ()) if c]
    if not cols:
        return pd.DataFrame()
    raw = df_client_filtered.pivot_table(
        index="품목명", columns="연도월_정렬", values="출고량", aggfunc="sum"
    ).fillna(0)
    data = {key: (raw[key] if key in raw.columns else 0) for key in cols}
    out = pd.DataFrame(data, index=raw.index)
    out = out.reindex(columns=cols, fill_value=0)
    out = out.loc[(out.fillna(0) != 0).any(axis=1)]
    return out


@st.cache_data
def cached_client_item_sales_pivot_two_years(df_client_filtered, col_keys_tuple):
    """거래처 품목×년월 매출(만원, VAT포함) — col_keys 순서 (당월→과거 역순). tab2 전용."""
    if df_client_filtered is None or df_client_filtered.empty:
        return pd.DataFrame()
    cols = [str(c) for c in (col_keys_tuple or ()) if c]
    if not cols:
        return pd.DataFrame()
    raw = df_client_filtered.pivot_table(
        index="품목명", columns="연도월_정렬", values="매출액", aggfunc="sum"
    ).fillna(0)
    data = {key: (raw[key] if key in raw.columns else 0) for key in cols}
    out = pd.DataFrame(data, index=raw.index)
    out = out.reindex(columns=cols, fill_value=0) * 1.1 / 10000
    out = out.loc[(out.fillna(0) != 0).any(axis=1)]
    return out


@st.cache_data
def cached_get_yearly_monthly_qty_pivot(data_df, all_months, years):
    """월 × 연도 출고량 합계 (전년·당해 비교용). years 순서 유지."""
    yrs = [str(y) for y in years if y]
    # 중복 제거하되 호출 측 순서 유지
    seen = set()
    yrs_ordered = []
    for y in yrs:
        if y not in seen:
            seen.add(y)
            yrs_ordered.append(y)
    if data_df is None or data_df.empty or not yrs_ordered:
        return pd.DataFrame(0, index=all_months, columns=yrs_ordered)
    pvt = data_df.pivot_table(
        index="월", columns="연도", values="출고량", aggfunc="sum"
    ).fillna(0)
    pvt = pvt.reindex(index=all_months, fill_value=0)
    pvt = pvt.reindex(columns=yrs_ordered, fill_value=0)
    pvt.columns = [str(c) for c in pvt.columns]
    return pvt


def create_item_share_hbar(share_series, title_text="당해 년평균 품목별 출고 비중"):
    """품목 비중(%) 가로 막대 — Blues 톤 (tab2 비교용)."""
    s = share_series.dropna().astype(float)
    s = s[s > 0].sort_values(ascending=True)
    if s.empty:
        return None
    colors = [
        f"rgb({int(30 + 180 * (1 - i / max(len(s) - 1, 1)))},"
        f"{int(58 + 150 * (1 - i / max(len(s) - 1, 1)))},"
        f"{int(138 + 100 * (1 - i / max(len(s) - 1, 1)))})"
        for i in range(len(s))
    ]
    fig = go.Figure(
        go.Bar(
            x=s.values,
            y=s.index.astype(str),
            orientation="h",
            marker=dict(color=colors, line=dict(width=0)),
            text=[f"{v:.1f}%" for v in s.values],
            textposition="outside",
            cliponaxis=False,
            hovertemplate="%{y}: %{x:.1f}%<extra></extra>",
        )
    )
    fig.update_layout(
        title=dict(text=title_text, font=dict(size=14, color="#1E3A8A")),
        xaxis=dict(
            title=None,
            ticksuffix="%",
            range=[0, max(s.values) * 1.25 if len(s) else 1],
            gridcolor="#E2E8F0",
            zeroline=False,
        ),
        yaxis=dict(title=None, automargin=True),
        margin=dict(l=10, r=40, t=40, b=20),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        height=max(280, 28 * len(s) + 80),
        showlegend=False,
    )
    return fig
# ==========================================
# 3. 데이터 로딩 & 메모리 캐싱 (최적화) - Error 무시(on_bad_lines) 적용
# ==========================================
@st.cache_data(show_spinner="주소록을 읽어오는 중입니다...")
def load_address_file(address_bytes):
    if not address_bytes:
        return {}
    try:
        for enc in ["utf-8-sig", "cp949", "euc-kr", "utf-8"]:
            try:
                temp_addr = pd.read_csv(io.BytesIO(address_bytes), encoding=enc, on_bad_lines='skip', engine='python')
                if len(temp_addr.columns) >= 2:
                    k_col = temp_addr.columns[0]
                    v_col = temp_addr.columns[1]
                    temp_addr = temp_addr.dropna(subset=[k_col])
                    keys = temp_addr[k_col].astype(str).str.strip()
                    vals = temp_addr[v_col].astype(str).str.strip()
                    # 빈/NaN 주소는 제외
                    ok = vals.ne("") & vals.str.lower().ne("nan") & vals.str.lower().ne("none")
                    return dict(zip(keys[ok], vals[ok]))
                break
            except UnicodeDecodeError:
                continue
    except Exception:
        pass
    return {}


def resolve_client_address(client_name, addr_dict):
    """거래처명 → 주소. 정확 일치 후 공백/법인표기 완화 매칭."""
    if not client_name or client_name == "전체 거래처" or not addr_dict:
        return None
    name = str(client_name).strip()
    if name in addr_dict:
        return addr_dict[name]
    # (주)/주식회사 표기 차이 완화
    def _norm(s):
        s = str(s).strip()
        s = s.replace("(주)", "").replace("주식회사", "").replace("(유)", "")
        return "".join(s.split())
    n = _norm(name)
    for k, v in addr_dict.items():
        if _norm(k) == n:
            return v
    return None
def _drop_debt_noise_rows(df):
    """ERP 내보내기 푸터·합계 행을 제거 (거래처/채권 데이터 오염 방지)."""
    if df.empty or "거래처" not in df.columns:
        return df
    client_raw = df["거래처"].astype(str).str.strip()
    gubun_raw = (
        df["구분"].astype(str).str.strip()
        if "구분" in df.columns
        else pd.Series("", index=df.index)
    )
    noise = (
        client_raw.str.match(r"^\d{4}-\d{2}-\d{2}\s", na=False)
        | gubun_raw.str.contains(r"총매출|총수금|합계|소계", case=False, na=False)
        | client_raw.isin(["", "nan", "None", "NaN"])
    )
    return df.loc[~noise].copy()

# 거래처명이 비어 있는 매출 행용 고정 키 (CSV 저장·재로드 시 NaN으로 깨지지 않음)
EMPTY_CLIENT_LABEL = "(거래처명 없음)"


def _is_valid_client_name(name) -> bool:
    """실거래처명 여부. 빈칸/NaN 마커는 False (이후 EMPTY_CLIENT_LABEL로 치환)."""
    if name is None or (isinstance(name, float) and pd.isna(name)):
        return False
    s = str(name).strip()
    if not s:
        return False
    if s == EMPTY_CLIENT_LABEL:
        return True
    return s.lower() not in {"nan", "none", "nat", "null"}


def _is_mappable_client_name(name) -> bool:
    """담당자 수동 매핑에 쓸 수 있는 거래처명(빈칸·미지정 표기 제외)."""
    if name is None or (isinstance(name, float) and pd.isna(name)):
        return False
    s = str(name).strip()
    if s == EMPTY_CLIENT_LABEL:
        return True
    if not _is_valid_client_name(name):
        return False
    return s not in {"미지정", "-", ".", "없음"} and s.lower() not in {"nan", "none"}


def _normalize_manual_client_key(name) -> str | None:
    """수동 매핑 CSV 키 정규화. 빈/NaN 키 → EMPTY_CLIENT_LABEL."""
    if name is None or (isinstance(name, float) and pd.isna(name)):
        return EMPTY_CLIENT_LABEL
    s = str(name).strip()
    if not s or s.lower() in {"nan", "none", "nat", "null"}:
        return EMPTY_CLIENT_LABEL
    if _is_mappable_client_name(s):
        return s
    return None


def _drop_sales_noise_rows(df):
    """매출 CSV 하단 타임스탬프·이월미수 행 제거(실거래처·빈거래처 매출 무손실)."""
    if df.empty:
        return df
    noise = pd.Series(False, index=df.index)
    if "거래처" in df.columns:
        noise |= df["거래처"].astype(str).str.match(r"^\d{4}-\d{2}-\d{2}\s", na=False)
    if "품목명" in df.columns:
        noise |= df["품목명"].astype(str).str.contains(r"이월\s*미수|\[이월", na=False)
    return df.loc[~noise].copy()
@st.cache_data(show_spinner="업종 분류 데이터를 읽어오는 중입니다...")
def load_industry_file(industry_bytes):
    if not industry_bytes:
        return {}
    try:
        for enc in ["utf-8-sig", "cp949", "euc-kr", "utf-8"]:
            try:
                temp_ind = pd.read_csv(io.BytesIO(industry_bytes), encoding=enc, on_bad_lines='skip', engine='python')
                temp_ind.columns = temp_ind.columns.astype(str).str.strip()
                
                c_client = next((c for c in temp_ind.columns if "거래처" in c or "상호" in c), None)
                c_ind = next((c for c in temp_ind.columns if "분류" in c or "업종" in c), None)
                
                if c_client and c_ind:
                    temp_ind = temp_ind.dropna(subset=[c_client, c_ind])
                    return temp_ind.astype(str).set_index(c_client)[c_ind].to_dict()
                break
            except UnicodeDecodeError:
                continue
    except Exception:
        pass
    return {}
def dedupe_debt_client_gubun(df):
    """동일 거래처·구분은 1행만 유지. 월 합계 절댓값이 큰 행 우선."""
    if df is None or df.empty:
        return df if df is not None else pd.DataFrame()
    if "거래처" not in df.columns or "구분" not in df.columns:
        return df
    out = df.copy()
    out["거래처"] = out["거래처"].astype(str).str.strip()
    out["구분"] = out["구분"].astype(str).str.strip()
    month_cols = [c for c in out.columns if c not in ("거래처", "구분")]
    if month_cols:
        out["_abs_sum"] = out[month_cols].apply(pd.to_numeric, errors="coerce").fillna(0).abs().sum(axis=1)
        out = out.sort_values(["거래처", "구분", "_abs_sum"], ascending=[True, True, False])
        out = out.drop_duplicates(subset=["거래처", "구분"], keep="first")
        out = out.drop(columns=["_abs_sum"])
    else:
        out = out.drop_duplicates(subset=["거래처", "구분"], keep="first")
    return out.reset_index(drop=True)
@st.cache_data(show_spinner=False)
def load_debt_file(debt_bytes):
    if not debt_bytes:
        return pd.DataFrame()
    try:
        for enc in ["utf-8-sig", "cp949", "euc-kr", "utf-8"]:
            try:
                df_direct = pd.read_csv(io.BytesIO(debt_bytes), encoding=enc, on_bad_lines='skip', engine='python')
                df_direct.columns = df_direct.columns.astype(str).str.strip()
                
                # ★ 채권 CSV 파일에 잘못 생성된 빈 껍데기 열(Unnamed) 무조건 제거
                df_direct = df_direct.loc[:, ~df_direct.columns.str.contains('^Unnamed')]
                
                if "거래처" in df_direct.columns and "구분" in df_direct.columns:
                    df_direct = _drop_debt_noise_rows(df_direct)
                    # 💡 거래처 이름 공백 제거로 누락 데이터 완벽 방지
                    df_direct["거래처"] = df_direct["거래처"].replace("", np.nan).ffill().astype(str).str.strip()
                    
                    def map_gubun(val):
                        val_clean = re.sub(r'[^가-힣a-zA-Z0-9]', '', str(val))
                        
                        if not val_clean or val_clean == 'nan': 
                            return "매출"
                            
                        if any(k in val_clean for k in ["이월", "전월", "기초"]): return "이월"
                        if any(k in val_clean for k in ["익월", "다음", "차월"]): return "익월"
                        if any(k in val_clean for k in ["잔액", "미수", "현재", "기말", "잔금"]): return "잔액"
                        if any(k in val_clean for k in ["수금", "입금", "결제", "회수", "대변", "감소"]): return "수금"
                        
                        return "매출"
                    df_direct["구분"] = df_direct["구분"].apply(map_gubun)
                    valid_types = ["익월", "이월", "매출", "수금", "잔액"]
                    df_filtered = df_direct[df_direct["구분"].isin(valid_types)].copy()
                    
                    if not df_filtered.empty:
                        month_cols = [c for c in df_filtered.columns if c not in ["거래처", "구분"]]
                        for m_col in month_cols:
                            df_filtered[m_col] = (
                                pd.to_numeric(
                                    df_filtered[m_col].astype(str).str.replace(r"[^\d.-]", "", regex=True),
                                    errors="coerce"
                                ).fillna(0)
                            )
                        return dedupe_debt_client_gubun(df_filtered)
                break
            except UnicodeDecodeError:
                continue
    except Exception:
        pass
    return pd.DataFrame()
def _parse_sales_filename_year_month(file_name: str):
    """파일명에서 (연도YYYY, 월MM|None) 추출. 2026.csv / 202608.csv / 202607.csv."""
    base = os.path.basename(str(file_name or ""))
    m = re.match(r"^(20\d{2})(\d{2})?(?:\.csv)?$", base, flags=re.IGNORECASE)
    if m:
        return m.group(1), m.group(2)
    m2 = re.search(r"(20\d{2})(\d{2})?", base)
    if m2:
        return m2.group(1), m2.group(2)
    return None, None


def _dedupe_sales_file_meta(file_meta):
    """매출 CSV 중복 제거.

    - 같은 연도·같은 용량(복사본: 2026.csv + 202607.csv) → 하나만
    - 같은 연도에 월별(YYYYMM)이 있으면 연간(YYYY.csv) 제외
    """
    if not file_meta:
        return file_meta
    rows = []
    for item in file_meta:
        try:
            name, path = item[0], item[1]
        except Exception:
            continue
        y, mm = _parse_sales_filename_year_month(name)
        try:
            sz = int(item[3]) if len(item) > 3 else int(os.path.getsize(path))
        except Exception:
            sz = 0
        try:
            mt = int(item[2]) if len(item) > 2 else 0
        except Exception:
            mt = 0
        rows.append(
            {"name": name, "item": item, "y": y, "mm": mm, "sz": sz, "mt": mt}
        )

    drop = set()
    # 1) 동일 연도+용량 → 월별명·최신 mtime 우선 1개
    winners = {}
    for r in rows:
        if not r["y"] or r["sz"] <= 0:
            continue
        key = (r["y"], r["sz"])
        prev = winners.get(key)
        if prev is None:
            winners[key] = r
            continue
        score = (1 if r["mm"] else 0, r["mt"], len(r["name"]))
        pscore = (1 if prev["mm"] else 0, prev["mt"], len(prev["name"]))
        if score >= pscore:
            winners[key] = r
    for r in rows:
        if not r["y"] or r["sz"] <= 0:
            continue
        key = (r["y"], r["sz"])
        w = winners.get(key)
        if w is not None and r["name"] != w["name"]:
            drop.add(r["name"])

    # 2) 같은 연도에 월별(YYYYMM)이 있으면 연간(YYYY.csv) 제외 — 용량 달라도 2배 집계 방지
    by_y = {}
    for r in rows:
        if r["name"] in drop or not r["y"]:
            continue
        by_y.setdefault(r["y"], []).append(r)
    for y, grp in by_y.items():
        annuals = [r for r in grp if not r["mm"]]
        monthlies = [r for r in grp if r["mm"]]
        if not annuals or not monthlies:
            continue
        for a in annuals:
            drop.add(a["name"])

    out = [r["item"] for r in rows if r["name"] not in drop]
    return out


# 가스코아산. 비고 → 종속 거래처명 (예: 동희 → 가스코아산(동희))
_GASCO_PARENT_CLIENT = "가스코아산."
_GASCO_NOTE_SKIP_RE = re.compile(
    r"이관|변경|출고|회수|반납|정리|차액|으로|부터|까지|취소|오류|수정|"
    r"미검|공병|대납|직접|고객사|용기|고압|일반으로|합계|이월"
)


def _is_gascocoasan_subclient_note(note):
    """비고가 실거래처명인지 (메모·업무문구 제외)."""
    s = str(note or "").strip()
    if not s or s.lower() in ("nan", "none", "-", "없음"):
        return False
    if _GASCO_NOTE_SKIP_RE.search(s):
        return False
    if len(s) < 2:
        return False
    if re.fullmatch(r"[\d\s./\-]+", s):
        return False
    return True


def expand_gascocoasan_remark_clients(df):
    """가스코아산. + 비고(종속처) → 거래처명 '가스코아산(종속처)'로 분리.

    담당자 필터에서 '가스코아산' 지정 시 종속처가 각각 검색·선택되도록
    해당 행 담당자를 '가스코아산'으로 묶는다. 비고 없는 행은 '가스코아산.' 유지.
    """
    if df is None or df.empty or "거래처" not in df.columns:
        return df if df is not None else pd.DataFrame()
    out = df.copy()
    client = out["거래처"].fillna("").astype(str).str.strip()
    parent_mask = client == _GASCO_PARENT_CLIENT
    if not parent_mask.any():
        # ERP 기등록 가스코아산(xxx)만 있는 경우: 미지정이면 담당자 가스코아산
        if "담당자" in out.columns:
            named = client.str.match(r"^가스코아산\(.+\)$", na=False)
            unassigned = out["담당자"].astype(str).str.strip().isin(
                ["", "nan", "None", "NaN", "미지정", "담당자없음"]
            )
            out.loc[named & unassigned, "담당자"] = "가스코아산"
        return out
    if "비고" in out.columns:
        note = out["비고"].fillna("").astype(str).str.strip()
        sub_mask = parent_mask & note.map(_is_gascocoasan_subclient_note)
        if sub_mask.any():
            out.loc[sub_mask, "거래처"] = "가스코아산(" + note.loc[sub_mask] + ")"
    if "담당자" in out.columns:
        client2 = out["거래처"].astype(str).str.strip()
        # 부모(가스코아산.) + 분리/기등록 가스코아산(xxx)
        gasco_mask = (client2 == _GASCO_PARENT_CLIENT) | client2.str.match(
            r"^가스코아산\(.+\)$", na=False
        )
        out.loc[gasco_mask, "담당자"] = "가스코아산"
    return out


def _parse_sales_uploaded_tuples(file_tuples):
    """매출 CSV 바이트 튜플 → DataFrame (캐시 없음, 순수 파싱)."""
    if not file_tuples:
        return pd.DataFrame()
    df_list = []
    for file_name, content in file_tuples:
        try:
            decoded_text = None
            for enc in ["utf-8-sig", "cp949", "euc-kr", "utf-8"]:
                try:
                    decoded_text = content.decode(enc)
                    break
                except UnicodeDecodeError:
                    continue
            if decoded_text is None:
                decoded_text = content.decode("utf-8", errors="replace")
            lines = [line for line in decoded_text.splitlines() if line.strip()]
            if not lines:
                continue
            header_idx = 0
            max_matches = 0
            for i, line in enumerate(lines[:30]):
                cells = re.split(r',|\t', line)
                matches = sum(1 for cell in cells if any(kw in cell for kw in ["거래처", "품목", "매출", "단가", "수량", "담당", "일자"]))
                
                if matches > max_matches:
                    max_matches = matches
                    header_idx = i
                
                if max_matches >= 3:
                    break
            df = pd.read_csv(io.StringIO("\n".join(lines[header_idx:])), on_bad_lines='skip', engine='python')
            df.columns = df.columns.astype(str).str.strip()
            cols = list(df.columns)
            def find_col(priority_keywords, exclude_keywords=[]):
                clean_cols = [c.replace(" ", "") for c in cols]
                
                for kw in priority_keywords:
                    kw_clean = kw.replace(" ", "")
                    for orig_c, clean_c in zip(cols, clean_cols):
                        if any(ex in clean_c for ex in exclude_keywords):
                            continue
                        if kw_clean == clean_c:
                            return orig_c
                            
                for kw in priority_keywords:
                    kw_clean = kw.replace(" ", "")
                    for orig_c, clean_c in zip(cols, clean_cols):
                        if any(ex in clean_c for ex in exclude_keywords):
                            continue
                        if kw_clean in clean_c:
                            return orig_c
                return None
            c_staff = find_col(["담당자", "담당자명", "영업담당", "영업사원", "담당"], ["코드", "ID", "번호"])
            c_client = find_col(["거래처", "거래처명", "상호명", "고객명", "회사명", "상호", "고객"], ["코드", "ID", "번호", "담당", "영업"])
            c_item = find_col(["품목명", "제품명", "상품명", "품목", "제품"], ["코드", "ID", "번호", "규격"])
            c_sales = find_col(["매출액", "금액", "매출"], ["일", "자", "수량", "량", "단가"])
            c_qty = find_col(["출고량", "수량", "출고"], ["액", "금액", "단가"])
            c_price = find_col(["단가", "단 가", "판매단가", "공급단가"], ["액", "금액", "수량", "량"])
            c_date = find_col(["매출일자", "매출일", "일자", "날짜", "출고일"])
            c_note = find_col(["비고", "적요", "메모", "특이사항"], ["코드", "ID"])
            rename_dict = {}
            if c_client: rename_dict[c_client] = "거래처"
            if c_item: rename_dict[c_item] = "품목명"
            if c_staff: rename_dict[c_staff] = "담당자"
            if c_sales: rename_dict[c_sales] = "매출액"
            if c_qty: rename_dict[c_qty] = "출고량"
            if c_price: rename_dict[c_price] = "단가"
            if c_date: rename_dict[c_date] = "매출일자_raw"
            if c_note: rename_dict[c_note] = "비고"
            df = df.rename(columns=rename_dict)
            df = _drop_sales_noise_rows(df)
            for req in ["거래처", "품목명", "담당자"]:
                if req not in df.columns:
                    df[req] = "미지정"
            file_year = _parse_sales_filename_year_month(file_name)[0] or "2026"
            date_col = "매출일자_raw" if "매출일자_raw" in df.columns else df.columns[0]
            df["매출일_dt"] = parse_date_series_robust(df[date_col], default_year=file_year)
            df["매출액"] = pd.to_numeric(
                df["매출액"].astype(str).str.replace(r"[^\d.-]", "", regex=True), errors="coerce"
            ).fillna(0) if "매출액" in df.columns else 0
            df["출고량"] = pd.to_numeric(
                df["출고량"].astype(str).str.replace(r"[^\d.-]", "", regex=True), errors="coerce"
            ).fillna(0) if "출고량" in df.columns else 0
            df["단가"] = pd.to_numeric(
                df["단가"].astype(str).str.replace(r"[^\d.-]", "", regex=True), errors="coerce"
            ).fillna(0) if "단가" in df.columns else 0
            df["거래처"] = df["거래처"].fillna("").astype(str).str.strip()
            df["담당자"] = df["담당자"].fillna("미지정").astype(str).str.strip()
            # 빈 거래처 → 고정 라벨(담당자 매핑·집계 가능, 매출 무손실)
            _blank_client = ~df["거래처"].map(_is_valid_client_name)
            if _blank_client.any():
                df.loc[_blank_client, "거래처"] = EMPTY_CLIENT_LABEL
            # 가스코아산. 비고 종속처 → 개별 거래처명으로 분리
            df = expand_gascocoasan_remark_clients(df)
            df = normalize_items_vectorized(df)
            df = df.dropna(subset=["매출일_dt"])
            df["연도"] = df["매출일_dt"].dt.year.astype(str)
            df["월"] = df["매출일_dt"].dt.strftime("%m월")
            df["연도월_정렬"] = df["연도"].astype(str).str[2:] + "년 " + df["월"].astype(str)
            if not df.empty:
                df_list.append(df)
        except Exception as e:
            st.sidebar.error(f"파일 읽기 오류 ({file_name}): {e}")
    result_df = pd.concat(df_list, ignore_index=True) if df_list else pd.DataFrame()
    # # 동일 행 중복(연간+월별 복사본 등) 제거 — 매출 2배 집계 방지
    # if not result_df.empty:
    #     subset = [
    #         c
    #         for c in ("거래처", "품목명", "매출일_dt", "매출액", "출고량", "단가", "담당자")
    #         if c in result_df.columns
    #     ]
    #     if len(subset) >= 4:
    #         result_df = result_df.drop_duplicates(subset=subset, keep="last")
    
    # ★ '미지정' 및 '담당자없음' 데이터 통합 및 최신 담당자 자동 추론 로직
    if not result_df.empty and "거래처" in result_df.columns and "담당자" in result_df.columns:
        result_df["담당자"] = result_df["담당자"].fillna("미지정").astype(str).str.strip()
        
        # '담당자없음'을 포함하여 비어있거나 누락된 표현들을 모두 '미지정'으로 통일
        invalid_staff_markers = ["", "nan", "None", "NAT", "NaN", "담당자없음", "지정안함", "없음"]
        result_df["담당자"] = result_df["담당자"].replace(invalid_staff_markers, "미지정")
        
        temp_df = result_df.dropna(subset=["매출일_dt"]).sort_values("매출일_dt")
        valid_staff_map = temp_df[~temp_df["담당자"].isin(["미지정"])].groupby("거래처")["담당자"].last().to_dict()
        
        mask_unassigned = result_df["담당자"] == "미지정"
        result_df.loc[mask_unassigned, "담당자"] = (
            result_df.loc[mask_unassigned, "거래처"]
            .map(valid_staff_map)
            .fillna("미지정")
        )
        result_df.loc[result_df["담당자"].isin(invalid_staff_markers), "담당자"] = "미지정"
    # 수동 지정 매핑 적용 (무손실 보존 · 빈 거래처 키는 EMPTY_CLIENT_LABEL로 복구)
    result_df = _apply_manual_staff_mapping(result_df)
    # ★ 거래처명 맨 앞이 z/Z → 담당자 '거래종료' (대시보드 전체 공통, 최종 확정)
    if not result_df.empty and "거래처" in result_df.columns and "담당자" in result_df.columns:
        _client_s = result_df["거래처"].astype(str).str.strip()
        mask_closed = _client_s.str.match(r"^[zZ]", na=False)
        if mask_closed.any():
            result_df.loc[mask_closed, "담당자"] = "거래종료"
    return result_df


def _apply_manual_staff_mapping(df):
    """수동 담당자 매핑을 캐시 밖에서도 강제 적용. 빈 거래처 → EMPTY_CLIENT_LABEL."""
    if df is None or df.empty or "거래처" not in df.columns:
        return df if df is not None else pd.DataFrame()
    out = df.copy()
    out["거래처"] = out["거래처"].fillna("").astype(str).str.strip()
    blank = ~out["거래처"].map(_is_valid_client_name)
    if blank.any():
        out.loc[blank, "거래처"] = EMPTY_CLIENT_LABEL
    manual_map_path = os.path.join(CACHE_DIR, "manual_staff_mapping.csv")
    if not os.path.exists(manual_map_path):
        return out
    if "담당자" not in out.columns:
        out["담당자"] = "미지정"
    try:
        manual_df = pd.read_csv(manual_map_path)
        if "거래처" not in manual_df.columns or "담당자" not in manual_df.columns:
            return out
        manual_dict = {}
        for _, mrow in manual_df.iterrows():
            ck = _normalize_manual_client_key(mrow["거래처"])
            if ck is None:
                continue
            staff = str(mrow["담당자"]).strip() if pd.notna(mrow["담당자"]) else ""
            if staff and staff not in ("nan", "None", "NaN", "미지정"):
                manual_dict[ck] = staff
        if manual_dict:
            mask_manual = out["거래처"].isin(manual_dict.keys())
            out.loc[mask_manual, "담당자"] = out.loc[mask_manual, "거래처"].map(manual_dict)
    except Exception:
        pass
    return out


@st.cache_data(show_spinner="데이터 파싱 중입니다...")
def load_uploaded_files_from_bytes(file_tuples, manual_map_token=None, parse_version=6):
    return _parse_sales_uploaded_tuples(file_tuples)


def _manual_staff_map_cache_token():
    """담당자 수동매핑 CSV 변경 시 매출 캐시를 무효화하기 위한 토큰."""
    path = os.path.join(CACHE_DIR, "manual_staff_mapping.csv")
    try:
        st_info = os.stat(path)
        return (int(st_info.st_mtime_ns), int(st_info.st_size))
    except OSError:
        return (0, 0)


@st.cache_data(show_spinner="데이터 파싱 중입니다...")
def load_uploaded_files_from_meta(file_meta, manual_map_token=None, parse_version=6):
    """파일명·경로·mtime·size만 캐시 키로 사용.
    새로고침마다 ~10MB CSV를 읽어 해싱하던 비용을 제거(파싱 결과는 동일).
    manual_map_token / parse_version: 매핑·파서 변경 시 캐시 무효화."""
    if not file_meta:
        return pd.DataFrame()
    file_tuples = []
    for item in file_meta:
        try:
            f_name, f_path = item[0], item[1]
            with open(f_path, "rb") as sf:
                file_tuples.append((f_name, sf.read()))
        except Exception:
            continue
    return _parse_sales_uploaded_tuples(file_tuples)


def _dash_pivot_cache_key(
    selected_client,
    selected_staff,
    selected_item,
    start_date,
    end_date,
    sales_meta,
    manual_token,
):
    return (
        selected_client,
        tuple(selected_staff or ()),
        tuple(selected_item or ()),
        str(start_date),
        str(end_date),
        sales_meta,
        manual_token,
    )


def _dash_compute_pivot_bundle(df_base, df_client_filtered, df_f, full_df, all_months):
    """필터 조합별 피벗·지표 (세션 캐시용 — 거래처 재선택 시 즉시 복원)."""
    raw_years = sorted(full_df["연도"].unique()) if "연도" in full_df.columns else ["2026"]
    years = sorted(raw_years, reverse=True)
    desired_order = [f"{y[2:]}년 {m}" for y in years for m in all_months]
    pivot_m_total = cached_get_yearly_monthly_pivot(df_base, all_months, years)
    client_item_qty_pivot = cached_client_item_qty_pivot(
        df_client_filtered, years, all_months
    )
    sales_p, qty_p, unit_price_p = cached_tab3_pivots(df_f, years, all_months)
    staff_pivot = cached_staff_pivot(df_base, desired_order)
    detail_cols = [
        "매출일_dt",
        "담당자",
        "거래처",
        "품목명",
        "출고량",
        "단가",
        "매출액",
    ]
    df_detail = (
        df_f[detail_cols] if not df_f.empty else pd.DataFrame(columns=detail_cols)
    )
    df_total_monthly = df_base.groupby(df_base["매출일_dt"].dt.to_period("M"))[
        "매출액"
    ].sum()
    if not df_total_monthly.empty:
        latest_period_total = df_total_monthly.index.max()
        cur_month_sales_total = df_total_monthly.loc[latest_period_total] * 1.1
        prev_period_total = latest_period_total - 1
        prev_month_sales_total = df_total_monthly.get(prev_period_total, 0.0) * 1.1
        mom_rate_total = (
            (cur_month_sales_total - prev_month_sales_total)
            / prev_month_sales_total
            * 100
            if prev_month_sales_total > 0
            else 0.0
        )
        avg_monthly_sales_total = df_total_monthly.mean() * 1.1
        avg_rate_total = (
            (cur_month_sales_total - avg_monthly_sales_total)
            / avg_monthly_sales_total
            * 100
            if avg_monthly_sales_total > 0
            else 0.0
        )
        latest_month_str_total = latest_period_total.strftime("%Y년 %m월")
    else:
        cur_month_sales_total = 0.0
        prev_month_sales_total = 0.0
        mom_rate_total = 0.0
        avg_monthly_sales_total = 0.0
        avg_rate_total = 0.0
        latest_month_str_total = "-"
    if not df_client_filtered.empty:
        df_client_monthly = (
            df_client_filtered.groupby(
                df_client_filtered["매출일_dt"].dt.to_period("M")
            )["매출액"].sum()
            * 1.1
        )
        latest_period_client = df_client_monthly.index.max()
        cur_month_sales_client = df_client_monthly.loc[latest_period_client]
        prev_period_client = latest_period_client - 1
        prev_month_sales_client = df_client_monthly.get(prev_period_client, 0.0)
        mom_rate_client = (
            (cur_month_sales_client - prev_month_sales_client)
            / prev_month_sales_client
            * 100
            if prev_month_sales_client > 0
            else 0.0
        )
        avg_monthly_sales_client = df_client_monthly.mean()
        avg_rate_client = (
            (cur_month_sales_client - avg_monthly_sales_client)
            / avg_monthly_sales_client
            * 100
            if avg_monthly_sales_client > 0
            else 0.0
        )
        latest_month_str_client = latest_period_client.strftime("%Y년 %m월")
    else:
        cur_month_sales_client = 0.0
        prev_month_sales_client = 0.0
        mom_rate_client = 0.0
        avg_monthly_sales_client = 0.0
        avg_rate_client = 0.0
        latest_month_str_client = "-"
    return {
        "years": years,
        "desired_order": desired_order,
        "pivot_m_total": pivot_m_total,
        "client_item_qty_pivot": client_item_qty_pivot,
        "sales_p": sales_p,
        "qty_p": qty_p,
        "unit_price_p": unit_price_p,
        "staff_pivot": staff_pivot,
        "df_detail": df_detail,
        "cur_month_sales_total": cur_month_sales_total,
        "prev_month_sales_total": prev_month_sales_total,
        "mom_rate_total": mom_rate_total,
        "avg_monthly_sales_total": avg_monthly_sales_total,
        "avg_rate_total": avg_rate_total,
        "latest_month_str_total": latest_month_str_total,
        "cur_month_sales_client": cur_month_sales_client,
        "prev_month_sales_client": prev_month_sales_client,
        "mom_rate_client": mom_rate_client,
        "avg_monthly_sales_client": avg_monthly_sales_client,
        "avg_rate_client": avg_rate_client,
        "latest_month_str_client": latest_month_str_client,
    }


# ==========================================
# 4. 피벗 및 지표 계산 연산 캐싱 (최적화)
# ==========================================
@st.cache_data
def cached_get_yearly_monthly_pivot(data_df, all_months, years):
    if data_df.empty:
        return pd.DataFrame(0, index=all_months, columns=[str(y) for y in years])
    
    pvt = data_df.pivot_table(
        index="월", columns="연도", values="매출액", aggfunc="sum"
    ).fillna(0) * 1.1 / 10000
    
    pvt = pvt.reindex(index=all_months, fill_value=0)
    all_yrs = [str(y) for y in years]
    pvt = pvt.reindex(columns=all_yrs, fill_value=0)
    return pvt
@st.cache_data
def cached_get_item_pivot(data_df, item_name, metric, all_months, years):
    if data_df.empty:
        return pd.DataFrame(0, index=all_months, columns=[str(y) for y in years])
        
    df_item = data_df[data_df["품목명"] == item_name]
    if df_item.empty:
        return pd.DataFrame(0, index=all_months, columns=[str(y) for y in years])
        
    if metric == "매출액 (만원)":
        pvt = df_item.pivot_table(index="월", columns="연도", values="매출액", aggfunc="sum").fillna(0) * 1.1 / 10000
    elif "출고량" in metric:
        pvt = df_item.pivot_table(index="월", columns="연도", values="출고량", aggfunc="sum").fillna(0)
        target_bulks = ["CO2 (kg, Bulk)", "N2 (kg, Bulk)", "O2 (kg, Bulk)", "AR (kg, Bulk)"]
        if item_name in target_bulks:
            pvt = pvt / 1000
    elif metric == "총매출 대비 비중 (%)":
        pvt_item = df_item.pivot_table(index="월", columns="연도", values="매출액", aggfunc="sum").fillna(0)
        pvt_total = data_df.pivot_table(index="월", columns="연도", values="매출액", aggfunc="sum").fillna(0)
        pvt = (pvt_item / pvt_total.replace(0, np.nan) * 100).fillna(0)
        
    pvt = pvt.reindex(index=all_months, fill_value=0)
    all_yrs = [str(y) for y in years]
    pvt = pvt.reindex(columns=all_yrs, fill_value=0)
    return pvt


def _calendar_week_of_month(day):
    """월 내 주차: 1~7→1주, 8~14→2주, 15~21→3주, 22~31→4주."""
    try:
        d = int(day)
    except (TypeError, ValueError):
        return "4주"
    if d <= 7:
        return "1주"
    if d <= 14:
        return "2주"
    if d <= 21:
        return "3주"
    return "4주"


@st.cache_data
def cached_get_item_week_pivot(data_df, item_name, metric, year, all_months):
    """품목·연도 기준 월(행) × 주차(열: 1~4주) 피벗. 사용량 비교용."""
    week_cols = ["1주", "2주", "3주", "4주"]
    empty = pd.DataFrame(0.0, index=all_months, columns=week_cols)
    if data_df is None or data_df.empty or not item_name or not year:
        return empty
    if "매출일_dt" not in data_df.columns:
        return empty
    yr = str(year)
    df = data_df[
        (data_df["품목명"] == item_name) & (data_df["연도"].astype(str) == yr)
    ].copy()
    if df.empty:
        return empty
    df = df.dropna(subset=["매출일_dt"])
    if df.empty:
        return empty
    df["주차"] = df["매출일_dt"].dt.day.map(_calendar_week_of_month)
    if metric == "매출액 (만원)":
        pvt = (
            df.pivot_table(index="월", columns="주차", values="매출액", aggfunc="sum")
            .fillna(0)
            * 1.1
            / 10000
        )
    elif "출고량" in str(metric):
        pvt = df.pivot_table(
            index="월", columns="주차", values="출고량", aggfunc="sum"
        ).fillna(0)
        target_bulks = [
            "CO2 (kg, Bulk)",
            "N2 (kg, Bulk)",
            "O2 (kg, Bulk)",
            "AR (kg, Bulk)",
        ]
        if item_name in target_bulks:
            pvt = pvt / 1000
    elif "비중" in str(metric):
        # 해당 월 품목 합 대비 주차별 비중(%) — 행 합 ≈ 100
        pvt_amt = df.pivot_table(
            index="월", columns="주차", values="매출액", aggfunc="sum"
        ).fillna(0)
        row_sum = pvt_amt.sum(axis=1).replace(0, np.nan)
        pvt = (pvt_amt.div(row_sum, axis=0) * 100).fillna(0)
    else:
        pvt = df.pivot_table(
            index="월", columns="주차", values="출고량", aggfunc="sum"
        ).fillna(0)
    pvt = pvt.reindex(index=all_months, fill_value=0)
    pvt = pvt.reindex(columns=week_cols, fill_value=0)
    return pvt


@st.cache_data
def cached_get_item_month_week_year(data_df, item_name, metric, all_months, years):
    """월×주차(행 MultiIndex) × 연도(열) — 월 행 펼침용 주간 상세."""
    week_labels = ["1주", "2주", "3주", "4주"]
    yr_cols = [str(y) for y in years]
    idx = pd.MultiIndex.from_product([all_months, week_labels], names=["월", "주차"])
    empty = pd.DataFrame(0.0, index=idx, columns=yr_cols)
    if data_df is None or data_df.empty or not item_name:
        return empty
    if "매출일_dt" not in data_df.columns:
        return empty
    df = data_df[data_df["품목명"] == item_name].copy()
    if df.empty:
        return empty
    df = df.dropna(subset=["매출일_dt"])
    if df.empty:
        return empty
    df["연도"] = df["연도"].astype(str)
    df["주차"] = df["매출일_dt"].dt.day.map(_calendar_week_of_month)
    target_bulks = [
        "CO2 (kg, Bulk)",
        "N2 (kg, Bulk)",
        "O2 (kg, Bulk)",
        "AR (kg, Bulk)",
    ]
    if metric == "매출액 (만원)":
        pvt = (
            df.pivot_table(
                index=["월", "주차"], columns="연도", values="매출액", aggfunc="sum"
            )
            .fillna(0)
            * 1.1
            / 10000
        )
    elif "출고량" in str(metric):
        pvt = df.pivot_table(
            index=["월", "주차"], columns="연도", values="출고량", aggfunc="sum"
        ).fillna(0)
        if item_name in target_bulks:
            pvt = pvt / 1000
    elif "비중" in str(metric):
        pvt_item = df.pivot_table(
            index=["월", "주차"], columns="연도", values="매출액", aggfunc="sum"
        ).fillna(0)
        df_all = data_df.dropna(subset=["매출일_dt"]).copy()
        df_all["연도"] = df_all["연도"].astype(str)
        df_all["주차"] = df_all["매출일_dt"].dt.day.map(_calendar_week_of_month)
        pvt_total = df_all.pivot_table(
            index=["월", "주차"], columns="연도", values="매출액", aggfunc="sum"
        ).fillna(0)
        pvt = (pvt_item / pvt_total.replace(0, np.nan) * 100).fillna(0)
    else:
        pvt = df.pivot_table(
            index=["월", "주차"], columns="연도", values="출고량", aggfunc="sum"
        ).fillna(0)
    pvt.columns = pvt.columns.astype(str)
    pvt = pvt.reindex(index=idx, fill_value=0)
    pvt = pvt.reindex(columns=yr_cols, fill_value=0)
    return pvt


def _styler_heatmap_colors(data_df, cmap_name="Blues"):
    """style_with_sum / st.dataframe 과 동일한 background_gradient 색·글자색.
    returns dict[(row_label, col_label)] -> (bg_hex, fg_hex)
    """
    out = {}
    if data_df is None or getattr(data_df, "empty", True):
        return out
    try:
        num = data_df.apply(pd.to_numeric, errors="coerce").fillna(0.0)
        if num.empty:
            return out
        name = cmap_name if cmap_name in ("Blues", "Greens", "Purples") else "Blues"
        styler = num.style.background_gradient(cmap=name, axis=None)
        styler._compute()
        ctx = getattr(styler, "ctx", None) or {}
        idxs = list(num.index)
        cols = list(num.columns)
        for (r, c), props in ctx.items():
            if r >= len(idxs) or c >= len(cols):
                continue
            bg, fg = "#FFFFFF", "#31333F"
            for item in props:
                if not isinstance(item, (list, tuple)) or len(item) < 2:
                    continue
                k, v = str(item[0]).lower(), str(item[1]).strip()
                if "background" in k:
                    bg = v
                elif k == "color":
                    fg = v
            try:
                if abs(float(num.iat[r, c])) < 1e-12:
                    bg, fg = "#FFFFFF", "#31333F"
            except Exception:
                pass
            if fg.lower() in ("#000", "#000000", "black", "rgb(0, 0, 0)"):
                fg = "#31333F"
            out[(idxs[r], cols[c])] = (bg, fg)
    except Exception:
        pass
    return out


def _heatmap_cell_bg(val, vmin, vmax, cmap_name="Blues"):
    """값 → 히트맵 배경색 (#rrggbb). 폴백용."""
    try:
        v = float(val)
    except (TypeError, ValueError):
        return "#FFFFFF"
    if vmax <= vmin or abs(v) < 1e-12:
        return "#FFFFFF"
    t = float(max(0.0, min(1.0, (v - vmin) / (vmax - vmin))))
    try:
        import matplotlib.pyplot as plt

        name = cmap_name if cmap_name in ("Blues", "Greens", "Purples") else "Blues"
        cmap = plt.get_cmap(name)
        r, g, b, _a = cmap(t)
        return f"#{int(round(r * 255)):02x}{int(round(g * 255)):02x}{int(round(b * 255)):02x}"
    except Exception:
        return "#FFFFFF"


def _heatmap_text_color(bg_hex):
    """어두운 칸은 흰 글씨."""
    try:
        h = str(bg_hex or "").lstrip("#")
        if len(h) != 6:
            return "#31333F"
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        lum = (0.299 * r + 0.587 * g + 0.114 * b) / 255.0
        return "#FFFFFF" if lum < 0.55 else "#31333F"
    except Exception:
        return "#31333F"


def render_month_expandable_week_table(
    month_pivot,
    week_year_pivot,
    fmt_kind="qty",
    cmap_name="Blues",
    height=460,
):
    """월 행 클릭 → 1~4주 펼침/접힘. tab1 4대품목 전용.
    그라데이션=상단 Blues와 동일, 글씨 크기는 고정(부모 복사 금지·레이아웃 안정).
    iPad: table-layout auto + 가로 스크롤로 연도·숫자 잘림 방지. 맥 레이아웃 무손실."""
    if month_pivot is None or month_pivot.empty:
        st.info("표시할 품목 데이터가 없습니다.")
        return
    week_labels = ["1주", "2주", "3주", "4주"]
    month_pivot = month_pivot.copy()
    month_pivot.columns = month_pivot.columns.astype(str)
    if week_year_pivot is not None and not week_year_pivot.empty:
        week_year_pivot = week_year_pivot.copy()
        week_year_pivot.columns = week_year_pivot.columns.astype(str)
    year_cols = sorted(list(month_pivot.columns), reverse=True)
    months = [m for m in month_pivot.index if str(m) != "연간 합계"]
    touch = False
    try:
        touch = bool(is_touch_ui())
    except Exception:
        touch = False

    def _fmt(v):
        try:
            x = float(v)
        except (TypeError, ValueError):
            return "0"
        if fmt_kind == "pct":
            return f"{x:,.0f}%"
        return f"{x:,.0f}"

    month_only = month_pivot.reindex(index=months, columns=year_cols).fillna(0.0)
    month_colors = _styler_heatmap_colors(month_only, cmap_name)
    week_colors = {}
    if week_year_pivot is not None and not week_year_pivot.empty:
        try:
            wdf = week_year_pivot.reindex(columns=year_cols).fillna(0.0)
            week_colors = _styler_heatmap_colors(wdf, cmap_name)
        except Exception:
            week_colors = {}

    month_nums = []
    for m in months:
        for y in year_cols:
            try:
                month_nums.append(float(month_only.at[m, y]))
            except Exception:
                month_nums.append(0.0)
    vmax = max(month_nums) if month_nums else 1.0
    vmin = 0.0

    def _cell_style(val, key, week=False):
        bg_fg = None
        if week and key in week_colors:
            bg_fg = week_colors[key]
        elif (not week) and key in month_colors:
            bg_fg = month_colors[key]
        if bg_fg:
            bg, fg = bg_fg
        else:
            bg = _heatmap_cell_bg(val, vmin, vmax, cmap_name)
            fg = _heatmap_text_color(bg)
        pad = "6px 8px" if not week else "5px 8px"
        fsz = "13px" if not week else "12px"
        # iPad: ellipsis 금지 · 최소폭 확보. 맥: 기존과 동일
        overflow = "visible" if touch else "hidden"
        ellipsis = "" if touch else "text-overflow:ellipsis;"
        minw = "min-width:3.4rem;" if touch else ""
        return (
            f"padding:{pad};text-align:right;vertical-align:middle;"
            f"background:{bg};color:{fg};font-size:{fsz};font-weight:400;"
            f"line-height:1.35;border-bottom:1px solid #E6EAF0;"
            f"border-right:1px solid #EEF2F6;white-space:nowrap;"
            f"overflow:{overflow};{ellipsis}{minw}"
        )

    sum_cells = []
    for y in year_cols:
        try:
            s = float(month_pivot[y].sum()) if y in month_pivot.columns else 0.0
        except Exception:
            s = 0.0
        sum_cells.append(s)

    th_min = "min-width:3.4rem;" if touch else ""
    th = "".join(
        f'<th style="position:sticky;top:0;z-index:2;background:#F0F2F6;padding:6px 8px;'
        f'border-bottom:1px solid #D0D7DE;border-right:1px solid #E6EAF0;'
        f'text-align:right;vertical-align:middle;font-weight:600;font-size:13px;'
        f'color:#31333F;line-height:1.35;white-space:nowrap;{th_min}">{html.escape(y)}</th>'
        for y in year_cols
    )
    body_parts = []
    for m in months:
        mid = html.escape(str(m))
        mcells = []
        for y in year_cols:
            try:
                val = float(month_only.at[m, y])
            except Exception:
                val = 0.0
            stl = _cell_style(val, (m, y), week=False)
            mcells.append(f'<td style="{stl}">{_fmt(val)}</td>')
        body_parts.append(
            f'<tr class="mrow" data-month="{mid}" style="cursor:pointer;">'
            f'<td class="mlabel" style="padding:6px 8px;font-weight:600;font-size:13px;'
            f'line-height:1.35;white-space:nowrap;color:#31333F;text-align:left;'
            f'position:sticky;left:0;z-index:1;background:#F0F2F6;'
            f'border-right:1px solid #D0D7DE;border-bottom:1px solid #E6EAF0;">'
            f'<span class="chev">▸</span> {mid}</td>{"".join(mcells)}</tr>'
        )
        for w in week_labels:
            wcells = []
            for y in year_cols:
                try:
                    val = float(week_year_pivot.loc[(m, w), y])
                except Exception:
                    val = 0.0
                stl = _cell_style(val, ((m, w), y), week=True)
                wcells.append(f'<td style="{stl}">{_fmt(val)}</td>')
            body_parts.append(
                f'<tr class="wrow" data-parent="{mid}" style="display:none;">'
                f'<td style="padding:5px 8px 5px 22px;color:#64748B;font-size:12px;'
                f'line-height:1.35;white-space:nowrap;text-align:left;'
                f'position:sticky;left:0;z-index:1;background:#F8FAFC;'
                f'border-right:1px solid #E6EAF0;border-bottom:1px solid #E6EAF0;">'
                f'{html.escape(w)}</td>{"".join(wcells)}</tr>'
            )
    sum_tds = []
    for s in sum_cells:
        sum_tds.append(
            f'<td style="padding:6px 8px;text-align:right;font-weight:700;font-size:13px;'
            f'line-height:1.35;background:#E2E8F0;color:#0F172A;white-space:nowrap;'
            f'border-top:1px solid #CBD5E1;">{_fmt(s)}</td>'
        )
    body_parts.append(
        f'<tr style="font-weight:700;">'
        f'<td style="padding:6px 8px;font-size:13px;line-height:1.35;text-align:left;'
        f'position:sticky;left:0;z-index:1;background:#E2E8F0;color:#0F172A;'
        f'border-right:1px solid #CBD5E1;border-top:1px solid #CBD5E1;white-space:nowrap;">'
        f'연간 합계</td>{"".join(sum_tds)}</tr>'
    )

    # 맥: fixed+ellipsis 유지 / iPad: auto + 가로스크롤로 숫자·연도 전부 표시
    if touch:
        table_css = (
            "border-collapse:separate;border-spacing:0;width:max-content;min-width:100%;"
            "table-layout:auto;font-family:inherit;font-size:13px;"
        )
        cell_overflow_css = "th,td{overflow:visible;text-overflow:clip;}"
        wrap_overflow = "overflow:auto;-webkit-overflow-scrolling:touch;"
    else:
        table_css = (
            "border-collapse:separate;border-spacing:0;width:100%;"
            "table-layout:fixed;font-family:inherit;font-size:13px;"
        )
        cell_overflow_css = "th,td{overflow:hidden;text-overflow:ellipsis;}"
        wrap_overflow = "overflow:auto;-webkit-overflow-scrolling:touch;"

    page_html = f"""
    <!DOCTYPE html>
    <html><head><meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <style>
      html,body{{
        margin:0;height:100%;overflow:hidden;
        font-family:"Source Sans Pro","Source Sans 3","Segoe UI",Roboto,"Helvetica Neue",Arial,sans-serif;
        font-size:13px;color:#31333F;background:#fff;
        -webkit-font-smoothing:antialiased;
      }}
      .hint{{
        padding:6px 10px;font-size:12px;color:#64748B;background:#F8FAFC;
        border:1px solid #E2E8F0;border-bottom:none;border-radius:0.5rem 0.5rem 0 0;
        font-family:inherit;
      }}
      .wrap{{
        height:calc(100% - 34px);{wrap_overflow}
        border:1px solid #E2E8F0;border-radius:0 0 0.5rem 0.5rem;background:#fff;
      }}
      table{{{table_css}}}
      {cell_overflow_css}
      .mrow:hover td{{filter:brightness(0.97);}}
      .mrow.open .chev{{display:inline-block;transform:rotate(90deg);}}
      .chev{{display:inline-block;width:1em;transition:transform .12s ease;color:#64748B;}}
    </style></head><body>
    <div class="hint">월 행 클릭 → 1~4주 펼침 · 다시 클릭 → 접기 · 1주=1~7일 · 2주=8~14일 · 3주=15~21일 · 4주=22~말일</div>
    <div class="wrap">
      <table>
        <thead><tr>
          <th style="position:sticky;top:0;left:0;z-index:3;background:#F0F2F6;padding:6px 8px;
            border-bottom:1px solid #D0D7DE;border-right:1px solid #D0D7DE;width:4.5rem;
            text-align:left;font-weight:600;font-size:13px;color:#31333F;line-height:1.35;">월</th>
          {th}
        </tr></thead>
        <tbody>{"".join(body_parts)}</tbody>
      </table>
    </div>
    <script>
    (function(){{
      document.querySelectorAll('tr.mrow').forEach(function(row){{
        row.addEventListener('click', function(){{
          var m = row.getAttribute('data-month');
          var open = row.classList.toggle('open');
          document.querySelectorAll('tr.wrow[data-parent="'+m+'"]').forEach(function(wr){{
            wr.style.display = open ? 'table-row' : 'none';
          }});
        }});
      }});
    }})();
    </script>
    </body></html>
    """
    components.html(page_html, height=height, scrolling=False)



@st.cache_data
def cached_get_industry_pivot(data_df, industry_name, metric, all_months, years):
    if data_df.empty:
        return pd.DataFrame(0, index=all_months, columns=[str(y) for y in years])
        
    df_ind = data_df[data_df["업종"] == industry_name]
    if df_ind.empty:
        return pd.DataFrame(0, index=all_months, columns=[str(y) for y in years])
        
    if metric == "매출액 (만원)":
        pvt = df_ind.pivot_table(index="월", columns="연도", values="매출액", aggfunc="sum").fillna(0) * 1.1 / 10000
    elif metric == "출고량":
        pvt = df_ind.pivot_table(index="월", columns="연도", values="출고량", aggfunc="sum").fillna(0) / 1000
    elif metric == "총매출 대비 비중 (%)":
        pvt_ind = df_ind.pivot_table(index="월", columns="연도", values="매출액", aggfunc="sum").fillna(0)
        pvt_total = data_df.pivot_table(index="월", columns="연도", values="매출액", aggfunc="sum").fillna(0)
        pvt = (pvt_ind / pvt_total.replace(0, np.nan) * 100).fillna(0)
        
    pvt = pvt.reindex(index=all_months, fill_value=0)
    all_yrs = [str(y) for y in years]
    pvt = pvt.reindex(columns=all_yrs, fill_value=0)
    return pvt
@st.cache_data
def cached_tab3_pivots(target_tab3_df, years, all_months):
    sales_p = pd.DataFrame()
    qty_p = pd.DataFrame()
    unit_price_p = pd.DataFrame()
    
    if target_tab3_df.empty:
        return sales_p, qty_p, unit_price_p
    sales_raw_p = target_tab3_df.pivot_table(index="품목명", columns="연도월_정렬", values="매출액", aggfunc="sum").fillna(0)
    qty_raw_p = target_tab3_df.pivot_table(index="품목명", columns="연도월_정렬", values="출고량", aggfunc="sum").fillna(0)
    sales_expanded_data = {}
    qty_expanded_data = {}
    for yr in years:
        yr_short = yr[2:]
        yr_sales_sum = 0
        yr_qty_sum = 0
        for m in all_months:
            col_key = f"{yr_short}년 {m}"
            s_val = sales_raw_p[col_key] if col_key in sales_raw_p.columns else 0
            q_val = qty_raw_p[col_key] if col_key in qty_raw_p.columns else 0
            sales_expanded_data[col_key] = s_val
            qty_expanded_data[col_key] = q_val
            yr_sales_sum += s_val
            yr_qty_sum += q_val
        sales_expanded_data[f"{yr_short}년 연간총합"] = yr_sales_sum
        qty_expanded_data[f"{yr_short}년 연간총합"] = yr_qty_sum
    sales_p = pd.DataFrame(sales_expanded_data, index=sales_raw_p.index)
    qty_p = pd.DataFrame(qty_expanded_data, index=qty_raw_p.index)
    latest_col = None
    for yr in years:
        for m in reversed(all_months):
            chk_c = f"{yr[2:]}년 {m}"
            if chk_c in qty_p.columns and qty_p[chk_c].sum() > 0:
                latest_col = chk_c
                break
        if latest_col: break
    if not latest_col and len(qty_p.columns) > 0:
        latest_col = qty_p.columns[0]
    if latest_col and latest_col in qty_p.columns:
        qty_p = qty_p.sort_values(by=latest_col, ascending=False)
        sales_p = sales_p.reindex(qty_p.index)
    df_for_price = target_tab3_df.copy()
    target_clients = ["두산판금", "드림맥", "모베이스전전", "지엔티테크", "태광기업", "경민산업", "동주산업"]
    mask_n2_special = (df_for_price["거래처"].isin(target_clients)) & (df_for_price["품목명"] == "N2 (kg, Bulk)")
    
    df_for_price.loc[mask_n2_special, "단가"] = df_for_price.loc[mask_n2_special, "단가"] * 1.238
    raw_up = df_for_price.pivot_table(index="품목명", columns="연도월_정렬", values="단가", aggfunc=get_exact_original_price)
    
    if not raw_up.empty:
        unit_price_p = raw_up.fillna(0)
        
        desired_price_cols = [f"{yr[2:]}년 {m}" for yr in years for m in reversed(all_months)]
        existing_cols = [c for c in desired_price_cols if c in unit_price_p.columns]
        unit_price_p = unit_price_p[existing_cols]
        unit_price_p = apply_forward_unit_price(unit_price_p, qty_p, years, all_months)
        
    return sales_p, qty_p, unit_price_p


# —— Tab3 4️⃣ 거래처 가스 사용량·재고관리 (신규 전용, 타 탭 미사용) ——
_TAB3_NON_GAS_ITEM_RE = re.compile(
    r"입금|이월|단가차액|잔액정리|기화기|공사|작업비|임대|운반비|회수|보증|수수료|운임"
)


def _tab3_is_bulk_item(name):
    s = str(name or "")
    return ("BULK" in s.upper()) or ("벌크" in s)


def _tab3_is_gas_item(name):
    s = str(name or "").strip()
    if not s or s.lower() in ("nan", "none", "-"):
        return False
    if _TAB3_NON_GAS_ITEM_RE.search(s):
        return False
    return True


@st.cache_data(show_spinner=False)
def cached_tab3_client_gas_usage(df_client, years_tuple, months_tuple):
    """거래처 납품(출고)량 → 벌크/그외 월·주·일 사용량 요약.

    years_tuple: ('2025','2026') 등
    months_tuple: ('01월',…) — 비우면 전체 월
    """
    empty = (
        pd.DataFrame(),
        pd.DataFrame(),
        {"bulk_month": 0.0, "other_month": 0.0, "n_months": 0, "n_days": 0},
    )
    if df_client is None or df_client.empty:
        return empty
    need = {"품목명", "출고량", "매출일_dt", "연도", "월"}
    if not need.issubset(set(df_client.columns)):
        return empty
    yrs = [str(y) for y in (years_tuple or ()) if y]
    mons = [str(m) for m in (months_tuple or ()) if m]
    d = df_client.copy()
    d["연도"] = d["연도"].astype(str)
    d = d[d["연도"].isin(yrs)]
    if mons:
        d = d[d["월"].astype(str).isin(mons)]
    d["출고량"] = pd.to_numeric(d["출고량"], errors="coerce").fillna(0.0)
    d = d[d["출고량"] > 0]
    d = d[d["품목명"].map(_tab3_is_gas_item)]
    if d.empty:
        return empty
    d["구분"] = d["품목명"].map(
        lambda x: "벌크(주요)" if _tab3_is_bulk_item(x) else "그외가스(부품목)"
    )
    d["매출일_dt"] = pd.to_datetime(d["매출일_dt"], errors="coerce")
    d = d[d["매출일_dt"].notna()]
    if d.empty:
        return empty
    # 기준기간 월수: 선택 연·월의 달력 기준(미래 월 제외). 납품 없는 달도 포함 → 평균 사용량 안정화
    latest_dt = d["매출일_dt"].max()
    latest_y, latest_m = int(latest_dt.year), int(latest_dt.month)
    cal_months = []
    for y in yrs:
        try:
            yi = int(y)
        except Exception:
            continue
        month_list = mons if mons else [f"{i:02d}월" for i in range(1, 13)]
        for mlab in month_list:
            try:
                mi = int(str(mlab).replace("월", "").strip())
            except Exception:
                continue
            if yi > latest_y or (yi == latest_y and mi > latest_m):
                continue
            cal_months.append((yi, mi))
    n_months = max(len(cal_months), 1)
    if cal_months:
        y0, m0 = min(cal_months)
        y1, m1 = max(cal_months)
        start = pd.Timestamp(year=y0, month=m0, day=1)
        end = pd.Timestamp(year=y1, month=m1, day=1) + pd.offsets.MonthEnd(0)
        n_days = max(int((end - start).days) + 1, 1)
    else:
        n_days = 30

    rows = []
    for item, g in d.groupby("품목명", sort=False):
        tot = float(g["출고량"].sum())
        n_deliv = int(len(g))
        active_m = int(g["매출일_dt"].dt.to_period("M").nunique())
        # 월사용량: 총납품 ÷ 기준기간 달력 월수
        month_u = tot / n_months
        week_u = month_u * 7.0 / 30.0
        day_u = month_u / 30.0
        dates = g["매출일_dt"].sort_values()
        last_dt = dates.iloc[-1]
        if len(dates) >= 2:
            gaps = dates.diff().dt.days.dropna()
            avg_gap = float(gaps.mean()) if not gaps.empty else 0.0
        else:
            avg_gap = 0.0
        avg_per = tot / n_deliv if n_deliv else 0.0
        # 예상 재고소진(일): 최근 회당 납품량 ÷ 일사용량
        cover_days = (avg_per / day_u) if day_u > 0 else 0.0
        rows.append(
            {
                "품목명": item,
                "구분": "벌크(주요)" if _tab3_is_bulk_item(item) else "그외가스(부품목)",
                "총납품량": tot,
                "납품횟수": n_deliv,
                "활성월수": active_m,
                "월사용량": month_u,
                "주사용량": week_u,
                "일사용량": day_u,
                "회당평균": avg_per,
                "평균납품간격(일)": avg_gap,
                "예상소진(일)": cover_days,
                "최근납품일": last_dt.strftime("%Y-%m-%d"),
            }
        )
    summary = pd.DataFrame(rows)
    if not summary.empty:
        summary = summary.sort_values(
            by=["구분", "월사용량"], ascending=[True, False]
        ).reset_index(drop=True)

    # 월별 피벗 (구분 합산용 원천)
    d["연월"] = d["매출일_dt"].dt.strftime("%y년 %m월")
    monthly = (
        d.pivot_table(
            index="품목명", columns="연월", values="출고량", aggfunc="sum"
        )
        .fillna(0)
    )
    # 열 시간순
    if not monthly.empty:
        def _ym_key(c):
            m = re.match(r"(\d{2})년\s*(\d{2})월", str(c))
            return (int(m.group(1)), int(m.group(2))) if m else (0, 0)

        monthly = monthly.reindex(columns=sorted(monthly.columns, key=_ym_key))

    meta = {
        "bulk_month": float(
            summary.loc[summary["구분"] == "벌크(주요)", "월사용량"].sum()
        )
        if not summary.empty
        else 0.0,
        "other_month": float(
            summary.loc[summary["구분"] == "그외가스(부품목)", "월사용량"].sum()
        )
        if not summary.empty
        else 0.0,
        "n_months": n_months,
        "n_days": n_days,
    }
    return summary, monthly, meta


def _tab3_usage_inv_card_html(title, subtitle, month_u, week_u, day_u, accent="#1D4ED8"):
    """인라인 카드 HTML (탭3 전용, 공유 CSS 미수정)."""
    return (
        f'<div style="background:#fff;border:1px solid #E2E8F0;border-left:4px solid {accent};'
        f'border-radius:10px;padding:14px 16px;margin-bottom:10px;'
        f'box-shadow:0 2px 6px rgba(15,23,42,.04);">'
        f'<div style="font-size:13px;font-weight:700;color:#0F172A;">{html.escape(title)}</div>'
        f'<div style="font-size:11px;color:#64748B;margin:2px 0 10px;">{html.escape(subtitle)}</div>'
        f'<div style="display:flex;gap:12px;flex-wrap:wrap;">'
        f'<div><div style="font-size:11px;color:#94A3B8;">월사용량</div>'
        f'<div style="font-size:18px;font-weight:700;color:{accent};">{month_u:,.1f}</div></div>'
        f'<div><div style="font-size:11px;color:#94A3B8;">주사용량</div>'
        f'<div style="font-size:18px;font-weight:700;color:#0F172A;">{week_u:,.1f}</div></div>'
        f'<div><div style="font-size:11px;color:#94A3B8;">일사용량</div>'
        f'<div style="font-size:18px;font-weight:700;color:#0F172A;">{day_u:,.1f}</div></div>'
        f"</div></div>"
    )


@st.cache_data
def cached_filter_tab3_year_columns(sales_p, qty_p, selected_detail_years, avail_years_short, all_months):
    """Tab3 연도별 컬럼 필터 (캐시로 탭 전환·재렌더 가속)."""
    def _filter_year_columns(df):
        if df.empty:
            return df
        cols = []
        for y in avail_years_short:
            if y in selected_detail_years:
                for m in reversed(all_months):
                    c = f"{y}년 {m}"
                    if c in df.columns:
                        cols.append(c)
            tot = f"{y}년 연간총합"
            if tot in df.columns:
                cols.append(tot)
        return df[[c for c in cols if c in df.columns]]
    sales_vat = sales_p * 1.1 / 10000
    return _filter_year_columns(sales_vat), _filter_year_columns(qty_p)
@st.cache_data
def cached_tab3_item_client_detail(
    df, item_name, metric, years, all_months, selected_years_short
):
    """품목 1개 → 거래처×월 상세 피벗 (sales=만원 VAT포함 / qty=출고량)."""
    if df is None or df.empty or not item_name:
        return pd.DataFrame()
    d = df[df["품목명"] == item_name]
    if d.empty:
        return pd.DataFrame()
    if metric == "sales":
        raw = (
            d.pivot_table(index="거래처", columns="연도월_정렬", values="매출액", aggfunc="sum")
            .fillna(0)
            * 1.1
            / 10000
        )
    else:
        raw = d.pivot_table(
            index="거래처", columns="연도월_정렬", values="출고량", aggfunc="sum"
        ).fillna(0)
        bulk = ["CO2 (kg, Bulk)", "N2 (kg, Bulk)", "O2 (kg, Bulk)", "AR (kg, Bulk)"]
        if item_name in bulk:
            raw = raw / 1000
    if raw.empty:
        return pd.DataFrame()
    expanded = {}
    for yr in years:
        yr_short = yr[2:]
        yr_sum = 0
        for m in all_months:
            col_key = f"{yr_short}년 {m}"
            val = raw[col_key] if col_key in raw.columns else 0
            expanded[col_key] = val
            yr_sum = yr_sum + val
        expanded[f"{yr_short}년 연간총합"] = yr_sum
    pvt = pd.DataFrame(expanded, index=raw.index)
    avail_years_short = [y[2:] for y in years]
    cols = []
    for y in avail_years_short:
        if y in selected_years_short:
            for m in reversed(all_months):
                c = f"{y}년 {m}"
                if c in pvt.columns:
                    cols.append(c)
        tot = f"{y}년 연간총합"
        if tot in pvt.columns:
            cols.append(tot)
    pvt = pvt[[c for c in cols if c in pvt.columns]]
    if pvt.empty:
        return pvt
    sort_col = None
    for c in pvt.columns:
        if "연간총합" in str(c):
            sort_col = c
            break
    if sort_col is None and len(pvt.columns):
        sort_col = pvt.columns[0]
    if sort_col:
        pvt = pvt.sort_values(by=sort_col, ascending=False)
    return pvt
def render_tab3_item_client_expanders(df_src, items, metric, years, all_months, selected_years_short):
    """상단에서 선택한 품목의 거래처별 상세 — 클릭 펼침 (매출/출고 각각)."""
    if not items:
        return
    label_metric = "매출" if metric == "sales" else "출고량"
    cmap = "Blues" if metric == "sales" else "Greens"
    bulk = ("CO2 (kg, Bulk)", "N2 (kg, Bulk)", "O2 (kg, Bulk)", "AR (kg, Bulk)")
    for item in items:
        with st.expander(
            f"📂 [{item}] 거래처별 {label_metric} 상세 데이터 파보기 (클릭하여 펼치기)",
            expanded=False,
        ):
            pvt = cached_tab3_item_client_detail(
                df_src,
                item,
                metric,
                tuple(years),
                tuple(all_months),
                tuple(sorted(selected_years_short)),
            )
            if pvt is None or pvt.empty:
                st.info(f"선택한 조건에서 [{item}] 거래처 상세 데이터가 없습니다.")
                continue
            pvt_disp = get_display_df_with_sum(pvt, "합계")
            fmt = "{:,.1f}" if metric == "qty" and item in bulk else "{:,.0f}"
            st.dataframe(
                style_with_sum(pvt_disp, fmt, cmap, axis=None),
                use_container_width=True,
                height=min(420, 80 + 28 * min(len(pvt_disp), 12)),
            )
# ==========================================
# ★ 누락되었던 필수 캐시 함수 복구 완료 ★
# ==========================================
@st.cache_data
def cached_client_item_qty_pivot(df_client_filtered, years, all_months):
    if df_client_filtered.empty:
        return pd.DataFrame()
    raw_ci_qty = df_client_filtered.pivot_table(
        index="품목명", columns="연도월_정렬", values="출고량", aggfunc="sum"
    ).fillna(0)
    ci_expanded_data = {}
    for yr in years:
        yr_short = yr[2:]
        for m in all_months:
            col_key = f"{yr_short}년 {m}"
            q_val = raw_ci_qty[col_key] if col_key in raw_ci_qty.columns else 0
            ci_expanded_data[col_key] = q_val
    return pd.DataFrame(ci_expanded_data, index=raw_ci_qty.index)
# ==========================================
# ★ 초고속 빈 껍데기 필터링 로직 ★
# ==========================================
def prepare_active_df_fast(df, target_col):
    if df is None or df.empty:
        return None, [], None
    
    df_active = df.loc[(df != 0).any(axis=1), (df != 0).any(axis=0)].copy()
    
    if df_active.empty:
        return None, [], None
    
    df_active.index.name = None
    df_active = df_active.reset_index().rename(columns={"index": "품목명"})
    df_active.columns.name = None 
    numeric_cols = [c for c in df_active.columns if c != "품목명"]
    highlight_col_name = None
    if target_col and target_col in df_active.columns:
        highlight_col_name = target_col
        
    return df_active, numeric_cols, highlight_col_name
@st.cache_data
def cached_prepare_active_df(df, target_col):
    """Tab3 표 렌더용 활성 행·열 필터 (캐시)."""
    return prepare_active_df_fast(df, target_col)
DEBT_PINK_SOFT = "#FFE4E6"
# 업체별 교차 구분색: 밝게 ↔ 다음 업체는 흐린 연한색 (금액 히트맵 그라디언트 아님)
DEBT_CLIENT_STRIPE_A = "#DCEEFF"
DEBT_CLIENT_STRIPE_B = "#F1F4F7"
DEBT_CLIENT_SEP = ""  # 구분선 없이 교차 배경색으로 거래처 구분
# 채권 히트맵(미사용·호환용) — 스크린샷 상한까지만
_DEBT_CMAP_BLUE = LinearSegmentedColormap.from_list(
    "debt_blue_cap", ["#FFFFFF", "#DEEBF7", "#9DC3E6", "#5B9BD5"]
)
_DEBT_CMAP_GREEN = LinearSegmentedColormap.from_list(
    "debt_green_cap", ["#FFFFFF", "#E2EFDA", "#A9D08E", "#70AD47"]
)


def _debt_heat_cell_css(val, vmin, vmax, family="blue"):
    """홀수거래처=파란 · 짝수=초록. 0=흰색 · 진하기 상한=스크린샷 수준."""
    try:
        v = abs(float(val))
    except (TypeError, ValueError):
        return "background-color:#FFFFFF;color:#31333F;"
    if vmax <= vmin or v < 1e-12:
        return "background-color:#FFFFFF;color:#31333F;"
    vs = float(np.sqrt(v))
    vmax_s = float(np.sqrt(max(vmax, 0.0)))
    t = float(max(0.0, min(1.0, vs / vmax_s))) if vmax_s > 0 else 0.0
    cmap = _DEBT_CMAP_BLUE if family != "green" else _DEBT_CMAP_GREEN
    try:
        r, g, b, _a = cmap(t)
        bg = f"#{int(round(r * 255)):02x}{int(round(g * 255)):02x}{int(round(b * 255)):02x}"
    except Exception:
        bg = "#FFFFFF"
    # 상한 색이 밝아 글자는 항상 진한 회색
    return f"background-color:{bg};color:#31333F;"


def _debt_blues_cell_css(val, vmin, vmax):
    """하위 호환: 파란 계열."""
    return _debt_heat_cell_css(val, vmin, vmax, family="blue")


def _debt_gubun_vmax(abs_vals):
    """구분별 스케일: 양의 값 95퍼센타일(없으면 max)."""
    pos = abs_vals[abs_vals > 1e-12]
    if pos.size == 0:
        return 0.0
    try:
        return float(max(np.percentile(pos, 95), pos.min()))
    except Exception:
        return float(np.nanmax(pos))


def _debt_client_color_family_map(clients_unique):
    """하위 호환: 홀수=파란, 짝수=초록."""
    out = {}
    n = 0
    for client in clients_unique:
        if client == "📌 [전체 합계]":
            out[client] = "blue"
            continue
        n += 1
        out[client] = "blue" if (n % 2 == 1) else "green"
    return out


def _debt_client_stripe_map(clients_unique):
    """거래처 교차 배경: 홀수=밝은색, 짝수=흐린 연한색."""
    out = {}
    n = 0
    for client in clients_unique:
        if client == "📌 [전체 합계]":
            out[client] = "#E2E8F0"
            continue
        n += 1
        out[client] = DEBT_CLIENT_STRIPE_A if (n % 2 == 1) else DEBT_CLIENT_STRIPE_B
    return out


def _debt_family_tint(family="blue", strength=0.30):
    """거래처/구분/결제/연체 열용 연한 배경 (같은 계열 그라디언트)."""
    try:
        t = float(max(0.0, min(1.0, strength)))
        cmap = _DEBT_CMAP_BLUE if family != "green" else _DEBT_CMAP_GREEN
        r, g, b, _a = cmap(t)
        return f"#{int(round(r * 255)):02x}{int(round(g * 255)):02x}{int(round(b * 255)):02x}"
    except Exception:
        return "#DEEBF7" if family != "green" else "#E2EFDA"


def payment_term_credit_months(term):
    """결제조건 → 정상채권으로 인정하는 최대 경과개월(age).
    age0=당월. 익월말 → age<=1(당월·전월) 정상, age>=2 연체.
    """
    if not term or not str(term).strip():
        return 1
    t = str(term).replace(" ", "").replace("　", "")
    if "선매출" in t:
        return 99
    if "당월" in t or "발행후" in t:
        return 0
    m = re.search(r"(\d+)일", t)
    if m:
        days = int(m.group(1))
        # 40일 ≈ 다다음달 10일 결제 → 2개월 유예 / 60일+ → 2개월 유지
        if days >= 40:
            return 2
        return 1
    if "2개월" in t or "이개월" in t:
        return 2
    if "익월말60" in t:
        return 2
    if "bulk-2" in t.lower() or "bulk-2개월" in t or "/bulk" in t.lower():
        return 2
    if "익월" in t:
        return 1
    return 1
def analyze_client_ar(sales_vals, bal_vals, credit_months, current_idx):
    """당월 잔액 배분 → (정상액, 연체액, 연체 매출열 set, 연체개월수).
    연체개월수 = 결제조건 초과 경과개월의 최대값. 잔액0 → 0.
    """
    try:
        cur_bal = float(bal_vals[current_idx]) if pd.notna(bal_vals[current_idx]) else 0.0
    except (TypeError, ValueError, IndexError):
        cur_bal = 0.0
    if cur_bal <= 0:
        return 0.0, 0.0, set(), 0
    remaining = cur_bal
    normal = 0.0
    overdue = 0.0
    pink_cols = set()
    overdue_months = 0
    for c in range(current_idx, -1, -1):
        if remaining <= 0:
            break
        try:
            sal = float(sales_vals[c]) if pd.notna(sales_vals[c]) else 0.0
        except (TypeError, ValueError):
            sal = 0.0
        sal = max(sal, 0.0)
        portion = min(sal, remaining) if sal > 0 else 0.0
        if portion > 0:
            age = current_idx - c
            if age <= credit_months:
                normal += portion
            else:
                overdue += portion
                pink_cols.add(c)
                overdue_months = max(overdue_months, age - credit_months)
            remaining -= portion
    if remaining > 0:
        overdue += remaining
        overdue_months = max(overdue_months, max(1, current_idx - credit_months))
    return normal, overdue, pink_cols, overdue_months
def format_debt_status_label(overdue_amt, overdue_months):
    """정상 / 연체 1~3개월 / 악성(연체 4개월+)."""
    if overdue_amt <= 0 or overdue_months <= 0:
        return "정상"
    if overdue_months >= 4:
        return "악성"
    return f"연체 {int(overdue_months)}개월"
@st.cache_data(show_spinner=False, max_entries=64)
def apply_debt_style_fast(df, highlight_debt=True, payment_terms_map=None):
    """업체별 교차 구분색(밝음/흐린연한) + 연체 분홍. 금액 히트맵 그라디언트 없음."""
    styles = np.full(df.shape, "", dtype=object)
    terms_map = payment_terms_map or {}
    clients = df.index.get_level_values("거래처")
    gubuns = df.index.get_level_values("구분")
    u_clients_fast = clients.unique()
    is_total = clients.to_numpy() == "📌 [전체 합계]"
    stripe_map = _debt_client_stripe_map(u_clients_fast)
    client_arr = clients.to_numpy()

    # 1) 합계=회색 · 그 외=업체 교차색 (금액 그라디언트 없음)
    for r in range(df.shape[0]):
        if is_total[r]:
            styles[r, :] = "background-color:#E2E8F0;font-weight:700;color:#0F172A;"
        else:
            bg = stripe_map.get(client_arr[r], DEBT_CLIENT_STRIPE_A)
            styles[r, :] = f"background-color:{bg};color:#31333F;"

    def apply_pink_cell(old_style, pink_bg=DEBT_PINK_SOFT):
        s = re.sub(r"background-color:\s*#[0-9a-fA-F]+;", "", old_style or "")
        s = re.sub(r"color:\s*#[0-9a-fA-F]+;", "", s)
        s = re.sub(r"font-weight:\s*\d+;", "", s)
        return s + f" background-color: {pink_bg};color:#9F1239;"

    if df.shape[1] == 0:
        return pd.DataFrame(styles, index=df.index, columns=df.columns)
    current_month_idx = df.shape[1] - 1
    for client in u_clients_fast:
        if client == "📌 [전체 합계]":
            continue
        client_rows = np.where(clients == client)[0]
        i_sal = -1
        i_bal = -1
        for r in client_rows:
            if gubuns[r] == "매출":
                i_sal = r
            elif gubuns[r] == "잔액":
                i_bal = r
        if i_sal == -1 or i_bal == -1:
            continue
        term = resolve_payment_term(client, terms_map)
        credit = payment_term_credit_months(term)
        sales_vals = [df.iat[i_sal, c] for c in range(df.shape[1])]
        bal_vals = [df.iat[i_bal, c] for c in range(df.shape[1])]
        _normal, overdue, pink_cols, _od_m = analyze_client_ar(
            sales_vals, bal_vals, credit, current_month_idx
        )
        if overdue > 0:
            styles[i_bal, current_month_idx] = apply_pink_cell(
                styles[i_bal, current_month_idx], DEBT_PINK_SOFT
            )
            for c in pink_cols:
                styles[i_sal, c] = apply_pink_cell(styles[i_sal, c], DEBT_PINK_SOFT)
    return pd.DataFrame(styles, index=df.index, columns=df.columns)
@st.cache_data(show_spinner=False, max_entries=64)
def compute_debt_status_by_client(disp_debt, payment_terms_map=None):
    """거래처 → 채권구분 라벨 (정상 / 연체 1~3개월 / 악성)."""
    terms_map = payment_terms_map or {}
    if disp_debt is None or disp_debt.empty:
        return {}
    clients = disp_debt.index.get_level_values("거래처")
    gubuns = disp_debt.index.get_level_values("구분")
    current_idx = disp_debt.shape[1] - 1
    out = {}
    for client in clients.unique():
        if client == "📌 [전체 합계]":
            continue
        rows = np.where(clients == client)[0]
        i_sal = i_bal = -1
        for r in rows:
            if gubuns[r] == "매출":
                i_sal = r
            elif gubuns[r] == "잔액":
                i_bal = r
        if i_sal < 0 or i_bal < 0 or current_idx < 0:
            out[client] = "정상"
            continue
        term = resolve_payment_term(client, terms_map)
        credit = payment_term_credit_months(term)
        sales_vals = [disp_debt.iat[i_sal, c] for c in range(disp_debt.shape[1])]
        bal_vals = [disp_debt.iat[i_bal, c] for c in range(disp_debt.shape[1])]
        _n, overdue, _p, od_m = analyze_client_ar(sales_vals, bal_vals, credit, current_idx)
        out[client] = format_debt_status_label(overdue, od_m)
    return out
def _debt_label_cell_style(client, gubun, color_map, compact=False):
    """채권표 거래처/구분 — 가독성 우선(14px·선명 대비)."""
    pad = "5px 5px" if compact else "7px 9px"
    base = (
        f"padding:{pad};border-bottom:1px solid #CBD5E1;white-space:nowrap;"
        "font-size:14px;font-weight:500;line-height:1.45;vertical-align:middle;"
        "font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Roboto,sans-serif;"
        "overflow:hidden;text-overflow:ellipsis;color:#0F172A;"
    )
    if client == "📌 [전체 합계]":
        return base + "background-color:#E2E8F0;font-weight:700;text-align:center;"
    bg = color_map.get(client, "#FFFFFF")
    return base + f"background-color:{bg};text-align:center;"
def _debt_num_cell_font(compact=False):
    """채권표 숫자 — 거래처명보다 한 단계 작게, 대비 강화."""
    pad = "5px 5px" if compact else "7px 9px"
    return (
        f"padding:{pad};border-bottom:1px solid #CBD5E1;white-space:nowrap;"
        "font-size:13px;font-weight:500;line-height:1.45;vertical-align:middle;text-align:center;"
        "font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Roboto,sans-serif;"
        "font-variant-numeric:tabular-nums;color:#0F172A;"
    )
PAYMENT_TERMS_PATH = os.path.join(CACHE_DIR, "payment_terms.csv")
PAYMENT_TERMS_FALLBACK = os.path.join(os.path.dirname(os.path.abspath(__file__)), "payment_terms.csv")
@st.cache_data(show_spinner=False)
def load_payment_terms_map():
    """거래처별 결제조건 로드 (입금기준표 반영)."""
    path = PAYMENT_TERMS_PATH if os.path.exists(PAYMENT_TERMS_PATH) else PAYMENT_TERMS_FALLBACK
    if not os.path.exists(path):
        return {}
    try:
        df = pd.read_csv(path, encoding="utf-8-sig")
    except Exception:
        try:
            df = pd.read_csv(path, encoding="cp949")
        except Exception:
            return {}
    if "거래처" not in df.columns or "결제조건" not in df.columns:
        return {}
    out = {}
    for _, row in df.iterrows():
        name = str(row["거래처"]).strip()
        term = "" if pd.isna(row["결제조건"]) else str(row["결제조건"]).strip()
        if name:
            out[name] = term
    return out
def resolve_payment_term(client_name, terms_map):
    """정확 매칭 → strip 매칭 → 접두(가스코아산/세진가스텍) 규칙."""
    if not client_name or client_name == "📌 [전체 합계]":
        return ""
    name = str(client_name)
    if name in terms_map:
        return terms_map[name]
    stripped = name.strip()
    if stripped in terms_map:
        return terms_map[stripped]
    for k, v in terms_map.items():
        if k.strip() == stripped:
            return v
    if name.startswith("가스코아산") or name.startswith("Z가스코아산"):
        return terms_map.get("가스코아산.", "익월말현금")
    if name.startswith("세진가스텍"):
        return "2개월"
    return ""
def render_interactive_html_table(
    headers,
    body_html,
    height=420,
    show_sum_popup=False,
    toolbar_hint=None,
    freeze_left_cols=0,
    freeze_left_widths=None,
    freeze_right_header=False,
    freeze_right_widths=None,
):
    hint = toolbar_hint or "셀 클릭 · 방향키 이동 · ⌘/Ctrl+클릭 또는 ⊕다중선택 · 복사 버튼/⌘C"
    sum_controls = ""
    if show_sum_popup:
        sum_controls = """
            <div class="dash-op-bar" id="dashOpBar" title="셀 선택 시 +/− 적용">
                <span class="dash-op-label">연산</span>
                <button type="button" class="dash-op-btn active" id="dashOpAdd">+</button>
                <button type="button" class="dash-op-btn" id="dashOpSub">−</button>
                <button type="button" class="dash-op-btn dash-op-reset" id="dashOpReset">초기화</button>
            </div>
            <div class="dash-sum-inline" id="dashSumInline">
                <span class="dash-sum-label">선택 합계</span>
                <span id="dashSumValue">0</span>
                <span class="dash-sum-meta" id="dashSumCount">0개</span>
            </div>
        """
    freeze_css = ""
    left_freeze_px = 0
    right_freeze_px = 0
    parts = [
        """
        thead th {
            position: sticky;
            top: 0;
            z-index: 5;
            background: #F0F2F6 !important;
            box-shadow: 0 1px 0 #CBD5E1;
        }
        """
    ]
    if freeze_left_cols > 0:
        widths = freeze_left_widths or ([150, 64] + [100] * max(0, freeze_left_cols - 2))
        left = 0
        for i in range(min(freeze_left_cols, len(widths))):
            w = widths[i]
            parts.append(
                f"""
        th.dash-freeze-{i}, td.dash-freeze-{i} {{
            position: sticky;
            left: {left}px;
            min-width: {w}px;
            max-width: {w}px;
            width: {w}px;
            box-shadow: 2px 0 0 #CBD5E1;
            overflow: hidden;
            text-overflow: ellipsis;
        }}
        th.dash-freeze-{i} {{
            top: 0;
            z-index: 8;
            background: #F0F2F6 !important;
        }}
        td.dash-freeze-{i} {{
            z-index: 4;
        }}
                """
            )
            left += w
        left_freeze_px = int(left)
    if freeze_right_header:
        # 우측 2열: r0=연체(right:0), r1=결제(right:r0폭)
        rw = freeze_right_widths or [64, 72]
        r0_w = int(rw[0]) if len(rw) > 0 else 64
        r1_w = int(rw[1]) if len(rw) > 1 else 72
        right_freeze_px = r0_w + r1_w
        parts.append(
            f"""
        th.dash-freeze-r0, td.dash-freeze-r0 {{
            position: sticky;
            right: 0;
            z-index: 6;
            min-width: {r0_w}px;
            max-width: {r0_w}px;
            width: {r0_w}px;
            box-shadow: -1px 0 0 #E2E8F0;
            background-clip: padding-box;
            word-break: keep-all;
        }}
        th.dash-freeze-r1, td.dash-freeze-r1 {{
            position: sticky;
            right: {r0_w}px;
            z-index: 6;
            min-width: {r1_w}px;
            max-width: {r1_w}px;
            width: {r1_w}px;
            box-shadow: -1px 0 0 #E2E8F0;
            background-clip: padding-box;
            word-break: keep-all;
            overflow: hidden;
        }}
        th.dash-freeze-r0, th.dash-freeze-r1 {{
            top: 0;
            z-index: 9;
            background: #F0F2F6 !important;
            color: #31333F;
            font-weight: 600;
        }}
        td.dash-freeze-r0, td.dash-freeze-r1 {{
            z-index: 5;
        }}
            """
        )
    freeze_css = "".join(parts)
    header_cells = []
    n_h = len(headers)
    for i, h in enumerate(headers):
        classes = []
        if i < freeze_left_cols:
            classes.append(f"dash-freeze-{i}")
        if freeze_right_header and n_h >= 2:
            if i == n_h - 1:
                classes.append("dash-freeze-r0")  # 연체개월수
            elif i == n_h - 2:
                classes.append("dash-freeze-r1")  # 결제조건
        freeze_cls = f' class="{" ".join(classes)}"' if classes else ""
        header_cells.append(
            f'<th{freeze_cls} style="padding:6px 8px;border-bottom:1px solid #E2E8F0;'
            f'background:#F0F2F6;text-align:center;vertical-align:middle;'
            f'font-size:13px;font-weight:600;white-space:nowrap;color:#31333F;">'
            f'{html.escape(str(h))}</th>'
        )
    page_html = f"""
    <!DOCTYPE html>
    <html><head><meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <style>
        html, body {{
            margin: 0;
            height: 100%;
            overflow: hidden;
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
            font-size: 13px;
            font-weight: 400;
            color: #31333F;
        }}
        .dash-shell {{
            display: flex;
            flex-direction: column;
            height: 100%;
            min-height: 100%;
        }}
        .dash-table-toolbar {{
            padding: 6px 10px;
            font-size: 12px;
            color: #64748B;
            background: #F8FAFC;
            border: 1px solid #E2E8F0;
            border-bottom: none;
            border-radius: 4px 4px 0 0;
            display: flex;
            gap: 8px;
            align-items: center;
            flex-wrap: wrap;
            flex-shrink: 0;
        }}
        .dash-action-btn, .dash-multi-btn {{
            padding: 4px 10px;
            font-size: 12px;
            border: 1px solid #CBD5E1;
            border-radius: 4px;
            background: #fff;
            cursor: pointer;
            color: #334155;
        }}
        .dash-multi-btn.active {{ background: #DBEAFE; border-color: #2563EB; color: #1D4ED8; }}
        .dash-op-bar {{
            display: inline-flex;
            align-items: center;
            gap: 4px;
            margin-left: 6px;
            padding: 2px 8px;
            border: 1px solid #E2E8F0;
            border-radius: 6px;
            background: #fff;
            flex: 0 0 auto;
        }}
        .dash-op-label {{
            font-size: 11px;
            color: #64748B;
            margin-right: 2px;
        }}
        .dash-op-btn {{
            min-width: 28px;
            padding: 3px 9px;
            font-size: 13px;
            font-weight: 700;
            line-height: 1.2;
            border: 1px solid #CBD5E1;
            border-radius: 4px;
            background: #F8FAFC;
            cursor: pointer;
            color: #334155;
        }}
        .dash-op-btn#dashOpAdd.active {{
            background: #DBEAFE;
            border-color: #2563EB;
            color: #1D4ED8;
        }}
        .dash-op-btn#dashOpSub.active {{
            background: #FEE2E2;
            border-color: #DC2626;
            color: #B91C1C;
        }}
        .dash-op-btn.dash-op-reset {{
            font-size: 11px;
            font-weight: 600;
            min-width: auto;
            padding: 3px 8px;
        }}
        .dash-sum-inline {{
            display: none;
            margin-left: auto;
            align-items: center;
            gap: 8px;
            padding: 4px 10px;
            border-radius: 6px;
            background: #1E293B;
            color: #fff;
            font-size: 13px;
            white-space: nowrap;
        }}
        .dash-sum-inline.show {{ display: inline-flex; }}
        .dash-sum-label {{ opacity: 0.85; font-size: 12px; }}
        #dashSumValue {{ font-size: 15px; font-weight: 700; }}
        .dash-sum-meta {{ opacity: 0.75; font-size: 11px; }}
        .dash-cell-selectable.selected.dash-op-minus {{
            outline-color: #DC2626;
            background-color: #FEE2E2 !important;
        }}
        .wrap {{
            flex: 1 1 auto;
            min-height: 0;
            overflow: auto;
            -webkit-overflow-scrolling: touch;
            overscroll-behavior: contain;
            touch-action: pan-x pan-y;
            border: 1px solid #E2E8F0;
            border-radius: 0 0 4px 4px;
            background: #fff;
        }}
        table {{ width: max-content; min-width: 100%; border-collapse: separate; border-spacing: 0; }}
        th, td {{ font-weight: 400; }}
        .dash-cell-selectable {{ cursor: cell; user-select: none; -webkit-user-select: none; }}
        .dash-cell-selectable.selected {{
            outline: 2px solid #2563EB;
            outline-offset: -2px;
            background-color: #DBEAFE !important;
        }}
        #dashTable tbody td.dash-cell-focus {{
            outline: 2px solid #1D4ED8;
            outline-offset: -2px;
            box-shadow: inset 0 0 0 1px #93C5FD;
        }}
        .wrap:focus {{ outline: none; }}
        .wrap:focus-visible {{ box-shadow: inset 0 0 0 2px #93C5FD; }}
        {freeze_css}
    </style></head><body>
    <div class="dash-shell">
    <div class="dash-table-toolbar">
        <span>{html.escape(hint)}</span>
        <button type="button" class="dash-multi-btn" id="dashMultiBtn">⊕ 다중선택</button>
        <button type="button" class="dash-action-btn" id="dashCopyBtn">📋 선택 복사</button>
        {sum_controls}
    </div>
    <div class="wrap" id="dashWrap" tabindex="0" aria-label="표 키보드 탐색">
        <table id="dashTable">
            <thead><tr>{''.join(header_cells)}</tr></thead>
            <tbody>{body_html}</tbody>
        </table>
    </div>
    </div>
    <script>
    (function() {{
        const showSum = {'true' if show_sum_popup else 'false'};
        const selected = new Set();
        const cellSign = new WeakMap();
        let multiMode = false;
        let opSign = 1; // +1 더하기, -1 빼기
        let focusCell = null;
        const popup = document.getElementById('dashSumInline');
        const sumValue = document.getElementById('dashSumValue');
        const sumCount = document.getElementById('dashSumCount');
        const wrap = document.getElementById('dashWrap');
        const table = document.getElementById('dashTable');
        const opAddBtn = document.getElementById('dashOpAdd');
        const opSubBtn = document.getElementById('dashOpSub');
        const opResetBtn = document.getElementById('dashOpReset');
        // sticky 고정열 폭 — scrollIntoView(nearest)는 가려진 칸을 보이도록 스크롤하지 않음
        const LEFT_FREEZE_PX = {int(left_freeze_px)};
        const RIGHT_FREEZE_PX = {int(right_freeze_px)};
        const HEADER_STICKY_PX = 34;
        function scrollCellIntoView(td) {{
            if (!wrap || !td) return;
            try {{
                const wr = wrap.getBoundingClientRect();
                const tr = td.getBoundingClientRect();
                const padL = LEFT_FREEZE_PX + 6;
                const padR = RIGHT_FREEZE_PX + 6;
                let dx = 0;
                if (tr.left < wr.left + padL) {{
                    dx = tr.left - (wr.left + padL);
                }} else if (tr.right > wr.right - padR) {{
                    dx = tr.right - (wr.right - padR);
                }}
                if (dx) wrap.scrollLeft += dx;
                let dy = 0;
                if (tr.top < wr.top + HEADER_STICKY_PX) {{
                    dy = tr.top - (wr.top + HEADER_STICKY_PX);
                }} else if (tr.bottom > wr.bottom - 4) {{
                    dy = tr.bottom - (wr.bottom - 4);
                }}
                if (dy) wrap.scrollTop += dy;
            }} catch (e) {{
                try {{ td.scrollIntoView({{ block: 'nearest', inline: 'nearest' }}); }} catch (e2) {{}}
            }}
        }}
        function fmt(n) {{
            return Math.round(n).toLocaleString('ko-KR');
        }}
        function setOpSign(s) {{
            opSign = s >= 0 ? 1 : -1;
            if (opAddBtn) opAddBtn.classList.toggle('active', opSign > 0);
            if (opSubBtn) opSubBtn.classList.toggle('active', opSign < 0);
        }}
        function clearSelection() {{
            selected.forEach(el => {{
                el.classList.remove('selected');
                el.classList.remove('dash-op-minus');
            }});
            selected.clear();
            updateSumUI();
        }}
        function applyCellSign(td, sign) {{
            if (!td) return;
            cellSign.set(td, sign >= 0 ? 1 : -1);
            if (sign < 0) td.classList.add('dash-op-minus');
            else td.classList.remove('dash-op-minus');
        }}
        function updateSumUI() {{
            if (!showSum) return;
            let total = 0, count = 0, plusN = 0, minusN = 0;
            selected.forEach(td => {{
                const raw = td.dataset.raw;
                if (raw === undefined || raw === '') return;
                const v = parseFloat(raw);
                if (isNaN(v)) return;
                const s = cellSign.get(td);
                const sign = (s === undefined || s === null) ? 1 : s;
                total += v * sign;
                count += 1;
                if (sign < 0) minusN += 1;
                else plusN += 1;
            }});
            if (count > 0) {{
                if (sumValue) sumValue.textContent = fmt(total);
                if (sumCount) {{
                    let meta = count + '개 숫자 셀';
                    if (minusN > 0) meta += ' · +' + plusN + '/−' + minusN;
                    sumCount.textContent = meta;
                }}
                if (popup) popup.classList.add('show');
            }} else {{
                if (popup) popup.classList.remove('show');
            }}
        }}
        function clearFocusMark() {{
            table.querySelectorAll('td.dash-cell-focus').forEach(el => el.classList.remove('dash-cell-focus'));
        }}
        function focusWrap() {{
            try {{ wrap && wrap.focus({{ preventScroll: true }}); }} catch (e) {{
                try {{ wrap && wrap.focus(); }} catch (e2) {{}}
            }}
        }}
        function selectCell(td, additive) {{
            if (!td) return;
            focusCell = td;
            clearFocusMark();
            td.classList.add('dash-cell-focus');
            if (!additive) {{
                selected.forEach(el => {{
                    el.classList.remove('selected');
                    el.classList.remove('dash-op-minus');
                }});
                selected.clear();
            }}
            if (td.classList.contains('dash-cell-selectable')) {{
                if (additive && td.classList.contains('selected')) {{
                    const cur = cellSign.get(td);
                    const curSign = (cur === undefined || cur === null) ? 1 : cur;
                    if (curSign !== opSign) {{
                        applyCellSign(td, opSign);
                    }} else {{
                        td.classList.remove('selected');
                        td.classList.remove('dash-op-minus');
                        selected.delete(td);
                    }}
                }} else {{
                    td.classList.add('selected');
                    selected.add(td);
                    applyCellSign(td, opSign);
                }}
            }}
            try {{
                scrollCellIntoView(td);
            }} catch (e) {{}}
            updateSumUI();
        }}
        function cellCoord(td) {{
            const tr = td.parentElement;
            const rows = Array.from(table.querySelectorAll('tbody tr'));
            const cells = Array.from(tr.children);
            return {{ r: rows.indexOf(tr), c: cells.indexOf(td), rows: rows }};
        }}
        function moveByArrow(dr, dc, additive) {{
            const rows = Array.from(table.querySelectorAll('tbody tr'));
            if (!rows.length) return;
            let cur = focusCell;
            if (!cur || !table.contains(cur)) {{
                cur = selected.size ? Array.from(selected)[selected.size - 1] : null;
            }}
            if (!cur || !table.contains(cur)) {{
                const first = rows[0].querySelector('td');
                if (first) selectCell(first, false);
                return;
            }}
            const {{ r, c }} = cellCoord(cur);
            const nr = Math.max(0, Math.min(rows.length - 1, r + dr));
            const tds = Array.from(rows[nr].children);
            if (!tds.length) return;
            const nc = Math.max(0, Math.min(tds.length - 1, c + dc));
            selectCell(tds[nc], additive);
        }}
        table.querySelectorAll('tbody td').forEach(td => {{
            td.addEventListener('click', function(e) {{
                const additive = multiMode || e.ctrlKey || e.metaKey;
                selectCell(td, additive);
                focusWrap();
            }});
        }});
        document.getElementById('dashMultiBtn').addEventListener('click', function() {{
            multiMode = !multiMode;
            this.classList.toggle('active', multiMode);
            focusWrap();
        }});
        if (opAddBtn) opAddBtn.addEventListener('click', function() {{
            setOpSign(1);
            focusWrap();
        }});
        if (opSubBtn) opSubBtn.addEventListener('click', function() {{
            setOpSign(-1);
            focusWrap();
        }});
        if (opResetBtn) opResetBtn.addEventListener('click', function() {{
            clearSelection();
            setOpSign(1);
            clearFocusMark();
            focusCell = null;
            focusWrap();
        }});
        if (showSum) setOpSign(1);
        function selectedText() {{
            const rows = new Map();
            selected.forEach(td => {{
                const tr = td.parentElement;
                if (!tr) return;
                const cells = Array.from(tr.querySelectorAll('td'));
                const idx = cells.indexOf(td);
                const key = Array.from(tr.parentElement.children).indexOf(tr);
                if (!rows.has(key)) rows.set(key, {{ order: key, parts: [] }});
                const label = (td.dataset.raw !== undefined && td.dataset.raw !== '')
                    ? td.dataset.raw
                    : (td.textContent || '').trim();
                rows.get(key).parts.push({{ idx: idx, text: label }});
            }});
            return Array.from(rows.values())
                .sort((a, b) => a.order - b.order)
                .map(r => r.parts.sort((a, b) => a.idx - b.idx).map(p => p.text).join('\\t'))
                .join('\\n');
        }}
        async function copySelected() {{
            const text = selectedText();
            const btn = document.getElementById('dashCopyBtn');
            if (!text) {{
                if (btn) {{
                    const prev = btn.textContent;
                    btn.textContent = '선택 없음';
                    setTimeout(() => {{ btn.textContent = prev; }}, 900);
                }}
                return;
            }}
            try {{
                if (navigator.clipboard && navigator.clipboard.writeText) {{
                    await navigator.clipboard.writeText(text);
                }} else {{
                    const ta = document.createElement('textarea');
                    ta.value = text;
                    ta.style.position = 'fixed';
                    ta.style.left = '-9999px';
                    document.body.appendChild(ta);
                    ta.select();
                    document.execCommand('copy');
                    document.body.removeChild(ta);
                }}
                if (btn) {{
                    const prev = btn.textContent;
                    btn.textContent = '✓ 복사됨';
                    setTimeout(() => {{ btn.textContent = prev; }}, 900);
                }}
            }} catch (err) {{
                if (btn) {{
                    const prev = btn.textContent;
                    btn.textContent = '복사 실패';
                    setTimeout(() => {{ btn.textContent = prev; }}, 900);
                }}
            }}
        }}
        const copyBtn = document.getElementById('dashCopyBtn');
        if (copyBtn) copyBtn.addEventListener('click', function(e) {{
            e.preventDefault();
            copySelected();
        }});
        document.addEventListener('keydown', function(e) {{
            const tag = (e.target && e.target.tagName) ? e.target.tagName.toLowerCase() : '';
            if (tag === 'input' || tag === 'textarea' || (e.target && e.target.isContentEditable)) return;
            if ((e.metaKey || e.ctrlKey) && (e.key === 'c' || e.key === 'C')) {{
                if (selected.size > 0) {{
                    e.preventDefault();
                    copySelected();
                }}
                return;
            }}
            const key = e.key;
            if (key !== 'ArrowUp' && key !== 'ArrowDown' && key !== 'ArrowLeft' && key !== 'ArrowRight') return;
            // 표 영역 포커스(또는 셀 선택) 있을 때만 — 사이드바/다른 위젯 방향키 방해 금지
            const inTable = wrap && (wrap === document.activeElement || wrap.contains(document.activeElement)
                || (focusCell && table.contains(focusCell)) || selected.size > 0);
            if (!inTable) return;
            e.preventDefault();
            e.stopPropagation();
            const additive = multiMode || e.shiftKey;
            if (key === 'ArrowUp') moveByArrow(-1, 0, additive);
            else if (key === 'ArrowDown') moveByArrow(1, 0, additive);
            else if (key === 'ArrowLeft') moveByArrow(0, -1, additive);
            else if (key === 'ArrowRight') moveByArrow(0, 1, additive);
            focusWrap();
        }});
    }})();
    </script>
    </body></html>
    """
    components.html(page_html, height=height, scrolling=True)
def render_tab3_dataframe_table(
    df, fmt, target_col, key_prefix="tab3", table_kind="sales", sort_order=None
):
    """Tab3 표 — 연한 파스텔 그라디언트 + 선택 행 복사.
    sort_order: 단가표용 품목명 순서(매출액 내림차순 인덱스).
    """
    df_active, numeric_cols, highlight_col_name = cached_prepare_active_df(df, target_col)
    if df_active is None or df_active.empty:
        return
    annual_cols = [c for c in numeric_cols if "연간총합" in str(c)]
    month_cols = [c for c in numeric_cols if "연간총합" not in str(c)]
    # 정렬: 단가표는 매출액 순(sort_order), 그 외는 해당 표 고실적 순
    if table_kind == "price" and sort_order is not None:
        rank = {name: i for i, name in enumerate(sort_order)}
        df_active = df_active.assign(
            _rk=df_active["품목명"].map(lambda x: rank.get(x, 10**9))
        ).sort_values(by="_rk").drop(columns=["_rk"])
    else:
        sort_col = None
        if table_kind != "price":
            if annual_cols:
                sort_col = annual_cols[0]
            elif highlight_col_name in numeric_cols:
                sort_col = highlight_col_name
            elif month_cols:
                sort_col = month_cols[0]
        if sort_col and sort_col in df_active.columns:
            df_active = df_active.sort_values(by=sort_col, ascending=False, na_position="last")
        elif numeric_cols and table_kind != "price":
            df_active = df_active.assign(_rk=df_active[numeric_cols].sum(axis=1)).sort_values(
                by="_rk", ascending=False
            ).drop(columns=["_rk"])
    styled = df_active.style.format(fmt, subset=numeric_cols)
    # 파스텔 히트맵 — 매출/출고만 (적용단가는 그라디언트 없음)
    if table_kind == "price":
        if highlight_col_name and highlight_col_name in numeric_cols:
            styled = styled.apply(
                lambda s: ["color:#B91C1C;background-color:#FEE2E2;"] * len(s),
                subset=[highlight_col_name],
                axis=0,
            )
        styled = styled.apply(
            lambda col: (
                ["border-right:2px solid #CBD5E1;background-color:#FAFAFA;"] * len(col)
                if col.name == "품목명"
                else [""] * len(col)
            ),
            axis=0,
        )
    else:
        cmap_month = _TAB3_CMAP_BLUE if table_kind == "sales" else _TAB3_CMAP_GREEN
        if month_cols:
            styled = styled.background_gradient(
                cmap=cmap_month, subset=month_cols, axis=0
            )
            styled = styled.set_properties(subset=month_cols, **{"color": "#0F172A"})
        if annual_cols:
            styled = styled.background_gradient(
                cmap=_TAB3_CMAP_ORANGE, subset=annual_cols, axis=0
            )
            styled = styled.set_properties(
                subset=annual_cols,
                **{
                    "font-weight": "700",
                    "border-left": "2px solid #FDBA74",
                    "color": "#0F172A",
                },
            )
        if highlight_col_name and highlight_col_name in numeric_cols:
            styled = styled.apply(
                lambda s: [
                    "color:#9F1239;font-weight:700;background-color:#FFE4E6;"
                    "border-left:2px solid #FB7185;border-right:2px solid #FB7185;"
                ]
                * len(s),
                subset=[highlight_col_name],
                axis=0,
            )
        styled = styled.apply(
            lambda col: (
                [
                    "font-weight:700;background-color:#F8FAFC;color:#0F172A;"
                    "border-right:2px solid #94A3B8;"
                ]
                * len(col)
                if col.name == "품목명"
                else [""] * len(col)
            ),
            axis=0,
        )
    event = st.dataframe(
        styled,
        use_container_width=True,
        height=420,
        hide_index=True,
        key=f"{key_prefix}_grid",
        on_select="rerun",
        selection_mode="multi-row",
        column_config={
            "품목명": st.column_config.TextColumn("품목명", width="medium"),
        },
    )
    rows = []
    try:
        rows = list(event.selection.rows or [])
    except Exception:
        rows = []
    if rows:
        subset = df_active.iloc[rows].copy()
        for c in numeric_cols:
            if c in subset.columns:
                subset[c] = pd.to_numeric(subset[c], errors="coerce").map(
                    lambda x: "" if pd.isna(x) else format(x, ",.0f")
                )
        tsv = subset.to_csv(sep="\t", index=False)
        enc = urllib.parse.quote(tsv, safe="")
        components.html(
            f"""
            <div style="font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;">
              <button id="tab3CopyBtn" type="button"
                style="padding:6px 12px;border:1px solid #CBD5E1;border-radius:6px;
                       background:#fff;color:#334155;font-size:13px;font-weight:600;cursor:pointer;">
                📋 선택 행 복사 ({len(rows)}행)
              </button>
              <span id="tab3CopyMsg" style="margin-left:8px;font-size:12px;color:#64748B;"></span>
            </div>
            <script>
            (function() {{
              var btn = document.getElementById("tab3CopyBtn");
              var msg = document.getElementById("tab3CopyMsg");
              if (!btn) return;
              btn.addEventListener("click", async function() {{
                var text = decodeURIComponent("{enc}");
                try {{
                  if (navigator.clipboard && navigator.clipboard.writeText) {{
                    await navigator.clipboard.writeText(text);
                  }} else {{
                    var ta = document.createElement("textarea");
                    ta.value = text;
                    document.body.appendChild(ta);
                    ta.select();
                    document.execCommand("copy");
                    document.body.removeChild(ta);
                  }}
                  if (msg) msg.textContent = "✓ 복사됨 — 엑셀/메모장에 붙여넣기 하세요";
                  btn.textContent = "✓ 복사됨 ({len(rows)}행)";
                }} catch (e) {{
                  if (msg) msg.textContent = "복사 실패 — 아래 텍스트를 길게 눌러 복사하세요";
                }}
              }});
            }})();
            </script>
            """,
            height=46,
        )
    else:
        st.caption("행을 클릭해 선택한 뒤 「선택 행 복사」로 복사할 수 있습니다. (⌘/Ctrl 클릭으로 여러 행)")
def render_debt_interactive_table(disp_debt, highlight_debt, height=700, payment_terms_map=None):
    """채권관리 표 — 우측: 결제조건 + 연체개월수(정상/연체1~3개월/악성)."""
    terms_map = payment_terms_map or {}
    style_df = apply_debt_style_fast(
        disp_debt, highlight_debt=highlight_debt, payment_terms_map=terms_map
    )
    status_map = compute_debt_status_by_client(disp_debt, terms_map)
    numeric_cols = list(disp_debt.columns)
    clients = disp_debt.index.get_level_values("거래처")
    u_clients = clients.unique()
    # 업체별 교차색 (밝음 / 흐린 연한색) — 금액 그라디언트 없음
    color_map = _debt_client_stripe_map(u_clients)
    long_cnt = sum(1 for v in status_map.values() if v == "악성")
    # 거래처분석 탭과 동일: 13px / 보통 굵기 · 중간정렬 (연체 분홍만 강조)
    compact = is_touch_ui()
    if compact:
        headers = ["거래처", "구분"] + numeric_cols + ["결제", "연체"]
        left_w, right_w = [108, 40], [52, 64]
        hint = "셀 클릭 · 업체 교차색(밝음/연한) · 분홍=연체금액"
    else:
        headers = ["거래처", "구분"] + numeric_cols + ["결제조건", "연체개월수"]
        left_w, right_w = [160, 68], [110, 120]
        hint = "셀 클릭 · 업체별 교차색(밝음↔흐린 연한) · 분홍=연체금액"
    cell_font = _debt_num_cell_font(compact=compact)
    body_rows = []
    prev_client = None
    for r, (idx, row) in enumerate(disp_debt.iterrows()):
        client, gubun = idx[0], idx[1]
        sep = DEBT_CLIENT_SEP if (prev_client is not None and client != prev_client) else ""
        prev_client = client
        cells = [
            f'<td class="dash-cell-selectable dash-freeze-0" style="{_debt_label_cell_style(client, "", color_map, compact=compact)}{sep}" title="{html.escape(str(client))}">{html.escape(str(client))}</td>',
            f'<td class="dash-cell-selectable dash-freeze-1" style="{_debt_label_cell_style(client, gubun, color_map, compact=compact)}{sep}">{html.escape(str(gubun))}</td>',
        ]
        for col in numeric_cols:
            val = row[col]
            try:
                num = float(val) if pd.notna(val) else 0.0
            except (TypeError, ValueError):
                num = 0.0
            display = f"{num:,.0f}"
            # 동일 (거래처,구분) 중복 시 .at 불가 → 행번호(iat)로 접근
            extra_style = ""
            if col in style_df.columns and 0 <= r < len(style_df):
                try:
                    extra_style = str(style_df.iat[r, style_df.columns.get_loc(col)] or "")
                except Exception:
                    extra_style = ""
            cells.append(
                f'<td class="dash-cell-selectable" style="{cell_font}{extra_style}{sep}" '
                f'data-raw="{num}">{display}</td>'
            )
        # 결제조건·연체도 거래처 계열 색 (잔액 행에만 문구)
        if client == "📌 [전체 합계]":
            term_bg = "#E2E8F0"
            if gubun == "잔액":
                term = ""
                status = f"악성 {long_cnt}곳" if long_cnt else "—"
                if compact and long_cnt:
                    status = f"악성{long_cnt}"
            else:
                term, status = "", ""
        else:
            term_bg = color_map.get(client, DEBT_CLIENT_STRIPE_A)
            if gubun == "잔액":
                term = resolve_payment_term(client, terms_map)
                status = status_map.get(client, "정상")
            else:
                term, status = "", ""
        status_disp = status
        if compact and status.startswith("연체 ") and status.endswith("개월"):
            status_disp = status.replace("연체 ", "").replace("개월", "M")
        meta_style = f"{cell_font}background-color:{term_bg};white-space:normal;{sep}"
        term_show = term if gubun == "잔액" else ""
        term_fallback = "—" if (client != "📌 [전체 합계]" and gubun == "잔액" and not term) else ""
        cells.append(
            f'<td class="dash-freeze-r1" style="{meta_style}" title="{html.escape(term)}">'
            f"{html.escape(term_show or term_fallback)}</td>"
        )
        if status == "악성" or (isinstance(status, str) and status.startswith("악성")):
            st_color = "#BE123C"
        elif status == "정상":
            st_color = "#047857"
        elif "연체" in str(status) and "개월" in str(status):
            st_color = "#C2410C"
        else:
            st_color = "#31333F"
        # 연체/악성이면 연체개월수 열에 분홍 배경 우선
        status_bg = term_bg
        if gubun == "잔액" and client != "📌 [전체 합계]":
            if status == "악성" or (isinstance(status, str) and ("연체" in status or status.startswith("악성"))):
                status_bg = DEBT_PINK_SOFT
        status_style = (
            f"{cell_font}background-color:{status_bg};white-space:normal;"
            f"color:{st_color};{sep}"
        )
        if status_bg == DEBT_PINK_SOFT:
            status_style += "font-weight:600;"
        cells.append(
            f'<td class="dash-freeze-r0" style="{status_style}" '
            f'title="{html.escape(status)}">{html.escape(status_disp)}</td>'
        )
        body_rows.append(f"<tr>{''.join(cells)}</tr>")
    render_interactive_html_table(
        headers,
        "".join(body_rows),
        height=height + 40,
        show_sum_popup=True,
        toolbar_hint=hint,
        freeze_left_cols=2,
        freeze_left_widths=left_w,
        freeze_right_header=True,
        freeze_right_widths=right_w,
    )
def build_debt_client_balance_matrix(debt_df, month_cols):
    """거래처 × 월 잔액 피벗 (잔액>0 업체만)."""
    if debt_df is None or debt_df.empty or not month_cols:
        return pd.DataFrame()
    bal = debt_df[debt_df["구분"] == "잔액"].copy()
    if bal.empty:
        return pd.DataFrame()
    cols = [c for c in month_cols if c in bal.columns]
    if not cols:
        return pd.DataFrame()
    mat = bal.groupby("거래처", as_index=True)[cols].sum()
    mat = mat[(mat.fillna(0) > 0).any(axis=1)]
    return mat
def compute_debt_status_from_raw(debt_df, month_cols, payment_terms_map=None):
    """원본 채권 df → 거래처별 연체 라벨 (정상 포함)."""
    meta = compute_debt_od_meta_from_raw(debt_df, month_cols, payment_terms_map)
    return {c: v["label"] for c, v in meta.items()}
@st.cache_data(show_spinner=False, max_entries=32)
def compute_debt_od_meta_from_raw(debt_df, month_cols, payment_terms_map=None):
    """거래처 → {label, months, overdue_amt, pink_months, od_month_amts, cur_bal}."""
    terms_map = payment_terms_map or {}
    if debt_df is None or debt_df.empty or not month_cols:
        return {}
    cols = [c for c in month_cols if c in debt_df.columns]
    if not cols:
        return {}
    out = {}
    for client, g in debt_df.groupby("거래처"):
        sal = g[g["구분"] == "매출"]
        bal = g[g["구분"] == "잔액"]
        if sal.empty or bal.empty:
            out[client] = {
                "label": "정상",
                "months": 0,
                "overdue_amt": 0.0,
                "pink_months": [],
                "od_month_amts": {},
                "cur_bal": 0.0,
            }
            continue
        sales_vals = [float(sal[c].sum()) if c in sal.columns else 0.0 for c in cols]
        bal_vals = [float(bal[c].sum()) if c in bal.columns else 0.0 for c in cols]
        credit = payment_term_credit_months(resolve_payment_term(client, terms_map))
        _n, overdue, pink_cols, od_m = analyze_client_ar(
            sales_vals, bal_vals, credit, len(cols) - 1
        )
        pink_months = [cols[i] for i in sorted(pink_cols)]
        od_month_amts = {
            cols[i]: float(sales_vals[i]) for i in pink_cols if sales_vals[i] > 0
        }
        try:
            cur_bal = float(bal_vals[-1]) if pd.notna(bal_vals[-1]) else 0.0
        except (TypeError, ValueError, IndexError):
            cur_bal = 0.0
        out[client] = {
            "label": format_debt_status_label(overdue, od_m),
            "months": int(od_m) if overdue > 0 else 0,
            "overdue_amt": float(overdue),
            "pink_months": pink_months,
            "od_month_amts": od_month_amts,
            "cur_bal": max(0.0, cur_bal),
        }
    return out
def _debt_od_filter_categories(status_map):
    """데이터에 존재하는 연체 필터 버튼 목록 (정상 제외)."""
    labels = {v for v in status_map.values() if v and v != "정상"}
    ordered = []
    for n in range(1, 12):
        lab = f"연체 {n}개월"
        if lab in labels:
            ordered.append(lab)
    if "악성" in labels:
        ordered.append("악성")
    for lab in sorted(labels):
        if lab not in ordered:
            ordered.append(lab)
    return ordered
@st.fragment
def render_debt_month_rank_panel(
    debt_df, month_cols, payment_terms_map=None, height=480, status_month_cols=None
):
    """화면 절반: 연체업체만 · 연체개월수 다중필터 · 채권액 내림차순.
    fragment: 연체 토글은 전체 앱 rerun 없이 패널만 갱신(거래처 원복 로딩과 분리).
    """
    terms_map = payment_terms_map or {}
    mat = build_debt_client_balance_matrix(debt_df, month_cols)
    if mat.empty:
        st.info("표시할 잔액 채권 데이터가 없습니다.")
        return
    cols = list(mat.columns)
    sort_m = cols[-1]
    status_cols = status_month_cols or cols
    od_meta = compute_debt_od_meta_from_raw(debt_df, status_cols, terms_map)
    status_map = {c: v["label"] for c, v in od_meta.items()}
    # 정상채권 제외
    overdue_idx = [c for c in mat.index if status_map.get(c, "정상") != "정상"]
    mat = mat.loc[overdue_idx] if overdue_idx else mat.iloc[0:0]
    if mat.empty:
        st.info("연체 거래처가 없습니다. (정상채권은 이 표에서 제외됩니다)")
        return
    filter_cats = _debt_od_filter_categories(
        {c: status_map[c] for c in mat.index if c in status_map}
    )
    if not filter_cats:
        st.info("연체개월수 필터 대상이 없습니다.")
        return
    if "debt_od_filters" not in st.session_state:
        st.session_state["debt_od_filters"] = list(filter_cats)
    else:
        kept = [x for x in st.session_state["debt_od_filters"] if x in filter_cats]
        _next = kept if kept else list(filter_cats)
        # 값 변경 시에만 기록 — 거래처 원복 rerun마다 session 쓰기로 추가 로딩 방지
        if list(st.session_state.get("debt_od_filters") or []) != _next:
            st.session_state["debt_od_filters"] = _next
    st.markdown(
        "<div style='font-size:15px;font-weight:700;color:#1E293B;margin:18px 0 8px;'>"
        "📊 연체개월수 기준 채권 현황 <span style='font-size:12px;font-weight:500;color:#64748B;'>"
        "(거래처 선택 무시·담당자 필터 · 정상 제외 · 연체개월 다중선택)</span></div>",
        unsafe_allow_html=True,
    )
    # 연체개월수 / 악성 다중 토글 (아담한 버튼 — 이 컨테이너에만 compact CSS)
    with st.container(key="debt_compact_btns_od"):
        bm = st.columns(len(filter_cats), gap="small")
        for i, cat in enumerate(filter_cats):
            with bm[i]:
                on = cat in st.session_state["debt_od_filters"]
                if st.button(
                    cat,
                    key=f"debt_od_filt_{cat}",
                    type="primary" if on else "secondary",
                    width="stretch",
                ):
                    cur = list(st.session_state["debt_od_filters"])
                    if on:
                        if len(cur) > 1:
                            cur = [x for x in cur if x != cat]
                    else:
                        cur.append(cat)
                        cur = [x for x in filter_cats if x in cur]
                    st.session_state["debt_od_filters"] = cur
                    try:
                        st.rerun(scope="fragment")
                    except TypeError:
                        st.rerun()
    selected = st.session_state["debt_od_filters"]
    mat_f = mat.loc[[c for c in mat.index if status_map.get(c) in selected]]
    if mat_f.empty:
        st.info("선택한 연체개월수에 해당하는 업체가 없습니다. 버튼을 하나 이상 켜 주세요.")
        return
    # 연체 있는 월만 열로 표시 (거래처별 해당 없으면 — / 상단 월토글과 무관)
    display_months = [
        m
        for m in status_cols
        if any(m in (od_meta.get(c, {}).get("pink_months") or []) for c in mat_f.index)
    ]
    # 연체합계 기준 내림차순 (동률이면 당월 잔액)
    clients_sorted = sorted(
        list(mat_f.index),
        key=lambda c: (
            float(od_meta.get(c, {}).get("overdue_amt") or 0.0),
            float(od_meta.get(c, {}).get("cur_bal") or 0.0),
        ),
        reverse=True,
    )
    left, right = st.columns([1, 1], gap="medium")
    with left:
        compact = is_touch_ui()
        bal_hdr = f"{sort_m}잔액"
        if compact:
            headers = ["거래처"] + display_months + ["연체합계", bal_hdr, "결제", "연체"]
            left_w, right_w = [100], [48, 60]
        else:
            headers = ["거래처"] + display_months + ["연체합계", bal_hdr, "결제조건", "연체개월수"]
            left_w, right_w = [140], [110, 120]
        cell = _debt_num_cell_font(compact=compact)
        label = _debt_label_cell_style  # 거래처명 13px 유지
        stripe_map = _debt_client_stripe_map(clients_sorted)
        body = []
        for i, client in enumerate(clients_sorted):
            bg = stripe_map.get(client, DEBT_CLIENT_STRIPE_A)
            sep = DEBT_CLIENT_SEP if i > 0 else ""
            meta = od_meta.get(client, {})
            od_amts = meta.get("od_month_amts") or {}
            pink_set = set(meta.get("pink_months") or [])
            tds = [
                f'<td class="dash-freeze-0" style="{label(client, "", {client: bg}, compact=compact)}{sep}'
                f'overflow:hidden;text-overflow:ellipsis;" '
                f'title="{html.escape(str(client))}">{html.escape(str(client))}</td>'
            ]
            for m in display_months:
                if m in pink_set and float(od_amts.get(m, 0) or 0) > 0:
                    v = float(od_amts[m])
                    # 연체월 금액: 분홍 배경 우선
                    tds.append(
                        f'<td class="dash-cell-selectable" style="{cell}'
                        f'background-color:{DEBT_PINK_SOFT};color:#9F1239;{sep}" '
                        f'data-raw="{v}">{v:,.0f}</td>'
                    )
                else:
                    tds.append(
                        f'<td class="dash-cell-selectable" style="{cell}'
                        f'background-color:{bg};color:#94A3B8;{sep}" '
                        f'data-raw="">—</td>'
                    )
            od_total = float(meta.get("overdue_amt") or 0.0)
            cur_bal = float(meta.get("cur_bal") or 0.0)
            # 연체합계도 분홍 우선
            tds.append(
                f'<td class="dash-cell-selectable" style="{cell}'
                f'background-color:{DEBT_PINK_SOFT};color:#9F1239;font-weight:600;{sep}" '
                f'data-raw="{od_total}">{od_total:,.0f}</td>'
            )
            tds.append(
                f'<td class="dash-cell-selectable" style="{cell}'
                f'background-color:{bg};color:#31333F;{sep}" '
                f'data-raw="{cur_bal}">{cur_bal:,.0f}</td>'
            )
            term = resolve_payment_term(client, terms_map) or "—"
            status = status_map.get(client, "악성")
            status_disp = status
            if compact and status.startswith("연체 ") and status.endswith("개월"):
                status_disp = status.replace("연체 ", "").replace("개월", "M")
            if status == "악성":
                sc = "#BE123C"
            elif "연체" in str(status) and "개월" in str(status):
                sc = "#C2410C"
            else:
                sc = "#31333F"
            tds.append(
                f'<td class="dash-freeze-r1" style="{cell}background-color:{bg};white-space:normal;{sep}" '
                f'title="{html.escape(term)}">{html.escape(term)}</td>'
            )
            # 연체개월수: 분홍 배경 최우선
            tds.append(
                f'<td class="dash-freeze-r0" style="{cell}'
                f'background-color:{DEBT_PINK_SOFT};color:{sc};font-weight:600;{sep}" '
                f'title="{html.escape(status)}">{html.escape(status_disp)}</td>'
            )
            body.append(f"<tr>{''.join(tds)}</tr>")
        sel_txt = ", ".join(selected)
        od_m_txt = ", ".join(display_months) if display_months else "없음"
        render_interactive_html_table(
            headers,
            "".join(body),
            height=height,
            show_sum_popup=True,
            toolbar_hint=f"필터: {sel_txt} · 분홍=연체개월/연체월/연체합계 · 연체합계↓ · 정상 제외",
            freeze_left_cols=1,
            freeze_left_widths=left_w,
            freeze_right_header=True,
            freeze_right_widths=right_w,
        )
    with right:
        top_series = pd.Series(
            {
                c: float(od_meta.get(c, {}).get("overdue_amt") or 0.0)
                for c in clients_sorted
            }
        )
        top_n = top_series[top_series > 0].head(15)
        if top_n.empty:
            st.info("선택한 연체 구간에 연체금액이 있는 거래처가 없습니다.")
        else:
            fig = go.Figure(
                go.Bar(
                    x=top_n.values,
                    y=[str(i) for i in top_n.index],
                    orientation="h",
                    marker_color="#FCA5A5",
                    text=[f"{v:,.0f}" for v in top_n.values],
                    textposition="outside",
                )
            )
            fig.update_layout(
                title=dict(
                    text=f"연체금액 Top {len(top_n)}",
                    font=dict(size=14, color="#334155"),
                ),
                yaxis=dict(autorange="reversed", tickfont=dict(size=11)),
                xaxis=dict(title="원", tickfont=dict(size=11)),
                margin=dict(l=10, r=40, t=40, b=30),
                height=height,
                plot_bgcolor="#FFFFFF",
                paper_bgcolor="#FFFFFF",
                showlegend=False,
            )
            chart_key = "debt_od_" + "_".join(selected)[:40]
            render_plotly_chart(fig, use_container_width=True, key=chart_key)
        # 선택 구간별 업체 수·연체합계
        rows_sum = []
        for cat in filter_cats:
            clients_c = [c for c in mat.index if status_map.get(c) == cat]
            if not clients_c:
                continue
            amt = float(sum(od_meta.get(c, {}).get("overdue_amt") or 0.0 for c in clients_c))
            on = "ON" if cat in selected else "OFF"
            rows_sum.append({"구분": cat, "업체수": len(clients_c), "연체합계": amt, "선택": on})
        if rows_sum:
            st.markdown(
                "<div style='font-size:12px;font-weight:600;color:#64748B;margin:8px 0 4px;'>"
                "연체개월수별 요약</div>",
                unsafe_allow_html=True,
            )
            sum_df = pd.DataFrame(rows_sum)
            st.dataframe(
                sum_df.style.format({"연체합계": "{:,.0f}"}),
                use_container_width=True,
                hide_index=True,
                height=min(280, 48 + 32 * len(sum_df)),
            )
@st.cache_data
def cached_staff_pivot(df_base, desired_order):
    if df_base.empty:
        return pd.DataFrame()
    staff_raw = (df_base.pivot_table(index="담당자", columns="연도월_정렬", values="매출액", aggfunc="sum").fillna(0) * 1.1 / 10000)
    staff_cols = [c for c in desired_order if c in staff_raw.columns]
    
    df_p = staff_raw.reindex(columns=staff_cols, fill_value=0)
    
    row_totals = df_p.sum(axis=1)
    # 합계 0인 담당자(특히 미지정 잔상)는 표에서 제거
    df_p = df_p.loc[row_totals > 1e-9].copy()
    if df_p.empty:
        return pd.DataFrame()
    row_totals = df_p.sum(axis=1)
    total_all = row_totals.sum()
    
    if total_all > 0:
        prop = row_totals / total_all * 100
    else:
        prop = 0.0
        
    df_p.insert(0, "총 매출 합계 (만원)", row_totals)
    df_p.insert(0, "매출 비중 (%)", prop)
    
    df_p = df_p.sort_values(by="총 매출 합계 (만원)", ascending=False)
    return df_p
@st.cache_data
def cached_ranking_pivot(df_base, current_year, sel_staff, all_months):
    df_ranking = df_base[(df_base["담당자"] == sel_staff) & (df_base["연도"] == current_year)]
    if df_ranking.empty:
        return pd.DataFrame()
    
    ranking_pivot = (df_ranking.pivot_table(index="거래처", columns="월", values="매출액", aggfunc="sum").fillna(0) * 1.1 / 10000)
    ranking_pivot = ranking_pivot.reindex(columns=all_months, fill_value=0)
    ranking_pivot["당해 누적 (만원)"] = ranking_pivot.sum(axis=1)
    ranking_pivot = ranking_pivot.sort_values(by="당해 누적 (만원)", ascending=False)
    return ranking_pivot
# ----------------------------------------------------
# 탭 전체 적용 업데이트 뱃지 렌더링 유틸
# ----------------------------------------------------
def render_update_badge(date_str):
    return f"<div style='text-align: right; margin-top: 0;'><span style='background-color: #FFFFFF; color: #475569; padding: 6px 12px; border-radius: 8px; font-size: 13px; font-weight: 600; border: 1px solid #CBD5E1; box-shadow: 0 1px 2px rgba(0,0,0,0.05);'>⏱️ 데이터 업데이트: {date_str}</span></div>"
def inject_sticky_tabs_script():
    """필터+탭 상단 fixed.
    - 맥: 현재 안정 코드 100% 무손실 유지 (양식 깨짐 절대 없음)
    - iPad: 필터/검색창 터치 시 키보드를 강제로 내리거나 입력을 방해하던 스크롤 고정 족쇄 제거 패치
    """
    components.html(
        """
        <script>
        (function() {
            var parentDoc = window.parent.document;
            var parentWin = window.parent;
            var SPACER_ID = 'dashboard-sticky-spacer';
            var SHIELD_ID = 'dashboard-top-shield';
            var STICKY_SCRIPT_VER_MAC = 12; 
            var STICKY_SCRIPT_VER_IPAD = 37; /* iPad: 검색창 방해 스크롤/포커스 족쇄 제거 */
            var syncTimer = null;
            var lastH = 0;
            function isTouchPadEarly() {
                var ua = navigator.userAgent || '';
                var ios = /iPad|iPhone|iPod/.test(ua);
                var ipadOs = navigator.platform === 'MacIntel' && navigator.maxTouchPoints > 1;
                return ios || ipadOs;
            }
            function ipadInjectTabHScrollCss() {
                var id = 'dashboard-ipad-tab-hscroll';
                var s = parentDoc.getElementById(id);
                if (!s) {
                    s = parentDoc.createElement('style');
                    s.id = id;
                    (parentDoc.head || parentDoc.documentElement).appendChild(s);
                }
                s.textContent = (
                    '#dashboard-ipad-h-tabs{display:block!important;overflow-x:auto!important;overflow-y:hidden!important;'
                    + 'white-space:nowrap!important;width:100%!important;max-width:100%!important;height:46px!important;'
                    + 'max-height:46px!important;-webkit-overflow-scrolling:touch!important;'
                    + 'border-top:1px solid #E2E8F0!important;background:#fff!important;}'
                    + '#dashboard-ipad-h-tabs button{display:inline-block!important;white-space:nowrap!important;'
                    + 'width:auto!important;height:46px!important;padding:0 12px!important;margin:0!important;'
                    + 'border:0!important;border-bottom:2px solid transparent!important;background:#fff!important;'
                    + 'font-size:13px!important;font-weight:700!important;color:#1E293B!important;vertical-align:top!important;}'
                    + '#dashboard-ipad-h-tabs button.dashboard-ipad-tab-on{color:#2563EB!important;border-bottom-color:#2563EB!important;}'
                    + '.dashboard-ipad-hide-tabs{display:none!important;height:0!important;max-height:0!important;'
                    + 'overflow:hidden!important;visibility:hidden!important;}'
                );
            }
            function ipadPatchScrollIntoView() {
                if (parentWin.__dashboardIpadSivPatched) return;
                parentWin.__dashboardIpadSivPatched = true;
                try {
                    var proto = parentWin.Element.prototype;
                    var orig = proto.scrollIntoView;
                    proto.scrollIntoView = function () {
                        try {
                            if (this && this.closest && (
                                this.closest('[role="tab"]') ||
                                this.closest('[role="tablist"]') ||
                                this.closest('[data-testid="stTab"]') ||
                                this.closest('[data-testid="stTabs"]') ||
                                this.closest('#dashboard-ipad-h-tabs') ||
                                this.closest('.dashboard-filter-sticky')
                            )) return;
                        } catch (eSiv) {}
                        return orig.apply(this, arguments);
                    };
                } catch (eP) {}
            }
            function ipadFreezeLayout() {
                parentWin.__dashboardIpadFreezeLayout = true;
                try {
                    if (parentWin.__dashboardStickyTouchInterval) {
                        clearInterval(parentWin.__dashboardStickyTouchInterval);
                        parentWin.__dashboardStickyTouchInterval = null;
                    }
                    if (parentWin.__dashboardStickyBootInterval) {
                        clearInterval(parentWin.__dashboardStickyBootInterval);
                        parentWin.__dashboardStickyBootInterval = null;
                    }
                } catch (eFz) {}
            }
            function ipadInstallScrollGuard() {
                if (parentWin.__dashboardIpadScrollGuard) return;
                parentWin.__dashboardIpadScrollGuard = true;
                
                /* [중요 패치 1] iOS 가상 키보드가 올라올 때 화면을 억지로 위로 당겨버려서 
                   드롭다운을 강제 종료시키던 과도한 스크롤 방어 로직 완전 무력화 */
                return;
            }
            if (isTouchPadEarly()) {
                try { ipadInjectTabHScrollCss(); } catch (eCss) {}
                try { ipadPatchScrollIntoView(); } catch (eSiv2) {}
                try { ipadInstallScrollGuard(); } catch (eSg) {}
            }
            if (isTouchPadEarly() && parentWin.__dashboardStickyTouchReady === STICKY_SCRIPT_VER_IPAD) {
                try {
                    if (typeof parentWin.__dashboardIpadPin === 'function') {
                        parentWin.__dashboardIpadPin();
                    }
                } catch (ePin) {}
                return;
            }
            if (!isTouchPadEarly() && parentWin.__dashboardStickyMacReady === STICKY_SCRIPT_VER_MAC) {
                try {
                    if (typeof parentWin.__dashboardFixDuplicateTabs === 'function') {
                        parentWin.__dashboardFixDuplicateTabs(false);
                    }
                } catch (eMacSkip) {}
                return;
            }
            if (parentWin.__dashboardStickyObserver) {
                parentWin.__dashboardStickyObserver.disconnect();
            }
            if (parentWin.__dashboardStickyBootInterval) {
                clearInterval(parentWin.__dashboardStickyBootInterval);
            }
            if (parentWin.__dashboardStickyTouchInterval) {
                clearInterval(parentWin.__dashboardStickyTouchInterval);
            }
            parentWin.__dashboardStickyRafLoop = false;
            if (parentWin.__dashboardIpadRaf) {
                parentWin.cancelAnimationFrame(parentWin.__dashboardIpadRaf);
                parentWin.__dashboardIpadRaf = null;
            }
            function isTouchPad() {
                var ua = navigator.userAgent || '';
                var ios = /iPad|iPhone|iPod/.test(ua);
                var ipadOs = navigator.platform === 'MacIntel' && navigator.maxTouchPoints > 1;
                return ios || ipadOs;
            }
            var touchMode = isTouchPad();
            function getTopOffsetMac() {
                var header = parentDoc.querySelector('[data-testid="stHeader"]');
                if (!header) return 46;
                var rect = header.getBoundingClientRect();
                if (rect.bottom > 0) {
                    return Math.round(rect.bottom);
                }
                if (header.offsetHeight > 0) {
                    return header.offsetHeight;
                }
                return 46;
            }
            function getMainRect() {
                var block = parentDoc.querySelector('section.main .block-container');
                return block ? block.getBoundingClientRect() : null;
            }
            function isMainTabList(el) {
                if (!el || !el.textContent) return false;
                return el.textContent.indexOf('📌 영업 종합 요약') !== -1;
            }
            function findMainTabsHost() {
                var hosts = parentDoc.querySelectorAll('div[data-testid="stTabs"]');
                for (var i = 0; i < hosts.length; i++) {
                    if (hosts[i].querySelector('[role="tabpanel"]')) return hosts[i];
                    if (hosts[i].textContent.indexOf('📌 영업 종합 요약') !== -1) return hosts[i];
                }
                return null;
            }
            function findMainTabList() {
                var hosts = parentDoc.querySelectorAll('div[data-testid="stTabs"]');
                var i, lists, j, el;
                for (i = 0; i < hosts.length; i++) {
                    lists = hosts[i].querySelectorAll('div[role="tablist"]');
                    for (j = 0; j < lists.length; j++) {
                        if (isMainTabList(lists[j])) return lists[j];
                    }
                }
                lists = parentDoc.querySelectorAll('div[role="tablist"]');
                var inFilter = null;
                for (i = 0; i < lists.length; i++) {
                    el = lists[i];
                    if (!isMainTabList(el)) continue;
                    if (el.closest('.dashboard-filter-sticky')) {
                        inFilter = inFilter || el;
                        continue;
                    }
                    return el;
                }
                return inFilter;
            }
            function findFilterBox() {
                var marker = parentDoc.getElementById('sticky-marker');
                if (!marker) return null;
                return marker.closest('div[data-testid="stVerticalBlockBorderWrapper"]') ||
                       marker.closest('div[data-testid="stVerticalBlock"]');
            }
            function collectMainTabLists() {
                var out = [];
                var lists = parentDoc.querySelectorAll('div[role="tablist"]');
                for (var i = 0; i < lists.length; i++) {
                    if (isMainTabList(lists[i])) out.push(lists[i]);
                }
                return out;
            }
            function hideOrphanTabList(el) {
                if (!el) return;
                try {
                    el.setAttribute('data-dashboard-orphan-tabs', '1');
                    el.classList.remove('dashboard-tabs-in-filter');
                    el.style.setProperty('display', 'none', 'important');
                    el.style.setProperty('visibility', 'hidden', 'important');
                    el.style.setProperty('height', '0', 'important');
                    el.style.setProperty('max-height', '0', 'important');
                    el.style.setProperty('overflow', 'hidden', 'important');
                    el.style.setProperty('pointer-events', 'none', 'important');
                    el.style.setProperty('margin', '0', 'important');
                    el.style.setProperty('padding', '0', 'important');
                } catch (eH) {}
            }
            function hideStaleTabListsInHost(tabsHost, keepList) {
                if (!tabsHost) return;
                Array.from(tabsHost.children).forEach(function(child) {
                    if (containsTabPanel(child)) return;
                    if (keepList && (child === keepList || child.contains(keepList))) return;
                    child.classList.add('dashboard-tabs-list-shell');
                    child.style.setProperty('display', 'none', 'important');
                    child.style.setProperty('height', '0', 'important');
                    child.style.setProperty('overflow', 'hidden', 'important');
                    child.style.setProperty('visibility', 'hidden', 'important');
                    child.style.setProperty('pointer-events', 'none', 'important');
                });
                var left = tabsHost.querySelectorAll('div[role="tablist"]');
                for (var k = 0; k < left.length; k++) {
                    if (keepList && left[k] === keepList) continue;
                    if (!isMainTabList(left[k])) continue;
                    hideOrphanTabList(left[k]);
                }
            }
            function purgeOrphanTabLists(filterBox, keepList) {
                if (!filterBox) return;
                var lists = filterBox.querySelectorAll('div[role="tablist"]');
                for (var i = 0; i < lists.length; i++) {
                    if (keepList && lists[i] === keepList) continue;
                    if (!isMainTabList(lists[i])) continue;
                    try { lists[i].remove(); } catch (eRm) {
                        hideOrphanTabList(lists[i]);
                    }
                }
            }
            function isTabMountHealthy(filterBox, all) {
                if (!filterBox || !all || all.length !== 1) return false;
                var t = all[0];
                return !!(filterBox.contains(t) && t.classList.contains('dashboard-tabs-in-filter'));
            }
            function fixDuplicateMainTabs(forceMove) {
                var filterBox = findFilterBox();
                var host = findMainTabsHost();
                var all = collectMainTabLists();
                if (!all.length) return false;
                if (!forceMove && isTabMountHealthy(filterBox, all)) {
                    parentWin.__dashboardFixDuplicateTabs = fixDuplicateMainTabs;
                    return true;
                }
                var keep = null;
                var i;
                if (host) {
                    for (i = 0; i < all.length; i++) {
                        if (host.contains(all[i])) { keep = all[i]; break; }
                    }
                }
                if (!keep) {
                    for (i = 0; i < all.length; i++) {
                        if (!all[i].closest('.dashboard-filter-sticky')) { keep = all[i]; break; }
                    }
                }
                if (!keep) keep = all[all.length - 1];
                for (i = 0; i < all.length; i++) {
                    if (all[i] === keep) continue;
                    try { all[i].remove(); } catch (eRm2) {
                        hideOrphanTabList(all[i]);
                    }
                }
                if (host) {
                    host.classList.add('dashboard-tabs-host-compact');
                    hideStaleTabListsInHost(host, keep);
                }
                var canMove = !!filterBox && (!isFilterDropdownOpen() || forceMove || all.length > 1);
                if (canMove) {
                    filterBox.classList.add('dashboard-filter-sticky');
                    if (touchMode) filterBox.classList.add('dashboard-filter-sticky-touch');
                    try {
                        var stuck = keep.closest('.dashboard-tabs-list-shell');
                        if (stuck && !filterBox.contains(keep)) {
                            stuck.style.removeProperty('display');
                            stuck.style.display = '';
                            stuck.style.removeProperty('visibility');
                            stuck.style.removeProperty('pointer-events');
                        }
                    } catch (eUh) {}
                    purgeOrphanTabLists(filterBox, keep);
                    if (!filterBox.contains(keep)) {
                        try { filterBox.appendChild(keep); } catch (eAp) {}
                    }
                    keep.classList.add('dashboard-tabs-in-filter');
                    keep.removeAttribute('data-dashboard-orphan-tabs');
                    keep.style.removeProperty('display');
                    keep.style.removeProperty('height');
                    keep.style.removeProperty('max-height');
                    keep.style.removeProperty('overflow');
                    keep.style.setProperty('visibility', 'visible', 'important');
                    keep.style.setProperty('pointer-events', 'auto', 'important');
                    if (touchMode) {
                        try { ipadScrollTabs(filterBox); } catch (eSc) {}
                    }
                    if (host) hideStaleTabListsInHost(host, keep);
                } else if (filterBox && filterBox.contains(keep)) {
                    keep.classList.add('dashboard-tabs-in-filter');
                    keep.removeAttribute('data-dashboard-orphan-tabs');
                } else if (!filterBox || !filterBox.contains(keep)) {
                    keep.classList.remove('dashboard-tabs-in-filter');
                    keep.removeAttribute('data-dashboard-orphan-tabs');
                    keep.style.removeProperty('display');
                    keep.style.setProperty('visibility', 'visible', 'important');
                    keep.style.setProperty('pointer-events', 'auto', 'important');
                }
                parentWin.__dashboardFixDuplicateTabs = fixDuplicateMainTabs;
                return true;
            }
            function remountLiveTabList(filterBox, tabList) {
                if (!filterBox || !tabList) return false;
                var mains = collectMainTabLists();
                if (mains.length > 1) {
                    return fixDuplicateMainTabs(true);
                }
                if (isTabMountHealthy(filterBox, mains)) {
                    return true;
                }
                if (isFilterDropdownOpen()) {
                    fixDuplicateMainTabs(false);
                    return false;
                }
                return fixDuplicateMainTabs(true);
            }
            function containsTabPanel(el) {
                if (!el || el.nodeType !== 1) return false;
                if (el.getAttribute && el.getAttribute('role') === 'tabpanel') return true;
                return !!(el.querySelector && el.querySelector('[role="tabpanel"]'));
            }
            function isFilterInteracting() {
                try {
                    if (isFilterDropdownOpen()) return true;
                    var ae = parentDoc.activeElement;
                    if (!ae || !ae.closest) return false;
                    return !!(ae.closest('[data-testid="stMultiSelect"]') ||
                        ae.closest('[data-testid="stSelectbox"]') ||
                        ae.closest('[data-baseweb="select"]') ||
                        ae.closest('[data-baseweb="popover"]') ||
                        ae.closest('[data-baseweb="menu"]'));
                } catch (eFi) {
                    return false;
                }
            }
            function isFilterDropdownOpen() {
                function visible(el) {
                    if (!el) return false;
                    var r = el.getBoundingClientRect();
                    if (r.width < 2 || r.height < 2) return false;
                    var st = parentWin.getComputedStyle(el);
                    if (st.display === 'none' || st.visibility === 'hidden' || st.opacity === '0') return false;
                    return true;
                }
                try {
                    var nodes = parentDoc.querySelectorAll(
                        '[data-baseweb="popover"], [data-baseweb="menu"], [data-baseweb="layer"], '
                        + '[data-baseweb="select"] [role="listbox"], ul[role="listbox"], '
                        + '[data-testid="stSelectboxVirtualDropdown"]'
                    );
                    for (var i = 0; i < nodes.length; i++) {
                        if (visible(nodes[i])) return true;
                    }
                } catch (eVis) {}
                return false;
            }
            function ensureSpacer(filterBox) {
                var spacer = parentDoc.getElementById(SPACER_ID);
                if (!spacer) {
                    spacer = parentDoc.createElement('div');
                    spacer.id = SPACER_ID;
                    filterBox.parentNode.insertBefore(spacer, filterBox);
                } else if (spacer.nextElementSibling !== filterBox) {
                    filterBox.parentNode.insertBefore(spacer, filterBox);
                }
                return spacer;
            }
            function syncTopShield(top, rect) {
                var shield = parentDoc.getElementById(SHIELD_ID);
                if (!shield) {
                    shield = parentDoc.createElement('div');
                    shield.id = SHIELD_ID;
                    parentDoc.body.appendChild(shield);
                }
                shield.style.setProperty('display', 'block', 'important');
                shield.style.setProperty('position', 'fixed', 'important');
                shield.style.setProperty('top', '0', 'important');
                shield.style.setProperty('left', rect.left + 'px', 'important');
                shield.style.setProperty('width', rect.width + 'px', 'important');
                shield.style.setProperty('height', top + 'px', 'important');
                shield.style.setProperty('background', '#F8FAFC', 'important');
                shield.style.setProperty('z-index', '989', 'important');
                shield.style.setProperty('pointer-events', 'none', 'important');
            }
            function mountTabs(filterBox, tabList) {
                return remountLiveTabList(filterBox, tabList);
            }
            
            /* ===== iPad 전용 로직 ===== */
            function isSidebarOpen() {
                var sidebar = parentDoc.querySelector('[data-testid="stSidebar"]');
                if (!sidebar) return false;
                if (sidebar.getAttribute('aria-expanded') === 'true') return true;
                return sidebar.getBoundingClientRect().width > 100;
            }
            function cleanupIpadPortal() {
                var portal = parentDoc.getElementById('dashboard-ipad-portal');
                if (!portal) return;
                var filterBox = findFilterBox();
                if (filterBox && portal.contains(filterBox)) {
                    var spacer = parentDoc.getElementById(SPACER_ID);
                    if (spacer && spacer.parentNode) {
                        spacer.parentNode.insertBefore(filterBox, spacer.nextSibling);
                    } else {
                        parentDoc.body.appendChild(filterBox);
                    }
                }
                portal.remove();
            }
            function lockIpadFilter(ms) {
                parentWin.__dashboardIpadFilterLockUntil = Date.now() + (ms || 10000);
            }
            function markIpadScrolling() {
                parentWin.__dashboardIpadScrollLockUntil = Date.now() + 1500;
            }
            function isIpadFilterLocked() {
                try {
                    if (Date.now() < (parentWin.__dashboardIpadScrollLockUntil || 0)) return true;
                    if (Date.now() < (parentWin.__dashboardIpadFilterLockUntil || 0)) return true;
                    return isFilterInteracting();
                } catch (eLk) {
                    return false;
                }
            }
            function applyIpad0804Hack() {
                if (parentWin.__dashboardIpadFreezeLayout) return;
                if (isIpadFilterLocked()) return;
                cleanupIpadPortal();
                var dropOpen = false;
                try { dropOpen = isFilterDropdownOpen(); } catch (eDo) {}
                if (!dropOpen) {
                    try { fixDuplicateMainTabs(collectMainTabLists().length > 1); } catch (eDup) {}
                }
                var targetBox = findFilterBox();
                if (!targetBox) return;
                if (!dropOpen) {
                    var tabList = findMainTabList();
                    if (!tabList) return;
                    if (!remountLiveTabList(targetBox, tabList)) return;
                    tabList.style.setProperty('padding', '0 10px 0px 10px', 'important');
                    tabList.style.setProperty('margin-top', '5px', 'important');
                    tabList.style.setProperty('border-bottom', 'none', 'important');
                    tabList.style.setProperty('background-color', 'transparent', 'important');
                    tabList.style.setProperty('pointer-events', 'auto', 'important');
                }
                ipadApplyBoxStyles(targetBox);
            }
            function ipadScrollTabs(box) {
                if (!box) return;
                try { ipadInjectTabHScrollCss(); } catch (eCss2) {}
                var list = null;
                var lists = box.querySelectorAll('[role="tablist"], [data-baseweb="tab-list"]');
                var i;
                for (i = 0; i < lists.length; i++) {
                    if (isMainTabList(lists[i])) { list = lists[i]; break; }
                }
                if (!list) list = findMainTabList();
                if (!list) {
                    var oneTab = parentDoc.querySelector('[data-testid="stTab"]');
                    if (oneTab) list = oneTab.parentElement;
                }
                if (!list) return;
                var origTabs = list.querySelectorAll('[data-testid="stTab"], [role="tab"]');
                if (!origTabs.length) return;
                var bar = parentDoc.getElementById('dashboard-ipad-h-tabs');
                if (!bar) {
                    bar = parentDoc.createElement('div');
                    bar.id = 'dashboard-ipad-h-tabs';
                }
                if (bar.parentNode !== box) {
                    box.appendChild(bar);
                }
                var labels = [];
                for (i = 0; i < origTabs.length; i++) {
                    labels.push((origTabs[i].textContent || '').replace(/\\s+/g, ' ').trim());
                }
                var sig = labels.join('|');
                if (bar.getAttribute('data-sig') !== sig || bar.childNodes.length !== origTabs.length) {
                    bar.setAttribute('data-sig', sig);
                    bar.innerHTML = '';
                    for (i = 0; i < origTabs.length; i++) {
                        (function (orig, label) {
                            var b = parentDoc.createElement('button');
                            b.type = 'button';
                            b.textContent = label;
                            
                            b.addEventListener('click', function (ev) {
                                ev.preventDefault();
                                var siblings = bar.querySelectorAll('button');
                                for(var k = 0; k < siblings.length; k++) {
                                    siblings[k].classList.remove('dashboard-ipad-tab-on');
                                }
                                b.classList.add('dashboard-ipad-tab-on');
                                try { orig.click(); } catch (eCl) {}
                            });
                            bar.appendChild(b);
                        })(origTabs[i], labels[i]);
                    }
                }
                var btns = bar.querySelectorAll('button');
                for (i = 0; i < origTabs.length && i < btns.length; i++) {
                    var on = origTabs[i].getAttribute('aria-selected') === 'true'
                        || origTabs[i].hasAttribute('data-selected')
                        || origTabs[i].getAttribute('data-selected') === 'true';
                    if (on) btns[i].classList.add('dashboard-ipad-tab-on');
                    else btns[i].classList.remove('dashboard-ipad-tab-on');
                }
                list.classList.add('dashboard-ipad-hide-tabs');
                list.style.setProperty('display', 'none', 'important');
                list.style.setProperty('height', '0', 'important');
                list.style.setProperty('max-height', '0', 'important');
                list.style.setProperty('overflow', 'hidden', 'important');
            }
            function ipadApplyBoxStyles(targetBox) {
                if (!targetBox) return;
                if (parentWin.__dashboardIpadFreezeLayout && targetBox.style.position === 'fixed') return;
                if (isFilterDropdownOpen()) return;
                ipadScrollTabs(targetBox);
                targetBox.classList.add('dashboard-filter-sticky');
                targetBox.classList.add('dashboard-filter-sticky-touch');
                var header = parentDoc.querySelector('[data-testid="stHeader"]');
                var topPx = 40;
                if (header) {
                    var hb = header.getBoundingClientRect().bottom;
                    if (hb > 20) topPx = Math.round(hb);
                }
                targetBox.style.setProperty('position', 'fixed', 'important');
                targetBox.style.setProperty('top', topPx + 'px', 'important');
                targetBox.style.setProperty('left', '8px', 'important');
                targetBox.style.setProperty('right', '8px', 'important');
                targetBox.style.setProperty('width', 'auto', 'important');
                targetBox.style.setProperty('max-width', 'calc(100vw - 16px)', 'important');
                targetBox.style.setProperty('z-index', '999999', 'important');
                targetBox.style.setProperty('height', 'auto', 'important');
                targetBox.style.setProperty('max-height', 'none', 'important');
                targetBox.style.setProperty('overflow', 'visible', 'important');
                
                targetBox.style.setProperty('-webkit-transform', 'translate3d(0, 0, 0)', 'important');
                targetBox.style.setProperty('transform', 'translate3d(0, 0, 0)', 'important');
                
                var spacer = parentDoc.getElementById(SPACER_ID);
                if (!spacer) {
                    spacer = parentDoc.createElement('div');
                    spacer.id = SPACER_ID;
                    if (targetBox.parentNode) targetBox.parentNode.insertBefore(spacer, targetBox);
                }
                var barH = Math.round(targetBox.offsetHeight || 160);
                var y = parentWin.pageYOffset || parentDoc.documentElement.scrollTop || 0;
                if (y < 24) {
                    try {
                        var pinH = Math.round(targetBox.getBoundingClientRect().bottom - spacer.getBoundingClientRect().top);
                        if (pinH >= 8 && pinH <= 280) barH = pinH;
                    } catch (ePinH) {}
                }
                spacer.style.setProperty('display', 'block', 'important');
                spacer.style.setProperty('height', barH + 'px', 'important');
                spacer.style.setProperty('min-height', '0', 'important');
                spacer.style.setProperty('margin', '0', 'important');
                parentDoc.documentElement.style.setProperty('--dashboard-fixed-bar-height', barH + 'px');
                parentWin.__dashboardIpadTarget = targetBox;
                parentWin.__dashboardIpadSpacer = spacer;
                ipadFreezeLayout();
            }
            function ipadPinFilterBox() {
                var box = findFilterBox();
                if (!box) return;
                if (box.style.position === 'fixed' && parentWin.__dashboardIpadFreezeLayout) return;
                if (isFilterDropdownOpen()) return;
                parentWin.__dashboardIpadFreezeLayout = false;
                try { fixDuplicateMainTabs(false); } catch (eDup2) {}
                ipadApplyBoxStyles(box);
            }
            parentWin.__dashboardIpadPin = ipadPinFilterBox;
            
            function syncIpadWidthLoop() {
                if (!parentWin.__dashboardIpadFreezeLayout) {
                    var targetBox = parentWin.__dashboardIpadTarget || findFilterBox();
                    var spacer = parentWin.__dashboardIpadSpacer || parentDoc.getElementById(SPACER_ID);
                    
                    if (targetBox && spacer && parentDoc.body.contains(spacer)) {
                        var currentHeight = Math.round(targetBox.getBoundingClientRect().height);
                        if (currentHeight > 0) {
                            spacer.style.setProperty('height', currentHeight + 'px', 'important');
                            parentDoc.documentElement.style.setProperty('--dashboard-fixed-bar-height', currentHeight + 'px');
                        }
                    }
                    
                    var now = Date.now();
                    if (!parentWin.__dashboardIpadLastWidthSync || now - parentWin.__dashboardIpadLastWidthSync >= 800) {
                        parentWin.__dashboardIpadLastWidthSync = now;
                        if (!isIpadFilterLocked()) {
                            if (!targetBox || !parentDoc.body.contains(targetBox)) {
                                applyIpad0804Hack();
                                targetBox = parentWin.__dashboardIpadTarget;
                                spacer = parentWin.__dashboardIpadSpacer;
                            }
                            if (targetBox && spacer && parentDoc.body.contains(spacer)) {
                                ipadApplyBoxStyles(targetBox);
                            }
                        }
                    }
                }
                
                try {
                    var bar = parentDoc.getElementById('dashboard-ipad-h-tabs');
                    if (bar) {
                        var list = parentDoc.querySelector('.dashboard-ipad-hide-tabs');
                        if (list) {
                            var origTabs = list.querySelectorAll('[data-testid="stTab"], [role="tab"]');
                            var btns = bar.querySelectorAll('button');
                            if (origTabs.length === btns.length) {
                                for (var i = 0; i < origTabs.length; i++) {
                                    var on = origTabs[i].getAttribute('aria-selected') === 'true'
                                        || origTabs[i].hasAttribute('data-selected')
                                        || origTabs[i].getAttribute('data-selected') === 'true';
                                    if (on) btns[i].classList.add('dashboard-ipad-tab-on');
                                    else btns[i].classList.remove('dashboard-ipad-tab-on');
                                }
                            }
                        }
                    }
                } catch(eTabSync) {}

                parentWin.__dashboardIpadRaf = parentWin.requestAnimationFrame(syncIpadWidthLoop);
            }

            function syncFixedBar() {
                try { fixDuplicateMainTabs(false); } catch (eSf) {}
                if (touchMode) {
                    if (isFilterDropdownOpen()) return;
                    applyIpad0804Hack();
                    return;
                }
                /* ===== Mac 환경: 원본 레이아웃 로직 100% 무손실 복구 ===== */
                var filterBox = findFilterBox();
                var tabList = findMainTabList();
                if (!filterBox || !tabList) {
                    parentWin.__dashTabRetry = (parentWin.__dashTabRetry || 0) + 1;
                    if (parentWin.__dashTabRetry < 25) scheduleSync(200);
                    return;
                }
                var mains = collectMainTabLists();
                if (!isTabMountHealthy(filterBox, mains)) {
                    var mounted = mountTabs(filterBox, tabList);
                    tabList = findMainTabList() || tabList;
                    if (!mounted || !filterBox.contains(tabList)) {
                        parentWin.__dashTabRetry = (parentWin.__dashTabRetry || 0) + 1;
                        if (parentWin.__dashTabRetry < 25) scheduleSync(180);
                        return;
                    }
                }
                parentWin.__dashTabRetry = 0;
                var rectMac = getMainRect();
                if (!rectMac) return;
                var topMac = getTopOffsetMac();
                var curLeft = parseFloat(filterBox.style.left) || -9999;
                var curTop = parseFloat(filterBox.style.top) || -9999;
                var curW = parseFloat(filterBox.style.width) || 0;
                var posOk = (
                    filterBox.style.position === 'fixed' &&
                    Math.abs(curLeft - rectMac.left) < 1.5 &&
                    Math.abs(curTop - topMac) < 1.5 &&
                    Math.abs(curW - rectMac.width) < 2.5
                );
                if (!posOk) {
                    filterBox.style.setProperty('position', 'fixed', 'important');
                    filterBox.style.setProperty('top', topMac + 'px', 'important');
                    filterBox.style.setProperty('left', rectMac.left + 'px', 'important');
                    filterBox.style.setProperty('width', rectMac.width + 'px', 'important');
                    filterBox.style.setProperty('max-width', rectMac.width + 'px', 'important');
                    filterBox.style.setProperty('z-index', '990', 'important');
                    parentDoc.documentElement.style.setProperty('--dashboard-bar-top', topMac + 'px');
                    parentDoc.documentElement.style.setProperty('--dashboard-bar-left', rectMac.left + 'px');
                    parentDoc.documentElement.style.setProperty('--dashboard-bar-width', rectMac.width + 'px');
                    syncTopShield(topMac, rectMac);
                }
                var barHMac = filterBox.offsetHeight + 4;
                if (Math.abs(barHMac - lastH) > 1) {
                    var spacerMac = ensureSpacer(filterBox);
                    spacerMac.style.height = barHMac + 'px';
                    spacerMac.style.display = 'block';
                    parentDoc.documentElement.style.setProperty('--dashboard-fixed-bar-height', barHMac + 'px');
                    lastH = barHMac;
                }
            }
            
            function mutationTouchesMainTabs(mutations) {
                function check(node) {
                    if (!node || node.nodeType !== 1) return false;
                    try {
                        if (node.id === 'sticky-marker' || node.id === SPACER_ID) return true;
                        if (node.getAttribute && node.getAttribute('data-testid') === 'stTabs') return true;
                        if (node.getAttribute && node.getAttribute('role') === 'tablist') return isMainTabList(node);
                        if (node.querySelector) {
                            if (node.querySelector('#sticky-marker')) return true;
                            var tl = node.querySelector('div[role="tablist"]');
                            if (tl && isMainTabList(tl)) return true;
                        }
                    } catch (eC) {}
                    return false;
                }
                for (var m = 0; m < mutations.length; m++) {
                    var rec = mutations[m];
                    if (check(rec.target)) return true;
                    var a, r;
                    for (a = 0; a < rec.addedNodes.length; a++) if (check(rec.addedNodes[a])) return true;
                    for (r = 0; r < rec.removedNodes.length; r++) if (check(rec.removedNodes[r])) return true;
                }
                return false;
            }
            function scheduleSync(ms) {
                try {
                    if (touchMode && (parentWin.__dashboardIpadFreezeLayout || isIpadFilterLocked())) return;
                    var all = collectMainTabLists();
                    if (all.length > 1) {
                        fixDuplicateMainTabs(true);
                        all = collectMainTabLists();
                    } else {
                        fixDuplicateMainTabs(false);
                    }
                    if (!touchMode) {
                        var fb = findFilterBox();
                        if (isTabMountHealthy(fb, all) && !isFilterDropdownOpen()) {
                            var hNow = fb ? (fb.offsetHeight + 4) : lastH;
                            if (Math.abs(hNow - lastH) <= 1) return;
                        }
                    }
                } catch (eFix) {}
                if (isFilterDropdownOpen() && collectMainTabLists().length <= 1) return;
                if (syncTimer) clearTimeout(syncTimer);
                var delay = ms || 40;
                delay = Math.max(delay, touchMode ? 350 : 280);
                syncTimer = setTimeout(syncFixedBar, delay);
            }

            if (touchMode) {
                parentDoc.documentElement.classList.add('dashboard-touch-mode');
                
                parentDoc.addEventListener('focusin', function(e) {
                    if (e.target.tagName === 'INPUT' || e.target.tagName === 'SELECT') {
                        setTimeout(function() {
                            var box = findFilterBox();
                            if(box) {
                                box.style.setProperty('transform', 'translate3d(0,0,0)', 'important');
                                /* [중요 패치 2] iOS 키보드 입력을 튕겨내던 scrollTo 강제고정 삭제 (자연스러운 스크롤 허용) */
                            }
                        }, 100);
                    }
                }, true);

                try {
                    try { parentDoc.cookie = 'dashboard_touch=1; path=/; max-age=31536000; SameSite=Lax'; } catch (eCk) {}
                    if (!parentWin.__dashboardTouchUiSynced) {
                        parentWin.__dashboardTouchUiSynced = true;
                        var tu = new URL(parentWin.location.href);
                        if (tu.searchParams.get('touch_ui') !== '1') {
                            tu.searchParams.set('touch_ui', '1');
                            try { parentWin.history.replaceState(null, '', tu.toString()); } catch (eHs) {}
                        }
                    }
                } catch (eTu) {}
                
                var boot = 0;
                parentWin.__dashboardStickyBootInterval = setInterval(function() {
                    if (!isIpadFilterLocked() && !isFilterDropdownOpen()) applyIpad0804Hack();
                    boot++;
                    if (boot > 40) {
                        clearInterval(parentWin.__dashboardStickyBootInterval);
                        parentWin.__dashboardStickyBootInterval = null;
                    }
                }, 200);
                
                parentWin.__dashboardStickyTouchInterval = setInterval(function() {
                    if (!isIpadFilterLocked() && !isFilterDropdownOpen()) applyIpad0804Hack();
                }, 5000);
                
                var observer = new MutationObserver(function() {
                    if (isIpadFilterLocked() || isFilterDropdownOpen()) return;
                    var box = findFilterBox();
                    if (!box) return;
                    if (box.style.position === 'fixed') return;
                    parentWin.__dashboardIpadFreezeLayout = false;
                    ipadPinFilterBox();
                });
                observer.observe(parentDoc.body, { childList: true, subtree: true });
                parentWin.__dashboardStickyObserver = observer;
                parentWin.__dashboardIpadScheduleSync = scheduleSync;
                parentWin.__dashboardFixDuplicateTabs = fixDuplicateMainTabs;
                parentWin.__dashboardStickyTouchReady = STICKY_SCRIPT_VER_IPAD;
                
                try {
                    var appScroll = parentDoc.querySelector('[data-testid="stAppViewContainer"]');
                    if (appScroll) {
                        appScroll.addEventListener('scroll', markIpadScrolling, { passive: true, capture: true });
                    }
                    parentDoc.addEventListener('scroll', markIpadScrolling, { passive: true, capture: true });
                    parentWin.addEventListener('scroll', markIpadScrolling, { passive: true, capture: true });
                    parentDoc.addEventListener('touchmove', markIpadScrolling, { passive: true, capture: true });
                } catch (eScr) {}
                parentDoc.addEventListener('pointerdown', function(e) {
                    try {
                        var t = e.target;
                        if (!t || !t.closest) return;
                        if (t.closest('[data-testid="stMultiSelect"]') ||
                            t.closest('[data-testid="stSelectbox"]') ||
                            t.closest('[data-baseweb="select"]') ||
                            t.closest('[data-baseweb="popover"]') ||
                            t.closest('[data-baseweb="menu"]') ||
                            t.closest('[data-testid="stSelectboxVirtualDropdown"]') ||
                            t.closest('ul[role="listbox"]')) {
                            lockIpadFilter(12000);
                        }
                    } catch (ePd) {}
                }, true);
                parentDoc.addEventListener('focusin', function(e) {
                    try {
                        var t = e.target;
                        if (!t || !t.closest) return;
                        if (t.closest('[data-testid="stMultiSelect"]') ||
                            t.closest('[data-testid="stSelectbox"]') ||
                            t.closest('[data-baseweb="select"]')) {
                            lockIpadFilter(12000);
                        }
                    } catch (eFi2) {}
                }, true);
                parentDoc.addEventListener('click', function(e) {
                    if (e.target.closest('[data-testid="collapsedControl"]') ||
                        e.target.closest('[data-testid="stSidebar"]') ||
                        e.target.closest('[role="tab"]')) {
                        scheduleSync(80);
                        scheduleSync(400);
                    }
                }, true);
                parentWin.addEventListener('resize', function() {
                    if (isIpadFilterLocked()) return;
                    scheduleSync(80);
                }, { passive: true });
                
                parentWin.addEventListener('orientationchange', function() {
                    parentWin.__dashboardIpadFreezeLayout = false;
                    parentWin.__dashboardIpadSpacerH = 0;
                    
                    var delays = [50, 150, 300, 500];
                    delays.forEach(function(ms) {
                        setTimeout(function() {
                            var box = findFilterBox();
                            var sp = parentDoc.getElementById(SPACER_ID);
                            if (box && sp) {
                                var h = Math.round(box.getBoundingClientRect().height);
                                if (h > 0) sp.style.setProperty('height', h + 'px', 'important');
                                box.style.setProperty('transform', 'translate3d(0,0,0)', 'important');
                            }
                        }, ms);
                    });
                    
                    scheduleSync(120);
                    scheduleSync(450);
                }, { passive: true });
                applyIpad0804Hack();
                syncIpadWidthLoop(); 
            } else {
                /* Mac 전용 오리지널 환경 유지 (무손실 복원) */
                var pollCount = 0;
                parentWin.__dashboardStickyBootInterval = setInterval(function() {
                    if (!isFilterDropdownOpen()) syncFixedBar();
                    pollCount++;
                    if (pollCount > 40) {
                        clearInterval(parentWin.__dashboardStickyBootInterval);
                        parentWin.__dashboardStickyBootInterval = null;
                    }
                }, 250);
                var observer = new MutationObserver(function(mutations) {
                    if (!mutationTouchesMainTabs(mutations)) return;
                    var n = collectMainTabLists().length;
                    if (n > 1) fixDuplicateMainTabs(true);
                    else scheduleSync(320);
                });
                observer.observe(parentDoc.body, {
                    childList: true,
                    subtree: true
                });
                parentWin.__dashboardStickyObserver = observer;
                parentWin.__dashboardMacScheduleSync = scheduleSync;
                parentWin.__dashboardFixDuplicateTabs = fixDuplicateMainTabs;
                parentWin.__dashboardStickyMacReady = STICKY_SCRIPT_VER_MAC;
                parentWin.addEventListener('resize', function() { scheduleSync(120); }, { passive: true });
                parentWin.addEventListener('pageshow', function() { scheduleSync(120); }, { passive: true });
                if (parentWin.visualViewport) {
                    parentWin.visualViewport.addEventListener('resize', function() { scheduleSync(100); }, { passive: true });
                }
                parentDoc.addEventListener('click', function(e) {
                    if (e.target.closest('[data-testid="collapsedControl"]') ||
                        e.target.closest('[role="tab"]')) {
                        scheduleSync(100);
                        scheduleSync(400);
                    }
                }, true);
                var sidebar = parentDoc.querySelector('[data-testid="stSidebar"]');
                if (sidebar) {
                    sidebar.addEventListener('transitionend', function() { scheduleSync(100); });
                }
                parentWin.__dashboardStickyTouchInterval = setInterval(function() {
                    if (!isFilterDropdownOpen()) syncFixedBar();
                }, 8000);
                syncFixedBar();
            }
        })();
        </script>
        """,
        height=0,
        width=0,
    )
# ----------------------------------------------------
# 날짜 유지용 로컬 캐시 함수 (7, 8번 탭 전용)
# ----------------------------------------------------
TAB7_DATE_FILE = os.path.join(CACHE_DIR, "tab7_date.txt")
TAB8_DATE_FILE = os.path.join(CACHE_DIR, "tab8_date.txt")
PROFIT_CACHE_FILE = os.path.join(CACHE_DIR, "profitability.json")
# 수익성분석.xlsx 기본값 (한국메티슨 시트 함수 그대로 이식)
PROFIT_DEFAULTS = {
    "project_name": "한국메티슨",
    "tank_gas": "질소",
    "tank_capacity_mode": "liters",  # liters | kg
    "tank_liters": 123763.0,  # 내용적 L (100000kg ÷ 0.808 ≈)
    "tank_spec": 100000.0,  # kg (환산값 또는 직접입력)
    "hourly_usage_mode": "kg",  # nm3 | kg
    "hourly_usage_kg": 50.0,  # 시간당 사용량 kg/h
    "hourly_usage_nm3": 39.98,  # 시간당 루베 Nm³/h (질소 50kg×0.7996)
    "operating_hours": 8.0,   # 일 가동시간 h
    "operating_days_per_month": 22.0,  # 월 가동일수 (일사용×가동일 = 월사용)
    "auto_monthly_from_cycle": False,  # 사용주기→월평균 공급량 자동반영
    "tank_price": 350_000_000.0,
    "monthly_usage_kg": 400_000.0,
    "vaporizer_capacity": 1500.0,
    "vaporizer_qty_note": "* 2ea",
    "vaporizer_price": 31_500_000.0,
    "construction_cost": 190_000_000.0,
    "purchase_unit": 120.0,
    "logistics_unit": 20.0,
    "supply_unit": 210.0,
    "interest_rate": 0.05,
    "mgmt_rate": 0.145,
    "depreciation_months": 120.0,
    "equipment_rent": 0.0,
    "rent_count": 0.0,
    "logi_km": 70.0,
    "logi_fuel_price": 1400.0,
    "logi_efficiency": 2.5,
    "logi_toll": 0.0,
    "logi_roundtrips": 10.0,
    "logi_supply_kg": 20_000.0,
    "logi_origin": "경기도 화성시 마도면 쌍송리 신일가스",
    "logi_dest": "",
}
def load_profit_inputs():
    """수익성분석 입력값 로드 (없으면 엑셀 기본값)."""
    data = dict(PROFIT_DEFAULTS)
    if os.path.exists(PROFIT_CACHE_FILE):
        try:
            with open(PROFIT_CACHE_FILE, "r", encoding="utf-8") as f:
                saved = json.load(f)
            if isinstance(saved, dict):
                data.update({k: saved[k] for k in PROFIT_DEFAULTS if k in saved})
        except Exception:
            pass
    return data
def save_profit_inputs(data):
    try:
        with open(PROFIT_CACHE_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        return True
    except Exception:
        return False
def compute_profitability(p):
    """엑셀 수익성분석 함수 동일 적용.
    C10=C11*12, C20=C9+C16+C18, C27=C26-(C24+C25),
    C30=H36*0.145, C35=F35*H35, C36=F36/120, H36=C11*C26,
    C37=F37*H37/12, C38=F38-H38-C30, C39=C38*3+(F39*H39),
    P18=(J18*K18/L18+M18)*N18/O18
    """
    monthly_usage = float(p.get("monthly_usage_kg") or 0)
    tank_price = float(p.get("tank_price") or 0)
    vap_price = float(p.get("vaporizer_price") or 0)
    const_cost = float(p.get("construction_cost") or 0)
    purchase = float(p.get("purchase_unit") or 0)
    logistics = float(p.get("logistics_unit") or 0)
    supply = float(p.get("supply_unit") or 0)
    rate = float(p.get("interest_rate") or 0)
    mgmt_rate = float(p.get("mgmt_rate") or 0)
    dep_m = float(p.get("depreciation_months") or 120) or 120.0
    rent = float(p.get("equipment_rent") or 0)
    rent_n = float(p.get("rent_count") or 0)
    yearly_usage = monthly_usage * 12  # C10
    total_invest = tank_price + vap_price + const_cost  # C20
    margin_kg = supply - (purchase + logistics)  # C27
    monthly_sales = monthly_usage * supply  # H36
    monthly_gross = monthly_usage * margin_kg  # C35
    depreciation = total_invest / dep_m  # C36
    finance = total_invest * rate / 12  # C37
    mgmt = monthly_sales * mgmt_rate  # C30
    invest_cost = depreciation + finance  # H38
    monthly_profit = monthly_gross - invest_cost - mgmt  # C38
    three_month = monthly_profit * 3 + (rent * rent_n)  # C39
    eff = float(p.get("logi_efficiency") or 0) or 1.0
    supply_kg = float(p.get("logi_supply_kg") or 0) or 1.0
    logi_per_kg = (
        (float(p.get("logi_km") or 0) * float(p.get("logi_fuel_price") or 0) / eff
         + float(p.get("logi_toll") or 0))
        * float(p.get("logi_roundtrips") or 0)
        / supply_kg
    )  # P18
    return {
        "yearly_usage": yearly_usage,
        "total_invest": total_invest,
        "margin_kg": margin_kg,
        "monthly_sales": monthly_sales,
        "monthly_gross": monthly_gross,
        "depreciation": depreciation,
        "finance": finance,
        "mgmt": mgmt,
        "invest_cost": invest_cost,
        "monthly_profit": monthly_profit,
        "three_month": three_month,
        "logi_per_kg": logi_per_kg,
    }
@st.cache_data(show_spinner=False, max_entries=8)
def _cached_profitability_report_excel(p_json, r_json, route_json, diesel_json):
    """동일 입력 반복 시 엑셀 재생성 생략 (Tab9 입력 체감 속도)."""
    return build_profitability_report_excel(
        json.loads(p_json),
        json.loads(r_json),
        route_info=json.loads(route_json),
        diesel_info=json.loads(diesel_json),
    )
def build_profitability_report_excel(p, r, route_info=None, diesel_info=None):
    """스크린샷형 수익성 분석 보고서 엑셀(bytes). Tab9 전용."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = "수익성분석"
    # 열 폭 (보고서 가독성)
    for col, w in enumerate([3, 38, 28, 28, 28, 3], start=1):
        ws.column_dimensions[get_column_letter(col)].width = w

    thin = Border(
        left=Side(style="thin", color="94A3B8"),
        right=Side(style="thin", color="94A3B8"),
        top=Side(style="thin", color="94A3B8"),
        bottom=Side(style="thin", color="94A3B8"),
    )
    dash = Border(
        left=Side(style="dashed", color="94A3B8"),
        right=Side(style="dashed", color="94A3B8"),
        top=Side(style="dashed", color="94A3B8"),
        bottom=Side(style="dashed", color="94A3B8"),
    )
    title_fill = PatternFill("solid", fgColor="F1F5F9")
    green_fill = PatternFill("solid", fgColor="DCFCE7")
    note_font = Font(name="맑은 고딕", size=8, color="64748B")
    body_font = Font(name="맑은 고딕", size=11, color="0F172A")
    bold_font = Font(name="맑은 고딕", size=11, bold=True, color="0F172A")
    head_font = Font(name="맑은 고딕", size=12, bold=True, color="0F172A")
    title_font = Font(name="맑은 고딕", size=16, bold=True, color="0F172A")
    green_font = Font(name="맑은 고딕", size=9, bold=True, color="166534")
    left_align = Alignment(horizontal="left", vertical="center", wrap_text=True)
    center = Alignment(horizontal="center", vertical="center", wrap_text=True)

    def _won(v):
        try:
            return f"{float(v):,.0f}"
        except Exception:
            return "0"

    def _num(v, d=1):
        try:
            return f"{float(v):,.{d}f}"
        except Exception:
            return "0"

    def put(row, col, text, font=None, align=None, fill=None, border=None, merge_to=None):
        cell = ws.cell(row=row, column=col, value=text)
        cell.font = font or body_font
        cell.alignment = align or left_align
        if fill:
            cell.fill = fill
        if border:
            cell.border = border
        if merge_to:
            ws.merge_cells(start_row=row, start_column=col, end_row=row, end_column=merge_to)
            for c in range(col, merge_to + 1):
                ws.cell(row=row, column=c).border = border or Border()
                ws.cell(row=row, column=c).fill = fill or PatternFill()
                ws.cell(row=row, column=c).alignment = align or left_align
        return cell

    name = str(p.get("project_name") or "프로젝트")
    dep_m = float(p.get("depreciation_months") or 120) or 120.0
    dep_years = dep_m / 12.0
    rate = float(p.get("interest_rate") or 0)
    mgmt_rate = float(p.get("mgmt_rate") or 0)
    route_info = route_info or {}
    diesel_info = diesel_info or {}

    # 1) 제목
    put(2, 2, "수익성 분석", font=title_font, align=center, fill=title_fill, border=thin, merge_to=5)
    ws.row_dimensions[2].height = 32

    # 2) 보고 헤더
    put(4, 2, f"◆ [{name}] 투자대비 수익성 분석 보고", font=head_font, merge_to=5)
    put(
        5, 2,
        f"TANK: {p.get('tank_gas','')} {float(p.get('tank_liters') or 0):,.0f} L → "
        f"{parse_tank_capacity_kg(p.get('tank_spec')):,.0f} kg · 기화기 {float(p.get('vaporizer_capacity') or 0):,.0f} Nm3/hr "
        f"{p.get('vaporizer_qty_note','')} · 년 사용량(자동) {_won(r['yearly_usage'])} kg",
        font=note_font, merge_to=5,
    )

    # 3) 장비 투자비
    row = 7
    put(row, 2, "○ 장비 투자비", font=bold_font, merge_to=5)
    row = 8
    _xlsx_kg = parse_tank_capacity_kg(p.get("tank_spec"))
    _xlsx_gas = str(p.get("tank_gas") or "질소")
    _xlsx_nm3 = tank_kg_to_nm3(_xlsx_kg, _xlsx_gas)
    put(
        row, 2,
        f"1. TANK 구매 : {_xlsx_gas} {float(p.get('tank_liters') or 0):,.0f} L "
        f"→ {_xlsx_kg:,.0f} kg · 기체환산 {_xlsx_nm3:,.1f} Nm³",
        font=body_font, merge_to=5,
    )
    row = 9
    put(row, 2, f"   구입가: {_won(p.get('tank_price'))} 원", font=body_font)
    put(row, 3, f"년 사용량: {_won(r['yearly_usage'])} kg", font=body_font)
    put(row, 4, f"월평균 공급량: {_won(p.get('monthly_usage_kg'))} kg", font=body_font)
    row = 10
    put(row, 2, f"2. 기화기 구매 : {float(p.get('vaporizer_capacity') or 0):,.0f} Nm3/hr {p.get('vaporizer_qty_note','')}", font=body_font, merge_to=5)
    row = 11
    put(row, 2, f"   구입가: {_won(p.get('vaporizer_price'))} 원", font=body_font, merge_to=5)
    row = 12
    put(row, 2, f"3. 공사비용: {_won(p.get('construction_cost'))} 원", font=body_font, merge_to=5)
    row = 13
    put(row, 2, f"총 투자금액: {_won(r['total_invest'])} 원", font=bold_font, merge_to=5)
    row = 14
    put(row, 2, "합계 = TANK + 기화기 + 공사비용", font=note_font, merge_to=5)

    # 4) 단가
    row = 16
    put(row, 2, "○ 단가", font=bold_font, merge_to=5)
    row = 17
    put(row, 2, f"1. 매입단가: {_num(p.get('purchase_unit'), 1)} 원/kg", font=body_font, merge_to=3)
    put(row, 4, f"2. 물류비: {_num(p.get('logistics_unit'), 2)} 원/kg", font=body_font, merge_to=5)
    row = 18
    put(row, 2, f"3. 공급단가: {_num(p.get('supply_unit'), 1)} 원/kg", font=body_font, merge_to=3)
    put(row, 4, f"4. 매출이익: {_num(r['margin_kg'], 1)} 원/kg", font=body_font, merge_to=5)
    row = 19
    put(row, 2, "매출이익 = 공급단가 − (매입단가 + 물류비)", font=note_font, merge_to=5)

    # 5) 일반관리비
    row = 21
    put(row, 2, "○ 일반관리비", font=bold_font)
    put(row, 3, f"{_won(r['mgmt'])} 원/월", font=bold_font)
    put(row, 4, f"(월 매출 {mgmt_rate*100:.1f}%)", font=note_font, merge_to=5)
    row = 22
    put(row, 2, "일반관리비 = 월매출 × 관리비율  ·  월매출 = 월평균공급량 × 공급단가", font=note_font, merge_to=5)

    # 6) 사용량 대비 영업이익 표
    row = 24
    put(row, 2, "○ 사용량 대비 영업이익", font=bold_font, merge_to=5)
    table = [
        [
            f"월평균 매출이익: {_won(r['monthly_gross'])} 원/월",
            f"월 사용량(kg): {_won(p.get('monthly_usage_kg'))}",
            f"매출마진: {_num(r['margin_kg'], 1)}",
        ],
        [
            "월사용량 × 매출이익",
            "프로젝트 월 평균 공급량",
            "공급 − (매입 + 물류)",
        ],
        [
            f"장비 감가상각비({dep_years:.0f}년): {_won(r['depreciation'])} 원/월",
            f"투자총액: {_won(r['total_invest'])}",
            f"월 매출: {_won(r['monthly_sales'])}",
        ],
        [
            f"투자합계 ÷ {dep_m:.0f}개월",
            "TANK+기화기+공사",
            "월사용량 × 공급단가",
        ],
        [
            f"금융비: {_won(r['finance'])} 원/월",
            f"원금: {_won(r['total_invest'])}",
            f"이자율: {rate*100:.1f}%",
        ],
        [
            "원금 × 이자율 ÷ 12",
            "총 투자금액",
            "입력 이자율",
        ],
        [
            f"월평균 이익금: {_won(r['monthly_profit'])} 원/월",
            f"투자비용: {_won(r['invest_cost'])}",
            f"일반관리비: {_won(r['mgmt'])}",
        ],
        [
            "매출이익 − 투자비용 − 관리비",
            "감가상각 + 금융비",
            f"월매출 × {mgmt_rate*100:.1f}%",
        ],
        [
            f"최근 3개월 매출 실이익: {_won(r['three_month'])} 원",
            f"장비 임대료: {_won(p.get('equipment_rent'))}",
            f"부가횟수: {_num(p.get('rent_count'), 0)}",
        ],
        [
            "월이익×3 + 임대료×횟수",
            "입력 임대료",
            "입력 부가횟수",
        ],
    ]
    t_row = 25
    for i, cells in enumerate(table):
        is_note = (i % 2 == 1)
        for j, text in enumerate(cells):
            c = put(
                t_row, 2 + j, text,
                font=note_font if is_note else body_font,
                border=dash,
                align=left_align,
            )
            if not is_note:
                c.font = bold_font if j == 0 and ("이익금" in text or "실이익" in text or "매출이익" in text) else body_font
        ws.row_dimensions[t_row].height = 18 if is_note else 22
        t_row += 1

    # 7) 물류비 계산 설명 (스크린샷 작은글씨 전부)
    t_row += 1
    put(t_row, 2, "○ 물류비 계산 상세 (20톤 벌크로리 · 경유)", font=bold_font, merge_to=5)
    t_row += 1
    put(
        t_row, 2,
        f"거리 { _num(p.get('logi_km'), 1) } km · 유류비 {_num(p.get('logi_fuel_price'), 2)} 원/L · "
        f"연비 {_num(p.get('logi_efficiency'), 1)} km/L · 통행료(편도·5종) {_won(p.get('logi_toll'))} 원 · "
        f"왕복 {_num(p.get('logi_roundtrips'), 0)} 회 · 월평균공급량 {_won(p.get('logi_supply_kg'))} kg",
        font=note_font, merge_to=5,
    )
    t_row += 1
    put(
        t_row, 2,
        f"물류비(원/kg) = ((거리 × 유류비 ÷ 연비) + 통행료) × 왕복 ÷ 월평균공급량  =  {_num(r['logi_per_kg'], 2)} 원/kg",
        font=note_font, merge_to=5,
    )
    t_row += 1
    _rt_xlsx = compute_roundtrips_from_tank(
        p.get("monthly_usage_kg"), p.get("tank_spec"), fill_ratio=0.8
    )
    put(
        t_row, 2,
        _rt_xlsx.get("message")
        or "왕복횟수 = ceil(월평균공급량 ÷ (탱크총용량kg × 80%))",
        font=note_font, merge_to=5,
    )
    t_row += 1
    route_msg = str(route_info.get("message") or "")
    if route_msg:
        put(t_row, 2, route_msg, font=green_font, fill=green_fill, border=thin, merge_to=5)
        ws.row_dimensions[t_row].height = 22
        t_row += 1
    o_lab = route_info.get("origin_label") or p.get("logi_origin") or ""
    d_lab = route_info.get("dest_label") or p.get("logi_dest") or ""
    if o_lab or d_lab:
        put(
            t_row, 2,
            f"출발: {o_lab}  →  도착: {d_lab}"
            + (f"  ·  {route_msg}" if route_msg else ""),
            font=note_font, merge_to=5,
        )
        t_row += 1
    if diesel_info.get("ok"):
        put(
            t_row, 2,
            f"경유 시장가 {_num(diesel_info.get('price'), 2)} 원/L  ·  "
            f"{diesel_info.get('source','')}  ·  기준일 {diesel_info.get('asof','')}  ·  "
            f"{diesel_info.get('month', datetime.date.today().strftime('%Y-%m'))} 자동반영",
            font=note_font, merge_to=5,
        )
        t_row += 1
    elif diesel_info.get("message"):
        put(t_row, 2, str(diesel_info.get("message")), font=note_font, merge_to=5)
        t_row += 1

    t_row += 1
    put(
        t_row, 2,
        "적용 함수: 년사용량=월×12 · 합계=탱크+기화기+공사 · 매출이익=공급−(매입+물류) · "
        "관리비=월매출×비율 · 감가=투자÷개월 · 금융=원금×이자÷12 · "
        "월이익=매출이익−(감가+금융)−관리비 · 3개월=월이익×3+(임대×횟수) · "
        "통행료=카카오1종×2.5(20톤벌크 5종)",
        font=note_font, merge_to=5,
    )
    t_row += 2
    put(t_row, 2, "1/1페이지", font=note_font, align=center, fill=title_fill, border=thin, merge_to=5)

    ws.print_title_rows = "1:2"
    ws.page_setup.fitToPage = True
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 1
    ws.sheet_properties.pageSetUpPr.fitToPage = True

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()
def get_saved_date(file_path):
    if os.path.exists(file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return datetime.datetime.strptime(f.read().strip(), "%Y-%m-%d").date()
        except:
            pass
    return datetime.date.today()
def set_saved_date(file_path, date_val):
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(date_val.strftime("%Y-%m-%d"))
# ==========================================
# ★ 탭 전체 적용 표 합계행 렌더링 유틸리티 (무손실 보존) ★
# ==========================================
def get_display_df_with_sum(df, sum_label="연간 합계", text_cols=None):
    if df is None or df.empty: return df
    disp = df.copy()
    numeric_cols = disp.select_dtypes(include=[np.number]).columns
    sum_series = disp[numeric_cols].sum()
    disp.loc[sum_label] = sum_series
    if text_cols:
        for col in text_cols:
            if col in disp.columns:
                disp.at[sum_label, col] = "총 합계"
    return disp
def style_with_sum(disp_df, fmt_str, cmap=None, subset_cols=None, axis=None):
    if disp_df.empty:
        return disp_df.style
    if subset_cols is None:
        subset_cols = disp_df.select_dtypes(include=[np.number]).columns
    grad_subset = pd.IndexSlice[disp_df.index[:-1], subset_cols]
    styled = disp_df.style.format(fmt_str)
    if cmap:
        styled = styled.background_gradient(cmap=cmap, axis=axis, subset=grad_subset)
    styled = styled.apply(
        lambda s: ['font-weight: 800; background-color: #E2E8F0; color: #0F172A;'] * len(s),
        subset=pd.IndexSlice[disp_df.index[-1], :],
        axis=1,
    )
    return styled
def inject_top30_month_bridge():
    """iframe(월 헤더 클릭) → 부모 URL query 갱신 + iPad Top30 도넛 스택. sticky/맥 로직 분리."""
    components.html(
        """
        <script>
        (function () {
            var parentWin = window.parent;
            if (!parentWin) return;
            var parentDoc = parentWin.document;
            function isTouchPad() {
                try {
                    var ua = parentWin.navigator.userAgent || "";
                    var ios = /iPad|iPhone|iPod/.test(ua);
                    var ipadOs = parentWin.navigator.platform === "MacIntel" && parentWin.navigator.maxTouchPoints > 1;
                    return ios || ipadOs || parentDoc.documentElement.classList.contains("dashboard-touch-mode");
                } catch (e) { return false; }
            }
            /* iPad touch_ui: 세션당 1회 reload만. 맥에서 touch_ui 삭제 reload 금지(루프 방지) */
            try {
                var touch = isTouchPad();
                if (touch) {
                    try { parentDoc.cookie = "dashboard_touch=1; path=/; max-age=31536000; SameSite=Lax"; } catch (eCk2) {}
                }
                if (!parentWin.__dashboardTouchUiSynced) {
                    parentWin.__dashboardTouchUiSynced = true;
                    var url = new URL(parentWin.location.href);
                    var cur = url.searchParams.get("touch_ui");
                    if (touch && cur !== "1") {
                        url.searchParams.set("touch_ui", "1");
                        try { parentWin.history.replaceState(null, "", url.toString()); } catch (eHs2) {}
                        try {
                            if (parentWin.sessionStorage.getItem("__dash_touch_boot") !== "1") {
                                parentWin.sessionStorage.setItem("__dash_touch_boot", "1");
                                parentWin.location.replace(url.toString());
                                return;
                            }
                        } catch (eSs2) {}
                    }
                }
            } catch (eSync) {}
            if (!parentWin.__dashboardTop30MonthBridge) {
                parentWin.__dashboardTop30MonthBridge = true;
                parentWin.addEventListener("message", function (e) {
                    try {
                        var d = e.data;
                        if (!d || d.type !== "dashboard-top30-month" || !d.month) return;
                        var url2 = new URL(parentWin.location.href);
                        url2.searchParams.set("top30_month", d.month);
                        parentWin.location.href = url2.toString();
                    } catch (err) {}
                });
            }
            /* iPad(touch-mode)만: Top30 표+도넛 행을 세로 100% 폭 — 맥은 즉시 return */
            function forceFullWidth(el) {
                if (!el || !el.style) return;
                el.style.setProperty("width", "100%", "important");
                el.style.setProperty("min-width", "100%", "important");
                el.style.setProperty("max-width", "100%", "important");
                el.style.setProperty("flex", "1 1 100%", "important");
                el.style.setProperty("box-sizing", "border-box", "important");
            }
            function applyTop30FullWidthRow(row) {
                row.classList.add("top30-touch-row");
                row.style.setProperty("display", "flex", "important");
                row.style.setProperty("flex-direction", "column", "important");
                row.style.setProperty("flex-wrap", "nowrap", "important");
                row.style.setProperty("align-items", "stretch", "important");
                row.style.setProperty("width", "100%", "important");
                row.style.setProperty("max-width", "100%", "important");
                Array.prototype.forEach.call(row.children, forceFullWidth);
                /* 직계(표|도넛) 칼럼만 — 중첩 칼럼에 min-width:100% 금지 */
                var topCols = row.querySelectorAll(':scope > [data-testid="column"], :scope > [data-testid="stColumn"], :scope > div > [data-testid="column"], :scope > div > [data-testid="stColumn"]');
                if (!topCols.length) {
                    topCols = [];
                    Array.prototype.forEach.call(row.children, function (ch) {
                        if (ch.matches && (ch.matches('[data-testid="column"]') || ch.matches('[data-testid="stColumn"]'))) {
                            topCols.push(ch);
                        } else if (ch.querySelector) {
                            var inner = ch.querySelector(':scope > [data-testid="column"], :scope > [data-testid="stColumn"]');
                            if (inner) topCols.push(inner);
                        }
                    });
                }
                Array.prototype.forEach.call(topCols, forceFullWidth);
                row.querySelectorAll("iframe").forEach(function (f) {
                    f.style.setProperty("width", "100%", "important");
                    f.style.setProperty("max-width", "100%", "important");
                    try { f.setAttribute("width", "100%"); } catch (e0) {}
                });
                var plots = row.querySelectorAll('[data-testid="stPlotlyChart"]');
                plots.forEach(function (p) {
                    p.style.setProperty("width", "100%", "important");
                    p.style.setProperty("min-height", "420px", "important");
                    p.style.setProperty("height", "420px", "important");
                });
                try { parentWin.dispatchEvent(new Event("resize")); } catch (e1) {}
                setTimeout(function () {
                    try { parentWin.dispatchEvent(new Event("resize")); } catch (e2) {}
                    plots.forEach(function (p) {
                        var ifr = p.querySelector("iframe");
                        if (!ifr) return;
                        try {
                            var pw = ifr.contentWindow;
                            if (pw && pw.Plotly) {
                                var gds = ifr.contentDocument.querySelectorAll(".js-plotly-plot");
                                for (var g = 0; g < gds.length; g++) {
                                    try { pw.Plotly.Plots.resize(gds[g]); } catch (e3) {}
                                }
                            }
                        } catch (e4) {}
                    });
                }, 250);
            }
            function stackTop30DonutRow() {
                try {
                    if (!isTouchPad()) return;
                    var flag = parentDoc.querySelector(".top30-section-flag");
                    if (!flag) return;
                    var scope = flag.closest('[data-testid="stVerticalBlock"]') || flag.parentElement;
                    if (!scope) return;
                    scope.classList.add("top30-touch-scope");
                    forceFullWidth(scope);
                    var rows = scope.querySelectorAll('[data-testid="stHorizontalBlock"]');
                    for (var i = 0; i < rows.length; i++) {
                        var row = rows[i];
                        var cols = row.querySelectorAll(':scope > [data-testid="column"], :scope > [data-testid="stColumn"], :scope > div > [data-testid="column"], :scope > div > [data-testid="stColumn"]');
                        if (!cols.length) {
                            cols = row.querySelectorAll('[data-testid="column"], [data-testid="stColumn"]');
                        }
                        /* 표|도넛(플롯 포함) 행만 풀폭 — 월 버튼 12칸 행은 제외 */
                        if (cols.length === 12) continue;
                        var hasPlot = !!row.querySelector('[data-testid="stPlotlyChart"]');
                        if (hasPlot) {
                            applyTop30FullWidthRow(row);
                        }
                    }
                } catch (err) {}
            }
            stackTop30DonutRow();
            if (!parentWin.__dashboardTop30StackTimer) {
                parentWin.__dashboardTop30StackTimer = parentWin.setInterval(stackTop30DonutRow, 600);
            }
        })();
        </script>
        """,
        height=0,
    )
def is_touch_ui():
    """iPad UI 분기. query touch_ui / cookie / UA / session. 맥 데스크톱은 False."""
    try:
        if st.session_state.get("force_touch_ui") is True:
            return True
        v = st.query_params.get("touch_ui", "")
        if isinstance(v, list):
            v = v[0] if v else ""
        if str(v) == "1":
            st.session_state["force_touch_ui"] = True
            return True
    except Exception:
        pass
    try:
        cookies = getattr(st.context, "cookies", None)
        if cookies is not None and str(cookies.get("dashboard_touch", "")) == "1":
            st.session_state["force_touch_ui"] = True
            return True
    except Exception:
        pass
    # Safari iPad는 쿠키/query 전에 UA로도 인식 (로딩 분기·감사 추출 지연에 사용)
    try:
        headers = getattr(st.context, "headers", None)
        ua = ""
        if headers is not None:
            ua = str(headers.get("User-Agent") or headers.get("user-agent") or "")
        if ua and (
            re.search(r"iPad|iPhone|iPod", ua)
            or ("Macintosh" in ua and "Mobile" in ua)
        ):
            st.session_state["force_touch_ui"] = True
            return True
    except Exception:
        pass
    return bool(st.session_state.get("force_touch_ui"))
def render_plotly_chart(fig, *, key=None, use_container_width=True, allow_drag=False, height=None, **kwargs):
    """맥: st.plotly_chart 무손실.
    iPad: components.html 직접 렌더 → 컨트롤바는 그래프 탭 시에만 표시.
    """
    if not is_touch_ui():
        st.plotly_chart(fig, use_container_width=use_container_width, key=key, **kwargs)
        return
    try:
        layout_h = fig.layout.height
    except Exception:
        layout_h = None
    try:
        h = int(height or layout_h or 450)
    except Exception:
        h = 450
    try:
        updates = {
            "clickmode": "none",
            "margin": dict(t=56, r=12, b=12, l=12),
            "modebar": dict(
                orientation="h",
                bgcolor="rgba(15, 23, 42, 0.92)",
                color="#E2E8F0",
                activecolor="#38BDF8",
            ),
        }
        if not allow_drag:
            updates["dragmode"] = False
        fig.update_layout(**updates)
    except Exception:
        pass
    config = {
        "displayModeBar": True,
        "displaylogo": False,
        "responsive": True,
        "scrollZoom": False,
        "doubleClick": "reset",
        "modeBarButtonsToRemove": ["lasso2d", "select2d"],
    }
    try:
        inner = fig.to_html(include_plotlyjs="cdn", full_html=False, config=config)
    except Exception:
        st.plotly_chart(fig, use_container_width=use_container_width, key=key, config=config, **kwargs)
        return
    page_html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<style>
  html, body {{ margin:0; padding:0; background:#fff; overflow:hidden; }}
  .wrap {{ position:relative; width:100%; }}
  .modebar-container, .modebar {{
    opacity:0 !important;
    visibility:hidden !important;
    pointer-events:none !important;
    transition: opacity .15s ease;
    z-index:9999 !important;
  }}
  body.mb-show .modebar-container,
  body.mb-show .modebar {{
    opacity:1 !important;
    visibility:visible !important;
    display:flex !important;
    pointer-events:auto !important;
  }}
  .modebar-container {{
    top:8px !important; right:8px !important; left:auto !important;
  }}
  .modebar {{
    background:rgba(15,23,42,.94) !important;
    border-radius:8px !important; padding:4px !important;
  }}
  .modebar-btn {{ min-width:32px !important; min-height:32px !important; }}
</style></head>
<body>
<div class="wrap" id="plot-wrap">
  {inner}
</div>
<script>
(function() {{
  var body = document.body;
  function showBar() {{ body.classList.add("mb-show"); }}
  function onTap(e) {{
    if (e.target && e.target.closest && e.target.closest(".modebar")) return;
    showBar();
  }}
  document.addEventListener("click", onTap, true);
  document.addEventListener("touchend", onTap, true);
}})();
</script>
</body></html>"""
    components.html(page_html, height=h + 8, scrolling=False)
def inject_ipad_plotly_controls():
    """iPad: touch cookie/query 동기화 + 예전 원복 플로팅 버튼 제거. 맥 무손실."""
    components.html(
        """
        <script>
        (function () {
            var parentWin = window.parent;
            if (!parentWin) return;
            var parentDoc = parentWin.document;
            function isTouchPad() {
                try {
                    var ua = parentWin.navigator.userAgent || "";
                    var ios = /iPad|iPhone|iPod/.test(ua);
                    var ipadOs = parentWin.navigator.platform === "MacIntel" && parentWin.navigator.maxTouchPoints > 1;
                    return ios || ipadOs || parentDoc.documentElement.classList.contains("dashboard-touch-mode");
                } catch (e) { return false; }
            }
            try {
                var old = parentDoc.getElementById("dashboard-ipad-plotly-reset");
                if (old) old.remove();
            } catch (e0) {}
            if (!isTouchPad()) return;
            try {
                parentDoc.cookie = "dashboard_touch=1; path=/; max-age=31536000; SameSite=Lax";
            } catch (e1) {}
            try {
                if (!parentWin.__dashboardTouchUiSynced) {
                    parentWin.__dashboardTouchUiSynced = true;
                    var url = new URL(parentWin.location.href);
                    if (url.searchParams.get("touch_ui") !== "1") {
                        url.searchParams.set("touch_ui", "1");
                        try { parentWin.history.replaceState(null, "", url.toString()); } catch (eHs3) {}
                        try {
                            if (parentWin.sessionStorage.getItem("__dash_touch_boot") !== "1") {
                                parentWin.sessionStorage.setItem("__dash_touch_boot", "1");
                                parentWin.location.replace(url.toString());
                                return;
                            }
                        } catch (eSs3) {}
                    }
                }
            } catch (e2) {}
        })();
        </script>
        """,
        height=0,
    )
def render_frozen_styler_html(
    styled,
    height=450,
    freeze_left_n=2,
    freeze_widths=None,
    clickable_cols=None,
    query_param=None,
    active_col=None,
):
    """Styler HTML 렌더 + 헤더/좌측열 틀고정 (데이터·스타일 무손실)."""
    widths = freeze_widths or ([44, 160] if freeze_left_n >= 2 else [160])
    left_css = []
    left = 0
    for i in range(freeze_left_n):
        w = widths[i] if i < len(widths) else 100
        n = i + 1
        left_css.append(
            f"""
            table thead tr th:nth-child({n}),
            table tbody tr th:nth-child({n}),
            table tbody tr td:nth-child({n}) {{
                position: sticky;
                left: {left}px;
                z-index: 4;
                min-width: {w}px;
                max-width: {w}px;
                box-shadow: 2px 0 0 #CBD5E1;
            }}
            table thead tr th:nth-child({n}) {{
                z-index: 8;
                background: #F0F2F6 !important;
            }}
            table tbody tr th:nth-child({n}),
            table tbody tr td:nth-child({n}) {{
                background-color: #FFFFFF;
                z-index: 5;
            }}
            table tbody tr:last-child th:nth-child({n}),
            table tbody tr:last-child td:nth-child({n}) {{
                background-color: #E2E8F0 !important;
            }}
            """
        )
        left += w
    table_html = styled.to_html()
    click_js = ""
    if clickable_cols and query_param:
        def _link_thead(match):
            block = match.group(0)
            for col in clickable_cols:
                active = "color:#B91C1C;font-weight:700;" if active_col == col else "color:inherit;font-weight:600;"
                block = re.sub(
                    rf"(<th[^>]*>)\s*{re.escape(str(col))}\s*(</th>)",
                    rf'\1<a href="#" role="button" data-top30-month="{col}" '
                    rf'style="text-decoration:none;cursor:pointer;display:block;{active}" '
                    rf'title="클릭: {col} 매출 순위·그래프 적용">{col}</a>\2',
                    block,
                )
            return block
        table_html = re.sub(r"<thead>.*?</thead>", _link_thead, table_html, count=1, flags=re.DOTALL)
        click_js = """
        <script>
        (function () {
            function notify(month) {
                try {
                    window.parent.postMessage({ type: "dashboard-top30-month", month: month }, "*");
                } catch (err) {}
            }
            document.querySelectorAll("a[data-top30-month]").forEach(function (a) {
                a.addEventListener("click", function (e) {
                    e.preventDefault();
                    e.stopPropagation();
                    notify(a.getAttribute("data-top30-month"));
                });
            });
        })();
        </script>
        """
    # 영업종합요약 st.dataframe(히트맵)과 동일 계열 글꼴·크기
    page_html = f"""
    <!DOCTYPE html>
    <html><head><meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <link rel="stylesheet" href="https://fonts.googleapis.com/css2?family=Source+Sans+3:wght@400;600;700&display=swap">
    <style>
        html, body {{
            margin: 0; height: 100%; overflow: hidden;
            font-family: "Source Sans 3", "Source Sans Pro", "Segoe UI", Roboto,
                "Helvetica Neue", Arial, "Apple SD Gothic Neo", "Noto Sans KR", sans-serif;
            font-size: 13px; color: #31333F;
            -webkit-font-smoothing: antialiased;
        }}
        .wrap {{
            height: 100%; overflow: auto;
            -webkit-overflow-scrolling: touch;
            border: 1px solid #E2E8F0; border-radius: 4px; background: #fff;
        }}
        table {{
            border-collapse: separate; border-spacing: 0;
            width: max-content; min-width: 100%;
            font-family: inherit;
            font-size: 13px;
        }}
        th, td {{
            padding: 6px 8px; white-space: nowrap;
            border-bottom: 1px solid #E2E8F0;
            font-family: inherit;
            font-size: 13px;
            font-weight: 400;
            font-variant-numeric: tabular-nums;
            text-align: center;
            vertical-align: middle;
            line-height: 1.35;
        }}
        thead th {{
            position: sticky; top: 0; z-index: 6;
            background: #F0F2F6 !important;
            box-shadow: 0 1px 0 #CBD5E1;
            font-weight: 600; text-align: center;
            font-size: 13px;
        }}
        table thead tr th:nth-child(1),
        table thead tr th:nth-child(2),
        table tbody tr th:nth-child(1),
        table tbody tr th:nth-child(2),
        table tbody tr td:nth-child(1),
        table tbody tr td:nth-child(2) {{
            text-align: left;
        }}
        thead th a:hover {{ text-decoration: underline !important; }}
        {''.join(left_css)}
    </style></head>
    <body><div class="wrap">{table_html}</div>
    <script>
    (function () {{
      /* 상단 st.dataframe과 동일 글꼴 우선 적용 */
      try {{
        var pDoc = window.parent && window.parent.document;
        if (pDoc) {{
          var pf = window.parent.getComputedStyle(pDoc.body).fontFamily;
          if (pf) document.body.style.fontFamily = '"Source Sans 3",' + pf;
        }}
      }} catch (e0) {{}}
      /* 어두운 히트맵 칸 → 흰 글씨 (상단 st.dataframe Blues와 동일 가독성) */
      function parseRgb(c) {{
        if (!c) return null;
        c = String(c).trim();
        if (c.indexOf("rgb") === 0) {{
          var m = c.match(/rgba?\\((\\d+),\\s*(\\d+),\\s*(\\d+)/);
          return m ? [+m[1], +m[2], +m[3]] : null;
        }}
        if (c.charAt(0) === "#" && (c.length === 7 || c.length === 4)) {{
          if (c.length === 4) c = "#" + c[1]+c[1]+c[2]+c[2]+c[3]+c[3];
          return [parseInt(c.slice(1,3),16), parseInt(c.slice(3,5),16), parseInt(c.slice(5,7),16)];
        }}
        return null;
      }}
      function applyContrast(el) {{
        /* 선택 월(빨강) 강조는 유지 */
        var inline = (el.getAttribute("style") || "").toLowerCase();
        if (inline.indexOf("#b91c1c") >= 0 || inline.indexOf("rgb(185, 28, 28)") >= 0) return;
        var bg = el.style.backgroundColor || window.getComputedStyle(el).backgroundColor;
        var rgb = parseRgb(bg);
        if (!rgb) return;
        var lum = (0.299*rgb[0] + 0.587*rgb[1] + 0.114*rgb[2]) / 255;
        if (lum < 0.55) el.style.color = "#FFFFFF";
      }}
      document.querySelectorAll("td, th").forEach(applyContrast);
    }})();
    </script>
    {click_js}</body></html>
    """
    components.html(page_html, height=height, scrolling=True)
# ==========================================
# 5. 메인 실행 흐름 및 영구 캐싱 관리
# ==========================================
inject_custom_css()
st.sidebar.header("📁 데이터 업로드 및 유지")
_render_cloud_sync_banner()
# Drive「dashboard 복사본」→ uploaded_cache (맥: Drive 마운트 / 클라우드: 배포 시드)
_drive_autoload_res = None
if sync_drive_copy_into_cache is not None:
    try:
        _drive_autoload_res = sync_drive_copy_into_cache(CACHE_DIR)
    except Exception as _dae:
        _drive_autoload_res = {"ok": False, "error": str(_dae)}
if isinstance(_drive_autoload_res, dict):
    if _drive_autoload_res.get("ok") and _drive_autoload_res.get("source"):
        _n = len(_drive_autoload_res.get("copied") or [])
        st.sidebar.caption(
            f"Drive 복사본 자동로드"
            + (f" · 갱신 {_n}개" if _n else " · 캐시 유지")
        )
    elif _drive_autoload_res.get("ok") and _drive_autoload_res.get("skipped"):
        st.sidebar.caption("캐시/시드 자동로드 (Drive 경로 없음)")
    elif not _drive_autoload_res.get("ok"):
        st.sidebar.warning(
            f"Drive 자동로드 실패: {_drive_autoload_res.get('error') or '알 수 없음'}"
        )
address_file_up = st.sidebar.file_uploader("거래처 주소록 (CSV)", type=["csv"])
industry_file_up = st.sidebar.file_uploader("🏢 거래처 업종 분류 (CSV)", type=["csv"])
debt_file_up = st.sidebar.file_uploader("채권 데이터 (채권.csv)", type=["csv"])
uploaded_files_up = st.sidebar.file_uploader("매출 데이터 (다중 업로드)", type=["csv"], accept_multiple_files=True)
st.sidebar.markdown("---")
st.sidebar.subheader("🏭 설비 재고 관리 (선택)")
tank_file_up = st.sidebar.file_uploader("탱크 재고현황 (CSV/Excel)", type=["csv", "xlsx"])
vaporizer_file_up = st.sidebar.file_uploader("기화기 재고현황 (CSV/Excel)", type=["csv", "xlsx"])
st.sidebar.markdown("---")
st.sidebar.subheader("🛢️ 통합 탱크 재고")
local_int_path = "통합탱크재고.csv"
int_bytes = None
int_name = ""
if os.path.exists(local_int_path):
    st.sidebar.success(f"✅ '{local_int_path}' 자동 로드됨")
    with open(local_int_path, "rb") as f:
        int_bytes = f.read()
    int_name = local_int_path
else:
    integrated_file_up = st.sidebar.file_uploader("통합 탱크 재고현황 (CSV)", type=["csv"])
    if integrated_file_up is not None:
        int_bytes = integrated_file_up.getvalue()
        int_name = integrated_file_up.name
@st.cache_data(show_spinner="설비 데이터를 읽어오는 중입니다...")
def load_equipment_file(file_bytes, file_name):
    if not file_bytes:
        return pd.DataFrame()
    try:
        if file_name.endswith('.csv'):
            decoded = None
            for enc in ['utf-8-sig', 'utf-8', 'cp949', 'euc-kr']:
                try:
                    decoded = file_bytes.decode(enc)
                    break
                except UnicodeDecodeError:
                    continue
            
            if decoded is None:
                decoded = file_bytes.decode('utf-8', errors='replace')
            
            lines = [line.strip().strip('"') for line in decoded.splitlines() if line.strip()]
            
            df = pd.read_csv(io.StringIO("\n".join(lines)), on_bad_lines='skip', engine='python')
            return df
        else:
            return pd.read_excel(io.BytesIO(file_bytes))
    except Exception as e:
        st.sidebar.error(f"파일 읽기 오류 ({file_name}): {e}")
        return pd.DataFrame()
st.sidebar.markdown("---")
st.sidebar.subheader("🔑 Open DART API 설정")
API_KEY_FILE = os.path.join(CACHE_DIR, "dart_api_key.txt")
API_KEY_DESKTOP = os.path.expanduser("~/Desktop/uploaded_cache/dart_api_key.txt")
DART_KEY_COOKIE = "dashboard_dart_api_key"


def _read_dart_key_file(path):
    try:
        if path and os.path.exists(path):
            with open(path, "r", encoding="utf-8") as f:
                return f.read().strip()
    except Exception:
        return ""
    return ""


def _write_dart_key_file(path, key):
    try:
        if not path or not key:
            return
        parent = os.path.dirname(path)
        if parent:
            os.makedirs(parent, exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            f.write(key.strip())
    except Exception:
        pass


def _read_dart_key_cookie():
    try:
        cookies = getattr(st.context, "cookies", None)
        if cookies is None:
            return ""
        return str(cookies.get(DART_KEY_COOKIE, "") or "").strip()
    except Exception:
        return ""


def _sync_dart_key_cookie(key):
    """브라우저 쿠키에 API 키 저장 — 맥/iPad 재시작·새로고침 후 자동 복원."""
    if not key:
        return
    safe = (
        str(key)
        .strip()
        .replace("\\", "\\\\")
        .replace("'", "\\'")
        .replace("\n", "")
        .replace("\r", "")
    )
    if not safe:
        return
    components.html(
        f"""
        <script>
        (function() {{
            try {{
                var k = '{safe}';
                var parentDoc = window.parent.document;
                parentDoc.cookie = "{DART_KEY_COOKIE}=" + encodeURIComponent(k)
                    + "; path=/; max-age=31536000; SameSite=Lax";
                try {{ window.parent.localStorage.setItem("{DART_KEY_COOKIE}", k); }} catch (e) {{}}
            }} catch (e) {{}}
        }})();
        </script>
        """,
        height=0,
    )


def _load_saved_dart_api_key():
    """파일 → 데스크톱 캐시 → 브라우저 쿠키 순으로 복원."""
    for path in (API_KEY_FILE, API_KEY_DESKTOP):
        v = _read_dart_key_file(path)
        if v:
            return v
    return _read_dart_key_cookie()


def _persist_dart_api_key(key):
    key = (key or "").strip()
    if not key:
        return
    _write_dart_key_file(API_KEY_FILE, key)
    _write_dart_key_file(API_KEY_DESKTOP, key)
    _sync_dart_key_cookie(key)


saved_api_key = _load_saved_dart_api_key()
# 쿠키/데스크톱에서만 온 키면 서버 파일에도 동기화
if saved_api_key and not _read_dart_key_file(API_KEY_FILE):
    _write_dart_key_file(API_KEY_FILE, saved_api_key)
# localStorage → 쿠키만 세팅 (reload 금지). iframe은 세션당 1회만 (iPad 새로고침 부담 완화)
if not saved_api_key and not st.session_state.get("_dart_ls_cookie_tried"):
    st.session_state["_dart_ls_cookie_tried"] = True
    components.html(
        f"""
        <script>
        (function() {{
            try {{
                var name = "{DART_KEY_COOKIE}";
                var parentWin = window.parent;
                var parentDoc = parentWin.document;
                var fromLs = "";
                try {{ fromLs = parentWin.localStorage.getItem(name) || ""; }} catch (e) {{}}
                if (!fromLs) return;
                var hasCookie = (parentDoc.cookie || "").indexOf(name + "=") >= 0;
                if (hasCookie) return;
                parentDoc.cookie = name + "=" + encodeURIComponent(fromLs)
                    + "; path=/; max-age=31536000; SameSite=Lax";
            }} catch (e) {{}}
        }})();
        </script>
        """,
        height=0,
    )
# text_input key 세션과 저장된 키 동기화 (재시작 자동 입력)
if "sidebar_dart_api_key" not in st.session_state:
    st.session_state["sidebar_dart_api_key"] = saved_api_key
elif saved_api_key and not str(st.session_state.get("sidebar_dart_api_key") or "").strip():
    st.session_state["sidebar_dart_api_key"] = saved_api_key
dart_api_key = st.sidebar.text_input(
    "DART API 키 (재무정보 연동용)",
    type="password",
    help="금융감독원 Open DART API 키. 파일+브라우저에 저장되어 맥/iPad 재시작 후에도 유지됩니다.",
    key="sidebar_dart_api_key",
)
# 키가 바뀌었을 때만 저장/쿠키 동기화 (매 rerun HTML 주입 방지)
if dart_api_key and dart_api_key != saved_api_key:
    _persist_dart_api_key(dart_api_key)
elif dart_api_key and not _read_dart_key_file(API_KEY_FILE):
    _write_dart_key_file(API_KEY_FILE, dart_api_key)
if OpenDartReader is None:
    if is_touch_ui():
        st.sidebar.caption("DART 재무연동은 선택 사항입니다. (미연결 시 해당 기능만 비활성)")
    else:     
        pass
elif dart_api_key:
    st.sidebar.caption("✓ DART 연동 준비됨 (거래처 분석 → 기업 재무정보)")
else:
    st.sidebar.caption("API 키 입력 시 매출액·영업이익 조회 가능")
addr_cache_path = os.path.join(CACHE_DIR, "address.csv")
industry_cache_path = os.path.join(CACHE_DIR, "industry.csv")
debt_cache_path = os.path.join(CACHE_DIR, "debt.csv")
tank_cache_path = os.path.join(CACHE_DIR, "tank_cache.dat")
vaporizer_cache_path = os.path.join(CACHE_DIR, "vaporizer_cache.dat")
integrated_cache_path = os.path.join(CACHE_DIR, "integrated_cache.dat")
sales_cache_dir = os.path.join(CACHE_DIR, "sales")
os.makedirs(sales_cache_dir, exist_ok=True)
if address_file_up is not None:
    addr_bytes = address_file_up.getvalue()
    with open(addr_cache_path, "wb") as f: f.write(addr_bytes)
elif os.path.exists(addr_cache_path):
    with open(addr_cache_path, "rb") as f: addr_bytes = f.read()
elif os.path.exists("주소.csv"):
    with open("주소.csv", "rb") as f: addr_bytes = f.read()
else:
    # Desktop 쪽 예전 캐시 경로 자동 복구
    _desktop_addr = os.path.expanduser("~/Desktop/uploaded_cache/address.csv")
    if os.path.exists(_desktop_addr):
        with open(_desktop_addr, "rb") as f:
            addr_bytes = f.read()
        with open(addr_cache_path, "wb") as f:
            f.write(addr_bytes)
    else:
        addr_bytes = None
if industry_file_up is not None:
    ind_bytes = industry_file_up.getvalue()
    with open(industry_cache_path, "wb") as f: f.write(ind_bytes)
elif os.path.exists(industry_cache_path):
    with open(industry_cache_path, "rb") as f: ind_bytes = f.read()
elif os.path.exists("업체대분류.csv"):
    with open("업체대분류.csv", "rb") as f: ind_bytes = f.read()
else:
    ind_bytes = None
if debt_file_up is not None:
    debt_bytes = debt_file_up.getvalue()
    # Google Drive 동기화 중 부분쓰기/충돌 완화: tmp → replace
    _debt_tmp = debt_cache_path + ".uploading"
    try:
        with open(_debt_tmp, "wb") as f:
            f.write(debt_bytes)
        os.replace(_debt_tmp, debt_cache_path)
    except Exception:
        with open(debt_cache_path, "wb") as f:
            f.write(debt_bytes)
        try:
            if os.path.exists(_debt_tmp):
                os.remove(_debt_tmp)
        except Exception:
            pass
    # 폴더의 채권.csv와도 맞춰 두면 Finder에서 바꾼 것과 사이드바 업로드가 어긋나지 않음
    try:
        with open("채권.csv", "wb") as f:
            f.write(debt_bytes)
    except Exception:
        pass
    try:
        load_debt_file.clear()
    except Exception:
        pass
    st.session_state["_debt_source"] = f"업로드:{getattr(debt_file_up, 'name', '채권.csv')}"
else:
    # ★ 핵심: uploaded_cache/debt.csv 가 있으면 예전엔 폴더 채권.csv를 영원히 무시했음
    # → mtime이 더 최신인 쪽을 사용 (폴더 파일 교체 반영)
    _debt_candidates = []
    if os.path.exists(debt_cache_path):
        try:
            st_ = os.stat(debt_cache_path)
            _debt_candidates.append(
                (st_.st_mtime, st_.st_size, debt_cache_path, "캐시")
            )
        except Exception:
            pass
    for f_name in os.listdir("."):
        if f_name.startswith("채권") and f_name.endswith(".csv"):
            try:
                st_ = os.stat(f_name)
                _debt_candidates.append(
                    (st_.st_mtime, st_.st_size, f_name, f"폴더:{f_name}")
                )
            except Exception:
                pass
    debt_bytes = None
    if _debt_candidates:
        _debt_candidates.sort(key=lambda x: (x[0], x[1]), reverse=True)
        _mtime, _size, _path, _label = _debt_candidates[0]
        with open(_path, "rb") as f:
            debt_bytes = f.read()
        st.session_state["_debt_source"] = (
            f"{_label} · {int(_size):,}B"
        )
        # 캐시가 오래됐으면 최신 폴더 파일로 캐시 갱신
        if _path != debt_cache_path and debt_bytes:
            try:
                _debt_tmp = debt_cache_path + ".uploading"
                with open(_debt_tmp, "wb") as f:
                    f.write(debt_bytes)
                os.replace(_debt_tmp, debt_cache_path)
                load_debt_file.clear()
            except Exception:
                try:
                    with open(debt_cache_path, "wb") as f:
                        f.write(debt_bytes)
                    load_debt_file.clear()
                except Exception:
                    pass
    else:
        st.session_state["_debt_source"] = "없음"
if tank_file_up is not None:
    tank_bytes = tank_file_up.getvalue()
    tank_name = tank_file_up.name
    with open(tank_cache_path, "wb") as f: f.write(tank_bytes)
    with open(tank_cache_path + "_name.txt", "w", encoding="utf-8") as f: f.write(tank_name)
elif os.path.exists(tank_cache_path) and os.path.exists(tank_cache_path + "_name.txt"):
    with open(tank_cache_path, "rb") as f: tank_bytes = f.read()
    with open(tank_cache_path + "_name.txt", "r", encoding="utf-8") as f: tank_name = f.read().strip()
elif os.path.exists("탱크.csv"):
    with open("탱크.csv", "rb") as f: tank_bytes = f.read()
    tank_name = "탱크.csv"
else:
    tank_bytes = None
    tank_name = ""
if vaporizer_file_up is not None:
    vaporizer_bytes = vaporizer_file_up.getvalue()
    vaporizer_name = vaporizer_file_up.name
    with open(vaporizer_cache_path, "wb") as f: f.write(vaporizer_bytes)
    with open(vaporizer_cache_path + "_name.txt", "w", encoding="utf-8") as f: f.write(vaporizer_name)
elif os.path.exists(vaporizer_cache_path) and os.path.exists(vaporizer_cache_path + "_name.txt"):
    with open(vaporizer_cache_path, "rb") as f: vaporizer_bytes = f.read()
    with open(vaporizer_cache_path + "_name.txt", "r", encoding="utf-8") as f: vaporizer_name = f.read().strip()
elif os.path.exists("기화기.csv"):
    with open("기화기.csv", "rb") as f: vaporizer_bytes = f.read()
    vaporizer_name = "기화기.csv"
else:
    vaporizer_bytes = None
    vaporizer_name = ""
if int_bytes is not None:
    with open(integrated_cache_path, "wb") as f: f.write(int_bytes)
    with open(integrated_cache_path + "_name.txt", "w", encoding="utf-8") as f: f.write(int_name)
elif os.path.exists(integrated_cache_path) and os.path.exists(integrated_cache_path + "_name.txt"):
    with open(integrated_cache_path, "rb") as f: int_bytes = f.read()
    with open(integrated_cache_path + "_name.txt", "r", encoding="utf-8") as f: int_name = f.read().strip()
else:
    int_bytes = None
    int_name = ""
sales_file_meta = []
if uploaded_files_up and len(uploaded_files_up) > 0:
    for f_name in os.listdir(sales_cache_dir):
        os.remove(os.path.join(sales_cache_dir, f_name))
    for f in uploaded_files_up:
        f_bytes = f.getvalue()
        f_path = os.path.join(sales_cache_dir, f.name)
        with open(f_path, "wb") as sf:
            sf.write(f_bytes)
        try:
            st_info = os.stat(f_path)
            sales_file_meta.append(
                (f.name, f_path, int(st_info.st_mtime), int(st_info.st_size))
            )
        except Exception:
            sales_file_meta.append((f.name, f_path, 0, len(f_bytes)))
else:
    # Drive 자동로드·시드 캐시 우선, 없으면 루트 20xx.csv
    if os.path.exists(sales_cache_dir):
        for f_name in sorted(os.listdir(sales_cache_dir)):
            if f_name.endswith(".csv"):
                f_path = os.path.join(sales_cache_dir, f_name)
                try:
                    st_info = os.stat(f_path)
                    sales_file_meta.append(
                        (f_name, f_path, int(st_info.st_mtime), int(st_info.st_size))
                    )
                except Exception:
                    continue
    if not sales_file_meta:
        for f_name in sorted(os.listdir(".")):
            if re.match(r"^20\d{2}.*\.csv$", f_name):
                f_path = os.path.abspath(f_name)
                try:
                    st_info = os.stat(f_path)
                    sales_file_meta.append(
                        (f_name, f_path, int(st_info.st_mtime), int(st_info.st_size))
                    )
                except Exception:
                    continue
sales_file_meta = _dedupe_sales_file_meta(sales_file_meta)
sales_file_meta = tuple(sales_file_meta)
if sales_file_meta:
    st.sidebar.caption(
        "매출 파일: " + ", ".join(x[0] for x in sales_file_meta)
    )

# 맥 사이드바 업로드 → Drive「dashboard 복사본」반영 (아이패드 Drive 폴더용)
_integrated_up = None
try:
    _integrated_up = integrated_file_up
except NameError:
    _integrated_up = None
_drive_upload_hit = bool(
    address_file_up
    or industry_file_up
    or debt_file_up
    or (uploaded_files_up and len(uploaded_files_up) > 0)
    or tank_file_up
    or vaporizer_file_up
    or _integrated_up
)

def _run_cache_to_drive_sync(*, force: bool = False):
    if sync_cache_to_drive_copy is None:
        return None
    if not force and st.session_state.get("_drive_synced_this_upload"):
        return None
    try:
        res = sync_cache_to_drive_copy(CACHE_DIR)
    except Exception as e:
        res = {"ok": False, "error": str(e)}
    st.session_state["_drive_synced_this_upload"] = True
    st.session_state["_drive_sync_out_msg"] = res
    return res

if _drive_upload_hit and sync_cache_to_drive_copy is not None:
    if not st.session_state.get("_drive_synced_this_upload"):
        _run_cache_to_drive_sync(force=False)
elif not _drive_upload_hit:
    st.session_state.pop("_drive_synced_this_upload", None)

if st.sidebar.button("☁️ Drive 복사본으로 동기화", help="맥 캐시 → Google Drive dashboard 복사본"):
    st.session_state.pop("_drive_synced_this_upload", None)
    _run_cache_to_drive_sync(force=True)
    st.rerun()

_drive_out = st.session_state.pop("_drive_sync_out_msg", None)
if isinstance(_drive_out, dict):
    if _drive_out.get("ok") and not _drive_out.get("skipped"):
        _nc = len([x for x in (_drive_out.get("copied") or []) if not str(x).startswith("-")])
        st.sidebar.success(f"Drive 복사본 반영 완료 · {_nc}개")
    elif _drive_out.get("skipped"):
        st.sidebar.caption("Drive 경로 없음 — 맥에서 Google Drive 앱 확인")
    elif not _drive_out.get("ok"):
        st.sidebar.warning(f"Drive 동기화 실패: {_drive_out.get('error') or '알 수 없음'}")

if st.sidebar.button("🗑️ 저장된 캐시 데이터 초기화"):
    for p in [addr_cache_path, industry_cache_path, debt_cache_path, 
              tank_cache_path, tank_cache_path + "_name.txt", 
              vaporizer_cache_path, vaporizer_cache_path + "_name.txt",
              integrated_cache_path, integrated_cache_path + "_name.txt",
              TAB7_DATE_FILE, TAB8_DATE_FILE]:
        if os.path.exists(p): os.remove(p)
    # DART API 키는 캐시 초기화에서 제외 (맥/iPad 재입력 방지)
    for f_name in os.listdir(sales_cache_dir):
        os.remove(os.path.join(sales_cache_dir, f_name))
    try:
        load_uploaded_files_from_meta.clear()
        load_uploaded_files_from_bytes.clear()
    except Exception:
        pass
    st.rerun()
addr_dict = load_address_file(addr_bytes) if addr_bytes else {}
industry_dict = load_industry_file(ind_bytes) if ind_bytes else {}
debt_df = load_debt_file(debt_bytes) if debt_bytes else pd.DataFrame()
if not debt_df.empty:
    debt_df = dedupe_debt_client_gubun(debt_df)
    _src = st.session_state.get("_debt_source") or ""
    st.sidebar.caption(
        f"채권 로드됨: {len(debt_df):,}행 · 거래처 {debt_df['거래처'].nunique():,}곳"
        + (f" · 출처 {_src}" if _src else "")
    )
elif debt_bytes:
    st.sidebar.warning(
        "채권 파일을 읽었지만 표로 변환하지 못했습니다. "
        "CSV에 '거래처','구분' 열이 있는지 확인하세요."
    )
df_tank = load_equipment_file(tank_bytes, tank_name) if tank_bytes else pd.DataFrame()
df_vaporizer = load_equipment_file(vaporizer_bytes, vaporizer_name) if vaporizer_bytes else pd.DataFrame()
df_integrated = load_equipment_file(int_bytes, int_name) if int_bytes else pd.DataFrame()
# ===== [초고속 데이터 전처리 캐시 엔진 (기존 로직 100% 무손실 보존)] =====
@st.cache_data(show_spinner=False)
def get_fast_processed_full_df(meta_data, staff_token, ind_dict):
    """필터 클릭 시 매번 발생하는 무거운 텍스트 연산(정규식, 매핑)을 1회로 압축"""
    # 1. 유저의 기존 고급 로드 기능 완벽 유지
    temp_df = (
        load_uploaded_files_from_meta(meta_data, staff_token, 6)
        if meta_data else pd.DataFrame()
    )
    
    # 2. 무거운 데이터 정제 작업을 캐시 안으로 이동 (한 번만 실행됨)
    if not temp_df.empty:
        temp_df = _apply_manual_staff_mapping(temp_df)
        _client_s = temp_df["거래처"].astype(str).str.strip()
        mask_closed = _client_s.str.match(r"^[zZ]", na=False)
        if mask_closed.any():
            temp_df.loc[mask_closed, "담당자"] = "거래종료"
            
        is_deposit_row = temp_df["품목명"].astype(str).str.contains("입금", na=False)
        temp_df = temp_df[~is_deposit_row].copy()
        temp_df["업종"] = temp_df["거래처"].map(ind_dict).fillna("미분류")
        
    return temp_df

# 단 0.01초 만에 메모리에서 정제 완료된 데이터를 즉시 꺼내옴
full_df = get_fast_processed_full_df(sales_file_meta, _manual_staff_map_cache_token(), industry_dict)
# ====================================================================
target_items = [
    "CO2 (kg, Bulk)",
    "N2 (kg, Bulk)",
    "O2 (kg, Bulk)",
    "AR (kg, Bulk)",
]
latest_update_str = "데이터 없음"
selected_staff = []
selected_item = []
selected_client = "전체 거래처"
df_base = pd.DataFrame()
df_client_filtered = pd.DataFrame()
df_f = pd.DataFrame()
pivot_m_total = pd.DataFrame()
client_item_qty_pivot = pd.DataFrame()
sales_p = pd.DataFrame()
qty_p = pd.DataFrame()
unit_price_p = pd.DataFrame()
staff_pivot = pd.DataFrame()
df_detail = pd.DataFrame()
all_months = [f"{i:02d}월" for i in range(1, 13)]
years = ["2026"]
if not full_df.empty:
    latest_dt_overall = full_df["매출일_dt"].max()
    if pd.notnull(latest_dt_overall):
        latest_update_str = latest_dt_overall.strftime("%Y-%m-%d")
    # ==============================================================
    # 필터 영역 (상단 고정은 inject_sticky_tabs_script에서 처리)
    # 맥·iPad 공통: 선택 즉시 반영. sticky 재기동 억제로 로딩감 완화.
    # ==============================================================
    st.markdown("<div id='dashboard-sticky-spacer'></div>", unsafe_allow_html=True)
    try:
        filter_container = st.container(border=True)
    except TypeError:
        filter_container = st.container()

    def _dash_parse_filter_dates(start_s, end_s):
        _sd = pd.to_datetime(start_s, format="%y%m%d", errors="coerce")
        _ed = pd.to_datetime(end_s, format="%y%m%d", errors="coerce")
        if pd.isna(_sd):
            _sd = pd.Timestamp("2000-01-01")
        if pd.isna(_ed):
            _ed = pd.Timestamp("2099-12-31")
        return _sd, _ed

    def _dash_staff_opts_from(df_src):
        _staff_raw = (
            sorted(df_src["담당자"].dropna().astype(str).unique())
            if not df_src.empty and "담당자" in df_src.columns
            else []
        )
        _staff_opts = [s for s in _staff_raw if s != "거래종료"]
        return ["거래종료"] + _staff_opts

    with filter_container:
        st.markdown(
            '<div id="dashboard-filter-bar" class="notranslate" translate="no" lang="ko"></div>',
            unsafe_allow_html=True,
        )
        st.markdown("<div id='sticky-marker' style='display:none;'></div>", unsafe_allow_html=True)
        fc1, fc2, fc3, fc4, fc5 = st.columns([1, 1, 1.15, 1.35, 1.15])
        start_date = fc1.text_input("📅 조회 시작", "200101", key="dash_filter_start")
        end_date = fc2.text_input("📅 조회 종료", "261231", key="dash_filter_end")
        start_dt, end_dt = _dash_parse_filter_dates(start_date, end_date)
        df_base_opts = full_df[
            (full_df["매출일_dt"] >= start_dt) & (full_df["매출일_dt"] <= end_dt)
        ]
        _staff_opts = _dash_staff_opts_from(df_base_opts)
        _prev_staff = st.session_state.get("dash_filter_staff", [])
        if isinstance(_prev_staff, list) and _prev_staff:
            _kept_staff = [x for x in _prev_staff if x in _staff_opts]
            if _kept_staff != _prev_staff:
                st.session_state["dash_filter_staff"] = _kept_staff
        # 👤 담당자 선택 (불편한 multiselect 버리고 selectbox로 교체 + 호환성 유지)
        _staff_opts_with_all = ["전체 담당자"] + _staff_opts
        _staff_picked = fc3.selectbox("👤 담당자", options=_staff_opts_with_all, index=0, key="dash_filter_staff_sb_new")
        selected_staff = [] if _staff_picked == "전체 담당자" else [_staff_picked]
        df_staff_for_opts = (
            df_base_opts[df_base_opts["담당자"].isin(selected_staff)]
            if selected_staff
            else df_base_opts
        )
        _client_opts_sig = (start_date, end_date, tuple(selected_staff or ()))
        if st.session_state.get("_dash_client_opts_sig") != _client_opts_sig:
            st.session_state["_dash_client_opts_sig"] = _client_opts_sig
            if df_staff_for_opts.empty:
                st.session_state["_dash_client_opts_tuple"] = ()
            else:
                st.session_state["_dash_client_opts_tuple"] = tuple(
                    df_staff_for_opts["거래처"].astype(str).unique()
                )
        all_clients = sorted(st.session_state.get("_dash_client_opts_tuple", ()))
        # 단일 선택 + 태그 X로 원복 (selectbox는 지우기 불가 → multiselect max 1)
        # 구 selectbox 키 문자열 → 새 리스트 키로 1회 이관
        if "dash_filter_client_ms" not in st.session_state:
            _old_c = st.session_state.get("dash_filter_client", "전체 거래처")
            if isinstance(_old_c, str) and _old_c and _old_c != "전체 거래처" and _old_c in all_clients:
                st.session_state["dash_filter_client_ms"] = [_old_c]
            else:
                st.session_state["dash_filter_client_ms"] = []
        _prev_clients = st.session_state.get("dash_filter_client_ms", [])
        if isinstance(_prev_clients, list) and _prev_clients:
            _kept_c = [x for x in _prev_clients if x in all_clients]
            if _kept_c != _prev_clients:
                st.session_state["dash_filter_client_ms"] = _kept_c
        # 🏢 거래처 선택 (기형적인 multiselect 버리고 1초 타자 가능한 selectbox로 순정 복구)
        client_options_with_all = ["전체 거래처"] + all_clients
        
        selected_client = fc4.selectbox(
            "🏢 거래처",
            options=client_options_with_all,
            index=0,
            key="dash_filter_client_selectbox"
        )
        # 하위 로직·엑셀 시그니처 호환용
        st.session_state["dash_filter_client"] = selected_client
        df_client_for_opts = (
            df_staff_for_opts[df_staff_for_opts["거래처"] == selected_client]
            if selected_client != "전체 거래처"
            else df_staff_for_opts
        )
        available_items = (
            sorted(df_client_for_opts["품목명"].astype(str).unique())
            if not df_client_for_opts.empty
            else []
        )
        _prev_items = st.session_state.get("dash_filter_items", [])
        if isinstance(_prev_items, list) and _prev_items:
            _kept = [x for x in _prev_items if x in available_items]
            if _kept != _prev_items:
                st.session_state["dash_filter_items"] = _kept
        # 📦 품목명 선택 (불편한 multiselect 버리고 selectbox로 교체 + 호환성 유지)
        _item_opts = ["전체 품목"] + available_items
        _item_picked = fc5.selectbox("📦 품목명", options=_item_opts, index=0, key="dash_filter_items_sb_new")
        selected_item = [] if _item_picked == "전체 품목" else [_item_picked]

        df_base = df_base_opts
        df_staff_filtered = (
            df_base[df_base["담당자"].isin(selected_staff)] if selected_staff else df_base
        )
        df_client_filtered = (
            df_staff_filtered[df_staff_filtered["거래처"] == selected_client]
            if selected_client != "전체 거래처"
            else df_staff_filtered
        )
        df_f = (
            df_client_filtered[df_client_filtered["품목명"].isin(selected_item)]
            if selected_item
            else df_client_filtered
        )
        _pivot_ck = _dash_pivot_cache_key(
            selected_client,
            selected_staff,
            selected_item,
            start_date,
            end_date,
            sales_file_meta,
            _manual_staff_map_cache_token(),
        )
        _pivot_store = st.session_state.setdefault("_dash_pivot_store", {})
        if _pivot_ck not in _pivot_store:
            _pivot_store[_pivot_ck] = _dash_compute_pivot_bundle(
                df_base, df_client_filtered, df_f, full_df, all_months
            )
            while len(_pivot_store) > 24:
                _pivot_store.pop(next(iter(_pivot_store)))
        _pb = _pivot_store[_pivot_ck]
        years = _pb["years"]
        desired_order = _pb["desired_order"]
        pivot_m_total = _pb["pivot_m_total"]
        client_item_qty_pivot = _pb["client_item_qty_pivot"]
        sales_p = _pb["sales_p"]
        qty_p = _pb["qty_p"]
        unit_price_p = _pb["unit_price_p"]
        staff_pivot = _pb["staff_pivot"]
        df_detail = _pb["df_detail"]
        cur_month_sales_total = _pb["cur_month_sales_total"]
        prev_month_sales_total = _pb["prev_month_sales_total"]
        mom_rate_total = _pb["mom_rate_total"]
        avg_monthly_sales_total = _pb["avg_monthly_sales_total"]
        avg_rate_total = _pb["avg_rate_total"]
        latest_month_str_total = _pb["latest_month_str_total"]
        cur_month_sales_client = _pb["cur_month_sales_client"]
        prev_month_sales_client = _pb["prev_month_sales_client"]
        mom_rate_client = _pb["mom_rate_client"]
        avg_monthly_sales_client = _pb["avg_monthly_sales_client"]
        avg_rate_client = _pb["avg_rate_client"]
        latest_month_str_client = _pb["latest_month_str_client"]
else:
    cur_month_sales_total = prev_month_sales_total = mom_rate_total = avg_monthly_sales_total = avg_rate_total = 0.0
    latest_month_str_total = "-"
    cur_month_sales_client = prev_month_sales_client = mom_rate_client = avg_monthly_sales_client = avg_rate_client = 0.0
    latest_month_str_client = "-"
# 담당자만 반영한 채권(거래처 선택 무관) — 연체개월수 요약표용
staff_debt_df = pd.DataFrame()
filtered_debt_df = pd.DataFrame()
if not debt_df.empty:
    if selected_staff:
        valid_staff_clients = full_df[full_df["담당자"].isin(selected_staff)]["거래처"].unique()
        staff_debt_df = debt_df[debt_df["거래처"].isin(valid_staff_clients)].copy()
    else:
        staff_debt_df = debt_df.copy()
    filtered_debt_df = staff_debt_df.copy()
    if selected_client != "전체 거래처":
        filtered_debt_df = filtered_debt_df[filtered_debt_df["거래처"] == selected_client].copy()
client_addr_raw = resolve_client_address(selected_client, addr_dict)
if not client_addr_raw:
    client_addr = "등록된 주소 정보가 없습니다."
else:
    client_addr = str(client_addr_raw)
st.sidebar.markdown("---")
st.sidebar.subheader("📥 엑셀 내보내기")
_excel_sig = (
    selected_client,
    tuple(selected_staff or []),
    tuple(selected_item or []),
    str(st.session_state.get("dash_filter_start", "")),
    str(st.session_state.get("dash_filter_end", "")),
    latest_update_str,
    int(len(df_base)),
    float(df_base["매출액"].sum()) if not df_base.empty else 0.0,
    int(len(filtered_debt_df)),
)
if st.session_state.get("_dash_excel_sig") != _excel_sig:
    st.session_state.pop("_dash_excel_bytes", None)
_excel_ready = bool(
    st.session_state.get("_dash_excel_sig") == _excel_sig
    and st.session_state.get("_dash_excel_bytes")
)
if st.sidebar.button(
    "📊 엑셀 파일 준비",
    key="dash_excel_prepare",
    use_container_width=True,
    help="현재 검색 조건으로 시트별 엑셀을 만듭니다.",
):
    with st.spinner("엑셀 생성 중…"):
        _sheets_dict = {
            "연도별_월매출(만원)": (pivot_m_total, True),
            "거래처별_품목별사용량": (client_item_qty_pivot, True),
            "선택거래처_품목별_매출액(만원)": (sales_p * 1.1 / 10000, True),
            "선택거래처_품목별_출고량": (qty_p, True),
            "선택거래처_품목별_적용단가": (unit_price_p, True),
            "담당자별_매출(만원)": (staff_pivot, True),
            "상세거래내역": (df_detail, False),
        }
        if not filtered_debt_df.empty:
            _sheets_dict["채권관리_현황"] = (filtered_debt_df, False)
        st.session_state["_dash_excel_sig"] = _excel_sig
        st.session_state["_dash_excel_bytes"] = convert_dfs_to_excel(_sheets_dict)
        _excel_ready = True
if _excel_ready:
    st.sidebar.download_button(
        label="⬇️ 전체 분석 시트별 엑셀 다운로드",
        data=st.session_state["_dash_excel_bytes"],
        file_name="통합영업분석_시트별보고서.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        use_container_width=True,
    )
elif st.session_state.get("_dash_excel_sig") and st.session_state.get("_dash_excel_sig") != _excel_sig:
    st.sidebar.caption("검색 조건이 바뀌었습니다. 「엑셀 파일 준비」를 다시 눌러 주세요.")
else:
    st.sidebar.caption("엑셀은 「엑셀 파일 준비」 후 다운로드 (거래처 지정 시 자동 생성 안 함)")
st.sidebar.subheader("📽️ PPT 내보내기")
_ppt_sig = (
    selected_client,
    tuple(sorted(selected_staff)),
    latest_update_str,
    latest_month_str_total,
)
if st.session_state.get("ppt_gen_sig") != _ppt_sig:
    st.session_state["ppt_ready_bytes"] = None
try:
    if st.sidebar.button("📽️ PPT 생성", key="gen_dashboard_ppt", use_container_width=True):
        with st.spinner("PPT 파일 생성 중..."):
            tot_sales_val_export = df_base["매출액"].sum() * 1.1 / 10000 if not df_base.empty else 0.0
            cur_sales_val_export = cur_month_sales_total / 10000
            pivot_m_client_export = cached_get_yearly_monthly_pivot(df_client_filtered, all_months, years)
            st.session_state["ppt_ready_bytes"] = convert_dashboard_to_ppt(
                latest_update_str,
                selected_client,
                tuple(selected_staff),
                latest_month_str_total,
                tot_sales_val_export,
                cur_sales_val_export,
                mom_rate_total,
                avg_rate_total,
                pivot_m_total,
                pivot_m_client_export,
                client_item_qty_pivot,
                sales_p,
                qty_p,
                staff_pivot,
                df_detail,
                filtered_debt_df,
                df_base,
                df_tank,
                df_vaporizer,
                df_integrated,
                tuple(all_months),
                tuple(years),
                tuple(target_items),
            )
            st.session_state["ppt_gen_sig"] = _ppt_sig
    if st.session_state.get("ppt_ready_bytes") and st.session_state.get("ppt_gen_sig") == _ppt_sig:
        st.sidebar.download_button(
            label="⬇️ PPT 다운로드",
            data=st.session_state["ppt_ready_bytes"],
            file_name=f"통합영업분석_전체탭_{datetime.date.today().strftime('%Y%m%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            key="dl_dashboard_ppt",
        )
    elif st.session_state.get("ppt_gen_sig") and st.session_state.get("ppt_gen_sig") != _ppt_sig:
        st.sidebar.caption("검색 조건이 변경되었습니다. PPT를 다시 생성해 주세요.")
except ImportError:
    st.sidebar.warning("PPT 내보내기: `pip install python-pptx kaleido` 설치 후 이용하세요.")
except Exception as exc:
    st.sidebar.error(f"PPT 생성 오류: {exc}")
# 탭 전환은 클라이언트 전환만 (rerun 없음). 필터 변경 시에만 전체 재계산.
tab1, tab2, tab3, tab4, tab5, tab6, tab7, tab8, tab9, tab10, tab11 = st.tabs(
    [
        "📌 영업 종합 요약",
        "🏢 거래처 분석",
        "📦 품목 및 단가 분석",
        "👤 담당자 & 상세내역",
        "📌 채권 관리",
        "📍 카카오맵",
        "🏭 설비 재고 현황",
        "🛢️ 통합 탱크 재고",
        "📈 수익성 분석",
        "📝 일일업무일지",
        "🔎 시장조사",
    ]
)
if is_touch_ui():
    # iPad: 거래처 변경 rerun마다 components.html 재주입하면 로딩이 되살아남
    if st.session_state.get("_ipad_sticky_ver") != 30:
        inject_sticky_tabs_script()
        inject_ipad_plotly_controls()
        st.session_state["_ipad_sticky_injected"] = True
        st.session_state["_ipad_sticky_ver"] = 30
else:
    inject_sticky_tabs_script()
    inject_ipad_plotly_controls()
# Tab 1: 📌 영업 종합 요약
with tab1:
    t1_c1, t1_c2 = st.columns([4, 1])
    t1_c1.markdown("<div class='sub-header dashboard-tab-panel-head'>📊 전체 영업 주요 실적 지표</div>", unsafe_allow_html=True)
    t1_c2.markdown(render_update_badge(latest_update_str), unsafe_allow_html=True)

    m1, m2, m3, m4 = st.columns(4)

    tot_sales_val = df_base["매출액"].sum() * 1.1 / 10000 if not df_base.empty else 0.0
    cur_sales_val = cur_month_sales_total / 10000
    m1.markdown(f"<div class='metric-box'><div class='metric-label'>총 누적 매출 (VAT포함)</div><div class='metric-value'>{tot_sales_val:,.0f} 만원</div></div>", unsafe_allow_html=True)
    m2.markdown(f"<div class='metric-box'><div class='metric-label'>최근 월 매출 ({latest_month_str_total})</div><div class='metric-value'>{cur_sales_val:,.0f} 만원</div></div>", unsafe_allow_html=True)
    m3.markdown(f"<div class='metric-box'><div class='metric-label'>전월 대비 (MoM)</div><div class='metric-value' style='color:{'#E11D48' if mom_rate_total < 0 else '#2563EB'};'>{mom_rate_total:+.0f}%</div></div>", unsafe_allow_html=True)
    m4.markdown(f"<div class='metric-box'><div class='metric-label'>월평균 대비 증감</div><div class='metric-value' style='color:{'#E11D48' if avg_rate_total < 0 else '#2563EB'};'>{avg_rate_total:+.0f}%</div></div>", unsafe_allow_html=True)
    st.markdown("<div class='sub-header dashboard-tab-panel-head'>📊 전체 영업 연도별 월 매출 추이</div>", unsafe_allow_html=True)
    col_left, col_right = st.columns([1, 1])

    with col_left:
        pivot_m_total_disp = get_display_df_with_sum(pivot_m_total, "연간 합계")
        st.dataframe(style_with_sum(pivot_m_total_disp, "{:,.0f}", "Blues", axis=None), use_container_width=True, height=460)
    with col_right:
        render_plotly_chart(
            create_stacked_bar_chart(pivot_m_total, title_text=""),
            use_container_width=True, key="tab1_total_chart"
        )
    
    st.markdown("---")
    st.markdown("<div class='sub-header dashboard-tab-panel-head'>📦 주요 4대 품목 상세 분석</div>", unsafe_allow_html=True)

    sel_col1, sel_col2 = st.columns([1, 1])
    with sel_col1:
        selected_target_item = st.radio("🔍 분석할 품목 선택", target_items, horizontal=True, key="overall_item_radio")
    with sel_col2:
        selected_metric = st.radio("📊 분석 지표 선택", ["매출액 (만원)", "출고량", "총매출 대비 비중 (%)"], horizontal=True, key="overall_metric_radio")
    
    item_pivot = cached_get_item_pivot(df_base, selected_target_item, selected_metric, all_months, years)
    week_year_pivot = cached_get_item_month_week_year(
        df_base, selected_target_item, selected_metric, all_months, years
    )

    if "비중" in selected_metric:
        y_suf, y_fmt, _fmt_kind, _cmap = "%", ",.1f", "pct", "Purples"
    elif "출고량" in selected_metric:
        y_suf, y_fmt, _fmt_kind, _cmap = " 천kg", ",.1f", "qty", "Greens"
    else:
        y_suf, y_fmt, _fmt_kind, _cmap = " 만원", ",.0f", "amt", "Blues"

    # 상단「연도별 월 매출 추이」와 동일: 좌 표 / 우 그래프 1:1
    # iPad 표는 render_month_expandable_week_table 내부에서 가로 스크롤로 숫자 표시
    i_col_left, i_col_right = st.columns([1, 1])
    with i_col_left:
        render_month_expandable_week_table(
            item_pivot,
            week_year_pivot,
            fmt_kind=_fmt_kind,
            cmap_name=_cmap,
            height=460,
        )
    with i_col_right:
        render_plotly_chart(
            create_stacked_bar_chart(
                item_pivot,
                title_text="",
                y_suffix=y_suf,
                y_format=y_fmt,
            ),
            use_container_width=True,
            key="tab1_item_chart",
        )

    st.markdown("---")
    st.markdown("<div class='sub-header dashboard-tab-panel-head'>🏭 업종별(분류별) 상세 분석</div>", unsafe_allow_html=True)

    if "업종" in df_base.columns:
        available_industries = sorted(list(df_base["업종"].unique()))
    
        if len(available_industries) == 1 and available_industries[0] == "미분류":
            st.info("💡 현재 모든 거래처가 '미분류' 상태입니다. 왼쪽 사이드바에서 '🏢 거래처 업종 분류 (CSV)' 파일을 업로드하시면 정확한 업종별 상세 분석이 가능합니다.")
        
        if available_industries:
            ind_col1, ind_col2 = st.columns([1, 1])
            with ind_col1:
                selected_industry = st.selectbox("🔍 분석할 업종(분류) 선택", available_industries, key="industry_selectbox")
            with ind_col2:
                selected_ind_metric = st.radio("📊 분석 지표 선택", ["매출액 (만원)", "출고량", "총매출 대비 비중 (%)"], horizontal=True, key="industry_metric_radio")
        
            ind_pivot = cached_get_industry_pivot(df_base, selected_industry, selected_ind_metric, all_months, years)
            
            i_col_left2, i_col_right2 = st.columns([1, 1])
            with i_col_left2:
                st.caption("💡 표의 행(월)을 클릭하면 아래 '업종별 매출 비중' 도넛 차트의 기준 월이 변경됩니다.")
                ind_pivot_disp = get_display_df_with_sum(ind_pivot, "연간 합계")
                
                # 표에 on_select 추가하여 클릭 감지
                if "비중" in selected_ind_metric:
                    ind_ev = st.dataframe(style_with_sum(ind_pivot_disp, "{:,.1f}%", "Purples", axis=None), use_container_width=True, height=460, on_select="rerun", selection_mode="single-row", key="ind_table_pct")
                    y_suf_i, y_fmt_i = "%", ",.1f"
                elif "출고량" in selected_ind_metric:
                    # 🟢 소수점 1자리(1f)와 '천' 단위가 적용된 부분입니다!
                    ind_ev = st.dataframe(style_with_sum(ind_pivot_disp, "{:,.1f}", "Greens", axis=None), use_container_width=True, height=460, on_select="rerun", selection_mode="single-row", key="ind_table_qty")
                    y_suf_i, y_fmt_i = " 천", ",.1f"
                else:
                    ind_ev = st.dataframe(style_with_sum(ind_pivot_disp, "{:,.0f}", "Blues", axis=None), use_container_width=True, height=460, on_select="rerun", selection_mode="single-row", key="ind_table_amt")
                    y_suf_i, y_fmt_i = " 만원", ",.0f"
                    
                # 행(월) 클릭 시 도넛 차트 월(top30_month) 강제 업데이트
                if ind_ev and ind_ev.selection.rows:
                    sel_idx = ind_ev.selection.rows[0]
                    sel_month = ind_pivot_disp.index[sel_idx]
                    
                    if sel_month in all_months and st.session_state.get("top30_month") != sel_month:
                        st.session_state["top30_month"] = sel_month
                        try:
                            st.query_params["top30_month"] = sel_month
                        except Exception:
                            pass
                        st.rerun()
                        
            with i_col_right2:
                render_plotly_chart(
                    create_stacked_bar_chart(
                        ind_pivot, 
                        title_text="", 
                        y_suffix=y_suf_i, 
                        y_format=y_fmt_i
                    ),
                    use_container_width=True, key="tab1_industry_chart"
                )
            st.markdown("<br>", unsafe_allow_html=True)
            with st.expander(f"📂 [{selected_industry}] 소속 거래처 상세 데이터 파보기 (클릭하여 펼치기)", expanded=False):
                df_ind_detail = df_base[df_base["업종"] == selected_industry]
            
                if not df_ind_detail.empty:
                    ind_client_pivot = df_ind_detail.pivot_table(index="거래처", columns="연도", values="매출액", aggfunc="sum").fillna(0) * 1.1 / 10000
                
                    avail_years_ind = sorted([y for y in ind_client_pivot.columns if str(y).isdigit()])
                    ind_client_pivot = ind_client_pivot.reindex(columns=avail_years_ind, fill_value=0)
                    ind_client_pivot["총 누적매출"] = ind_client_pivot.sum(axis=1)
                    ind_client_pivot = ind_client_pivot.sort_values(by="총 누적매출", ascending=False)
                
                    ind_client_pivot_disp = get_display_df_with_sum(ind_client_pivot, "합계")
                    st.dataframe(style_with_sum(ind_client_pivot_disp, "{:,.0f}", "Blues", axis=None), use_container_width=True, height=280)
                
                    st.markdown("<hr style='margin: 15px 0px; border-top: 1px dashed #E2E8F0;'>", unsafe_allow_html=True)
                
                    ind_clients = sorted(df_ind_detail["거래처"].unique())
                
                    c1, c2, c3 = st.columns([1, 1, 1])
                    with c1:
                        sel_ind_client = st.selectbox(f"🏢 [{selected_industry}] 거래처 선택", ind_clients, key="ind_client_sel")
                    with c2:
                        client_items = sorted(df_ind_detail[df_ind_detail["거래처"] == sel_ind_client]["품목명"].unique())
                        sel_ind_item = st.selectbox("📦 품목 선택", client_items if client_items else ["없음"], key="ind_item_sel")
                    with c3:
                        sel_ind_sub_metric = st.radio("📊 지표 선택", ["매출액 (만원)", "출고량", "총매출 대비 비중 (%)"], horizontal=True, key="ind_sub_metric")
                
                    if client_items:
                        df_target_client = df_ind_detail[df_ind_detail["거래처"] == sel_ind_client]
                        sub_item_pivot = cached_get_item_pivot(df_target_client, sel_ind_item, sel_ind_sub_metric, all_months, years)
                    
                        sc1, sc2 = st.columns([1, 1])
                        with sc1:
                            sub_item_pivot_disp = get_display_df_with_sum(sub_item_pivot, "연간 합계")
                            if "비중" in sel_ind_sub_metric:
                                st.dataframe(style_with_sum(sub_item_pivot_disp, "{:,.1f}%", "Purples", axis=None), use_container_width=True, height=410)
                                y_suf_sub, y_fmt_sub = "%", ",.1f"
                            elif "출고량" in sel_ind_sub_metric:
                                if sel_ind_item in target_items:
                                    y_suf_sub, y_fmt_sub = " 천kg", ",.1f"
                                elif "LPG" in str(sel_ind_item).upper():
                                    y_suf_sub, y_fmt_sub = " kg", ",.0f"
                                else:
                                    y_suf_sub, y_fmt_sub = " 개(병)", ",.0f"
                                st.dataframe(style_with_sum(sub_item_pivot_disp, f"{{:{y_fmt_sub}}}", "Greens", axis=None), use_container_width=True, height=410)
                            else:
                                st.dataframe(style_with_sum(sub_item_pivot_disp, "{:,.0f}", "Blues", axis=None), use_container_width=True, height=410)
                                y_suf_sub, y_fmt_sub = " 만원", ",.0f"
                            
                        with sc2:
                            render_plotly_chart(
                                create_stacked_bar_chart(sub_item_pivot, title_text="", y_suffix=y_suf_sub, y_format=y_fmt_sub),
                                use_container_width=True, key="ind_client_sub_chart"
                            )
    st.markdown("---")
    st.markdown("<div class='sub-header dashboard-tab-panel-head'>🏆 당해년도 상위 30위 거래처 실적 (1월~12월) 및 업종 비중 · 월 헤더 클릭</div>", unsafe_allow_html=True)

    if not df_base.empty:
        current_year_str = str(df_base["연도"].max())
        df_curr_year = df_base[df_base["연도"] == current_year_str]
        if not df_curr_year.empty:
            latest_dt = df_curr_year["매출일_dt"].max()
            latest_m = latest_dt.strftime("%m월")
            # 월 선택 → 해당 월 순위·도넛만 변경 (집계 로직 무손실)
            inject_top30_month_bridge()
            if "top30_month" not in st.session_state:
                st.session_state["top30_month"] = latest_m
            try:
                qp_m = st.query_params.get("top30_month", None)
                if isinstance(qp_m, list):
                    qp_m = qp_m[0] if qp_m else None
                if qp_m is not None:
                    qp_m = urllib.parse.unquote(str(qp_m))
                    if qp_m in all_months:
                        st.session_state["top30_month"] = qp_m
            except Exception:
                pass
            if st.session_state.get("top30_month") not in all_months:
                st.session_state["top30_month"] = latest_m
            rank_m = st.session_state["top30_month"]
            rank_month_label = f"{current_year_str}년 {rank_m}"
            # top30-section-flag: iPad CSS 스택용 마커 (맥 레이아웃/데이터 무손실)
            # 상단 연도별 월매출 표·그래프와 동일: 1:1 열, 높이 460
            _TOP30_H = 460
            with st.container():
                st.markdown("<div class='top30-section-flag' style='display:none'></div>", unsafe_allow_html=True)
                p_col1, p_col2 = st.columns([1, 1])
                with p_col1:
                    st.markdown(
                        f"<div style='font-size: 14px; font-weight: 600; color: #334155; margin: 0 0 8px; min-height: 22px; line-height: 1.4;'>"
                        f"🥇 [{rank_month_label} 기준] 상위 30위 거래처 월별 실적 (VAT포함, 만원)"
                        f"<span style='font-size:12px;font-weight:500;color:#64748B;margin-left:8px;'>"
                        f"{'← 월 선택 또는 표 헤더 클릭' if is_touch_ui() else '← 월 버튼 또는 표 헤더 클릭'}"
                        f"</span></div>",
                        unsafe_allow_html=True,
                    )
                    # 맥: 이전 가로 월버튼 / iPad(touch_ui=1): 현재 selectbox
                    if is_touch_ui():
                        _m_idx = all_months.index(rank_m) if rank_m in all_months else 0
                        picked_m = st.selectbox(
                            "기준 월",
                            all_months,
                            index=_m_idx,
                            key="top30_month_select",
                            label_visibility="collapsed",
                        )
                        if picked_m != st.session_state.get("top30_month"):
                            st.session_state["top30_month"] = picked_m
                            try:
                                st.query_params["top30_month"] = picked_m
                            except Exception:
                                pass
                            st.rerun()
                    else:
                        m_cols = st.columns(12, gap="small")
                        _clicked_m = None
                        for _i, _m in enumerate(all_months):
                            with m_cols[_i]:
                                if st.button(
                                    _m,
                                    key=f"top30_month_btn_{_m}",
                                    type="primary" if _m == rank_m else "secondary",
                                    width="stretch",
                                ):
                                    _clicked_m = _m
                        if _clicked_m:
                            st.session_state["top30_month"] = _clicked_m
                            try:
                                st.query_params["top30_month"] = _clicked_m
                            except Exception:
                                pass
                            st.rerun()
                    rank_m = st.session_state["top30_month"]
                    rank_month_label = f"{current_year_str}년 {rank_m}"
                
                    pvt_curr = df_curr_year.pivot_table(index="거래처", columns="월", values="매출액", aggfunc="sum").fillna(0) * 1.1 / 10000
                    pvt_curr = pvt_curr.reindex(columns=all_months, fill_value=0)
                
                    if rank_m in pvt_curr.columns:
                        pvt_curr = pvt_curr.sort_values(by=rank_m, ascending=False)
                    top30_pvt = pvt_curr.head(30).reset_index()
                
                    top30_pvt.index = range(1, len(top30_pvt) + 1)
                
                    top30_pvt_disp = get_display_df_with_sum(top30_pvt, sum_label="합계", text_cols=["거래처"])
                
                    fmt_dict = {m: "{:,.0f}" for m in all_months}
                    styled_top30 = style_with_sum(top30_pvt_disp, fmt_dict, "Blues", subset_cols=all_months, axis=0)
                
                    if rank_m in all_months:
                        styled_top30 = styled_top30.apply(
                            lambda s: ['color: #B91C1C; font-weight: bold; background-color: #DBEAFE;'] * len(s),
                            subset=[rank_m],
                            axis=0
                        )
                
                    render_frozen_styler_html(
                        styled_top30,
                        height=_TOP30_H,
                        freeze_left_n=2,
                        freeze_widths=[44, 160],
                        clickable_cols=all_months,
                        query_param="top30_month",
                        active_col=rank_m,
                    )
                with p_col2:
                    st.markdown(
                        f"<div style='font-size: 14px; font-weight: 600; color: #334155; margin: 0 0 8px; min-height: 22px; line-height: 1.4;'>"
                        f"🍩 [{rank_month_label}] 업종별 매출 비중</div>",
                        unsafe_allow_html=True,
                    )
                    # 왼쪽 월 선택 줄과 높이를 맞춰 표·그래프 상단 정렬
                    _ctrl_h = 42 if is_touch_ui() else 40
                    st.markdown(
                        f"<div style='height:{_ctrl_h}px;margin:0 0 8px;' aria-hidden='true'></div>",
                        unsafe_allow_html=True,
                    )
                
                    df_rank_month = df_curr_year[df_curr_year["월"] == rank_m]
                    ind_sales = df_rank_month.groupby("업종")["매출액"].sum().reset_index()
                    ind_sales = ind_sales[ind_sales["매출액"] > 0]
                
                    if ind_sales.empty:
                        st.info(f"{rank_month_label} 업종별 매출 데이터가 없습니다.")
                    else:
                        fig_donut = px.pie(
                            ind_sales, 
                            values='매출액', 
                            names='업종', 
                            hole=0.4,
                            color_discrete_sequence=px.colors.qualitative.Pastel
                        )
                        fig_donut.update_traces(
                            textposition='inside',
                            textinfo='percent+label',
                            textfont_size=12,
                        )
                        fig_donut.update_layout(
                            showlegend=True,
                            legend=dict(
                                orientation="h",
                                yanchor="top",
                                y=-0.02,
                                xanchor="center",
                                x=0.5,
                                font=dict(size=11),
                            ),
                            margin=dict(l=8, r=8, t=8, b=72),
                            height=_TOP30_H,
                            autosize=True,
                            paper_bgcolor="rgba(0,0,0,0)",
                            plot_bgcolor="rgba(0,0,0,0)",
                        )
                        render_plotly_chart(
                            fig_donut,
                            key=f"top30_donut_{rank_m}",
                            height=_TOP30_H,
                        )
        else:
            st.info("당해년도 매출 데이터가 없습니다.")
    st.markdown("---")
    if not df_base.empty and '매출일_dt' in df_base.columns:
        latest_period = df_base["매출일_dt"].dt.to_period("M").max()
        prev_period = latest_period - 1
    
        curr_m_label = latest_period.strftime("%m월")
        prev_m_label = prev_period.strftime("%m월")
    
        st.markdown(f"<div class='sub-header dashboard-tab-panel-head'>📊 전월 대비 실적 증감 및 신규/이탈 분석 ({prev_m_label} vs {curr_m_label})</div>", unsafe_allow_html=True)
    
        df_prev = df_base[df_base["매출일_dt"].dt.to_period("M") == prev_period].groupby("거래처")["매출액"].sum().reset_index().rename(columns={"매출액": f"{prev_m_label} 매출"})
        df_curr = df_base[df_base["매출일_dt"].dt.to_period("M") == latest_period].groupby("거래처")["매출액"].sum().reset_index().rename(columns={"매출액": f"{curr_m_label} 매출"})
        df_diff = pd.merge(df_prev, df_curr, on="거래처", how="outer").fillna(0)
        df_diff["매출 증감액"] = df_diff[f"{curr_m_label} 매출"] - df_diff[f"{prev_m_label} 매출"]
    
        df_diff[[f"{prev_m_label} 매출", f"{curr_m_label} 매출", "매출 증감액"]] = (df_diff[[f"{prev_m_label} 매출", f"{curr_m_label} 매출", "매출 증감액"]] * 1.1) / 10000
        top_gains = df_diff[(df_diff[f"{prev_m_label} 매출"] > 0) & (df_diff["매출 증감액"] > 0)].sort_values(by="매출 증감액", ascending=False).head(10)
        top_drops = df_diff[(df_diff[f"{prev_m_label} 매출"] > 0) & (df_diff[f"{curr_m_label} 매출"] > 0) & (df_diff["매출 증감액"] < 0)].sort_values(by="매출 증감액", ascending=True).head(10)
        new_clients = df_diff[(df_diff[f"{prev_m_label} 매출"] == 0) & (df_diff[f"{curr_m_label} 매출"] > 0)].sort_values(by=f"{curr_m_label} 매출", ascending=False)
        lost_clients = df_diff[(df_diff[f"{prev_m_label} 매출"] > 0) & (df_diff[f"{curr_m_label} 매출"] == 0)].sort_values(by=f"{prev_m_label} 매출", ascending=False)
        mom_view = st.radio(
            "전월 대비 분석 보기",
            ["🚀 상승 Top 10", "📉 하락 Top 10", "🎉 신규/재개 거래처", "⚠️ 미거래/이탈 의심"],
            horizontal=True,
            key="tab1_mom_view",
            label_visibility="collapsed",
        )
        if mom_view == "🚀 상승 Top 10":
            st.markdown(f"**🔥 기존 거래처 중 매출이 가장 많이 [상승]한 10곳 (단위: 만원, VAT 포함)**")
            if top_gains.empty:
                st.info("해당 조건에 맞는 상승 거래처가 없습니다.")
            else:
                st.dataframe(
                    top_gains.style.format({f"{prev_m_label} 매출": "{:,.0f}", f"{curr_m_label} 매출": "{:,.0f}", "매출 증감액": "{:,.0f}"})
                    .apply(lambda s: ['color: #2563EB; font-weight: bold;' if v > 0 else '' for v in s], subset=['매출 증감액']),
                    use_container_width=True, hide_index=True, height=min(420, 38 + len(top_gains) * 35)
                )
        elif mom_view == "📉 하락 Top 10":
            st.markdown(f"**📉 기존 거래처 중 매출이 가장 많이 [하락]한 10곳 (단위: 만원, VAT 포함)**")
            if top_drops.empty:
                st.info("해당 조건에 맞는 하락 거래처가 없습니다.")
            else:
                st.dataframe(
                    top_drops.style.format({f"{prev_m_label} 매출": "{:,.0f}", f"{curr_m_label} 매출": "{:,.0f}", "매출 증감액": "{:,.0f}"})
                    .apply(lambda s: ['color: #B91C1C; font-weight: bold;' if v < 0 else '' for v in s], subset=['매출 증감액']),
                    use_container_width=True, hide_index=True, height=min(420, 38 + len(top_drops) * 35)
                )
        elif mom_view == "🎉 신규/재개 거래처":
            st.markdown(f"**🎉 {prev_m_label}엔 거래가 없었으나 {curr_m_label}에 새롭게 매출이 발생한 곳 (총 {len(new_clients)}곳, 단위: 만원, VAT 포함)**")
            if new_clients.empty:
                st.info("신규/재개 거래처가 없습니다.")
            else:
                st.dataframe(
                    new_clients[["거래처", f"{curr_m_label} 매출"]].style.format({f"{curr_m_label} 매출": "{:,.0f}"}),
                    use_container_width=True, hide_index=True, height=min(480, 38 + len(new_clients) * 35)
                )
        else:
            st.markdown(f"**⚠️ {prev_m_label}엔 매출이 있었으나 {curr_m_label}엔 거래가 없는 곳 (총 {len(lost_clients)}곳, 단위: 만원, VAT 포함)**")
            if lost_clients.empty:
                st.info("미거래/이탈 의심 거래처가 없습니다.")
            else:
                st.dataframe(
                    lost_clients[["거래처", f"{prev_m_label} 매출"]].style.format({f"{prev_m_label} 매출": "{:,.0f}"}),
                    use_container_width=True, hide_index=True, height=min(480, 38 + len(lost_clients) * 35)
                )

# Tab 2: 🏢 거래처 분석
with tab2:
    t2_c1, t2_c2 = st.columns([4, 1])
    t2_c1.markdown(f"<div class='sub-header dashboard-tab-panel-head'>🏢 [{selected_client}] 영업 실적 및 요약</div>", unsafe_allow_html=True)
    t2_c2.markdown(render_update_badge(latest_update_str), unsafe_allow_html=True)

    if "show_corp_info" not in st.session_state:
        st.session_state.show_corp_info = False

    # [핵심 패치] 거래처 필터를 바꿀 때 기업정보 창이 열려있으면 자동으로 닫아서 API 무한 로딩(프리징) 완벽 방지
    if "last_opened_client" not in st.session_state:
        st.session_state.last_opened_client = selected_client
    if st.session_state.show_corp_info and st.session_state.last_opened_client != selected_client:
        st.session_state.show_corp_info = False
    st.session_state.last_opened_client = selected_client

    # 가로 넓은 직사각형 버튼 — 글씨 한 줄로 박스 안에
    with st.container(key="tab2_action_btns"):
        btn_c1, btn_c2, btn_c3 = st.columns([2.4, 2.0, 1.8], gap="medium")
        with btn_c1:
            _notes_addr = (
                client_addr
                if client_addr and client_addr != "등록된 주소 정보가 없습니다."
                else None
            )
            _notes_disabled = selected_client == "전체 거래처"
            if _is_local_macos():
                if st.button(
                    "📝 macOS 메모에서 노트 열기/생성",
                    key="btn_notes",
                    width="stretch",
                    disabled=_notes_disabled,
                    help="특정 거래처를 선택하세요." if _notes_disabled else "메모「거래처」폴더에서 같은 거래처명 노트를 엽니다.",
                ):
                    with st.spinner("메모「거래처」폴더에서 같은 거래처명을 찾는 중..."):
                        _notes_res = open_macos_notes_folder(
                            selected_client,
                            dart_api_key,
                            df_integrated,
                            address=_notes_addr,
                        )
                    st.session_state["_tab2_loaded_note"] = {
                        **_notes_res,
                        "client": selected_client,
                    }
                _loaded = st.session_state.get("_tab2_loaded_note") or {}
                if _loaded.get("client") == selected_client:
                    if _loaded.get("ok"):
                        st.success(_loaded.get("msg") or "메모를 열었습니다.")
                        _body = str(_loaded.get("body") or "").strip()
                        _title = _loaded.get("matched_name") or selected_client
                        with st.expander(f"불러온 메모 · {_title}", expanded=True):
                            if _body:
                                st.text_area(
                                    "메모 본문",
                                    value=_body,
                                    height=220,
                                    key=f"tab2_loaded_note_body_{selected_client}",
                                    label_visibility="collapsed",
                                )
                            else:
                                st.caption("노트는 열렸지만 본문을 읽지 못했습니다.")
                    elif _loaded.get("msg"):
                        st.error(_loaded.get("msg"))
            else:
                _notes_label = "📝 거래처 메모 · 내보내기"
                with st.popover(
                    _notes_label,
                    width="stretch",
                    disabled=_notes_disabled,
                    help="특정 거래처를 선택하세요." if _notes_disabled else None,
                ):
                    if not _notes_disabled:
                        # [핵심 패치] 팝오버를 열자마자 API를 긁어와 무한 로딩에 빠지는 현상 방지
                        if st.button("🚀 메모 데이터 수집/생성 (클릭)", key="btn_gen_memo_export", use_container_width=True):
                            with st.spinner("DART/팩토리온/네이버 조회 중..."):
                                _, _n_plain, _n_html, _n_fname = prepare_client_note_export(
                                    selected_client,
                                    dart_api_key,
                                    df_integrated,
                                    address=_notes_addr,
                                )
                                st.session_state["_ready_note_export"] = {
                                    "client": selected_client,
                                    "plain": _n_plain,
                                    "html": _n_html,
                                    "fname": _n_fname
                                }
                        
                        _note_cache = st.session_state.get("_ready_note_export", {})
                        if _note_cache.get("client") == selected_client:
                            st.success("✅ 메모가 준비되었습니다!")
                            st.caption(
                                "Cloud·iPad에서는 아래 **다운로드·복사·공유**로 메모 앱에 넣을 수 있습니다."
                            )
                            st.download_button(
                                "HTML 파일 다운로드",
                                data=_note_cache["html"].encode("utf-8"),
                                file_name=_note_cache["fname"],
                                mime="text/html",
                                key="tab2_notes_download",
                                use_container_width=True,
                            )
                            st.caption("Mac: 다운로드 후 Safari로 열어 전체 선택 → 메모에 붙여넣기.")
                            _render_tab2_note_share_html(_note_cache["plain"], selected_client)
                        else:
                            st.info("👆 위 버튼을 눌러 기업 정보를 먼저 수집하세요.")
        with btn_c2:
            btn_label = "🏢 기업정보 닫기" if st.session_state.show_corp_info else "🏢 기업 기본/재무정보 보기"
            if st.button(btn_label, key="btn_dart_info", width="stretch"):
                st.session_state.show_corp_info = not st.session_state.show_corp_info
        with btn_c3:
            # 주소록 주소 우선, 없으면 거래처명으로 카카오맵 검색
            if selected_client and selected_client != "전체 거래처":
                kakao_q = client_addr if client_addr != "등록된 주소 정보가 없습니다." else selected_client
                kakao_url = f"https://map.kakao.com/link/search/{urllib.parse.quote(kakao_q)}"
                st.link_button("🗺️ 카카오맵에서 주소 보기", kakao_url, width="stretch")
            else:
                st.button(
                    "🗺️ 카카오맵에서 주소 보기",
                    disabled=True,
                    key="btn_kakao_disabled",
                    width="stretch",
                    help="사이드바에서 특정 거래처를 선택하세요.",
                )
            # 버튼 바로 아래 주소 표시 (tab2 전용, 버튼과 동일 폭·글자크기)
            _addr_color = "#64748B" if client_addr == "등록된 주소 정보가 없습니다." else "#334155"
            st.markdown(
                f"<div class='tab2-kakao-addr' style='color:{_addr_color};'>"
                f"📍 {html.escape(client_addr)}</div>",
                unsafe_allow_html=True,
            )

    # 목록에 없는 거래처도 기업정보만 조회 가능 (상단 필터·매출 집계는 변경 없음)
    _ov_c1, _ov_c2 = st.columns([1.4, 1], gap="small")
    with _ov_c1:
        _corp_name_override = st.text_input(
            "목록 외 상호 (기업정보 조회용)",
            key="tab2_corp_name_override",
            placeholder="예: OO산업 — 비우면 상단 선택 거래처 사용",
            help="매출 목록에 없어도 상호를 입력하면 기업 기본/재무·공장등록을 조회합니다.",
        )
    with _ov_c2:
        _corp_addr_override = st.text_input(
            "주소 힌트 (선택)",
            key="tab2_corp_addr_override",
            placeholder="동명 구분 · 예: 평택시 서탄면",
        )
    _corp_query_name = str(_corp_name_override or "").strip()
    if not _corp_query_name and selected_client and selected_client != "전체 거래처":
        _corp_query_name = str(selected_client).strip()
    _corp_query_is_override = bool(str(_corp_name_override or "").strip())

    if st.session_state.show_corp_info:
        if not _corp_query_name:
            st.info(
                "상단에서 거래처를 선택하거나, 위에 **목록 외 상호**를 입력한 뒤 "
                "기업정보를 다시 열어 주세요."
            )
        else:
            if _corp_query_is_override:
                st.caption(
                    f"목록 외 조회: **{_corp_query_name}** "
                    "(상단 매출 필터와 별개 · 기업정보만)"
                )
            _addr_for_lookup = None
            if str(_corp_addr_override or "").strip():
                _addr_for_lookup = str(_corp_addr_override).strip()
            elif not _corp_query_is_override:
                _addr_for_lookup = (
                    client_addr
                    if client_addr and client_addr != "등록된 주소 정보가 없습니다."
                    else None
                )
            else:
                # 목록 외 상호: 주소록에 같은 이름이 있으면 활용
                _ov_addr = resolve_client_address(_corp_query_name, addr_dict)
                if _ov_addr and _ov_addr != "등록된 주소 정보가 없습니다.":
                    _addr_for_lookup = _ov_addr

            _memo_key = (
                str(_corp_query_name),
                str(dart_api_key or ""),
                str(_addr_for_lookup or ""),
                "ov" if _corp_query_is_override else "sel",
            )
            _memo = st.session_state.get("_tab2_corp_memo")
            _use_memo = (
                isinstance(_memo, dict)
                and _memo.get("key") == _memo_key
                and isinstance(_memo.get("c_info"), dict)
            )
            _touch_corp = is_touch_ui()
            if _use_memo:
                c_info = dict(_memo["c_info"])
                _latest_audit = _memo.get("latest_audit")
                _audit_sum = dict(_memo.get("audit_sum") or {})
                _matched = (
                    c_info.get("matched_name")
                    or c_info.get("clean_name")
                    or _corp_query_name
                )
                _ccode = c_info.get("corp_code") or ""
                _lookup = _ccode or _matched
            else:
                # 1차: 기업개요·재무만 (거래처명+주소로 동명 오매칭 완화)
                with st.spinner("기업 정보 불러오는 중…"):
                    c_info = get_company_info_hybrid(
                        _corp_query_name, dart_api_key, address=_addr_for_lookup
                    )
                    _matched = (
                        c_info.get("matched_name")
                        or c_info.get("clean_name")
                        or _corp_query_name
                    )
                    _ccode = c_info.get("corp_code") or ""
                    _lookup = _ccode or _matched
                    _latest_audit = None
                    _audit_sum = {}
                    if dart_api_key and _lookup and OpenDartReader is not None:
                        _years_back = 2 if _touch_corp else 4
                        _audits = list_dart_audit_reports(
                            _lookup, dart_api_key, years_back=_years_back
                        )
                        if _audits:
                            _latest_audit = _audits[0]
                    st.session_state["_tab2_corp_memo"] = {
                        "key": _memo_key,
                        "c_info": dict(c_info),
                        "latest_audit": _latest_audit,
                        "audit_sum": {},
                    }
                    # 📊 DART 재무제표 표 화면에 출력하기
                    st.markdown("##### 📊 DART 재무제표 요약")
                    st.error(f"🕵️‍♂️ 파이썬 검색 단어: {_lookup} / 기업코드: {c_info.get('corp_code')}")
                    if _latest_audit is not None and not _latest_audit.empty:
                     st.dataframe(_latest_audit, use_container_width=True)
                    else:
                     st.info("💡 다트에 등록된 재무제표가 없는 기업(비상장 등)입니다.")
            # 감사 본문 추출: 맥은 자동, iPad는 버튼(동일 데이터·무손실)
            _want_audit_parse = bool(st.session_state.get("_tab2_force_audit_parse"))
            if (
                _latest_audit
                and not _audit_sum
                and dart_api_key
                and OpenDartReader is not None
                and (not _touch_corp or _want_audit_parse)
            ):
                with st.spinner("감사 주석 개요·계속기업 이슈 추출 중…"):
                    _audit_sum = parse_dart_audit_report_summary(
                        _latest_audit["rcept_no"], dart_api_key
                    )
                    if _audit_sum.get("revenue") and c_info.get("revenue") == "정보 없음":
                        c_info["revenue"] = _audit_sum["revenue"] + " (감사보고서 추정)"
                    if _audit_sum.get("profit") and c_info.get("profit") == "정보 없음":
                        c_info["profit"] = _audit_sum["profit"] + " (감사보고서 추정)"
                    # 주석 개요로 대표/업종 보강 (기존 값 있을 때는 덮지 않음)
                    if _audit_sum.get("ceo_note") and c_info.get("ceo") in (
                        "",
                        "정보 없음",
                        None,
                    ):
                        c_info["ceo"] = _audit_sum["ceo_note"]
                    if _audit_sum.get("business") and c_info.get("industry") in (
                        "",
                        "정보 없음",
                        None,
                    ):
                        c_info["industry"] = _audit_sum["business"]
                    if (
                        _audit_sum.get("revenue")
                        or _audit_sum.get("profit")
                        or _audit_sum.get("overview_ok")
                    ):
                        if "DART" not in str(c_info.get("source")):
                            c_info["source"] = "DART 감사보고서 주석·본문"
                    st.session_state["_tab2_corp_memo"] = {
                        "key": _memo_key,
                        "c_info": dict(c_info),
                        "latest_audit": _latest_audit,
                        "audit_sum": dict(_audit_sum) if _audit_sum else {},
                    }
                    st.session_state.pop("_tab2_force_audit_parse", None)
                      # 📊 DART 재무제표 표 화면에 출력하기
                    st.markdown("##### 📊 DART 재무제표 요약")
                    
                    # 👉 파이썬이 무슨 단어로 검색했는지 화면에 박제하기!
                    st.error(f"🕵️‍♂️ 파이썬 검색 단어: {_lookup} / 기업코드: {c_info.get('corp_code')}")
                    
                    if _latest_audit is not None and not _latest_audit.empty:
                        st.dataframe(_latest_audit, use_container_width=True)
                    else:
                        st.info("💡 다트에 등록된 재무제표가 없는 기업(비상장 등)입니다.")
            def _autosave_factory_api_key():
                v = str(st.session_state.get("tab2_factory_api_key_input") or "").strip()
                if not v:
                    return
                old = _load_factory_api_key()
                if v == old:
                    return
                _persist_factory_api_key(v)
                try:
                    fetch_factory_registry.clear()
                except Exception:
                    pass
                st.session_state.pop("_tab2_factory_memo", None)

            _factory_key = _load_factory_api_key()
            if "tab2_factory_api_key_input" not in st.session_state:
                st.session_state["tab2_factory_api_key_input"] = _factory_key

            # 공장등록 조회 (표시 전에 취합)
            _f_memo_key = (
                str(_corp_query_name),
                str(_matched or ""),
                str(_addr_for_lookup or ""),
                str(_load_factory_api_key() or ""),
            )
            _f_memo = st.session_state.get("_tab2_factory_memo")
            _f_use = (
                isinstance(_f_memo, dict)
                and _f_memo.get("key") == _f_memo_key
                and isinstance(_f_memo.get("info"), dict)
            )
            if _f_use:
                _f_info = dict(_f_memo["info"])
            elif not _load_factory_api_key():
                _f_info = {"ok": False, "error": "공장등록 API 키 없음"}
            else:
                with st.spinner("공장등록 정보 조회 중…"):
                    _f_info = fetch_factory_registry(
                        _matched or _corp_query_name,
                        address=_addr_for_lookup,
                        api_key=_load_factory_api_key(),
                    )
                st.session_state["_tab2_factory_memo"] = {
                    "key": _f_memo_key,
                    "info": dict(_f_info) if isinstance(_f_info, dict) else {},
                }
            if not isinstance(_f_info, dict):
                _f_info = {"ok": False, "error": "공장등록 조회 실패"}

            def _corp_val(*vals):
                for v in vals:
                    s = str(v or "").strip()
                    if s and s not in ("정보 없음", "-", "None", "nan"):
                        return s
                return ""

            _ceo = _corp_val(
                c_info.get("ceo"),
                _f_info.get("ceo"),
                (_audit_sum or {}).get("ceo_note"),
            )
            _industry = _corp_val(
                c_info.get("industry"),
                _f_info.get("industry"),
                (_audit_sum or {}).get("business"),
            )
            _addr_show = _corp_val(
                _f_info.get("address"),
                _addr_for_lookup,
                (_audit_sum or {}).get("hq"),
            )
            _product = _corp_val(_f_info.get("product"))
            _tel = _corp_val(_f_info.get("tel"))
            _hp = _corp_val(_f_info.get("homepage"))
            if _hp and not _hp.startswith("http"):
                _hp_href = "https://" + _hp
            else:
                _hp_href = _hp
            _rev = _corp_val(c_info.get("revenue")) or "정보 없음"
            _prf = _corp_val(c_info.get("profit")) or "정보 없음"
            _op = _corp_val((_audit_sum or {}).get("opinion"))
            _op_color = {
                "적정의견": "#166534",
                "한정의견": "#A16207",
                "부적정의견": "#B91C1C",
                "의견거절": "#9F1239",
            }.get(_op, "#334155")
            _gc_issue = ((_audit_sum or {}).get("going_concern_issue") or "").strip()
            _gc_flag = bool((_audit_sum or {}).get("going_concern_flag"))

            _toolbar = st.columns([1.2, 1, 3], gap="small")
            with _toolbar[0]:
                if st.button("🔄 다시 조회", key="btn_refresh_corp", width="stretch"):
                    get_company_info_hybrid.clear()
                    list_dart_audit_reports.clear()
                    parse_dart_audit_report_summary.clear()
                    try:
                        fetch_factory_registry.clear()
                    except Exception:
                        pass
                    try:
                        _make_opendart_reader.clear()
                    except Exception:
                        pass
                    st.session_state.pop("_opendart_last_error", None)
                    st.session_state.pop("_tab2_corp_memo", None)
                    st.session_state.pop("_tab2_factory_memo", None)
                    st.session_state.pop("_tab2_force_audit_parse", None)
                    st.rerun()
            with _toolbar[1]:
                if _touch_corp and _latest_audit and not _audit_sum:
                    if st.button("📄 감사추출", key="btn_parse_audit_sum", width="stretch"):
                        st.session_state["_tab2_force_audit_parse"] = True
                        st.rerun()

            # 통합 카드 (중복 제거: 대표/업종/주소 1회)
            _rows_basic = [
                ("상호", html.escape(_matched or _corp_query_name)),
                ("대표", html.escape(_ceo or "-")),
                ("업종", html.escape(_industry or "-")),
                ("주소", html.escape(_addr_show or "-")),
                ("전화", html.escape(_tel or "-")),
                (
                    "홈페이지",
                    (
                        f"<a href='{html.escape(_hp_href)}' target='_blank' rel='noopener'>"
                        f"{html.escape(_hp)}</a>"
                        if _hp_href
                        else "-"
                    ),
                ),
            ]
            _rows_fin = [
                ("매출액", html.escape(_rev)),
                ("영업이익", html.escape(_prf)),
            ]
            _rows_fac = []
            if _f_info.get("ok"):
                for lab, key in (
                    ("주생산품", "product"),
                    ("용지면적", "land_area"),
                    ("건축면적", "bldg_area"),
                    ("용도지역", "zone"),
                    ("행정기관", "admin"),
                    ("등록일자", "reg_date"),
                    ("고용인원", "employees"),
                    ("산업단지", "complex"),
                ):
                    vv = _corp_val(_f_info.get(key))
                    if vv:
                        _rows_fac.append((lab, html.escape(vv)))

            def _grid_html(rows):
                parts = ['<div class="tab2-corp-grid">']
                for k, v in rows:
                    parts.append(
                        f'<div class="row"><span class="k">{html.escape(k)}</span>'
                        f'<span class="v">{v}</span></div>'
                    )
                parts.append("</div>")
                return "".join(parts)

            _audit_html = ""
            if _latest_audit:
                _audit_html += (
                    f'<div class="tab2-corp-sec"><div class="sec-title">감사 · 리스크</div>'
                    f'<div style="font-size:13px;margin-bottom:4px;">'
                    f'<a href="{html.escape(_latest_audit["url"])}" target="_blank" rel="noopener">'
                    f'{html.escape(_latest_audit["date"])} · {html.escape(_latest_audit["name"])}'
                    f"</a></div>"
                )
                if _op:
                    _audit_html += (
                        f'<span class="tab2-corp-op" style="color:{_op_color};">'
                        f"감사의견: {html.escape(_op)}</span>"
                    )
                if _gc_flag and _gc_issue:
                    _audit_html += (
                        f'<div style="margin-top:8px;padding:10px 12px;border:1px solid #FECACA;'
                        f'border-radius:8px;background:#FEF2F2;color:#7F1D1D;font-size:13px;'
                        f'line-height:1.45;">{html.escape(_gc_issue)}</div>'
                    )
                elif _audit_sum:
                    _gc_cap = _corp_val((_audit_sum or {}).get("going_concern")) or "관련 문구 없음"
                    _audit_html += (
                        f'<div style="margin-top:6px;font-size:12px;color:#64748B;">'
                        f"계속기업: {html.escape(_gc_cap)}</div>"
                    )
                _audit_html += "</div>"
            elif dart_api_key and OpenDartReader is not None:
                _audit_html = (
                    '<div class="tab2-corp-sec"><div class="sec-title">감사 · 리스크</div>'
                    '<div style="font-size:12px;color:#64748B;">최근 감사보고서 공시 없음</div></div>'
                )

            _fac_html = ""
            if _rows_fac:
                _fac_html = (
                    '<div class="tab2-corp-sec"><div class="sec-title">공장등록 (팩토리온)</div>'
                    + _grid_html(_rows_fac)
                    + "</div>"
                )
            elif _f_info.get("error") and "키 없음" not in str(_f_info.get("error")):
                _fac_html = (
                    '<div class="tab2-corp-sec"><div class="sec-title">공장등록 (팩토리온)</div>'
                    f'<div style="font-size:12px;color:#64748B;">{html.escape(str(_f_info.get("error")))}'
                    "</div></div>"
                )

            _src_bits = []
            if c_info.get("source"):
                _src_bits.append(str(c_info.get("source")))
            if _f_info.get("ok"):
                _src_bits.append("팩토리온")
            _src_line = " · ".join(dict.fromkeys(_src_bits)) if _src_bits else ""

            st.markdown(
                f"""
    <div class="tab2-corp-card">
      <h4>🏢 {html.escape(str(_corp_query_name))}
        <span style="font-size:12px;font-weight:500;color:#94A3B8;margin-left:8px;">
          {html.escape(_src_line)}</span>
      </h4>
      <div class="sec-title">기본 · 재무</div>
      {_grid_html(_rows_basic + _rows_fin)}
      {_fac_html}
      {_audit_html}
    </div>
    """,
                unsafe_allow_html=True,
            )

            if (_rev == "정보 없음" or _prf == "정보 없음") and c_info.get("dart_error"):
                st.caption(f"DART: {c_info.get('dart_error')}")

            _sh = (_audit_sum or {}).get("shareholders") or []
            if _sh:
                with st.expander("주요 주주 · 지분율", expanded=False):
                    st.dataframe(
                        pd.DataFrame(_sh),
                        width="stretch",
                        hide_index=True,
                        height=min(160, 38 + 28 * len(_sh)),
                    )

            with st.expander(
                "공장등록 API 키 (자동저장)",
                expanded=not bool(_load_factory_api_key()),
            ):
                st.caption(
                    "입력 후 Enter(또는 포커스 이동) 시 자동 저장됩니다. "
                    "Decoding 키(끝 `==`) 권장. "
                    "([생산정보](https://www.data.go.kr/data/15087611/openapi.do) · "
                    "[필지](https://www.data.go.kr/data/15087615/openapi.do))"
                )
                st.text_input(
                    "공공데이터 일반 인증키",
                    type="password",
                    key="tab2_factory_api_key_input",
                    placeholder="data.go.kr 일반 인증키",
                    on_change=_autosave_factory_api_key,
                    label_visibility="collapsed",
                )
                if _load_factory_api_key():
                    st.caption("✓ 키가 저장되어 있습니다.")

            _loc_bits = _loc_tokens_from_address(_addr_for_lookup)
            _q_base = str(_matched)
            _q_merged = f"{_q_base} {' '.join(_loc_bits[:2])}".strip() if _loc_bits else _q_base
            _q = urllib.parse.quote(_q_merged)
            _links = (c_info.get("job_links") or {}) if isinstance(c_info, dict) else {}
            _saramin = _links.get("saramin_company") or (
                "https://www.saramin.co.kr/zf_user/search/company?searchword=" + _q
            )
            st.markdown(
                f"[DART](https://dart.fss.or.kr/) · "
                f"[네이버 기업정보](https://search.naver.com/search.naver?query={_q}%20기업정보) · "
                f"[사람인]({_saramin})"
                + (
                    f" · [팩토리온 원문]({_f_info.get('source_url')})"
                    if _f_info.get("ok")
                    else ""
                )
            )
    m1, m2, m3, m4 = st.columns(4)
    tot_sales_c = df_client_filtered["매출액"].sum() * 1.1 / 10000 if not df_client_filtered.empty else 0.0

    cur_sales_c = cur_month_sales_client / 10000
    m1.markdown(f"<div class='metric-box'><div class='metric-label'>총 누적 매출 (VAT포함)</div><div class='metric-value'>{tot_sales_c:,.0f} 만원</div></div>", unsafe_allow_html=True)
    m2.markdown(f"<div class='metric-box'><div class='metric-label'>최근 월 매출 ({latest_month_str_client})</div><div class='metric-value'>{cur_sales_c:,.0f} 만원</div></div>", unsafe_allow_html=True)
    m3.markdown(f"<div class='metric-box'><div class='metric-label'>전월 대비 (MoM)</div><div class='metric-value' style='color:{'#E11D48' if mom_rate_client < 0 else '#2563EB'};'>{mom_rate_client:+.0f}%</div></div>", unsafe_allow_html=True)
    m4.markdown(f"<div class='metric-box'><div class='metric-label'>월평균 대비 증감</div><div class='metric-value' style='color:{'#E11D48' if avg_rate_client < 0 else '#2563EB'};'>{avg_rate_client:+.0f}%</div></div>", unsafe_allow_html=True)
    if not df_client_filtered.empty:
        pivot_m_client = cached_get_yearly_monthly_pivot(df_client_filtered, all_months, years)
        cl, cr = st.columns([1, 1])
        with cl:
            pivot_m_client_disp = get_display_df_with_sum(pivot_m_client, "연간 합계")
            st.dataframe(style_with_sum(pivot_m_client_disp, "{:,.0f}", "Blues", axis=None), use_container_width=True, height=460)
        with cr:
            render_plotly_chart(
                create_stacked_bar_chart(pivot_m_client, title_text=""),
                use_container_width=True, key="tab2_client_total_chart"
            )
        # —— 매출 비교를 품목별 상세 분석보다 위에 배치 (tab2 전용 순서) ——
        st.markdown("---")
        _client_years = sorted(
            {str(y) for y in df_client_filtered["연도"].dropna().unique()},
            key=lambda y: int(y) if str(y).isdigit() else 0,
        )
        _cur_y = _client_years[-1] if _client_years else (str(years[0]) if years else None)
        _prev_y = None
        if _cur_y:
            try:
                _prev_cand = str(int(_cur_y) - 1)
            except Exception:
                _prev_cand = None
            if _prev_cand and _prev_cand in _client_years:
                _prev_y = _prev_cand
            elif len(_client_years) >= 2:
                _prev_y = _client_years[-2]
        _yr_label = f"{_prev_y}·{_cur_y}" if _prev_y else str(_cur_y)
        # 당월 기준: 미래 월 제외, 당월→과거 역순 (26년 08·07·…·01 → 25년 12·…·01)
        _latest_dt_sales = df_client_filtered["매출일_dt"].max()
        _cur_month = (
            _latest_dt_sales.strftime("%m월")
            if pd.notnull(_latest_dt_sales)
            else all_months[0]
        )
        _mi = all_months.index(_cur_month) if _cur_month in all_months else 0
        _sales_col_keys = []
        if _cur_y:
            _ys = str(_cur_y)[2:]
            for i in range(_mi, -1, -1):
                _sales_col_keys.append(f"{_ys}년 {all_months[i]}")
        if _prev_y:
            _ps = str(_prev_y)[2:]
            for i in range(len(all_months) - 1, -1, -1):
                _sales_col_keys.append(f"{_ps}년 {all_months[i]}")
        # 오른쪽 월 그래프: 당월→01월 역순만 (미도래 월 제외)
        _months_back = list(reversed(all_months[: _mi + 1]))
        st.markdown(
            f"<div class='sub-header dashboard-tab-panel-head'>"
            f"📊 [{selected_client}] 매출 비교 ({_yr_label}, {_cur_month}→과거)</div>",
            unsafe_allow_html=True,
        )
        st.caption(
            f"왼쪽: 품목별 매출(만원·VAT포함) · {_cur_y}년 {_cur_month}→과거 → {_prev_y or '전년'}년 "
            f"(매출 많은 순) · 오른쪽: 월별 매출 비교 그래프 · 하단: 당해 년평균 매출 비중"
        )
        sales_two_y = cached_client_item_sales_pivot_two_years(
            df_client_filtered, tuple(_sales_col_keys)
        )
        sales_month_two_y = cached_get_yearly_monthly_pivot(
            df_client_filtered,
            _months_back,
            [y for y in (_cur_y, _prev_y) if y],
        )
        q_left, q_right = st.columns([1, 1])
        with q_left:
            if sales_two_y.empty:
                st.info("전년·당해 매출 데이터가 없습니다.")
            else:
                # 품목별 총 매출 많은 순 내림차순 — tab2 매출 비교만
                sales_two_y = sales_two_y.loc[
                    sales_two_y.sum(axis=1).sort_values(ascending=False).index
                ]
                sales_disp = get_display_df_with_sum(sales_two_y, "합계")
                st.dataframe(
                    style_with_sum(sales_disp, "{:,.0f}", "Blues", axis=None),
                    use_container_width=True,
                    height=460,
                )
        with q_right:
            if sales_month_two_y.empty or (sales_month_two_y.fillna(0) == 0).all().all():
                st.info("전년·당해 월별 매출 데이터가 없습니다.")
            else:
                # 월별(당해·전년) 매출 비교 그룹 막대 — tab2 전용
                _fig_sales_cmp = create_grouped_bar_chart(
                    sales_month_two_y,
                    title_text=f"월별 매출 비교 ({_yr_label}, 만원)",
                    y_suffix="만원",
                    y_format=",.0f",
                )
                render_plotly_chart(
                    _fig_sales_cmp,
                    use_container_width=True,
                    key="tab2_sales_yoy_grouped",
                )
        # 당해 년평균(월평균) 매출 기준 품목별 비중 — 가로 막대(블루 톤)
        if _cur_y:
            _df_cy = df_client_filtered[
                df_client_filtered["연도"].astype(str) == str(_cur_y)
            ]
            if not _df_cy.empty and "매출액" in _df_cy.columns:
                _m_pvt = _df_cy.pivot_table(
                    index="품목명", columns="월", values="매출액", aggfunc="sum"
                ).fillna(0)
                _yr_avg = _m_pvt.mean(axis=1)
                _yr_avg = _yr_avg[_yr_avg > 0]
                if not _yr_avg.empty:
                    _share = (_yr_avg / _yr_avg.sum() * 100).sort_values(ascending=False)
                    _fig_share = create_item_share_hbar(
                        _share,
                        title_text=f"당해({_cur_y}) 년평균 매출 기준 품목별 비중",
                    )
                    if _fig_share is not None:
                        render_plotly_chart(
                            _fig_share,
                            use_container_width=True,
                            key="tab2_cur_year_item_share",
                        )
        st.markdown("---")
        st.markdown(f"<div class='sub-header dashboard-tab-panel-head'>📦 [{selected_client}] 품목별 상세 분석</div>", unsafe_allow_html=True)
    
        client_available_items = sorted(df_client_filtered["품목명"].unique())
        if client_available_items:
            item_ratios = {}
            if not df_client_filtered.empty:
                latest_dt_c = df_client_filtered["매출일_dt"].max()
                if pd.notnull(latest_dt_c):
                    latest_ym = latest_dt_c.strftime("%Y-%m")
                    df_cm = df_client_filtered[df_client_filtered["매출일_dt"].dt.strftime("%Y-%m") == latest_ym]
                    tot_sales = df_cm["매출액"].sum()
                    if tot_sales > 0:
                        grp = df_cm.groupby("품목명")["매출액"].sum()
                        for item, val in grp.items():
                            item_ratios[item] = (val / tot_sales) * 100
            def format_item_with_ratio(item_name):
                pct = item_ratios.get(item_name, 0.0)
                return f"{item_name} (당월 {pct:.1f}%)"
            sel_col1_c, sel_col2_c = st.columns([1, 1])
            with sel_col1_c:
                selected_target_item_c = st.selectbox(
                    "🔍 분석할 품목 선택 (전체 거래 품목)", 
                    options=client_available_items, 
                    format_func=format_item_with_ratio,
                    key="client_item_selectbox"
                )
            with sel_col2_c:
                selected_metric_c = st.radio("📊 분석 지표 선택", ["매출액 (만원)", "출고량", "총매출 대비 비중 (%)"], horizontal=True, key="client_metric_radio")
            
            client_item_pivot = cached_get_item_pivot(df_client_filtered, selected_target_item_c, selected_metric_c, all_months, years)
        
            i_col_left_c, i_col_right_c = st.columns([1, 1])
            with i_col_left_c:
                client_item_pivot_disp = get_display_df_with_sum(client_item_pivot, "연간 합계")
                # tab2 표 색상 통일: Blues 그라데이션
                if "비중" in selected_metric_c:
                    st.dataframe(style_with_sum(client_item_pivot_disp, "{:,.1f}%", "Blues", axis=None), use_container_width=True, height=460)
                    y_suf_c, y_fmt_c = "%", ",.1f"
                elif "출고량" in selected_metric_c:
                    if selected_target_item_c in target_items:
                        y_suf_c, y_fmt_c = " 천kg", ",.1f"
                    elif "LPG" in str(selected_target_item_c).upper():
                        y_suf_c, y_fmt_c = " kg", ",.0f"
                    else:
                        y_suf_c, y_fmt_c = " 개(병)", ",.0f"
                    st.dataframe(style_with_sum(client_item_pivot_disp, f"{{:{y_fmt_c}}}", "Blues", axis=None), use_container_width=True, height=460)
                else:
                    st.dataframe(style_with_sum(client_item_pivot_disp, "{:,.0f}", "Blues", axis=None), use_container_width=True, height=460)
                    y_suf_c, y_fmt_c = " 만원", ",.0f"
                
            with i_col_right_c:
                render_plotly_chart(
                    create_stacked_bar_chart(
                        client_item_pivot, 
                        title_text="", 
                        y_suffix=y_suf_c, 
                        y_format=y_fmt_c
                    ),
                    use_container_width=True, key="tab2_client_item_chart"
                )
# Tab 3: 📦 품목 및 단가 분석
with tab3:
    t3_c1, t3_c2 = st.columns([4, 1])
    t3_c1.markdown(f"<div class='sub-header dashboard-tab-panel-head'>📦 [{selected_client}] 품목별 실적 분석</div>", unsafe_allow_html=True)
    t3_c2.markdown(render_update_badge(latest_update_str), unsafe_allow_html=True)

    latest_dt_overall = df_base["매출일_dt"].max() if not df_base.empty else None
    target_month_col = latest_dt_overall.strftime("%y년 %m월") if pd.notnull(latest_dt_overall) else None
    avail_years_short = [y[2:] for y in years]
    current_year_short = str(df_base["연도"].max())[2:] if not df_base.empty else (avail_years_short[0] if avail_years_short else "26")

    selected_detail_years = st.multiselect(
        "📅 월별 상세 내역을 펼쳐볼 연도 선택 (단가표 제외)",
        options=avail_years_short,
        default=[current_year_short] if current_year_short in avail_years_short else avail_years_short[:1],
        format_func=lambda x: f"20{x}년",
        key="tab3_detail_years",
    )
    
    sales_p_filtered, qty_p_filtered = cached_filter_tab3_year_columns(
        sales_p,
        qty_p,
        tuple(sorted(selected_detail_years)),
        tuple(avail_years_short),
        tuple(all_months),
    )
    st.markdown("<div style='font-size: 14px; font-weight: 600; color: #334155; margin-bottom: 10px;'>1️⃣ 매출액 (VAT 포함, 만원)</div>", unsafe_allow_html=True)
    render_tab3_dataframe_table(
        sales_p_filtered, "{:,.0f}", target_month_col, key_prefix="tab3_sales", table_kind="sales"
    )
    # 상단 담당자·품목 선택 시 → 품목별 거래처 매출 상세 (클릭 펼침)
    if selected_staff and selected_item:
        render_tab3_item_client_expanders(
            df_f,
            list(selected_item),
            "sales",
            years,
            all_months,
            selected_detail_years,
        )
    elif selected_item and not selected_staff:
        st.caption("💡 거래처별 상세를 보려면 상단 고정바에서 담당자를 먼저 선택한 뒤 품목을 선택하세요.")
    
    st.markdown("<div style='font-size: 14px; font-weight: 600; color: #334155; margin-bottom: 10px;'>2️⃣ 출고량</div>", unsafe_allow_html=True)
    render_tab3_dataframe_table(
        qty_p_filtered, "{:,.0f}", target_month_col, key_prefix="tab3_qty", table_kind="qty"
    )
    if selected_staff and selected_item:
        render_tab3_item_client_expanders(
            df_f,
            list(selected_item),
            "qty",
            years,
            all_months,
            selected_detail_years,
        )
    
    st.markdown("<div style='font-size: 14px; font-weight: 600; color: #334155; margin-bottom: 10px;'>3️⃣ 적용 단가 (실제 원본 단가) - 전체 기간 월별 고정 표시</div>", unsafe_allow_html=True)
    _price_sort_order = list(sales_p.index) if not sales_p.empty else None
    _price_items = []
    if not unit_price_p.empty:
        if _price_sort_order:
            _price_items = [i for i in _price_sort_order if i in unit_price_p.index]
            _price_items += [i for i in unit_price_p.index if i not in _price_items]
        else:
            _price_items = list(unit_price_p.index)
    _price_view = st.selectbox(
        "단가 보기 품목",
        options=["전체 품목 (월별 고정)"] + _price_items,
        key="tab3_price_item_select",
        help="특정 품목을 고르면 단가가 처음 적용·변동된 연월만 표시합니다.",
    )
    if _price_view == "전체 품목 (월별 고정)":
        st.caption("전체 기간 월별 고정(이월) 단가입니다. 품목을 선택하면 변동 연월만 봅니다.")
        render_tab3_dataframe_table(
            unit_price_p,
            "{:,.0f}",
            target_month_col,
            key_prefix="tab3_price",
            table_kind="price",
            sort_order=_price_sort_order,
        )
    else:
        st.caption(
            f"[{_price_view}] · 열=단가 최초 적용·변동 연월만 · 값=그 시점 단가 "
            "(월별 고정/이월 표시가 아님)"
        )
        _chg = build_unit_price_change_pivot(
            unit_price_p, years, all_months, item_names=[_price_view]
        )
        if _chg.empty:
            st.info("이 품목의 단가 변동(또는 최초 적용) 이력이 없습니다.")
        else:
            _chg_fmt = {c: "{:,.0f}" for c in _chg.columns}
            st.dataframe(
                _chg.style.format(_chg_fmt, na_rep=""),
                use_container_width=True,
                height=120,
            )

    # —— 4️⃣ 거래처 선택 시: 벌크(주요) / 그외가스(부품목) 납품량 기준 사용·재고 ——
    st.markdown(
        "<div style='font-size: 14px; font-weight: 600; color: #334155; margin: 18px 0 8px;'>"
        "4️⃣ 가스 사용량 · 재고관리 (납품량 기준 · 벌크 / 그외)</div>",
        unsafe_allow_html=True,
    )
    if selected_client == "전체 거래처":
        st.info("상단에서 거래처를 선택하면 해당 거래처의 벌크·그외 가스 월/주/일 사용량을 볼 수 있습니다.")
    elif df_client_filtered.empty:
        st.warning("선택한 거래처의 매출·출고 데이터가 없습니다.")
    else:
        _u_years_all = sorted(
            {str(y) for y in df_client_filtered["연도"].dropna().unique()},
            key=lambda y: int(y) if str(y).isdigit() else 0,
        )
        _u_cur = _u_years_all[-1] if _u_years_all else (
            str(years[-1]) if years else None
        )
        _u_prev = None
        if _u_cur:
            try:
                _u_prev_cand = str(int(_u_cur) - 1)
            except Exception:
                _u_prev_cand = None
            if _u_prev_cand and _u_prev_cand in _u_years_all:
                _u_prev = _u_prev_cand
            elif len(_u_years_all) >= 2:
                _u_prev = _u_years_all[-2]
        _u_default_years = [y for y in (_u_prev, _u_cur) if y]
        if not _u_default_years and _u_years_all:
            _u_default_years = _u_years_all[-1:]

        st.caption(
            "납품(출고)량으로 사용량을 산출합니다. 기본 기준기간은 **전년도 + 당해년도**이며, "
            "연도·월을 바꿔 납품 참고 구간을 조정할 수 있습니다. "
            "월사용량 = 총납품 ÷ 기준기간 달력 월수(미래월 제외) · 주 = 월×7/30 · 일 = 월/30."
        )
        _uc1, _uc2 = st.columns([1.2, 1.8])
        with _uc1:
            _u_sel_years = st.multiselect(
                "📅 납품 기준 연도",
                options=_u_years_all,
                default=[y for y in _u_default_years if y in _u_years_all],
                format_func=lambda x: f"{x}년",
                key="tab3_usage_years",
                help="기본: 전년도 + 당해년도",
            )
        with _uc2:
            _u_sel_months = st.multiselect(
                "📆 납품 기준 월 (비우면 선택 연도의 전체 월)",
                options=all_months,
                default=[],
                key="tab3_usage_months",
                help="특정 월만 보고 싶을 때 선택. 비우면 선택 연도 전체.",
            )
        if not _u_sel_years:
            st.warning("기준 연도를 하나 이상 선택하세요.")
        else:
            _u_sum, _u_monthly, _u_meta = cached_tab3_client_gas_usage(
                df_client_filtered,
                tuple(_u_sel_years),
                tuple(_u_sel_months),
            )
            _yr_lbl = "·".join(_u_sel_years)
            _mo_lbl = (
                ",".join(_u_sel_months) if _u_sel_months else "전체 월"
            )
            if _u_sum.empty:
                st.info(
                    f"[{selected_client}] {_yr_lbl} / {_mo_lbl} 구간에 "
                    "가스 납품(출고) 실적이 없습니다."
                )
            else:
                _k1, _k2, _k3, _k4 = st.columns(4)
                _k1.markdown(
                    f"<div class='metric-box'><div class='metric-label'>"
                    f"기준기간</div><div class='metric-value' style='font-size:16px;'>"
                    f"{html.escape(_yr_lbl)} · {_u_meta['n_months']}개월"
                    f"</div></div>",
                    unsafe_allow_html=True,
                )
                _k2.markdown(
                    f"<div class='metric-box'><div class='metric-label'>"
                    f"🛢️ 벌크 월사용량 합</div><div class='metric-value' style='color:#1D4ED8;'>"
                    f"{_u_meta['bulk_month']:,.0f}</div></div>",
                    unsafe_allow_html=True,
                )
                _k3.markdown(
                    f"<div class='metric-box'><div class='metric-label'>"
                    f"🧪 그외가스 월사용량 합</div><div class='metric-value' style='color:#0F766E;'>"
                    f"{_u_meta['other_month']:,.0f}</div></div>",
                    unsafe_allow_html=True,
                )
                _k4.markdown(
                    f"<div class='metric-box'><div class='metric-label'>"
                    f"품목 수</div><div class='metric-value'>"
                    f"{len(_u_sum):,} 종</div></div>",
                    unsafe_allow_html=True,
                )

                _bulk_df = _u_sum[_u_sum["구분"] == "벌크(주요)"]
                _other_df = _u_sum[_u_sum["구분"] == "그외가스(부품목)"]
                _left, _right = st.columns(2)
                with _left:
                    st.markdown(
                        "<div style='font-size:13px;font-weight:700;color:#1E3A8A;"
                        "margin:4px 0 8px;'>🛢️ 주요품목 · 벌크</div>",
                        unsafe_allow_html=True,
                    )
                    if _bulk_df.empty:
                        st.caption("이 기간 벌크 납품 없음")
                    else:
                        for _, _r in _bulk_df.iterrows():
                            _sub = (
                                f"총 {_r['총납품량']:,.0f} · 납품 {_r['납품횟수']}회 · "
                                f"간격 {_r['평균납품간격(일)']:.0f}일 · "
                                f"최근 {_r['최근납품일']} · "
                                f"회당≈{_r['회당평균']:,.0f} (≈{_r['예상소진(일)']:.0f}일분)"
                            )
                            st.markdown(
                                _tab3_usage_inv_card_html(
                                    str(_r["품목명"]),
                                    _sub,
                                    float(_r["월사용량"]),
                                    float(_r["주사용량"]),
                                    float(_r["일사용량"]),
                                    accent="#1D4ED8",
                                ),
                                unsafe_allow_html=True,
                            )
                with _right:
                    st.markdown(
                        "<div style='font-size:13px;font-weight:700;color:#115E59;"
                        "margin:4px 0 8px;'>🧪 부품목 · 그외 가스</div>",
                        unsafe_allow_html=True,
                    )
                    if _other_df.empty:
                        st.caption("이 기간 그외 가스 납품 없음")
                    else:
                        for _, _r in _other_df.head(12).iterrows():
                            _sub = (
                                f"총 {_r['총납품량']:,.0f} · 납품 {_r['납품횟수']}회 · "
                                f"간격 {_r['평균납품간격(일)']:.0f}일 · "
                                f"최근 {_r['최근납품일']}"
                            )
                            st.markdown(
                                _tab3_usage_inv_card_html(
                                    str(_r["품목명"]),
                                    _sub,
                                    float(_r["월사용량"]),
                                    float(_r["주사용량"]),
                                    float(_r["일사용량"]),
                                    accent="#0F766E",
                                ),
                                unsafe_allow_html=True,
                            )
                        if len(_other_df) > 12:
                            st.caption(f"외 {len(_other_df) - 12}개 품목 → 아래 표 참고")

                _disp_cols = [
                    "구분",
                    "품목명",
                    "월사용량",
                    "주사용량",
                    "일사용량",
                    "총납품량",
                    "납품횟수",
                    "회당평균",
                    "평균납품간격(일)",
                    "예상소진(일)",
                    "최근납품일",
                    "활성월수",
                ]
                _tbl = _u_sum[[c for c in _disp_cols if c in _u_sum.columns]].copy()
                st.markdown(
                    "<div style='font-size:13px;font-weight:600;color:#334155;"
                    "margin:12px 0 6px;'>📋 사용량 상세표</div>",
                    unsafe_allow_html=True,
                )
                st.dataframe(
                    _tbl.style.format(
                        {
                            "월사용량": "{:,.1f}",
                            "주사용량": "{:,.1f}",
                            "일사용량": "{:,.1f}",
                            "총납품량": "{:,.0f}",
                            "회당평균": "{:,.0f}",
                            "평균납품간격(일)": "{:,.1f}",
                            "예상소진(일)": "{:,.0f}",
                        }
                    ),
                    use_container_width=True,
                    height=min(420, 56 + 28 * max(len(_tbl), 1)),
                )

                # 월별 납품 추이 (벌크 우선, 없으면 상위 그외)
                _chart_items = list(_bulk_df["품목명"]) if not _bulk_df.empty else []
                if len(_chart_items) < 4 and not _other_df.empty:
                    _chart_items += list(_other_df["품목명"].head(4 - len(_chart_items)))
                if (
                    _chart_items
                    and not _u_monthly.empty
                    and any(i in _u_monthly.index for i in _chart_items)
                ):
                    _plot = _u_monthly.reindex(
                        [i for i in _chart_items if i in _u_monthly.index]
                    )
                    if not _plot.empty and (_plot.fillna(0) != 0).any().any():
                        _plot_t = _plot.T.copy()
                        _plot_t.index.name = "연월"
                        _melt = _plot_t.reset_index().melt(
                            id_vars="연월", var_name="품목명", value_name="납품량"
                        )
                        _fig_u = px.bar(
                            _melt,
                            x="연월",
                            y="납품량",
                            color="품목명",
                            barmode="group",
                            title=f"[{selected_client}] 월별 납품량 추이 ({_yr_lbl})",
                        )
                        _fig_u.update_layout(
                            margin=dict(l=10, r=10, t=40, b=10),
                            paper_bgcolor="rgba(0,0,0,0)",
                            plot_bgcolor="rgba(0,0,0,0)",
                            legend=dict(
                                orientation="h",
                                yanchor="bottom",
                                y=-0.35,
                                x=0.5,
                                xanchor="center",
                            ),
                            height=360,
                            xaxis_title=None,
                            yaxis_title=None,
                        )
                        render_plotly_chart(
                            _fig_u,
                            use_container_width=True,
                            key="tab3_usage_monthly_chart",
                        )

# Tab 4: 👤 담당자 & 상세내역
with tab4:
    t4_c1, t4_c2 = st.columns([4, 1])
    t4_c1.markdown("<div class='sub-header dashboard-tab-panel-head'>👤 담당자별 월 매출 실적 (만원)</div>", unsafe_allow_html=True)
    t4_c2.markdown(render_update_badge(latest_update_str), unsafe_allow_html=True)

    if not staff_pivot.empty:
        format_dict = {col: "{:,.0f}" for col in staff_pivot.columns if col != "매출 비중 (%)"}
        format_dict["매출 비중 (%)"] = "{:,.1f}%"
    
        monthly_cols = [c for c in staff_pivot.columns if c not in ["매출 비중 (%)", "총 매출 합계 (만원)"]]
    
        styled_staff = (
            staff_pivot.style.format(format_dict)
            .background_gradient(cmap="Purples", subset=["매출 비중 (%)"])
            .background_gradient(cmap="Oranges", subset=["총 매출 합계 (만원)"])
            .background_gradient(cmap="Blues", subset=monthly_cols)
        )
        st.dataframe(styled_staff, use_container_width=True, height=350)
    st.markdown("<div class='sub-header dashboard-tab-panel-head'>🏆 담당자별 거래처 매출 순위 (당해년도)</div>", unsafe_allow_html=True)
    if not df_base.empty:
        current_year = str(df_base["연도"].max())
        all_staffs = sorted(df_base["담당자"].unique())
    
        sel_staff = st.selectbox("👤 순위를 조회할 담당자 선택", all_staffs, key="ranking_staff_select")
        ranking_pivot = cached_ranking_pivot(df_base, current_year, sel_staff, all_months)
    
        if not ranking_pivot.empty:
            r_col1, r_col2 = st.columns([1.2, 1])
        
            with r_col1:
                st.dataframe(
                    ranking_pivot.style.format("{:,.0f}")
                    .background_gradient(cmap="Blues", subset=all_months)
                    .background_gradient(cmap="Oranges", subset=["당해 누적 (만원)"]),
                    use_container_width=True, height=380
                )
            
            with r_col2:
                top_clients = ranking_pivot.head(10).sort_values(by="당해 누적 (만원)", ascending=True)
            
                fig_ranking = px.bar(
                    top_clients.reset_index(),
                    x="당해 누적 (만원)",
                    y="거래처",
                    orientation='h',
                    text="당해 누적 (만원)",
                    color="당해 누적 (만원)",
                    color_continuous_scale="Blues"
                )
                fig_ranking.update_traces(texttemplate='%{text:,.0f}', textposition='outside')
                fig_ranking.update_layout(
                    xaxis_title=None,
                    yaxis_title=None,
                    margin=dict(l=10, r=40, t=10, b=10),
                    plot_bgcolor='rgba(0,0,0,0)',
                    paper_bgcolor='rgba(0,0,0,0)',
                    coloraxis_showscale=False,
                    height=380
                )
                render_plotly_chart(fig_ranking, use_container_width=True, key=f"ranking_chart_{sel_staff}")
        else:
            st.info(f"💡 {current_year}년에 선택한 담당자({sel_staff})의 거래처 매출 실적 데이터가 없습니다.")
    if not df_base.empty:
        with st.expander("⚠️ 담당자 미지정 신규/누락 거래처 (직접 지정) 열기/닫기", expanded=True):
            unassigned_df = df_base[df_base["담당자"] == "미지정"]
            # 거래처명 없는 노이즈 행은 지정 불가 → 목록에서 제외(실거래처만 표시)
            if not unassigned_df.empty and "거래처" in unassigned_df.columns:
                unassigned_df = unassigned_df[unassigned_df["거래처"].map(_is_mappable_client_name)]
        
            custom_staffs = ["가스코아산", "거래종료"]
            existing_staffs = [s for s in full_df["담당자"].unique() if s not in ("미지정",)]
            combined_staffs = sorted(list(set(existing_staffs + custom_staffs)))
        
            all_staff_options = ["미지정"] + combined_staffs
        
            if not unassigned_df.empty:
                unassigned_summary = unassigned_df.groupby("거래처").agg(
                    최근매출일=("매출일_dt", "max"),
                    총매출액_만원=("매출액", lambda x: sum(x) * 1.1 / 10000)
                ).reset_index()
            
                unassigned_summary = unassigned_summary.sort_values(by="총매출액_만원", ascending=False)
                unassigned_summary["최근매출일"] = unassigned_summary["최근매출일"].dt.strftime("%Y-%m-%d")
                unassigned_summary["담당자지정"] = "미지정"
            
                st.warning(f"💡 자동 추론으로도 담당자를 찾을 수 없는 거래처가 총 **{len(unassigned_summary)}곳** 있습니다. 표의 **'담당자지정'** 열을 클릭하여 담당자를 선택하고 아래 저장 버튼을 누르세요.")
            
                edited_unassigned = st.data_editor(
                    unassigned_summary,
                    column_config={
                        "거래처": st.column_config.TextColumn("거래처", disabled=True),
                        "최근매출일": st.column_config.TextColumn("최근매출일", disabled=True),
                        "총매출액_만원": st.column_config.NumberColumn("총매출액(만원)", disabled=True, format="%d"),
                        "담당자지정": st.column_config.SelectboxColumn(
                            "👤 담당자 지정 (클릭하여 변경)",
                            help="이 거래처의 담당자를 선택하세요.",
                            options=all_staff_options,
                            required=True
                        )
                    },
                    use_container_width=True,
                    hide_index=True,
                    key="unassigned_editor"
                )
            
                if st.button("💾 변경된 담당자 저장 및 전체 대시보드 적용", type="primary"):
                    changed_rows = edited_unassigned[edited_unassigned["담당자지정"] != "미지정"]
                    if not changed_rows.empty:
                        manual_map_path = os.path.join(CACHE_DIR, "manual_staff_mapping.csv")
                        existing_map = {}
                        if os.path.exists(manual_map_path):
                            try:
                                _em = pd.read_csv(manual_map_path)
                                for _, mrow in _em.iterrows():
                                    ck = _normalize_manual_client_key(mrow["거래처"])
                                    if ck is None:
                                        continue
                                    staff = str(mrow["담당자"]).strip() if pd.notna(mrow["담당자"]) else ""
                                    if staff:
                                        existing_map[ck] = staff
                            except Exception:
                                existing_map = {}
                        
                        for _, row in changed_rows.iterrows():
                            ck = _normalize_manual_client_key(row["거래처"])
                            if ck is None:
                                continue
                            existing_map[ck] = row["담당자지정"]
                        
                        save_df = pd.DataFrame(list(existing_map.items()), columns=["거래처", "담당자"])
                        save_df.to_csv(manual_map_path, index=False, encoding="utf-8-sig")
                    
                        st.success("✅ 담당자 지정이 완료되었습니다! 대시보드를 새로고침합니다.")
                        load_uploaded_files_from_bytes.clear()
                        load_uploaded_files_from_meta.clear()
                        st.rerun()
                    else:
                        st.warning("저장할 담당자 지정이 없습니다.")
            else:
                st.success("🎉 모든 거래처에 담당자가 완벽하게 지정되어 있습니다!")
            
        with st.expander("🔄 이미 지정된 기존 거래처 담당자 수정/강제 변경하기 열기/닫기"):
            assigned_df = df_base[df_base["담당자"] != "미지정"]
            if not assigned_df.empty:
                assigned_summary = assigned_df.groupby("거래처").agg(
                    현재담당자=("담당자", "first"),
                    최근매출일=("매출일_dt", "max")
                ).reset_index()
            
                assigned_summary["새담당자변경"] = assigned_summary["현재담당자"]
                assigned_summary["최근매출일"] = assigned_summary["최근매출일"].dt.strftime("%Y-%m-%d")
            
                all_staffs_for_edit_assign = combined_staffs + ["미지정"]
            
                st.info("💡 잘못 지정된 거래처나 인수인계된 거래처의 담당자를 새롭게 변경할 수 있습니다.")
                edited_assigned = st.data_editor(
                    assigned_summary,
                    column_config={
                        "거래처": st.column_config.TextColumn("거래처", disabled=True),
                        "현재담당자": st.column_config.TextColumn("현재 담당자", disabled=True),
                        "최근매출일": st.column_config.TextColumn("최근매출일", disabled=True),
                        "새담당자변경": st.column_config.SelectboxColumn(
                            "👤 새 담당자로 변경 (클릭)",
                            options=all_staffs_for_edit_assign,
                            required=True
                        )
                    },
                    use_container_width=True,
                    hide_index=True,
                    key="assigned_editor"
                )
            
                if st.button("💾 변경된 기존 거래처 담당자 저장", type="primary", key="save_assigned_btn"):
                    changed_assigned = edited_assigned[edited_assigned["새담당자변경"] != edited_assigned["현재담당자"]]
                    if not changed_assigned.empty:
                        manual_map_path = os.path.join(CACHE_DIR, "manual_staff_mapping.csv")
                        existing_map = {}
                        if os.path.exists(manual_map_path):
                            try:
                                _em = pd.read_csv(manual_map_path)
                                for _, mrow in _em.iterrows():
                                    ck = _normalize_manual_client_key(mrow["거래처"])
                                    if ck is None:
                                        continue
                                    staff = str(mrow["담당자"]).strip() if pd.notna(mrow["담당자"]) else ""
                                    if staff:
                                        existing_map[ck] = staff
                            except Exception:
                                existing_map = {}
                        
                        for _, row in changed_assigned.iterrows():
                            ck = _normalize_manual_client_key(row["거래처"])
                            if ck is None:
                                continue
                            existing_map[ck] = row["새담당자변경"]
                        
                        save_df = pd.DataFrame(list(existing_map.items()), columns=["거래처", "담당자"])
                        save_df.to_csv(manual_map_path, index=False, encoding="utf-8-sig")
                    
                        st.success("✅ 담당자 변경이 완료되었습니다! 대시보드를 새로고침합니다.")
                        load_uploaded_files_from_bytes.clear()
                        load_uploaded_files_from_meta.clear()
                        st.rerun()
    st.markdown("<div class='sub-header dashboard-tab-panel-head'>📋 거래 상세 내역 (최신순 800건)</div>", unsafe_allow_html=True)
    if not df_detail.empty:
        view_detail_df = df_detail.sort_values(by="매출일_dt", ascending=False).head(800)
    
        styled_detail = (
            view_detail_df
            .style.format({
                "출고량": "{:,.0f}",
                "단가": "{:,.0f}",
                "매출액": "{:,.0f}",
                "매출일_dt": lambda t: t.strftime("%Y-%m-%d") if pd.notnull(t) else ""
            })
            .background_gradient(subset=["매출액"], cmap="Blues")
        )
        st.dataframe(styled_detail, use_container_width=True, height=600, hide_index=True)
# Tab 5: 📌 채권 관리
with tab5:
    latest_month = None
    if not filtered_debt_df.empty:
        # 데이터가 0(없는) 달 제거 로직 추가
        numeric_cols_temp = [c for c in filtered_debt_df.columns if c not in ["거래처", "구분"]]
        valid_numeric_cols = [c for c in numeric_cols_temp if filtered_debt_df[c].abs().sum() > 0]

        filtered_debt_df = filtered_debt_df[["거래처", "구분"] + valid_numeric_cols]
        numeric_cols_debt = valid_numeric_cols

        if numeric_cols_debt:
            latest_month = numeric_cols_debt[-1]

    debt_update_str = f"{latest_month} 기준" if latest_month else "데이터 없음"

    t5_c1, t5_c2 = st.columns([4, 1])
    t5_c1.markdown("<div class='sub-header dashboard-tab-panel-head'>💰 채권(외상대금) 관리 현황 및 연령 분석</div>", unsafe_allow_html=True)
    t5_c2.markdown(render_update_badge(debt_update_str), unsafe_allow_html=True)

    _debt_chip = (
        f"담당자 {', '.join(selected_staff)}" if selected_staff else "담당자 전체"
    )
    _debt_chip += (
        f" · 거래처 {selected_client}"
        if selected_client != "전체 거래처"
        else " · 거래처 전체"
    )
    st.markdown(
        f"<div class='dashboard-debt-filter-chip'>현재 필터: {_debt_chip}</div>",
        unsafe_allow_html=True,
    )

    if not debt_df.empty:
        if not filtered_debt_df.empty:
            numeric_cols = [c for c in filtered_debt_df.columns if c not in ["거래처", "구분"]]

            total_outstanding = 0
            warning_count = 0

            if latest_month:
                # 거래처별 집계 — 반복 loc 대신 groupby로 원복/지정 시 부하 완화
                _bal = filtered_debt_df[filtered_debt_df["구분"] == "잔액"]
                _sal = filtered_debt_df[filtered_debt_df["구분"] == "매출"]
                if latest_month in _bal.columns and not _bal.empty:
                    bal_by = _bal.groupby("거래처", sort=False)[latest_month].sum()
                    sal_by = (
                        _sal.groupby("거래처", sort=False)[latest_month].sum()
                        if not _sal.empty and latest_month in _sal.columns
                        else pd.Series(dtype=float)
                    )
                    for uc, b_val in bal_by.items():
                        b_val = float(b_val) if pd.notna(b_val) else 0.0
                        s_val = float(sal_by.get(uc, 0.0) or 0.0)
                        total_outstanding += max(0.0, b_val)
                        if b_val > 0 and b_val > s_val:
                            warning_count += 1

            m1, m2 = st.columns(2)
            m1.markdown(
                f"<div class='metric-box dashboard-debt-metric'><div class='metric-label'>총 미수금 잔액 ({latest_month} 기준)</div>"
                f"<div class='metric-value'>{total_outstanding:,.0f} 원</div></div>",
                unsafe_allow_html=True,
            )
            m2.markdown(
                f"<div class='metric-box dashboard-debt-metric'><div class='metric-label'>매출 초과 악성/지연 채권 업체 수</div>"
                f"<div class='metric-value' style='color:#E11D48;'>{warning_count} 곳</div></div>",
                unsafe_allow_html=True,
            )

            st.markdown("<div style='height:8px'></div>", unsafe_allow_html=True)
            summary_rows = []
            for gubun in ["이월", "익월", "매출", "수금", "잔액", "합계"]:
                if gubun in filtered_debt_df["구분"].values:
                    sum_vals = filtered_debt_df[filtered_debt_df["구분"] == gubun][numeric_cols].sum()
                    row_data = {"거래처": "📌 [전체 합계]", "구분": gubun}
                    for col in numeric_cols:
                        row_data[col] = sum_vals[col]
                    summary_rows.append(row_data)

            if summary_rows:
                summary_df = pd.DataFrame(summary_rows)
                disp_debt = pd.concat([filtered_debt_df, summary_df], ignore_index=True)
            else:
                disp_debt = filtered_debt_df.copy()

            gubun_order = {"이월": 1, "매출": 2, "수금": 3, "잔액": 4, "합계": 5}
            disp_debt["구분순위"] = disp_debt["구분"].map(gubun_order).fillna(99)
            client_order = {c: i for i, c in enumerate(disp_debt["거래처"].unique())}
            disp_debt["거래처순위"] = disp_debt["거래처"].map(client_order)

            disp_debt = disp_debt.sort_values(by=["거래처순위", "구분순위"]).drop(columns=["거래처순위", "구분순위"])

            disp_debt = disp_debt.set_index(["거래처", "구분"])
            debt_highlight = selected_client != "전체 거래처"
            df_height = 520 if selected_client != "전체 거래처" else 720
            show_cols = list(numeric_cols)
            # 입금기준표 결제조건 → 표 맨 우측 고정 표시
            if os.path.exists(PAYMENT_TERMS_FALLBACK) and not os.path.exists(PAYMENT_TERMS_PATH):
                try:
                    shutil.copy2(PAYMENT_TERMS_FALLBACK, PAYMENT_TERMS_PATH)
                except Exception:
                    pass
            payment_terms_map = load_payment_terms_map()
            st.markdown(
                "<div style='font-size:14px;font-weight:700;color:#1E293B;margin:4px 0 6px;'>"
                "📋 거래처별 채권 상세</div>",
                unsafe_allow_html=True,
            )
            render_debt_interactive_table(
                disp_debt[show_cols],
                debt_highlight,
                height=df_height,
                payment_terms_map=payment_terms_map,
            )
            # 상세표 아래: 연체개월수 요약 — 거래처(상위검색) 무시, 담당자 필터만 적용
            # staff 기준 메타는 cache → 거래처만 바꿔도 연체패널 재계산 부담 감소
            if not staff_debt_df.empty:
                _staff_num = [c for c in staff_debt_df.columns if c not in ("거래처", "구분")]
                _staff_months = [c for c in _staff_num if staff_debt_df[c].abs().sum() > 0]
                _rank_months = [c for c in show_cols if c in _staff_months] or _staff_months
                render_debt_month_rank_panel(
                    staff_debt_df,
                    _rank_months,
                    payment_terms_map=payment_terms_map,
                    height=480,
                    status_month_cols=_staff_months,
                )
            else:
                render_debt_month_rank_panel(
                    filtered_debt_df,
                    show_cols,
                    payment_terms_map=payment_terms_map,
                    height=480,
                    status_month_cols=numeric_cols,
                )
# Tab 6: 📍 대한민국 V-World 고해상도 한글/위성 지도 적용
with tab6:
    t6_c1, t6_c2 = st.columns([4, 1])
    t6_c1.markdown("<div class='sub-header dashboard-tab-panel-head'>📍 담당자별 거래처 지도 분포 (대한민국 V-World 지도)</div>", unsafe_allow_html=True)
    t6_c2.markdown(render_update_badge(latest_update_str), unsafe_allow_html=True)

    rest_api_key = "21a8c4d7312051598c2e05dba0b9c0c7"

    map_col1, map_col2, map_col3 = st.columns([1, 1, 1])
    with map_col1:
        map_style_choice = st.radio(
            "🗺️ 지도 배경 스타일 선택",
            ["일반 지도 (V-World 한글 기본도)", "위성 지도 (V-World 고해상도 위성)"],
            horizontal=True,
            key="map_style_radio"
        )
    with map_col2:
        all_staff_list = sorted(df_base["담당자"].unique()) if not df_base.empty else []
        map_selected_staff = st.multiselect(
            "👤 지도 전용 담당자 선택", 
            options=all_staff_list, 
            default=all_staff_list,
            key="map_staff_multiselect"
        )
    with map_col3:
        all_map_clients = sorted(df_base["거래처"].unique()) if not df_base.empty else []
        map_selected_client = st.multiselect(
            "🏢 특정 거래처 위치 검색", 
            options=all_map_clients,
            placeholder="검색할 거래처명을 입력하세요...",
            key="map_client_multiselect"
        )

    if map_selected_client:
        addr_display_html = "<div style='background-color: #F1F5F9; padding: 8px 12px; border-radius: 6px; border: 1px solid #CBD5E1; margin-top: 5px; margin-bottom: 15px; font-size: 13px; color: #334155;'>"
        for sc in map_selected_client:
            raw_a = resolve_client_address(sc, addr_dict)
            clean_a = raw_a if raw_a else "등록된 주소 정보가 없습니다."
            addr_display_html += f"<div>📍 <b>{sc}:</b> {clean_a}</div>"
        addr_display_html += "</div>"
        st.markdown(addr_display_html, unsafe_allow_html=True)
    st.markdown("<br>", unsafe_allow_html=True)
    ctrl_space, ctrl_c1, ctrl_c2, ctrl_c3, ctrl_c4 = st.columns([5, 1.2, 1.2, 1.2, 1.2])

    with ctrl_space:
        st.empty()
    with ctrl_c1:
        btn_load_map = st.button("🗺️ 지도 새로고침/조회", type="primary", use_container_width=True)
    with ctrl_c2:
        btn_zoom_in = st.button("➕ 확대 (+)", use_container_width=True)
    with ctrl_c3:
        btn_zoom_out = st.button("➖ 축소 (-)", use_container_width=True)
    with ctrl_c4:
        btn_reset_map = st.button("🏠 기본 위치", use_container_width=True)

    # 재시작·다른 탭 조작 시 전체 지오코딩이 돌지 않도록: 조회 버튼 후에만 로드
    if "show_map" not in st.session_state:
        st.session_state.show_map = False
    if btn_load_map:
        st.session_state.show_map = True
        st.session_state.map_force_rebuild = True
    if btn_zoom_in or btn_zoom_out or btn_reset_map:
        st.session_state.show_map = True

    map_filter_fp = (
        tuple(sorted(map_selected_staff or [])),
        tuple(sorted(map_selected_client or [])),
    )

    if not st.session_state.show_map:
        st.info("담당자·거래처를 선택한 뒤 **지도 새로고침/조회**를 누르면 지도를 불러옵니다. (재시작 시 자동 조회하지 않아 앱이 빨라집니다)")
    else:
        need_rebuild = (
            st.session_state.pop("map_force_rebuild", False)
            or st.session_state.get("tab6_map_fp") != map_filter_fp
            or "tab6_map_df" not in st.session_state
        )

        if need_rebuild:
            target_map_df = df_base.copy()
            if map_selected_client:
                target_map_df = target_map_df[target_map_df["거래처"].isin(map_selected_client)]
            elif map_selected_staff:
                target_map_df = target_map_df[target_map_df["담당자"].isin(map_selected_staff)]

            map_data = []
            invalid_clients = []
            if not target_map_df.empty:
                unique_clients_df = target_map_df[["거래처", "담당자"]].drop_duplicates(subset=["거래처"])
                total_cnt = len(unique_clients_df)
                disk_cache = _load_kakao_geocode_disk()
                dirty = [False]
                progress_text = "주소 좌표 변환 중 (디스크 캐시 우선) 🚀"
                my_bar = st.progress(0, text=progress_text)

                for i, (_, row) in enumerate(unique_clients_df.iterrows()):
                    c_name = row["거래처"]
                    c_staff = row["담당자"]
                    c_addr_raw = resolve_client_address(c_name, addr_dict)
                    c_addr = c_addr_raw if c_addr_raw else "등록된 주소 정보가 없습니다."
                    lat, lon = get_lat_lon_kakao_disk(c_name, c_addr, rest_api_key, disk_cache, dirty)
                    if lat is not None and lon is not None:
                        map_data.append(
                            {
                                "거래처": c_name,
                                "담당자": c_staff,
                                "주소": c_addr,
                                "lat": lat,
                                "lon": lon,
                            }
                        )
                    else:
                        invalid_clients.append(c_name)
                    my_bar.progress((i + 1) / total_cnt, text=f"{progress_text} ({i + 1}/{total_cnt})")
                my_bar.empty()
                if dirty[0]:
                    _save_kakao_geocode_disk(disk_cache)

            st.session_state.tab6_map_df = pd.DataFrame(map_data) if map_data else pd.DataFrame()
            st.session_state.tab6_invalid_clients = invalid_clients
            st.session_state.tab6_map_fp = map_filter_fp

        map_df = st.session_state.get("tab6_map_df", pd.DataFrame())
        invalid_clients = st.session_state.get("tab6_invalid_clients", [])

        if map_df is not None and not map_df.empty:
            center_lat = float(map_df["lat"].mean())
            center_lon = float(map_df["lon"].mean())

            default_zoom = 13 if map_selected_client and len(map_selected_client) <= 3 else 8

            if "map_zoom" not in st.session_state or btn_reset_map or btn_load_map:
                st.session_state.map_zoom = default_zoom

            if btn_zoom_in:
                st.session_state.map_zoom = min(st.session_state.map_zoom + 2, 20)
            elif btn_zoom_out:
                st.session_state.map_zoom = max(st.session_state.map_zoom - 2, 2)

            vworld_base = "https://xdworld.vworld.kr/2d/Base/service/{z}/{x}/{y}.png"
            vworld_sat = "https://xdworld.vworld.kr/2d/Satellite/service/{z}/{x}/{y}.jpeg"
            vworld_hybrid = "https://xdworld.vworld.kr/2d/Hybrid/service/{z}/{x}/{y}.png"
            dynamic_key = f"map_chart_{hash(str(map_selected_staff))}_{hash(str(map_selected_client))}"

            # iPad 전용: Plotly Mapbox WebGL이 Safari에서 마커/범례를 검정으로 그림
            # → Leaflet 원형 마커(명시 HEX)로만 우회. 맥 경로(아래 else)는 일절 변경 없음.
            if is_touch_ui():
                _palette = [
                    "#636EFA", "#EF553B", "#00CC96", "#AB63FA", "#FFA15A",
                    "#19D3F3", "#FF6692", "#B6E880", "#FF97FF", "#FECB52",
                    "#1F77B4", "#D62728", "#2CA02C", "#9467BD", "#8C564B",
                ]
                _staffs = sorted(map_df["담당자"].astype(str).unique())
                _cmap = {s: _palette[i % len(_palette)] for i, s in enumerate(_staffs)}
                _pts = []
                for _, _r in map_df.iterrows():
                    _staff = str(_r["담당자"])
                    _pts.append({
                        "lat": float(_r["lat"]),
                        "lon": float(_r["lon"]),
                        "name": str(_r["거래처"]),
                        "staff": _staff,
                        "addr": str(_r.get("주소") or ""),
                        "color": _cmap.get(_staff, "#636EFA"),
                    })
                _legend_html = "".join(
                    f'<span style="display:inline-flex;align-items:center;margin:0 10px 6px 0;'
                    f'font-size:13px;color:#334155;">'
                    f'<span style="width:12px;height:12px;border-radius:50%;background:{_cmap[s]};'
                    f'display:inline-block;margin-right:5px;border:1px solid #94A3B8;"></span>'
                    f"{html.escape(s)}</span>"
                    for s in _staffs
                )
                _use_sat = "일반" not in map_style_choice
                _tiles_js = (
                    f'L.tileLayer("{vworld_sat}", {{maxZoom:19, attribution:"VWorld"}}).addTo(map);'
                    f'L.tileLayer("{vworld_hybrid}", {{maxZoom:19, attribution:"VWorld"}}).addTo(map);'
                    if _use_sat
                    else f'L.tileLayer("{vworld_base}", {{maxZoom:19, attribution:"VWorld"}}).addTo(map);'
                )
                _pts_json = json.dumps(_pts, ensure_ascii=False)
                _leaflet_html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"/>
<meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0"/>
<link rel="stylesheet" href="https://unpkg.com/leaflet@1.9.4/dist/leaflet.css"/>
<script src="https://unpkg.com/leaflet@1.9.4/dist/leaflet.js"></script>
<style>
  html, body {{ margin:0; height:100%; }}
  #map {{ width:100%; height:560px; }}
  .legend {{
    padding:8px 10px; font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",sans-serif;
    background:#fff; border-top:1px solid #E2E8F0;
  }}
</style></head>
<body>
<div id="map"></div>
<div class="legend"><b>담당자</b><div style="margin-top:6px;">{_legend_html}</div></div>
<script>
(function() {{
  var map = L.map("map", {{ zoomControl: true }}).setView(
    [{float(center_lat)}, {float(center_lon)}], {int(st.session_state.map_zoom)}
  );
  {_tiles_js}
  var pts = {_pts_json};
  pts.forEach(function(p) {{
    var m = L.circleMarker([p.lat, p.lon], {{
      radius: 8,
      color: "#ffffff",
      weight: 1.5,
      fillColor: p.color,
      fillOpacity: 0.95
    }});
    m.bindPopup("<b>" + p.name + "</b><br/>담당자: " + p.staff + "<br/>" + (p.addr || ""));
    m.addTo(map);
  }});
}})();
</script>
</body></html>"""
                components.html(_leaflet_html, height=620, scrolling=False)
            else:
                fig_map = px.scatter_mapbox(
                    map_df,
                    lat="lat",
                    lon="lon",
                    color="담당자",
                    hover_name="거래처",
                    hover_data={"주소": True, "lat": False, "lon": False, "담당자": False},
                    zoom=st.session_state.map_zoom,
                    center={"lat": center_lat, "lon": center_lon},
                    height=600
                )
                fig_map.update_traces(marker=dict(size=14, opacity=0.9))
                if "일반" in map_style_choice:
                    mapbox_layers = [
                        {"below": 'traces', "sourcetype": "raster", "source": [vworld_base]}
                    ]
                else:
                    mapbox_layers = [
                        {"below": 'traces', "sourcetype": "raster", "source": [vworld_sat]},
                        {"below": 'traces', "sourcetype": "raster", "source": [vworld_hybrid]}
                    ]
                fig_map.update_layout(
                    mapbox_style="white-bg",
                    mapbox_layers=mapbox_layers,
                    margin={"r": 0, "t": 10, "l": 0, "b": 0},
                    legend=dict(
                        orientation="h",
                        yanchor="bottom",
                        y=-0.15,
                        xanchor="center",
                        x=0.5
                    )
                )
                render_plotly_chart(fig_map, use_container_width=True, key=dynamic_key, allow_drag=True)
            if invalid_clients:
                with st.expander("⚠️ 지도에 표시되지 않은 거래처 (주소 정보 없음 또는 좌표 변환 실패)"):
                    st.write(", ".join(invalid_clients))
        else:
            st.info("조건에 맞는 거래처 데이터가 없습니다.")
# Tab 7: 🏭 설비 재고 현황
with tab7:
    t7_c1, t7_c2 = st.columns([4, 1])
    t7_c1.markdown("<div class='sub-header dashboard-tab-panel-head'>🏭 고압가스 탱크 및 기화기 재고 현황</div>", unsafe_allow_html=True)
    with t7_c2:
        st.markdown("<div style='margin-top: 15px;'></div>", unsafe_allow_html=True)
        tab7_default = get_saved_date(TAB7_DATE_FILE)
        tab7_date = st.date_input("기준일", value=tab7_default, key="tab7_date", label_visibility="collapsed")
        if tab7_date != tab7_default:
            set_saved_date(TAB7_DATE_FILE, tab7_date)
    if not df_tank.empty or not df_vaporizer.empty:
        eq_col1, eq_col2, eq_col3, eq_col4 = st.columns(4)
    
        with eq_col1:
            available_branches = []
            if not df_tank.empty and '지사' in df_tank.columns:
                available_branches.extend(df_tank['지사'].dropna().unique())
            if not df_vaporizer.empty and '지사' in df_vaporizer.columns:
                available_branches.extend(df_vaporizer['지사'].dropna().unique())
            selected_branch = st.selectbox("📍 지사 선택 (전체 조회)", ["전체 지사"] + sorted(list(set(available_branches))))
    
        with eq_col2:
            selected_equip_type = st.selectbox("🛢️ 설비 종류 선택", ["전체 보기", "탱크 재고", "기화기 재고"])
        with eq_col3:
            selected_status = st.selectbox("📌 사용구분 필터", ["전체 상태", "유휴 장비", "거래처 사용중"])
        with eq_col4:
            eq_items = []
            if not df_tank.empty and '품목' in df_tank.columns:
                eq_items.extend(df_tank['품목'].dropna().astype(str).tolist())
            if not df_vaporizer.empty and '기화형식' in df_vaporizer.columns:
                eq_items.extend(df_vaporizer['기화형식'].dropna().astype(str).tolist())
            unique_eq_items = sorted(list(set([i.strip() for i in eq_items if i.strip() != ''])))
            selected_eq_item = st.selectbox("📦 품목/형식 선택", ["전체 품목/형식"] + unique_eq_items)
        st.markdown("---")
        if not df_tank.empty:
            if '사용구분' in df_tank.columns:
                mask_idle = df_tank['사용구분'].astype(str).str.contains('유휴')
                df_tank.loc[mask_idle, '사용구분'] = '🟢 ' + df_tank.loc[mask_idle, '사용구분'].astype(str).str.replace('🟢 ', '').str.replace('🏢 ', '')
                df_tank.loc[~mask_idle, '사용구분'] = '🏢 ' + df_tank.loc[~mask_idle, '사용구분'].astype(str).str.replace('🟢 ', '').str.replace('🏢 ', '')
        if not df_vaporizer.empty:
            if '사용구분' in df_vaporizer.columns:
                mask_idle_v = df_vaporizer['사용구분'].astype(str).str.contains('유휴')
                df_vaporizer.loc[mask_idle_v, '사용구분'] = '🟢 ' + df_vaporizer.loc[mask_idle_v, '사용구분'].astype(str).str.replace('🟢 ', '').str.replace('🏢 ', '')
                df_vaporizer.loc[~mask_idle_v, '사용구분'] = '🏢 ' + df_vaporizer.loc[~mask_idle_v, '사용구분'].astype(str).str.replace('🟢 ', '').str.replace('🏢 ', '')
        if selected_equip_type in ["전체 보기", "탱크 재고"]:
            st.markdown("<div style='font-size: 16px; font-weight: 700; color: #1E3A8A; margin-bottom: 10px;'>🛢️ 초저온 탱크 재고 현황</div>", unsafe_allow_html=True)
            if not df_tank.empty:
                filtered_tank = df_tank.copy()
                if selected_branch != "전체 지사" and '지사' in filtered_tank.columns:
                    filtered_tank = filtered_tank[filtered_tank['지사'].astype(str).str.contains(selected_branch)]
                if selected_status != "전체 상태" and '사용구분' in filtered_tank.columns:
                    filtered_tank = filtered_tank[filtered_tank['사용구분'].astype(str).str.contains(selected_status)]
                if selected_eq_item != "전체 품목/형식" and '품목' in filtered_tank.columns:
                    filtered_tank = filtered_tank[filtered_tank['품목'].astype(str).str.strip() == selected_eq_item]
            
                st.dataframe(filtered_tank, use_container_width=True, height=350, hide_index=True)
        if selected_equip_type in ["전체 보기", "기화기 재고"]:
            if selected_equip_type == "전체 보기":
                st.markdown("<br>", unsafe_allow_html=True)
        
            st.markdown(f"<div style='font-size: 16px; font-weight: 700; color: #1E3A8A; margin-bottom: 10px;'>♨️ 기화기 재고 현황</div>", unsafe_allow_html=True)
            if not df_vaporizer.empty:
                filtered_vap = df_vaporizer.copy()
                if selected_branch != "전체 지사" and '지사' in filtered_vap.columns:
                    filtered_vap = filtered_vap[filtered_vap['지사'].astype(str).str.contains(selected_branch)]
                if selected_status != "전체 상태" and '사용구분' in filtered_vap.columns:
                    filtered_vap = filtered_vap[filtered_vap['사용구분'].astype(str).str.contains(selected_status)]
                if selected_eq_item != "전체 품목/형식" and '기화형식' in filtered_vap.columns:
                    filtered_vap = filtered_vap[filtered_vap['기화형식'].astype(str).str.strip() == selected_eq_item]
            
                st.dataframe(filtered_vap, use_container_width=True, height=350, hide_index=True)
# Tab 8: 🛢️ 통합 탱크 재고
with tab8:
    t8_c1, t8_c2 = st.columns([4, 1])
    t8_c1.markdown("<div class='sub-header dashboard-tab-panel-head'>🛢️ 통합 고압가스 탱크 재고 현황</div>", unsafe_allow_html=True)
    with t8_c2:
        st.markdown("<div style='margin-top: 15px;'></div>", unsafe_allow_html=True)
        tab8_default = get_saved_date(TAB8_DATE_FILE)
        tab8_date = st.date_input("기준일", value=tab8_default, key="int_date", label_visibility="collapsed")
        if tab8_date != tab8_default:
            set_saved_date(TAB8_DATE_FILE, tab8_date)
    
    if not df_integrated.empty:
        int_col1, int_col2 = st.columns(2)
        with int_col1:
            items = ["전체 품목"] + sorted([str(x) for x in df_integrated['품목'].dropna().unique() if str(x).strip()])
            sel_item = st.selectbox("📦 품목 선택", items, key="int_item")
        with int_col2:
            statuses = ["전체 상태"] + sorted([str(x) for x in df_integrated['사용구분'].dropna().unique() if str(x).strip()])
            sel_status = st.selectbox("📌 사용구분", statuses, key="int_status")
        st.markdown("---")
        df_int_filtered = df_integrated.copy()
        if sel_item != "전체 품목":
            df_int_filtered = df_int_filtered[df_int_filtered['품목'].astype(str) == sel_item]
        if sel_status != "전체 상태":
            df_int_filtered = df_int_filtered[df_int_filtered['사용구분'].astype(str) == sel_status]
        total_tanks = len(df_int_filtered)
        idle_tanks = len(df_int_filtered[df_int_filtered['사용구분'].astype(str).str.contains('유휴', na=False)])
        inuse_tanks = total_tanks - idle_tanks
    
        k1, k2, k3 = st.columns(3)
        k1.markdown(f"<div class='metric-box'><div class='metric-label'>총 탱크 수량</div><div class='metric-value'>{total_tanks:,} 기</div></div>", unsafe_allow_html=True)
        k2.markdown(f"<div class='metric-box'><div class='metric-label'>🟢 유휴 장비 (대기중)</div><div class='metric-value' style='color:#059669;'>{idle_tanks:,} 기</div></div>", unsafe_allow_html=True)
        k3.markdown(f"<div class='metric-box'><div class='metric-label'>🏢 사용/충전중</div><div class='metric-value' style='color:#2563EB;'>{inuse_tanks:,} 기</div></div>", unsafe_allow_html=True)
        st.markdown("<div style='font-size: 14px; font-weight: 600; color: #334155; margin-bottom: 10px; margin-top: 20px;'>📋 상세 재고 데이터</div>", unsafe_allow_html=True)
    
        df_display = df_int_filtered.copy()
        if '사용구분' in df_display.columns:
            mask_idle = df_display['사용구분'].astype(str).str.contains('유휴')
            df_display.loc[mask_idle, '사용구분'] = '🟢 ' + df_display.loc[mask_idle, '사용구분'].astype(str).str.replace('🟢 ', '').str.replace('🏢 ', '')
            df_display.loc[~mask_idle, '사용구분'] = '🏢 ' + df_display.loc[~mask_idle, '사용구분'].astype(str).str.replace('🟢 ', '').str.replace('🏢 ', '')
        st.dataframe(df_display, use_container_width=True, height=600, hide_index=True)
    else:
        st.warning("통합 탱크 재고 데이터가 없습니다. 폴더에 '통합탱크재고.csv'를 넣거나 왼쪽 사이드바에서 업로드해주세요.")
# Tab 9: 📈 수익성 분석 (엑셀 함수 동일 적용)
# 입력 변경 시 Tab9만 부분 재실행 (다른 탭·상단 로딩 생략)
@st.fragment
def _render_profitability_analysis_tab(latest_update_str):
    t9_c1, t9_c2 = st.columns([4, 1])
    t9_c1.markdown(
        "<div class='sub-header dashboard-tab-panel-head'>📈 투자대비 수익성 분석</div>",
        unsafe_allow_html=True,
    )
    t9_c2.markdown(render_update_badge(latest_update_str), unsafe_allow_html=True)
    st.caption("엑셀「수익성분석.xlsx」함수를 그대로 적용합니다. 입력값을 바꾸면 결과가 즉시 재계산됩니다.")
    if "profit_inputs" not in st.session_state:
        st.session_state["profit_inputs"] = load_profit_inputs()
    p0 = st.session_state["profit_inputs"]
    _pf_keys = [
        "pf_name", "pf_tank_gas", "pf_tank_cap_mode",
        "pf_tank_liters", "pf_tank_liters__comma",
        "pf_tank_kg", "pf_tank_kg__comma", "pf_tank_spec",
        "pf_hourly_mode", "pf_hourly_usage", "pf_hourly_nm3", "pf_operating_hours",
        "pf_operating_days", "pf_auto_monthly",
        "pf_tank_price", "pf_tank_price__comma", "pf_usage", "pf_usage__comma",
        "pf_const", "pf_const__comma", "pf_vap_cap", "pf_vap_cap__comma",
        "pf_vap_note", "pf_vap_price", "pf_vap_price__comma",
        "pf_buy", "pf_logi", "pf_supply",
        "pf_rate", "pf_mgmt", "pf_dep", "pf_rent", "pf_rent_n",
        "pf_origin", "pf_dest", "pf_origin_cands", "pf_dest_cands",
        "pf_origin_q", "pf_dest_q", "pf_origin_pick", "pf_dest_pick",
        "pf_lkm", "pf_lfuel", "pf_leff", "pf_ltoll", "pf_lrt", "pf_lkg",
    ]
    # —— 입력 (물류비 계산 → 단가 순, 글씨·위젯 스타일은 단가란과 동일) ——
    st.markdown("##### ◆ 프로젝트 / 장비 투자비")
    c_name, c_gas = st.columns([1, 1])
    project_name = c_name.text_input("거래처/프로젝트명", value=str(p0.get("project_name", "")), key="pf_name")
    _gas0 = str(p0.get("tank_gas") or PROFIT_DEFAULTS["tank_gas"])
    if _gas0 not in GAS_OPTIONS:
        _gas0 = GAS_OPTIONS[0]
    tank_gas = c_gas.selectbox(
        "탱크 가스 종류",
        GAS_OPTIONS,
        index=GAS_OPTIONS.index(_gas0),
        key="pf_tank_gas",
        help="질소·알곤·산소·탄산·수소·헬륨. 내용적(L) 입력 시 kg 환산에 사용됩니다.",
    )
    _dens = float(GAS_DENSITY_KG_PER_L.get(tank_gas, 0.808))
    _mode0 = str(p0.get("tank_capacity_mode") or "liters")
    if _mode0 not in ("liters", "kg"):
        _mode0 = "liters"
    tank_capacity_mode = st.radio(
        "탱크 용량 입력 방식",
        options=["liters", "kg"],
        index=0 if _mode0 == "liters" else 1,
        format_func=lambda m: "내용적(L) → kg 환산" if m == "liters" else "용량(kg) 직접 입력",
        horizontal=True,
        key="pf_tank_cap_mode",
        help="L로 넣거나, kg를 바로 넣을 수 있습니다.",
    )
    _kg0 = parse_tank_capacity_kg(p0.get("tank_spec", PROFIT_DEFAULTS["tank_spec"]))
    _liters0 = p0.get("tank_liters")
    if _liters0 is None or float(_liters0 or 0) <= 0:
        _liters0 = (_kg0 / _dens) if _dens > 0 and _kg0 > 0 else float(PROFIT_DEFAULTS["tank_liters"])
    t_l, t_k = st.columns([1, 1])
    if tank_capacity_mode == "liters":
        with t_l:
            tank_liters = profit_int_comma_input(
                "TANK 내용적 (L)",
                key="pf_tank_liters",
                value=_liters0,
                help=f"{tank_gas} 밀도 {_dens:g} kg/L × 내용적(L) = 용량(kg)",
            )
        tank_kg = round(liters_to_tank_kg(tank_liters, tank_gas))
        st.session_state["pf_tank_kg"] = int(tank_kg)
        st.session_state["pf_tank_kg__comma"] = f"{int(tank_kg):,}"
        t_k.markdown(
            f"<div style='padding-top:0.2rem;'>"
            f"<div style='font-size:0.875rem;color:#31333F;margin-bottom:0.25rem;'>TANK 용량 (kg) 환산</div>"
            f"<div style='font-size:1rem;font-weight:600;color:#0F172A;'>{tank_kg:,.0f}</div>"
            f"<div style='font-size:0.75rem;color:#64748B;'>"
            f"{tank_gas} {_dens:g} kg/L × {float(tank_liters):,.0f} L</div>"
            f"</div>",
            unsafe_allow_html=True,
        )
    else:
        with t_k:
            tank_kg = profit_int_comma_input(
                "TANK 용량 (kg) 직접입력",
                key="pf_tank_kg",
                value=_kg0 if _kg0 > 0 else 4900,
                help="용량을 kg로 바로 입력합니다. 왕복횟수 계산에 사용됩니다.",
            )
        tank_liters = round(float(tank_kg) / _dens) if _dens > 0 else 0.0
        st.session_state["pf_tank_liters"] = int(tank_liters)
        st.session_state["pf_tank_liters__comma"] = f"{int(tank_liters):,}"
        t_l.markdown(
            f"<div style='padding-top:0.2rem;'>"
            f"<div style='font-size:0.875rem;color:#31333F;margin-bottom:0.25rem;'>TANK 내용적 (L) 환산</div>"
            f"<div style='font-size:1rem;font-weight:600;color:#0F172A;'>{tank_liters:,.0f}</div>"
            f"<div style='font-size:0.75rem;color:#64748B;'>"
            f"{float(tank_kg):,.0f} kg ÷ {tank_gas} {_dens:g} kg/L</div>"
            f"</div>",
            unsafe_allow_html=True,
        )
    tank_spec = float(tank_kg)  # 저장·계산용 (단위: kg)
    _nm3 = tank_kg_to_nm3(tank_kg, tank_gas)
    _nm3_per_kg = float(GAS_NM3_PER_KG.get(tank_gas, 0))
    st.caption(
        f"기체환산({tank_gas}): {float(tank_kg):,.0f} kg × {_nm3_per_kg:g} Nm³/kg = "
        f"**{_nm3:,.1f} Nm³**  ·  내용적 {float(tank_liters):,.0f} L 기준 "
        f"{tank_liters_to_nm3(tank_liters, tank_gas):,.1f} Nm³"
    )
    with st.expander("각 가스별 밀도·기체환산 비교", expanded=False):
        st.caption("동일 내용적(L)으로 가스별 kg·Nm³를 비교합니다. ★ = 현재 선택 가스. (0℃·1atm 대략값)")
        _gas_df = pd.DataFrame(gas_conversion_rows(tank_liters, tank_kg, tank_gas))
        _gas_df["밀도(kg/L)"] = _gas_df["밀도(kg/L)"].map(lambda x: f"{float(x):.4g}")
        _gas_df["Nm³/kg"] = _gas_df["Nm³/kg"].map(lambda x: f"{float(x):.4g}")
        _gas_df["내용적 기준 kg"] = _gas_df["내용적 기준 kg"].map(lambda x: f"{float(x):,.1f}")
        _gas_df["내용적 기준 Nm³"] = _gas_df["내용적 기준 Nm³"].map(lambda x: f"{float(x):,.1f}")
        _gas_df["현재탱크 Nm³"] = _gas_df["현재탱크 Nm³"].map(
            lambda x: f"{float(x):,.1f}" if isinstance(x, (int, float)) else str(x)
        )
        st.dataframe(_gas_df, hide_index=True, width="stretch")
    # 별도 작은 타일: 탱크 사용주기 (화면 복잡도 ↓ — 접힌 expander)
    _nm3pkg = float(GAS_NM3_PER_KG.get(tank_gas, 0) or 0)
    with st.expander("⏱ 탱크 사용주기 (시간당 사용량 · 가동시간)", expanded=False):
        st.caption(
            f"충전기준 = 탱크×80% · 사용주기 = 충전기준 ÷ (시간당kg × 일가동) · "
            f"{tank_gas} 환산 {_nm3pkg:g} Nm³/kg"
        )
        # 기존 float 세션값 → 정수 (format=%d 호환)
        for _pk in ("pf_hourly_nm3", "pf_hourly_usage", "pf_operating_hours", "pf_operating_days"):
            if _pk in st.session_state:
                try:
                    st.session_state[_pk] = int(round(float(st.session_state[_pk])))
                except Exception:
                    pass
        _hmode0 = str(p0.get("hourly_usage_mode") or "kg")
        if _hmode0 not in ("nm3", "kg"):
            _hmode0 = "kg"
        hourly_usage_mode = st.radio(
            "시간당 사용량 입력",
            options=["nm3", "kg"],
            index=0 if _hmode0 == "nm3" else 1,
            format_func=lambda m: "루베(Nm³/h) → kg/h 환산" if m == "nm3" else "kg/h 직접 입력",
            horizontal=True,
            key="pf_hourly_mode",
        )
        _kg_h0 = float(p0.get("hourly_usage_kg", PROFIT_DEFAULTS["hourly_usage_kg"]) or 0)
        _nm3_h0 = p0.get("hourly_usage_nm3")
        if _nm3_h0 is None or float(_nm3_h0 or 0) <= 0:
            _nm3_h0 = kg_per_h_to_nm3_per_h(_kg_h0, tank_gas)
        u_a, u_b, u_o = st.columns(3)
        if hourly_usage_mode == "nm3":
            with u_a:
                hourly_usage_nm3 = st.number_input(
                    "시간당 사용량 (Nm³/h, 루베)",
                    min_value=0,
                    step=1,
                    value=int(round(float(_nm3_h0))),
                    format="%d",
                    key="pf_hourly_nm3",
                    help=f"{tank_gas}: kg/h = Nm³/h ÷ {_nm3pkg:g}",
                )
            hourly_usage_kg = nm3_per_h_to_kg_per_h(hourly_usage_nm3, tank_gas)
            st.session_state["pf_hourly_usage"] = int(round(hourly_usage_kg))
            u_b.markdown(
                f"<div style='padding-top:0.2rem;'>"
                f"<div style='font-size:0.875rem;color:#31333F;margin-bottom:0.25rem;'>시간당 사용량 (kg/h) 환산</div>"
                f"<div style='font-size:1rem;font-weight:600;color:#0F172A;'>{hourly_usage_kg:,.0f}</div>"
                f"<div style='font-size:0.75rem;color:#64748B;'>{tank_gas} {float(hourly_usage_nm3):,.0f} Nm³/h ÷ {_nm3pkg:g}</div>"
                f"</div>",
                unsafe_allow_html=True,
            )
        else:
            with u_b:
                hourly_usage_kg = st.number_input(
                    "시간당 사용량 (kg/h)",
                    min_value=0,
                    step=1,
                    value=int(round(float(_kg_h0))),
                    format="%d",
                    key="pf_hourly_usage",
                )
            hourly_usage_nm3 = kg_per_h_to_nm3_per_h(hourly_usage_kg, tank_gas)
            st.session_state["pf_hourly_nm3"] = int(round(hourly_usage_nm3))
            u_a.markdown(
                f"<div style='padding-top:0.2rem;'>"
                f"<div style='font-size:0.875rem;color:#31333F;margin-bottom:0.25rem;'>시간당 사용량 (Nm³/h) 환산</div>"
                f"<div style='font-size:1rem;font-weight:600;color:#0F172A;'>{hourly_usage_nm3:,.0f}</div>"
                f"<div style='font-size:0.75rem;color:#64748B;'>{tank_gas} {float(hourly_usage_kg):,.0f} kg/h × {_nm3pkg:g}</div>"
                f"</div>",
                unsafe_allow_html=True,
            )
        operating_hours = u_o.number_input(
            "일 가동시간 (h/일)",
            min_value=0,
            max_value=24,
            step=1,
            value=int(round(float(p0.get("operating_hours", PROFIT_DEFAULTS["operating_hours"])))),
            format="%d",
            key="pf_operating_hours",
        )
        d_days, d_auto = st.columns([1, 2])
        operating_days = d_days.number_input(
            "월 가동일수 (일/월)",
            min_value=1,
            max_value=31,
            step=1,
            value=int(round(float(p0.get("operating_days_per_month", PROFIT_DEFAULTS["operating_days_per_month"])))),
            format="%d",
            key="pf_operating_days",
            help="월 사용량 = 일 사용량(시간당×일가동) × 월 가동일수. 예: 주5일≈22일, 연중무휴≈30일",
        )
        _cycle = compute_tank_usage_cycle(
            tank_kg,
            hourly_usage_kg,
            operating_hours,
            fill_ratio=0.8,
            days_per_month=operating_days,
        )
        _monthly_est = float(_cycle.get("monthly_kg") or 0)
        auto_monthly_from_cycle = d_auto.checkbox(
            f"월 평균 공급량에 자동 반영 (일사용량 × {int(operating_days)}일)",
            value=bool(p0.get("auto_monthly_from_cycle", False)),
            key="pf_auto_monthly",
            help="체크 시 아래 「월 평균 공급량」= 일사용량 × 월가동일수. 해제하면 직접 입력합니다.",
        )
        if _cycle.get("ok"):
            c1, c2, c3, c4 = st.columns(4)
            c1.markdown(
                f"<div class='metric-box'><div class='metric-label'>일 사용량</div>"
                f"<div class='metric-value' style='font-size:18px;'>{_cycle['daily_kg']:,.0f} kg/일</div></div>",
                unsafe_allow_html=True,
            )
            c2.markdown(
                f"<div class='metric-box'><div class='metric-label'>사용주기</div>"
                f"<div class='metric-value' style='font-size:18px;color:#2563EB;'>{_cycle['cycle_days']:,.0f} 일</div></div>",
                unsafe_allow_html=True,
            )
            c3.markdown(
                f"<div class='metric-box'><div class='metric-label'>가동시간 기준</div>"
                f"<div class='metric-value' style='font-size:18px;'>{_cycle['cycle_hours']:,.0f} h</div></div>",
                unsafe_allow_html=True,
            )
            c4.markdown(
                f"<div class='metric-box'><div class='metric-label'>월 사용량({int(operating_days)}일)</div>"
                f"<div class='metric-value' style='font-size:18px;'>{_monthly_est:,.0f} kg</div></div>",
                unsafe_allow_html=True,
            )
            st.caption(
                f"탱크 {float(tank_kg):,.0f} kg · 충전기준 {_cycle['charge_kg']:,.0f} kg(80%) · "
                f"월가동 {int(operating_days)}일 · 월충전 약 {_cycle['fills_per_month']:,.0f}회 · {_cycle['message']}"
            )
            if auto_monthly_from_cycle:
                _mu = int(round(_monthly_est))
                st.session_state["pf_usage"] = _mu
                st.session_state["pf_usage__comma"] = f"{_mu:,}"
                st.caption(
                    f"→ 월 평균 공급량 {_mu:,} kg "
                    f"(일 {_cycle['daily_kg']:,.0f} × {int(operating_days)}일) 자동 반영 중"
                )
        else:
            st.caption(_cycle.get("message") or "입력값을 확인하세요.")
    i1, i2, i3 = st.columns(3)
    with i1:
        tank_price = profit_int_comma_input(
            "1. TANK 구입가 (원)", key="pf_tank_price", value=p0["tank_price"]
        )
    with i2:
        if auto_monthly_from_cycle and _cycle.get("ok"):
            monthly_usage = float(st.session_state.get("pf_usage", round(_monthly_est)))
            st.markdown(
                f"<div style='padding-top:0.2rem;'>"
                f"<div style='font-size:0.875rem;color:#31333F;margin-bottom:0.25rem;'>월 평균 공급량 (kg)</div>"
                f"<div style='font-size:1rem;font-weight:600;color:#0F172A;'>{monthly_usage:,.0f}</div>"
                f"<div style='font-size:0.75rem;color:#64748B;'>사용주기 자동반영 (직접입력은 위 체크 해제)</div>"
                f"</div>",
                unsafe_allow_html=True,
            )
        else:
            monthly_usage = profit_int_comma_input(
                "월 평균 공급량 (kg)",
                key="pf_usage",
                value=p0["monthly_usage_kg"],
                help="직접 입력. 사용주기에서 자동반영하려면 위 체크박스를 켜세요.",
            )
    with i3:
        construction = profit_int_comma_input(
            "3. 공사비용 (원)", key="pf_const", value=p0["construction_cost"]
        )
    v1, v2, v3 = st.columns(3)
    with v1:
        vap_cap = profit_int_comma_input(
            "2. 기화기 용량 (Nm3/hr)", key="pf_vap_cap", value=p0["vaporizer_capacity"]
        )
    vap_note = v2.text_input("기화기 수량 메모", value=str(p0.get("vaporizer_qty_note", "")), key="pf_vap_note")
    with v3:
        vap_price = profit_int_comma_input(
            "기화기 구입가 (원)", key="pf_vap_price", value=p0["vaporizer_price"]
        )
    # —— 물류비 계산 (단가 바로 위) ——
    st.markdown("##### ◎ 물류비 계산")
    st.caption(
        "20톤 벌크로리 · 경유 · 통행료 5종  ·  "
        "왕복=탱크용량×80%충전 기준 자동  ·  "
        "((거리km × 유류비/L ÷ 연비) + 통행료) × 왕복 ÷ 월평균공급량 → 「2. 물류비」반영"
    )
    # 경유 시장가(오피넷) — 매월 자동 반영
    force_diesel = st.session_state.pop("pf_diesel_force", False)
    diesel_info = get_diesel_price_monthly(force_refresh=force_diesel)
    month_now = datetime.date.today().strftime("%Y-%m")
    if diesel_info.get("ok"):
        if st.session_state.get("pf_diesel_applied_month") != month_now or force_diesel:
            st.session_state["pf_lfuel"] = float(diesel_info["price"])
            st.session_state["pf_diesel_applied_month"] = month_now
    a1, a2 = st.columns(2)
    logi_origin = a1.text_input(
        "출발지 (주소·상호·명칭)",
        value=str(p0.get("logi_origin", PROFIT_DEFAULTS["logi_origin"])),
        key="pf_origin",
        placeholder="예: 신일가스 화성공장 / 경기도 화성시 …",
        help="상호·지점명이 여러 개면 아래에서 선택할 수 있습니다.",
    )
    logi_dest = a2.text_input(
        "도착지 (주소·상호·명칭)",
        value=str(p0.get("logi_dest", "")),
        key="pf_dest",
        placeholder="예: 제이아이금속 / 경북 구미시 …",
        help="상호·지점명이 여러 개면 아래에서 선택할 수 있습니다.",
    )
    _oq = str(logi_origin or "").strip()
    _dq = str(logi_dest or "").strip()
    _o_cands = st.session_state.get("pf_origin_cands") if st.session_state.get("pf_origin_q") == _oq else []
    _d_cands = st.session_state.get("pf_dest_cands") if st.session_state.get("pf_dest_q") == _dq else []
    _o_cands = _o_cands or []
    _d_cands = _d_cands or []
    # 지점이 여러 개면 선택 UI
    if _o_cands or _d_cands:
        p1, p2 = st.columns(2)
        with p1:
            if len(_o_cands) > 1:
                st.selectbox(
                    f"출발 지점 선택 ({len(_o_cands)}곳)",
                    options=list(range(len(_o_cands))),
                    format_func=lambda i: _o_cands[i]["label"],
                    key="pf_origin_pick",
                )
            elif len(_o_cands) == 1:
                st.caption(f"출발: {_o_cands[0]['label']}")
        with p2:
            if len(_d_cands) > 1:
                st.selectbox(
                    f"도착 지점 선택 ({len(_d_cands)}곳)",
                    options=list(range(len(_d_cands))),
                    format_func=lambda i: _d_cands[i]["label"],
                    key="pf_dest_pick",
                )
            elif len(_d_cands) == 1:
                st.caption(f"도착: {_d_cands[0]['label']}")
    btn_r, btn_f = st.columns(2)
    if btn_r.button("📍 거리·통행료 조회", key="pf_route_btn", type="secondary", width="stretch"):
        with st.spinner("통합검색·거리·통행료 조회 중…"):
            need_search = (
                st.session_state.get("pf_origin_q") != _oq
                or st.session_state.get("pf_dest_q") != _dq
                or not st.session_state.get("pf_origin_cands")
                or not st.session_state.get("pf_dest_cands")
            )
            if need_search:
                o_cands = kakao_place_search(_oq)
                d_cands = kakao_place_search(_dq)
                st.session_state["pf_origin_cands"] = o_cands
                st.session_state["pf_dest_cands"] = d_cands
                st.session_state["pf_origin_q"] = _oq
                st.session_state["pf_dest_q"] = _dq
                st.session_state.pop("pf_origin_pick", None)
                st.session_state.pop("pf_dest_pick", None)
                if not o_cands or not d_cands:
                    st.session_state["pf_route_info"] = {
                        "ok": False,
                        "message": "출발/도착을 찾지 못했습니다. 주소·상호·명칭을 확인해 주세요.",
                    }
                    st.error(st.session_state["pf_route_info"]["message"])
                elif len(o_cands) > 1 or len(d_cands) > 1:
                    st.warning("지점이 여러 개입니다. 아래에서 지점을 고른 뒤 다시 「거리·통행료 조회」를 눌러주세요.")
                    st.rerun()
                else:
                    oi, di = 0, 0
                    route = kakao_route_from_coords(
                        o_cands[oi]["lat"], o_cands[oi]["lon"],
                        d_cands[di]["lat"], d_cands[di]["lon"],
                        o_cands[oi]["label"], d_cands[di]["label"],
                    )
                    st.session_state["pf_route_info"] = route
                    if route.get("ok"):
                        st.session_state["pf_lkm"] = round(float(route["km"]), 1)
                        st.session_state["pf_ltoll"] = float(route.get("toll") or 0)
                        st.success(route.get("message") or f"거리 {route['km']:.1f} km")
                    else:
                        st.error(route.get("message") or "거리 조회 실패")
            else:
                o_cands = st.session_state.get("pf_origin_cands") or []
                d_cands = st.session_state.get("pf_dest_cands") or []
                if not o_cands or not d_cands:
                    st.error("출발/도착을 찾지 못했습니다. 주소·상호·명칭을 확인해 주세요.")
                else:
                    oi = int(st.session_state.get("pf_origin_pick", 0) or 0) if len(o_cands) > 1 else 0
                    di = int(st.session_state.get("pf_dest_pick", 0) or 0) if len(d_cands) > 1 else 0
                    oi = max(0, min(oi, len(o_cands) - 1))
                    di = max(0, min(di, len(d_cands) - 1))
                    route = kakao_route_from_coords(
                        o_cands[oi]["lat"], o_cands[oi]["lon"],
                        d_cands[di]["lat"], d_cands[di]["lon"],
                        o_cands[oi]["label"], d_cands[di]["label"],
                    )
                    st.session_state["pf_route_info"] = route
                    if route.get("ok"):
                        st.session_state["pf_lkm"] = round(float(route["km"]), 1)
                        st.session_state["pf_ltoll"] = float(route.get("toll") or 0)
                        st.success(route.get("message") or f"거리 {route['km']:.1f} km")
                    else:
                        st.error(route.get("message") or "거리 조회 실패")
    if btn_f.button("⛽ 경유 시세 새로고침", key="pf_diesel_btn", width="stretch"):
        st.session_state["pf_diesel_force"] = True
        st.rerun()
    route_info = st.session_state.get("pf_route_info") or {}
    if route_info.get("ok"):
        st.caption(
            f"출발: {route_info.get('origin_label','')} → 도착: {route_info.get('dest_label','')}  ·  "
            f"{route_info.get('message','')}"
        )
    if diesel_info.get("ok"):
        st.caption(
            f"경유 시장가 {float(diesel_info['price']):,.2f} 원/L  ·  "
            f"{diesel_info.get('source','')}  ·  기준일 {diesel_info.get('asof','')}  ·  "
            f"{month_now} 자동반영"
        )
    else:
        st.caption(diesel_info.get("message") or "경유 시세 조회 실패 — 유류비를 수동 입력하세요.")
    if "pf_lkm" not in st.session_state:
        st.session_state["pf_lkm"] = float(p0.get("logi_km", 70.0))
    if "pf_lfuel" not in st.session_state:
        st.session_state["pf_lfuel"] = float(
            diesel_info["price"] if diesel_info.get("ok") else p0.get("logi_fuel_price", 1400.0)
        )
    l1, l2, l3 = st.columns(3)
    logi_km = l1.number_input("거리 (km)", min_value=0.0, step=0.1, key="pf_lkm")
    logi_fuel = l2.number_input(
        "유류비 (원/L, 경유)", min_value=0.0, step=1.0, key="pf_lfuel",
    )
    logi_eff = l3.number_input(
        "연비 (km/L, 20톤벌크)", min_value=0.1, step=0.1,
        value=float(p0.get("logi_efficiency", 2.5)), key="pf_leff",
    )
    l4, l5, l6 = st.columns(3)
    if "pf_ltoll" not in st.session_state:
        st.session_state["pf_ltoll"] = float(p0.get("logi_toll", 0.0))
    logi_toll = l4.number_input(
        "통행료 (원, 편도·20톤벌크)", min_value=0.0, step=100.0, key="pf_ltoll",
    )
    # 왕복횟수 = ceil(월평균공급량 ÷ (탱크총용량kg × 80%))
    _rt_info = compute_roundtrips_from_tank(monthly_usage, tank_spec, fill_ratio=0.8)
    logi_rt = float(_rt_info["roundtrips"]) if _rt_info.get("ok") else float(p0.get("logi_roundtrips", 0) or 0)
    st.session_state["pf_lrt"] = logi_rt
    l5.markdown(
        f"<div style='padding-top:0.2rem;'>"
        f"<div style='font-size:0.875rem;color:#31333F;margin-bottom:0.25rem;'>왕복 횟수 (회/월)</div>"
        f"<div style='font-size:1rem;font-weight:600;color:#0F172A;'>{logi_rt:,.0f}</div>"
        f"<div style='font-size:0.75rem;color:#64748B;'>탱크 80% 소모시 충전</div>"
        f"</div>",
        unsafe_allow_html=True,
    )
    # 물류비 공급량 = 프로젝트 「월 평균 공급량」 그대로 (별도 입력 없음)
    logi_kg = float(monthly_usage) if float(monthly_usage or 0) > 0 else 1.0
    st.session_state["pf_lkg"] = logi_kg
    l6.markdown(
        f"<div style='padding-top:0.2rem;'>"
        f"<div style='font-size:0.875rem;color:#31333F;margin-bottom:0.25rem;'>월평균 공급량 (kg)</div>"
        f"<div style='font-size:1rem;font-weight:600;color:#0F172A;'>{logi_kg:,.0f}</div>"
        f"<div style='font-size:0.75rem;color:#64748B;'>프로젝트 월 평균 공급량 연동</div>"
        f"</div>",
        unsafe_allow_html=True,
    )
    logi_calc = compute_logistics_unit_cost(
        logi_km, logi_fuel, logi_eff, logi_toll, logi_rt, logi_kg
    )
    logi_per = int(round(float(logi_calc["per_kg"])))
    # 계산값을 단가 「물류비」위젯에 즉시 반영 (단가 렌더 전에 세팅)
    st.session_state["pf_logi"] = int(logi_per)
    st.caption(_rt_info.get("message") or "")
    st.caption(
        f"계산 물류비 {logi_per:,} 원/kg  ·  "
        f"편도연료 {logi_calc['fuel_one_way']:,.0f}원  ·  "
        f"왕복총비용 {logi_calc['round_trip_cost']:,.0f}원"
    )
    st.markdown("##### ◎ 단가")
    for _uk in ("pf_buy", "pf_logi", "pf_supply", "pf_dep", "pf_rent", "pf_rent_n"):
        if _uk in st.session_state:
            try:
                st.session_state[_uk] = int(round(float(st.session_state[_uk])))
            except Exception:
                pass
    u1, u2, u3 = st.columns(3)
    with u1:
        purchase_unit = st.number_input(
            "1. 매입단가 (원/kg)",
            min_value=0,
            step=1,
            value=int(round(float(p0["purchase_unit"]))),
            format="%d",
            key="pf_buy",
        )
    with u2:
        logistics_unit = st.number_input(
            "2. 물류비 (원/kg)",
            min_value=0,
            step=1,
            format="%d",
            key="pf_logi",
        )
        st.caption(
            f"적용: ((거리 {float(logi_km):,.0f}km × 유류 {float(logi_fuel):,.0f}원/L ÷ 연비 {float(logi_eff):g}) "
            f"+ 통행료 {float(logi_toll):,.0f}원) × 왕복 {float(logi_rt):,.0f}회 "
            f"÷ 월공급 {float(logi_kg):,.0f}kg = **{logi_per:,} 원/kg** → 위 칸에 자동반영"
        )
    with u3:
        supply_unit = st.number_input(
            "3. 공급단가 (원/kg)",
            min_value=0,
            step=1,
            value=int(round(float(p0["supply_unit"]))),
            format="%d",
            key="pf_supply",
        )
    st.markdown("##### ◎ 금융 / 관리 / 임대")
    f1, f2, f3, f4 = st.columns(4)
    interest = f1.number_input(
        "이자율 (예: 0.05=5%)", min_value=0.0, max_value=1.0, step=0.005, format="%.3f",
        value=float(p0["interest_rate"]), key="pf_rate",
    )
    mgmt_rate = f2.number_input(
        "일반관리비 비율 (월매출)", min_value=0.0, max_value=1.0, step=0.005, format="%.3f",
        value=float(p0["mgmt_rate"]), key="pf_mgmt",
    )
    dep_months = f3.number_input(
        "감가상각 개월 (10년=120)",
        min_value=1,
        step=12,
        value=int(round(float(p0["depreciation_months"]))),
        format="%d",
        key="pf_dep",
    )
    rent = f4.number_input(
        "장비 임대료 (원)",
        min_value=0,
        step=10000,
        value=int(round(float(p0["equipment_rent"]))),
        format="%d",
        key="pf_rent",
    )
    rent_count = st.number_input(
        "부가횟수",
        min_value=0,
        step=1,
        value=int(round(float(p0["rent_count"]))),
        format="%d",
        key="pf_rent_n",
    )
    b_save, b_reset = st.columns(2)
    submitted = b_save.button("💾 계산 / 저장", type="primary", width="stretch", key="pf_save_btn")
    reset = b_reset.button("↺ 엑셀 기본값으로 초기화", width="stretch", key="pf_reset_btn")
    if reset:
        st.session_state["profit_inputs"] = dict(PROFIT_DEFAULTS)
        save_profit_inputs(st.session_state["profit_inputs"])
        for k in _pf_keys:
            st.session_state.pop(k, None)
        st.session_state.pop("pf_route_info", None)
        st.session_state.pop("pf_diesel_applied_month", None)
        st.session_state.pop("pf_diesel_force", None)
        st.rerun()
    if submitted:
        new_p = {
            "project_name": project_name,
            "tank_gas": tank_gas,
            "tank_capacity_mode": tank_capacity_mode,
            "tank_liters": float(tank_liters),
            "tank_spec": tank_spec,
            "hourly_usage_mode": str(hourly_usage_mode),
            "hourly_usage_kg": float(hourly_usage_kg),
            "hourly_usage_nm3": float(hourly_usage_nm3),
            "operating_hours": float(operating_hours),
            "operating_days_per_month": float(operating_days),
            "auto_monthly_from_cycle": bool(auto_monthly_from_cycle),
            "tank_price": tank_price,
            "monthly_usage_kg": monthly_usage,
            "vaporizer_capacity": vap_cap,
            "vaporizer_qty_note": vap_note,
            "vaporizer_price": vap_price,
            "construction_cost": construction,
            "purchase_unit": purchase_unit,
            "logistics_unit": float(logistics_unit),
            "supply_unit": supply_unit,
            "interest_rate": interest,
            "mgmt_rate": mgmt_rate,
            "depreciation_months": dep_months,
            "equipment_rent": rent,
            "rent_count": rent_count,
            "logi_km": float(logi_km),
            "logi_fuel_price": float(logi_fuel),
            "logi_efficiency": float(logi_eff),
            "logi_toll": float(logi_toll),
            "logi_roundtrips": float(logi_rt),
            "logi_supply_kg": float(logi_kg),
            "logi_origin": str(logi_origin),
            "logi_dest": str(logi_dest),
        }
        st.session_state["profit_inputs"] = new_p
        save_profit_inputs(new_p)
        st.success("저장되었습니다. 아래 결과가 갱신됩니다.")
    # 화면 결과도 현재 입력(자동 반영된 물류비 포함) 기준으로 표시
    p = {
        "project_name": project_name,
        "tank_gas": tank_gas,
        "tank_capacity_mode": tank_capacity_mode,
        "tank_liters": float(tank_liters),
        "tank_spec": tank_spec,
        "hourly_usage_mode": str(hourly_usage_mode),
        "hourly_usage_kg": float(hourly_usage_kg),
        "hourly_usage_nm3": float(hourly_usage_nm3),
        "operating_hours": float(operating_hours),
        "operating_days_per_month": float(operating_days),
        "auto_monthly_from_cycle": bool(auto_monthly_from_cycle),
        "tank_price": tank_price,
        "monthly_usage_kg": monthly_usage,
        "vaporizer_capacity": vap_cap,
        "vaporizer_qty_note": vap_note,
        "vaporizer_price": vap_price,
        "construction_cost": construction,
        "purchase_unit": purchase_unit,
        "logistics_unit": float(logistics_unit),
        "supply_unit": supply_unit,
        "interest_rate": interest,
        "mgmt_rate": mgmt_rate,
        "depreciation_months": dep_months,
        "equipment_rent": rent,
        "rent_count": rent_count,
        "logi_km": float(logi_km),
        "logi_fuel_price": float(logi_fuel),
        "logi_efficiency": float(logi_eff),
        "logi_toll": float(logi_toll),
        "logi_roundtrips": float(logi_rt),
        "logi_supply_kg": float(logi_kg),
        "logi_origin": str(logi_origin),
        "logi_dest": str(logi_dest),
    }
    st.session_state["profit_inputs"] = p
    r = compute_profitability(p)
    st.markdown("---")
    st.markdown(
        f"<div class='sub-header dashboard-tab-panel-head'>◆ [{html.escape(str(p.get('project_name') or '프로젝트'))}] 투자대비 수익성 분석 보고</div>",
        unsafe_allow_html=True,
    )
    st.caption(
        f"TANK: {p.get('tank_gas','')} {float(p.get('tank_liters') or 0):,.0f} L → "
        f"{parse_tank_capacity_kg(p.get('tank_spec')):,.0f} kg · "
        f"기화기 {p.get('vaporizer_capacity',0):,.0f} Nm3/hr {p.get('vaporizer_qty_note','')} · "
        f"년 사용량(자동) {r['yearly_usage']:,.0f} kg"
    )
    m1, m2, m3, m4 = st.columns(4)
    m1.markdown(
        f"<div class='metric-box'><div class='metric-label'>합계 투자액 (C9+C16+C18)</div>"
        f"<div class='metric-value'>{r['total_invest']:,.0f} 원</div></div>",
        unsafe_allow_html=True,
    )
    m2.markdown(
        f"<div class='metric-box'><div class='metric-label'>매출이익 (원/kg)</div>"
        f"<div class='metric-value'>{r['margin_kg']:,.1f} 원/kg</div></div>",
        unsafe_allow_html=True,
    )
    m3.markdown(
        f"<div class='metric-box'><div class='metric-label'>월평균 이익금</div>"
        f"<div class='metric-value' style='color:{'#059669' if r['monthly_profit']>=0 else '#E11D48'};'>"
        f"{r['monthly_profit']:,.0f} 원/월</div></div>",
        unsafe_allow_html=True,
    )
    m4.markdown(
        f"<div class='metric-box'><div class='metric-label'>최근 3개월 매출 실이익</div>"
        f"<div class='metric-value'>{r['three_month']:,.0f} 원</div></div>",
        unsafe_allow_html=True,
    )
    left, right = st.columns(2)
    def _pf_i(v):
        try:
            return f"{int(round(float(v or 0))):,}"
        except Exception:
            return "0"
    with left:
        st.markdown("**◎ 단가 · 투자 요약**")
        summary_df = pd.DataFrame(
            [
                {"항목": "TANK 구입가", "값": _pf_i(p["tank_price"]), "단위": "원"},
                {"항목": "기화기 구입가", "값": _pf_i(p["vaporizer_price"]), "단위": "원"},
                {"항목": "공사비용", "값": _pf_i(p["construction_cost"]), "단위": "원"},
                {"항목": "합계 금액", "값": _pf_i(r["total_invest"]), "단위": "원"},
                {"항목": "년 사용량", "값": _pf_i(r["yearly_usage"]), "단위": "kg"},
                {"항목": "월 평균 공급량", "값": _pf_i(p["monthly_usage_kg"]), "단위": "kg"},
                {"항목": "매입단가", "값": _pf_i(p["purchase_unit"]), "단위": "원/kg"},
                {"항목": "물류비", "값": _pf_i(p["logistics_unit"]), "단위": "원/kg"},
                {"항목": "공급단가", "값": _pf_i(p["supply_unit"]), "단위": "원/kg"},
                {"항목": "매출이익", "값": _pf_i(r["margin_kg"]), "단위": "원/kg"},
                {"항목": "물류비 계산(P18)", "값": _pf_i(r["logi_per_kg"]), "단위": "원/kg"},
            ]
        )
        st.dataframe(summary_df, hide_index=True, width="stretch", height=420)
    with right:
        st.markdown("**◎ 사용량 대비 영업이익 (월)**")
        result_df = pd.DataFrame(
            [
                {"항목": "월평균 매출이익", "계산": "월사용량 × 매출이익", "값": _pf_i(r["monthly_gross"]), "단위": "원/월"},
                {"항목": "장비 감가상각", "계산": f"투자합계 ÷ {p['depreciation_months']:.0f}", "값": _pf_i(r["depreciation"]), "단위": "원/월"},
                {"항목": "금융비", "계산": "원금 × 이자율 ÷ 12", "값": _pf_i(r["finance"]), "단위": "원/월"},
                {"항목": "월 매출", "계산": "월사용량 × 공급단가", "값": _pf_i(r["monthly_sales"]), "단위": "원/월"},
                {"항목": "일반관리비", "계산": f"월매출 × {p['mgmt_rate']*100:.1f}%", "값": _pf_i(r["mgmt"]), "단위": "원/월"},
                {"항목": "투자비용(상각+금융)", "계산": "감가상각 + 금융비", "값": _pf_i(r["invest_cost"]), "단위": "원/월"},
                {"항목": "월평균 이익금", "계산": "매출이익 − 투자비용 − 관리비", "값": _pf_i(r["monthly_profit"]), "단위": "원/월"},
                {"항목": "최근 3개월 실이익", "계산": "월이익×3 + 임대료×횟수", "값": _pf_i(r["three_month"]), "단위": "원"},
            ]
        )
        st.dataframe(result_df, hide_index=True, width="stretch", height=420)
    st.info(
        "적용 함수: `년사용량=월×12` · `합계=탱크+기화기+공사` · `매출이익=공급−(매입+물류)` · "
        "`관리비=월매출×14.5%` · `감가=투자÷120` · `금융=원금×이자÷12` · "
        "`월이익=매출이익−(감가+금융)−관리비` · `3개월=월이익×3+(임대×횟수)` · "
        "`물류원/kg=(KM×유류비/연비+통행료)×왕복÷KG`"
    )
    # Tab9 전용 — 스크린샷형 수익성 보고서 엑셀 (동일 입력이면 캐시 재사용)
    _route_for_xlsx = st.session_state.get("pf_route_info") or {}
    _diesel_for_xlsx = diesel_info if isinstance(diesel_info, dict) else {}
    try:
        _profit_xlsx = _cached_profitability_report_excel(
            json.dumps(p, ensure_ascii=False, sort_keys=True, default=str),
            json.dumps(r, ensure_ascii=False, sort_keys=True, default=str),
            json.dumps(_route_for_xlsx, ensure_ascii=False, sort_keys=True, default=str),
            json.dumps(_diesel_for_xlsx, ensure_ascii=False, sort_keys=True, default=str),
        )
        st.download_button(
            "📥 수익성 분석 보고서 엑셀 내보내기",
            data=_profit_xlsx,
            file_name=f"수익성분석_{str(p.get('project_name') or '보고서')}_{datetime.date.today().strftime('%Y%m%d')}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            width="stretch",
            key="pf_report_xlsx_btn",
            help="스크린샷과 같은 보고서 형태. 거리·통행료·경유시세 설명은 작은 글씨로 포함됩니다.",
        )
    except Exception as _xlsx_err:
        st.warning(f"보고서 엑셀 생성 실패: {_xlsx_err}")

with tab9:
    _render_profitability_analysis_tab(latest_update_str)

with tab10:
    # 업무일지 탭 전용 — 다른 탭과 공유 상태/헬퍼를 쓰지 않음.
    # 파일 mtime 변경 시에만 reload (매번 reload 금지 → 달력/저장 로딩 감소).
    # 로드 실패 시에도 다른 탭은 유지.
    try:
        import importlib
        import os
        import sys

        import worklog_tab as _worklog_tab

        _wl_path = getattr(_worklog_tab, "__file__", None) or ""
        _wl_mtime = os.path.getmtime(_wl_path) if _wl_path and os.path.exists(_wl_path) else 0
        if "_wl_mod_mtime" not in st.session_state:
            st.session_state["_wl_mod_mtime"] = _wl_mtime
        elif st.session_state.get("_wl_mod_mtime") != _wl_mtime:
            _worklog_tab = importlib.reload(_worklog_tab)
            st.session_state["_wl_mod_mtime"] = _wl_mtime
            sys.modules["worklog_tab"] = _worklog_tab
        _worklog_tab.render_worklog_tab(latest_update_str)
    except ModuleNotFoundError:
        st.error(
            "일일업무일지 모듈(`worklog_tab.py`)을 찾을 수 없습니다. "
            "배포 파일에 포함되었는지 확인해 주세요."
        )
        st.info("다른 탭은 정상 이용 가능합니다.")
    except Exception as _wl_err:
        if is_touch_ui():
            st.error("일일업무일지 탭을 표시하지 못했습니다. 잠시 후 새로고침해 주세요.")
            st.info("다른 탭은 정상 이용 가능합니다.")
        else:
            st.error(f"일일업무일지 탭 오류: {_wl_err}")
            st.info("다른 탭은 정상 이용 가능합니다.")

with tab11:
    # 시장조사 탭 전용 — 다른 탭과 공유 상태/헬퍼를 쓰지 않음.
    try:
        import importlib
        import os
        import sys

        import market_research_tab as _mr_tab

        _mr_path = getattr(_mr_tab, "__file__", None) or ""
        _mr_mtime = os.path.getmtime(_mr_path) if _mr_path and os.path.exists(_mr_path) else 0
        if "_mr_mod_mtime" not in st.session_state:
            st.session_state["_mr_mod_mtime"] = _mr_mtime
        elif st.session_state.get("_mr_mod_mtime") != _mr_mtime:
            _mr_tab = importlib.reload(_mr_tab)
            st.session_state["_mr_mod_mtime"] = _mr_mtime
            sys.modules["market_research_tab"] = _mr_tab
        _mr_tab.render_market_research_tab(latest_update_str)
    except ModuleNotFoundError:
        st.error(
            "시장조사 모듈(`market_research_tab.py`)을 찾을 수 없습니다. "
            "배포 파일에 포함되었는지 확인해 주세요."
        )
        st.info("다른 탭은 정상 이용 가능합니다.")
    except Exception as _mr_err:
        st.error(f"시장조사 탭 오류: {_mr_err}")
        st.info("다른 탭은 정상 이용 가능합니다. 새로고침 후에도 같으면 관리자에게 오류 문구를 보내 주세요.")
